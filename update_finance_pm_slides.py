#!/usr/bin/env python3
"""
Update File 07 — Finance & Performance Management: Accounting Operations
Enriches 46 slides with Pakistan banking finance context, adds 3 new slides,
removes Teradata branding, updates dates, adds speaker notes.
"""

import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.shapes.group import GroupShape

INPUT_FILE = './pptout/07_Finance_PM_Accounting_Operations.pptx'
CLEAN_FILE = './pptout/07_Finance_PM_Accounting_Operations_CLEAN.pptx'
OUTPUT_FILE = './pptout/07_Finance_PM_Accounting_Operations_UPDATED.pptx'


# ─────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────

def set_notes(slide, text):
    """Set or replace the speaker notes on a slide."""
    if not slide.has_notes_slide:
        slide.notes_slide  # creates notes slide
    notes_tf = slide.notes_slide.notes_text_frame
    # Clear existing
    for i in range(len(notes_tf.paragraphs) - 1, 0, -1):
        p = notes_tf.paragraphs[i]._element
        p.getparent().remove(p)
    notes_tf.paragraphs[0].text = text


def replace_cell_text(cell, new_text):
    """Replace cell text while preserving first run formatting."""
    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
        fmt = cell.text_frame.paragraphs[0].runs[0].font
        font_size = fmt.size
        font_bold = fmt.bold
        font_name = fmt.name
    else:
        font_size = Pt(10)
        font_bold = False
        font_name = None

    # Clear all paragraphs
    for i in range(len(cell.text_frame.paragraphs) - 1, 0, -1):
        p = cell.text_frame.paragraphs[i]._element
        p.getparent().remove(p)

    lines = new_text.split('\n')
    for idx, line in enumerate(lines):
        if idx == 0:
            p = cell.text_frame.paragraphs[0]
            p.clear()
        else:
            p = cell.text_frame.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.bold = font_bold
        if font_name:
            run.font.name = font_name


def append_cell_text(cell, additional_text):
    """Append text to existing cell content."""
    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
        fmt = cell.text_frame.paragraphs[0].runs[0].font
        font_size = fmt.size
        font_name = fmt.name
    else:
        font_size = Pt(10)
        font_name = None

    for line in additional_text.split('\n'):
        p = cell.text_frame.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        if font_name:
            run.font.name = font_name


def remove_engagement_text(slide):
    """Remove 'For use in Maturity Assessment & Roadmap Engagements' text from groups."""
    removed = False
    for shape in list(slide.shapes):
        if isinstance(shape, GroupShape):
            for s in shape.shapes:
                if s.has_text_frame:
                    txt = s.text_frame.text.strip()
                    if 'Engagements' in txt and 'Maturity Assessment' in txt:
                        for para in s.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = ''
                        removed = True
    return removed


def get_table(slide):
    """Return the first table shape on a slide."""
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def replace_h1_2018(table):
    """Replace 'H1 2018' with 'H1 2025 — H2 2026' in table cells."""
    replaced = False
    for row in table.rows:
        for cell in row.cells:
            if 'H1 2018' in cell.text:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if 'H1 2018' in run.text:
                            run.text = run.text.replace('H1 2018', 'H1 2025 —\nH2 2026')
                            replaced = True
    return replaced


# ─────────────────────────────────────────────
# Capability enrichments for 10 detail slides
# Slides 22,24,26,28,30,32,34,36,38,40 (0-indexed: 21,23,25,27,29,31,33,35,37,39)
# Table structure: 4x3
#   [0,0]=Business Objectives, [1,0]=content
#   [2,0]=Data & Solution, [2,2]=Outcome
#   [3,0]=data content, [3,2]=outcome content
# ─────────────────────────────────────────────

CAPABILITY_ENRICHMENTS = {
    # Slide 22: Billing & Collections
    21: {
        'biz_obj': [
            "Reduce NPL ratio from 7.5% industry average through early collection intervention",
            "Optimize collection strategies by segment: mass (SMS), affluent (RM call), corporate (legal)",
            "Integrate ECIB credit bureau data for collection prioritization",
            "Track SBP classification migration: Current > OAEM > Substandard > Doubtful > Loss",
        ],
        'data_sol': [
            "ECIB credit bureau scores and payment history",
            "SBP classification data (OAEM/Substandard/Doubtful/Loss)",
            "Court decree data for legal recovery tracking",
            "Restructured loan tracking and performance monitoring",
        ],
        'outcome': [
            "Reduce days past due (DPD) migration across Pakistan loan portfolio",
            "Improve recovery rates on PKR 900B+ classified assets industry-wide",
            "Automate ECIB reporting and SBP classification compliance",
            "Early analytical intervention to avoid 3-7 year legal recovery process",
        ],
    },
    # Slide 24: Procurement & Payment
    23: {
        'biz_obj': [
            "Centralize procurement analytics across group entities (conventional + Islamic + Takaful + brokerage)",
            "Comply with FBR WHT requirements on vendor payments (variable rates by vendor type and filer status)",
            "Automate FBR Active Taxpayer List (ATL) lookup for WHT rate determination",
        ],
        'data_sol': [
            "FBR filer/non-filer status of vendors (determines WHT rate)",
            "NTN (National Tax Number) for all vendors",
            "Vendor banking details for RAAST payment processing",
            "WHT rates: Services 8%/16%, Supplies 4%/8%, Contracts 7%/14% (filer/non-filer)",
        ],
        'outcome': [
            "Automated FBR WHT calculation and e-filing integration",
            "Vendor performance analytics across all banking group entities",
            "Procurement consolidation reducing costs 10-15%",
            "Full compliance with Pakistan FBR withholding tax regulations",
        ],
    },
    # Slide 26: Sub-ledger Accounting
    25: {
        'biz_obj': [
            "Unify sub-ledger from 4-5 core banking systems (CTL, Temenos, Oracle, Finacle) into FSDM warehouse",
            "Support dual accounting: IFRS entries for conventional + AAOIFI entries for Islamic products",
            "Maintain individual loan-level sub-ledger for IFRS 9 ECL calculations",
            "Enable drill-down from GL balance to individual transactions across all source systems",
        ],
        'data_sol': [
            "Core banking transaction data from all systems (CTL, Temenos, Oracle, Finacle)",
            "Islamic banking sub-system data (parallel IFRS + AAOIFI entries)",
            "Card system data (Visa/Mastercard settlements)",
            "Trade finance and treasury system transaction data",
        ],
        'outcome': [
            "Single source of truth for SBP returns, IFRS 9, management reporting, audit",
            "Consistent chart of accounts mapping across all core systems",
            "Full transaction lineage from GL to individual operational entries",
            "Dual IFRS/AAOIFI accounting for Islamic banking operations",
        ],
    },
    # Slide 28: Fair Value & Hedge Accounting
    27: {
        'biz_obj': [
            "Manage fair value for PKR 10T+ government securities (PIBs, T-Bills, Sukuk, Ijara) per SBP classification",
            "Hedge FX exposure — Pakistan banks hold significant USD, GBP, EUR positions for trade/remittance",
            "Comply with IFRS 9 hedge effectiveness requirements — SBP expects robust documentation",
            "Islamic Sukuk valuation under both IFRS and AAOIFI measurement approaches",
        ],
        'data_sol': [
            "SBP government securities rates (PKR-denominated PIBs, T-Bills)",
            "KIBOR rates for yield curve construction",
            "FX rates (SBP daily fixing) for foreign currency positions",
            "Islamic Sukuk and Ijara valuations from MUFAP/SBP",
        ],
        'outcome': [
            "Accurate P&L impact from HTM/AFS/HFT reclassification under IFRS 9",
            "Automated mark-to-market for AFS and HFT government securities",
            "SBP-compliant FX open position limit monitoring",
            "Hedge effectiveness documentation meeting external audit requirements",
        ],
    },
    # Slide 30: Asset Valuations
    29: {
        'biz_obj': [
            "Value bank physical assets: 16,000+ branch properties, ATM network, IT infrastructure, vehicle fleet",
            "IFRS 9 impairment assessment for investment portfolio",
            "Comply with SBP Fixed Asset regulations (revaluation surplus for Tier 2 capital)",
            "Collateral valuations for secured lending portfolio (property, vehicle, gold, shares)",
        ],
        'data_sol': [
            "Branch property valuations updated per SBP requirements (every 3-5 years)",
            "IT asset register across all bank locations",
            "Investment portfolio market values (daily for AFS/HFT)",
            "Collateral valuations for PKR 13T+ lending portfolio",
        ],
        'outcome': [
            "Accurate regulatory capital from revaluation surplus (Tier 2)",
            "Automated impairment testing per IFRS 9 and IAS 36",
            "Collateral coverage monitoring for secured lending risk management",
            "Complete fixed asset lifecycle tracking across 16,000+ locations",
        ],
    },
    # Slide 32: Reserve Analytics
    31: {
        'biz_obj': [
            "Calculate IFRS 9 ECL reserves at individual loan level (PD x LGD x EAD)",
            "Maintain general provisions (1-2% per SBP) plus specific provisions (25-100%) on classified assets",
            "Islamic banking provisioning for Diminishing Musharaka, Ijarah, Murabaha under IFRS 9 + AAOIFI FAS 30",
            "Forward-looking ECL with Pakistan macroeconomic scenarios (GDP, inflation, KIBOR, default rates)",
        ],
        'data_sol': [
            "ECIB credit bureau data for PD estimation",
            "Historical default and recovery data (5+ years) for LGD calibration",
            "Collateral values for secured loan LGD reduction",
            "Restructured loan data and customer financial statements",
        ],
        'outcome': [
            "IFRS 9 ECL automation: Stage 1/2/3 classification with SICR monitoring",
            "SBP-compliant provisioning with regulatory overlay capability",
            "Reduced external audit findings on provisioning adequacy",
            "Pakistan macro-scenario-driven ECL for stress testing and SBP reporting",
        ],
    },
    # Slide 34: Tax Management & Optimisation
    33: {
        'biz_obj': [
            "Manage Pakistan banking tax: corporate tax (39%), super tax (10%), WHT on deposits (15%/30%), advance tax",
            "Automate FBR withholding tax on customer deposits and vendor payments",
            "Track tax credits, brought-forward losses, group entity tax planning",
            "Calculate effective tax rate (~49% for Pakistan banks — highest in industry)",
        ],
        'data_sol': [
            "FBR Active Taxpayer List (ATL) — updated monthly for filer status",
            "Customer CNIC/NTN for filer status determination on deposit WHT",
            "Deposit profit payment data across all accounts (millions of transactions)",
            "Vendor payment records with WHT rates by category and filer status",
        ],
        'outcome': [
            "Automated WHT calculation on millions of deposit profit payments",
            "FBR ATL auto-lookup eliminating manual filer status verification",
            "Consolidated tax provisioning across all group entities",
            "Tax optimization within legal framework saving PKR hundreds of millions",
        ],
    },
    # Slide 36: Reconciliation, Validation & Adjustments
    35: {
        'biz_obj': [
            "Automate reconciliation between 4-5 core banking systems and data warehouse",
            "GL-to-sub-ledger reconciliation for SBP reporting accuracy",
            "Automated nostro/vostro reconciliation for correspondent banking",
            "Reduce finance team reconciliation effort from 40-50% to <15% of time",
        ],
        'data_sol': [
            "GL balances from multiple ERPs/core banking systems",
            "RAAST settlement data and 1Link switch data",
            "Card network settlement data (Visa/Mastercard)",
            "SBP clearing house (NIFT) and nostro/vostro account data",
        ],
        'outcome': [
            "Automated daily reconciliation across all payment systems",
            "Exception-based processing — only manual review for breaks",
            "Month-end close acceleration from 15-30 days to <10 days",
            "Audit-ready reconciliation documentation with full trail",
        ],
    },
    # Slide 38: Financial Consolidation
    37: {
        'biz_obj': [
            "Consolidate across group: conventional bank + Islamic window/subsidiary + Takaful + brokerage + AMC + modaraba",
            "Eliminate inter-company transactions between group entities",
            "Produce consolidated SBP returns AND consolidated IFRS financial statements",
            "Handle different ERPs and chart of accounts across group entities",
        ],
        'data_sol': [
            "GL data from all group entities (often different ERP/core systems)",
            "Islamic banking separate books (IFRS + AAOIFI)",
            "Subsidiary financial data with different chart of accounts",
            "Inter-company transaction register for elimination entries",
        ],
        'outcome': [
            "Consolidated SBP returns (solo + consolidated) within regulatory deadlines",
            "SECP-compliant consolidated IFRS statements for listed banking groups",
            "Automated IC elimination reducing manual consolidation effort by 60%",
            "Drill-through from consolidated view to any entity-level detail",
        ],
    },
    # Slide 40: Functional Profit & Loss Statement
    39: {
        'biz_obj': [
            "Branch-level P&L for all 16,000+ branches across Pakistan",
            "Customer-level profitability (building on UBL Customer Profitability Engine / FSDM)",
            "Product-level P&L across conventional and Islamic product lines",
            "Business segment P&L per SBP/IFRS 8: Corporate, Commercial, Retail, Treasury, Islamic, Digital",
        ],
        'data_sol': [
            "FSDM star schema: FACT_CUSTOMER_PROFITABILITY (35+ measures)",
            "Activity-based costing allocations across branch, digital, ATM, agent channels",
            "Funds Transfer Pricing (FTP) rates for NII allocation",
            "Direct and allocated costs by cost center, branch, product, segment",
        ],
        'outcome': [
            "360-degree profitability: branch × customer × product × segment × channel × region",
            "FTP-adjusted Net Interest Income at individual customer level",
            "Islamic P&L with correct line items (profit on financing, return on deposits)",
            "Data-driven decisions on branch viability, product pricing, customer retention",
        ],
    },
}


# ─────────────────────────────────────────────
# Maturity enrichments for 10 maturity slides
# Slides 23,25,27,29,31,33,35,37,39,41 (0-indexed: 22,24,26,28,30,32,34,36,38,40)
# Table structure: 10x7
#   [7,3] = key areas for improvement, [7,4] = benefits, [7,5] = timeframe
# ─────────────────────────────────────────────

MATURITY_ENRICHMENTS = {
    # Slide 23: Billing & Collections maturity
    22: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "ECIB integration for collection prioritization and early intervention\n"
            "Segment-based collection strategy (mass/affluent/corporate/Islamic)\n"
            "SBP classification automation: OAEM/Substandard/Doubtful/Loss tracking\n"
            "Analytical early-warning system to avoid 3-7 year legal recovery process\n"
            "Automated provisioning per SBP general (1-2%) and specific (25-100%) rules"
        ),
        'benefits': (
            "Reduced NPL ratio through early analytical intervention\n"
            "Lower provisioning charges from faster recovery and DPD management\n"
            "SBP compliance for classification and provisioning requirements\n"
            "Reduced legal costs from court-based recovery process"
        ),
    },
    # Slide 25: Procurement & Payment maturity
    24: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Automated FBR WHT calculation on vendor payments (variable rates by category)\n"
            "FBR ATL auto-lookup for filer/non-filer status determination\n"
            "Centralized procurement across banking group entities\n"
            "RAAST integration for vendor payments\n"
            "Group-wide spend analytics across conventional + Islamic + subsidiaries"
        ),
        'benefits': (
            "Full FBR WHT compliance — automated calculation and e-filing\n"
            "10-15% procurement cost reduction from group-wide consolidation\n"
            "Vendor performance visibility across all banking group entities\n"
            "Reduced processing cost through RAAST vendor payments"
        ),
    },
    # Slide 27: Sub-ledger Accounting maturity
    26: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Unified sub-ledger from 4-5 core banking systems (CTL, Temenos, Oracle, Finacle)\n"
            "Dual IFRS + AAOIFI accounting entries for Islamic banking operations\n"
            "Transaction-level detail for IFRS 9 ECL at individual loan level\n"
            "Consistent chart of accounts mapping across all source systems\n"
            "Full audit trail: GL balance to source system transaction"
        ),
        'benefits': (
            "Single source of truth for all financial reporting (SBP, IFRS, management)\n"
            "IFRS 9 ECL calculation enabled by loan-level sub-ledger data\n"
            "External audit efficiency: drill-down from any GL balance\n"
            "Foundation for branch P&L, customer profitability, product profitability"
        ),
    },
    # Slide 29: Fair Value & Hedge Accounting maturity
    28: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "PKR 10T+ government securities portfolio fair value management\n"
            "SBP HTM/AFS/HFT classification impact on P&L and OCI\n"
            "FX hedge documentation meeting IFRS 9 and SBP requirements\n"
            "Islamic Sukuk/Ijara fair valuation under both IFRS and AAOIFI\n"
            "Daily mark-to-market automation for AFS and HFT portfolios"
        ),
        'benefits': (
            "Accurate P&L impact from government securities portfolio movements\n"
            "SBP-compliant FX open position limit monitoring\n"
            "Reduced earnings volatility through proper hedge accounting\n"
            "External audit sign-off on hedge effectiveness documentation"
        ),
    },
    # Slide 31: Asset Valuations maturity
    30: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "16,000+ branch property revaluation management (SBP 3-5 year cycle)\n"
            "IFRS 9 impairment testing for investment portfolio\n"
            "Collateral valuation for PKR 13T+ secured lending portfolio\n"
            "Fixed asset lifecycle tracking across all bank locations\n"
            "Revaluation surplus management for SBP Tier 2 capital calculation"
        ),
        'benefits': (
            "Regulatory capital optimization from accurate revaluation surplus\n"
            "Timely impairment recognition per IFRS 9 and IAS 36\n"
            "Improved collateral coverage monitoring for credit risk management\n"
            "Complete asset register for insurance and depreciation management"
        ),
    },
    # Slide 33: Reserve Analytics maturity
    32: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "IFRS 9 ECL automation: PD × LGD × EAD at individual loan level\n"
            "SBP provisioning compliance: general + specific + regulatory overlay\n"
            "Forward-looking macro scenarios (GDP, inflation, KIBOR, default rates)\n"
            "Islamic financing provisioning under AAOIFI FAS 30\n"
            "Stage 1/2/3 classification with SICR monitoring automation"
        ),
        'benefits': (
            "Accurate ECL reducing external audit provisioning adjustments\n"
            "SBP stress testing capability with multiple Pakistan macro scenarios\n"
            "Reduced manual effort in quarterly IFRS 9 ECL recalculation\n"
            "Islamic and conventional provisioning from unified data source"
        ),
    },
    # Slide 35: Tax Management & Optimisation maturity
    34: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Automated WHT on deposit profit (millions of transactions, 15%/30% rates)\n"
            "FBR ATL integration for real-time filer/non-filer determination\n"
            "Consolidated tax provisioning across all banking group entities\n"
            "Super tax (10%), advance tax, minimum tax calculation automation\n"
            "Tax planning optimization within Pakistan legal framework"
        ),
        'benefits': (
            "Full FBR compliance — automated WHT and advance tax calculation\n"
            "Eliminated manual filer status verification (ATL auto-lookup)\n"
            "Effective tax rate optimization from ~49% towards legal minimum\n"
            "Consolidated tax reporting across group entities"
        ),
    },
    # Slide 37: Reconciliation, Validation & Adjustments maturity
    36: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Automated reconciliation across 4-5 core banking systems and GL\n"
            "Payment system settlement reconciliation (RAAST, 1Link, SWIFT, NIFT)\n"
            "Card network settlement matching (Visa/Mastercard)\n"
            "Nostro/vostro daily reconciliation for correspondent banking\n"
            "Inter-branch transaction reconciliation across 16,000+ branches"
        ),
        'benefits': (
            "Finance team time freed: reconciliation from 40-50% to <15% of effort\n"
            "Month-end close acceleration from 15-30 days to <10 days\n"
            "Exception-based processing — manual review only for breaks\n"
            "Audit-ready reconciliation trail for external auditors"
        ),
    },
    # Slide 39: Financial Consolidation maturity
    38: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Consolidation across 6-8 group entities with different ERPs/core systems\n"
            "Islamic window/subsidiary dual IFRS + AAOIFI consolidation\n"
            "Automated inter-company elimination entries\n"
            "Solo + consolidated SBP regulatory returns\n"
            "SECP-compliant consolidated IFRS statements for listed groups"
        ),
        'benefits': (
            "Consolidated close within SBP and SECP regulatory deadlines\n"
            "60% reduction in manual consolidation effort\n"
            "Drill-through from consolidated to entity-level detail\n"
            "December 31 year-end consolidation compliance (SBP mandate)"
        ),
    },
    # Slide 41: Functional P&L Statement maturity
    40: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Branch-level P&L for all 16,000+ branches (currently available at <15% of banks)\n"
            "Customer profitability using FSDM FACT_CUSTOMER_PROFITABILITY (35+ measures)\n"
            "FTP engine using KIBOR yield curve (not simplified flat spreads)\n"
            "Activity-based costing for branch, digital, ATM, agent channels\n"
            "Islamic P&L with correct line items (profit on financing, not interest income)"
        ),
        'benefits': (
            "Data-driven branch viability decisions (open/close/resize)\n"
            "Customer-level profitability driving retention and cross-sell\n"
            "Product pricing optimization based on true cost allocation\n"
            "SBP/IFRS 8 segment reporting from integrated profitability data"
        ),
    },
}


# ─────────────────────────────────────────────
# Speaker notes for all 46 existing slides + 3 new
# ─────────────────────────────────────────────

SPEAKER_NOTES = {
    # Part A: FPM Overview (slides 1-15, 0-indexed 0-14)
    0: """FINANCE & PERFORMANCE MANAGEMENT — Domain Overview

This domain covers 4 sub-domains with 35 BVF capabilities:
1. Accounting Operations & Close (11 capabilities) — covered in this file
2. Enterprise Performance Management (11 capabilities) — detailed in File 08
3. Treasury Management (7 capabilities)
4. Financial Reporting (6 capabilities)

PAKISTAN CONTEXT: Aligned to SBP regulatory requirements, IFRS 9/IFRS 17, Basel III capital adequacy, and AAOIFI Islamic accounting standards. Pakistan banks face dual accounting: conventional IFRS + AAOIFI for Islamic banking windows/subsidiaries.

FSDM: GL_ACCT (General Ledger Account), GL_BLNC (GL Balance), ACCTNG_EVNT (Accounting Event), FIN_INSTRMT (Financial Instrument), CST_CNTR (Cost Center), PRDCT (Product), PRTY (Party), ORGN_UNT (Organization Unit)""",

    1: """FPM — Accounting Operations Capabilities (11 listed)

The 11 Accounting Operations capabilities cover the end-to-end financial close process. In Pakistan banking context:
- Sub-ledger integration across 4-5 core banking systems is foundational
- IFRS 9 ECL has made close significantly more complex since 2018
- Islamic banking windows/subsidiaries require parallel IFRS + AAOIFI entries
- SBP statutory returns require granular data that most banks cannot produce efficiently

FSDM: SUB_LDGR (Sub-Ledger), GL_PSTNG (GL Posting), COA (Chart of Accounts), RCNCLTN (Reconciliation), CNSLDTN (Consolidation), PL_STMNT (P&L Statement)""",

    2: """FPM — EPM Capabilities (11 listed)

EPM capabilities (detailed in File 08) include: Revenue Analytics, Activity-Based Costing, Profitability Modeling, Funds Transfer Pricing, Pricing Analytics, Performance Management, Budget & Forecasting, and more.

Pakistan banking EPM context:
- Branch-level P&L not available at most banks (vs standard practice globally)
- FTP engines use simplified KIBOR spreads, not yield-curve based
- Budget cycles take 3-4 months (manual process)
- Customer profitability rarely calculated at individual level

FSDM: PRFTBLTY (Profitability), FTP (Funds Transfer Pricing), ABC (Activity-Based Costing), BDGT (Budget), FRCST (Forecast)""",

    3: """FPM — Treasury Capabilities (7 listed)

Treasury capabilities cover: Investment Analytics, Liquidity Management, ALM (Asset-Liability Management), Capital Management, FX Management, and Money Market operations.

Pakistan treasury context:
- Banks hold PKR 10T+ in government securities — investment portfolio management is critical
- SBP Basel III requirements: LCR (min 100%), NSFR, Capital Adequacy (min 11.5%)
- FX management significant: USD 30B+ annual remittances, trade finance positions
- KIBOR-based pricing dominates — yield curve construction is key for FTP and treasury

FSDM: INVSTMNT (Investment), SCRTY (Security), LQD (Liquidity), ALM (Asset-Liability Management), FX_PSITN (FX Position)""",

    4: """FPM — Financial Reporting Capabilities (6 listed)

Financial reporting capabilities cover: Regulatory Reporting, Management Reporting, External Reporting (SECP/PSX), Board Reporting, Statutory Returns.

Pakistan financial reporting context:
- SBP requires 30+ statutory returns at varying frequencies (weekly, monthly, quarterly, annual)
- SECP requires consolidated IFRS statements for listed banking groups
- PSX disclosure requirements for listed banks
- External auditors (Big 4 in Pakistan) require detailed drill-down capability

FSDM: RPT (Report), RGLTY_RPT (Regulatory Report), FIN_STMNT (Financial Statement), DSCL (Disclosure)""",

    5: """Finance BVF — Use Cases (1 of 3)

Use cases mapped to Finance & Performance Management capabilities. Pakistan priority use cases include:
- IFRS 9 ECL Calculation and Reporting
- Automated SBP Statutory Returns
- Branch-Level Profitability Analysis
- Customer Profitability Engine (UBL FSDM reference)
- Islamic Accounting Hub (dual IFRS + AAOIFI)

Note: Original triangle color legend removed. Use cases are prioritized based on Pakistan banking regulatory and business requirements.

FSDM: Use cases span all FPM entities across Accounting, Performance, Treasury, and Reporting domains.""",

    6: """Finance BVF — Use Cases (2 of 3)

Continued use case mapping. Key Pakistan-relevant use cases:
- Funds Transfer Pricing (KIBOR yield-curve based)
- Activity-Based Costing (branch, digital, ATM, agent channels)
- Tax Management Automation (WHT, super tax, advance tax)
- Financial Consolidation across Group Entities
- Reconciliation Automation (core systems, payment networks, nostro/vostro)

FSDM: FTP (Funds Transfer Pricing), ABC (Activity-Based Costing), TAX (Tax), CNSLDTN (Consolidation), RCNCLTN (Reconciliation)""",

    7: """Finance BVF — Use Cases (3 of 3)

Continued use case mapping. Additional Pakistan use cases:
- Reserve Analytics (IFRS 9 ECL with Pakistan macro scenarios)
- Fair Value Management (Government Securities, FX)
- Asset Valuation (16,000+ branch properties, investment portfolio)
- Procurement Analytics (FBR WHT compliance)
- Billing & Collections (NPL management, ECIB integration)

FSDM: ECL (Expected Credit Loss), FV_ADJSTMNT (Fair Value), ASST_VLTN (Asset Valuation), VNDR (Vendor), CLLCTN (Collection)""",

    8: """FPM — Why Important

GLOBAL: Finance function transformation is a top-3 priority for bank CFOs globally. The shift from backward-looking reporting to forward-looking analytics and real-time decisioning is accelerating.

PAKISTAN CONTEXT:
- Dual accounting challenge: conventional IFRS + AAOIFI for Islamic banking windows/subsidiaries
- SBP requires 30+ statutory returns at varying frequencies — data integration is critical
- IFRS 9 ECL provisioning requires granular loan-level data across all core banking systems
- Monthly close takes 15-30 days (vs 3-5 day global benchmark) — competitive disadvantage
- Branch-level P&L not available — cannot make data-driven decisions on 16,000+ branches

FSDM: All FPM entities — foundational for regulatory compliance, profitability analysis, and management decision-making.""",

    9: """FPM — Challenges

PAKISTAN-SPECIFIC CHALLENGES:
1. 4-5 core banking systems (CTL, Temenos, Oracle, Finacle, Symbols) with incompatible chart of accounts
2. Islamic banking windows require parallel GL entries under both IFRS and AAOIFI
3. FBR tax complexity: WHT on deposits (15%/30%), super tax (10%), advance tax, minimum tax
4. Monthly close takes 15-30 days due to manual reconciliation across fragmented systems
5. Branch-level P&L not feasible without integrated sub-ledger data
6. IFRS 9 ECL requires individual loan-level data — most banks have only summary GL
7. External auditors (Big 4) increasingly demand transaction-level drill-down capability

FSDM: Integration challenges span all domains — GL, Sub-Ledger, Product, Party, Financial Instrument, Cost Center.""",

    10: """FPM — Challenges (continued)

Additional Pakistan challenges:
- SBP statutory return deadlines are strict — penalties for late/inaccurate filing
- SECP/PSX disclosure requirements increasing for listed banking groups
- Basel III implementation requires integrated data for LCR, NSFR, CAR calculations
- Talent shortage: Pakistan lacks sufficient IFRS 9 / data analytics professionals
- Legacy system retirement: Some banks still run 20+ year old core systems

FSDM: RGLTY_RPT (Regulatory Report), BSL (Basel), LCR (Liquidity Coverage Ratio), CAR (Capital Adequacy Ratio)""",

    11: """FPM — Are You Able To...

PAKISTAN ASSESSMENT QUESTIONS:
- Can you produce branch-level P&L for all 16,000+ branches?
- Can you calculate customer profitability across conventional AND Islamic product lines?
- Can you automate SBP statutory returns from a single data source?
- Can you produce IFRS 9 ECL reports with full drill-down to individual loan data?
- Can you calculate FTP using KIBOR yield curve rather than flat spreads?
- Can you consolidate across all group entities (bank + Islamic + Takaful + AMC + brokerage)?
- Can you close monthly books in under 10 days?
- Can you automate reconciliation across 4-5 core banking systems?

These questions define the gap between current Pakistan banking finance capability and global best practice.

FSDM: Assessment questions map to all 35 FPM sub-capabilities across 4 sub-domains.""",

    12: """FPM — How (Capabilities)

PAKISTAN SOLUTION APPROACH:
FSDM-based data warehouse provides the integrated financial data foundation required for all downstream finance analytics. The approach:

1. DATA INTEGRATION: Unified sub-ledger from all core banking systems
2. CHART OF ACCOUNTS: Consistent COA mapping across all source systems
3. DUAL ACCOUNTING: IFRS + AAOIFI entries from same operational transactions
4. GRANULARITY: Transaction-level detail enabling any aggregation
5. TIMELINESS: Near-real-time data feeds reducing close cycle
6. AUDITABILITY: Full lineage from report to source transaction

UBL CONTEXT: Customer Profitability Engine demonstrates how FSDM star schema enables PKR-denominated profitability across all product lines.

FSDM: FSDM_DW (FSDM Data Warehouse), COA (Chart of Accounts), SUB_LDGR (Sub-Ledger), ACCTNG_HUB (Accounting Hub)""",

    13: """FPM — Value Proposition

VALUE FOR PAKISTAN BANKS:
1. COMPLIANCE: Timely, accurate SBP statutory returns — avoid penalties
2. EFFICIENCY: Close cycle from 15-30 days to <5 days
3. PROFITABILITY: Branch/customer/product P&L enabling data-driven decisions
4. RISK: IFRS 9 ECL automation reducing provisioning errors
5. TAX: Automated WHT and tax calculations — FBR compliance
6. CONSOLIDATION: Group-wide view across all entities

ESTIMATED VALUE: PKR 5B+ annual cost savings from optimized finance operations, reduced provisioning errors, and improved capital allocation across a top-5 Pakistan bank.

FSDM: Value delivered through integrated data across all FPM domains.""",

    14: """FPM — Value (continued)

Additional value dimensions:
- REGULATORY: Reduced risk of SBP penalties for delayed/inaccurate returns
- STRATEGIC: Data-driven decision making for branch network optimization
- CAPITAL: Optimized capital allocation through accurate risk-weighted asset calculation
- INVESTOR: Improved disclosure quality for PSX-listed banking groups
- ISLAMIC: Proper Shariah-compliant accounting builds trust with Islamic banking customers

ROI: Finance transformation investments typically deliver 5-8x ROI over 3 years through efficiency gains, reduced regulatory risk, and improved profitability analytics.

FSDM: All FPM entities contribute to comprehensive value realization.""",

    15: """SECTION DIVIDER: Finance and Performance Management Capabilities

This section transitions from the FPM domain overview to the detailed capability slides. The following 25 slides cover 10 Accounting Operations capabilities in detail, each with a Capability Detail slide and a Maturity Assessment slide.

Pakistan banks should prioritize these capabilities based on:
1. SBP regulatory requirements (IFRS 9, Basel III, statutory returns)
2. Quick wins (reconciliation automation, WHT automation)
3. Strategic value (branch P&L, customer profitability)
4. Islamic banking requirements (dual IFRS + AAOIFI accounting)""",

    16: """ACCOUNTING OPERATIONS AND CLOSE — Why Important

Pakistan banking context: Quality of financial close is under SBP scrutiny. Banks with delayed or restated financials face regulatory action. IFRS 9 implementation (since 2018) has made the close process significantly more complex — ECL recalculation required at each reporting date.

Key importance for Pakistan banks:
- SBP demands increasing granularity in regulatory returns
- External auditors require transaction-level drill-down
- IFRS 9 ECL adds 5-10 days to close cycle at most banks
- Islamic banking parallel books add additional reconciliation layer

FSDM: ACCTNG_OPS (Accounting Operations), CLS (Close), PRCSS (Process), WRKFLW (Workflow)""",

    17: """ACCOUNTING OPERATIONS AND CLOSE — Business Challenges

PAKISTAN-SPECIFIC CHALLENGES:
1. 4-5 core banking systems with different chart of accounts — manual mapping and reconciliation
2. Islamic banking window sub-ledger creates additional reconciliation layer
3. SBP requires increasing granularity in regulatory returns — summary-level GL data is insufficient
4. External auditors (Big 4 in Pakistan) require detailed drill-down capability
5. IFRS 9 ECL recalculation at each period-end adds complexity and time
6. Lack of automated reconciliation between payment networks (RAAST, 1Link, SWIFT) and GL
7. Month-end close takes 15-30 days — finance team overwhelmed with manual processes

FSDM: ACCTNG_OPS (Accounting Operations), CLS (Close), RCNCLTN (Reconciliation), SUB_LDGR (Sub-Ledger)""",

    18: """ACCOUNTING OPERATIONS AND CLOSE — Are You Able To...

Pakistan assessment questions for Accounting Operations:
- Can you reconcile all core banking system GL balances within 24 hours?
- Can you produce IFRS 9 ECL at individual loan level with full audit trail?
- Can you generate dual IFRS + AAOIFI entries from the same operational transactions?
- Can you automate nostro/vostro reconciliation with correspondent banks?
- Can you match RAAST, 1Link, and card network settlements to GL within same day?
- Can you close monthly books in under 10 days?
- Can you provide external auditors with drill-down from GL to source transaction?

FSDM: Assessment questions map to all 10 Accounting Operations sub-capabilities.""",

    19: """ACCOUNTING OPERATIONS AND CLOSE — How (Capabilities)

Solution approach for Pakistan banks:
1. UNIFIED SUB-LEDGER: Capture transaction-level data from all core banking systems
2. COA MAPPING: Consistent chart of accounts across all source systems
3. DUAL ACCOUNTING: Automated IFRS + AAOIFI entry generation for Islamic transactions
4. RECONCILIATION: Automated daily reconciliation with exception-based processing
5. IFRS 9 ENGINE: ECL calculation at individual loan level with stage monitoring
6. CLOSE ORCHESTRATION: Automated close workflow with dependency management
7. REPORTING: Automated SBP statutory returns from integrated data source

FSDM: SUB_LDGR, COA, RCNCLTN, ECL, CLS_WRKFLW (Close Workflow), RGLTY_RPT (Regulatory Report)""",

    20: """ACCOUNTING OPERATIONS AND CLOSE — Value

Value proposition for Pakistan banks:
- Close cycle: From 15-30 days to <5 days (global benchmark)
- Reconciliation: From 40-50% of finance time to <15%
- IFRS 9: From semi-manual ECL to automated loan-level calculation
- Audit: From weeks of manual drill-down requests to self-service
- Islamic: Dual accounting from parallel manual entry to automated generation
- Compliance: From late/restated returns to on-time, accurate SBP filing

Estimated PKR 1-3B annual value from improved finance operations efficiency across a top-5 Pakistan bank.

FSDM: Value delivered through integrated Accounting Operations capabilities.""",

    # Part C: 10 Capability pairs (slides 22-41, 0-indexed 21-40)
    # Detail slides
    21: """BILLING & COLLECTIONS — Capability Detail

PAKISTAN COLLECTIONS CONTEXT:
NPL landscape: PKR 900B+ in non-performing loans across industry (~7.5% ratio).
SBP classification: Current > OAEM (90 days) > Substandard (180 days) > Doubtful (1 year) > Loss (5 years)
Provisioning: General 1-2% + Specific (25%-100% depending on classification)
Recovery channels: Call center, field recovery, legal (banking courts), write-off + recovery agents
Challenge: Pakistan's legal system averages 3-7 years for loan recovery — incentivizes early analytical intervention.

FSDM: AR_BLNC (Accounts Receivable Balance), INVCE (Invoice), PMT (Payment), CLLCTN (Collection), DPD (Days Past Due), NPL (Non-Performing Loan), CLSFCTN (Classification), PRVSN (Provision)""",

    22: """BILLING & COLLECTIONS — Maturity Assessment

Current State (Developing): Basic collection tracking. SBP classification managed semi-manually. Limited ECIB integration for collection prioritization.

Target State (Innovating): Automated ECIB-driven collection strategy, real-time SBP classification tracking, segment-based collection optimization, predictive early-warning system.

Pakistan Context: With PKR 900B+ NPLs and 3-7 year legal recovery timelines, analytical early intervention is the highest-value capability.

FSDM: CLLCTN (Collection), DPD (Days Past Due), NPL (Non-Performing Loan), CLSFCTN (Classification), PRVSN (Provision), ECIB (Credit Bureau)""",

    23: """PROCUREMENT & PAYMENT — Capability Detail

PAKISTAN PROCUREMENT CONTEXT:
Pakistan banks are significant purchasers: IT infrastructure, branch construction/maintenance, security, stationery, marketing services, consulting, audit fees.

Tax complexity: Every vendor payment requires WHT deduction at source. Rates vary:
- Services: 8% (filer) / 16% (non-filer)
- Supplies: 4% (filer) / 8% (non-filer)
- Contracts: 7% (filer) / 14% (non-filer)
FBR filer status must be verified for each payment — active taxpayer list (ATL) changes monthly.

FSDM: VNDR (Vendor), PO (Purchase Order), INVCE (Invoice), PMT (Payment), WHT (Withholding Tax), AP_BLNC (Accounts Payable Balance)""",

    24: """PROCUREMENT & PAYMENT — Maturity Assessment

Current State (Developing): Manual FBR filer status verification. WHT calculated manually for many vendor categories. No centralized procurement across group entities.

Target State (Innovating): Automated ATL lookup, automated WHT calculation across all vendor categories, centralized procurement analytics across all group entities.

Pakistan Context: FBR ATL changes monthly — manual verification for thousands of vendors is error-prone and non-compliant.

FSDM: VNDR (Vendor), WHT (Withholding Tax), ATL (Active Taxpayer List), AP_BLNC (Accounts Payable Balance)""",

    25: """SUB-LEDGER ACCOUNTING — Capability Detail

PAKISTAN SUB-LEDGER CONTEXT:
Pakistan banks typically operate "thick" GLs because sub-ledger data is not centralized:
- Core banking system generates summary GL postings (not transaction-level)
- Cards system posts net settlements to GL
- Islamic banking window posts separately
- Trade finance posts separately
- Investment/treasury system posts separately

Result: GL has 50,000+ accounts but limited drill-down to transaction detail.

FSDM solution: Create sub-ledger accounting hub that captures individual transaction-level entries from each source, applies consistent COA mapping, generates dual IFRS/AAOIFI entries for Islamic transactions, and enables full drill-down.

FSDM: ACCTNG_EVNT (Accounting Event), SUB_LDGR (Sub-Ledger), GL_PSTNG (GL Posting), COA (Chart of Accounts), FIN_INSTRMT (Financial Instrument), ACCTNG_RULE (Accounting Rule)""",

    26: """SUB-LEDGER ACCOUNTING — Maturity Assessment

Current State (Developing): Summary GL postings from core systems. No centralized sub-ledger. Limited drill-down. Manual COA mapping between systems.

Target State (Innovating): Unified transaction-level sub-ledger across all core systems with consistent COA, dual IFRS/AAOIFI entries, and full audit trail.

Pakistan Context: This is the foundational capability — without unified sub-ledger, IFRS 9 ECL, branch P&L, customer profitability, and regulatory reporting are all compromised.

FSDM: SUB_LDGR (Sub-Ledger), GL_PSTNG (GL Posting), COA (Chart of Accounts), ACCTNG_EVNT (Accounting Event)""",

    27: """FAIR VALUE & HEDGE ACCOUNTING — Capability Detail

PAKISTAN FAIR VALUE CONTEXT:
Government securities portfolio: Pakistan banks hold PKR 10T+ in government securities (PIBs, T-Bills, Sukuk).
SBP classification impacts P&L: HTM (amortized cost) vs. AFS (OCI) vs. HFT (P&L). Reclassification under IFRS 9 created significant one-time impacts.

FX hedging: Banks hedge trade finance, remittance (USD 30B+/year), and treasury positions. SBP limits open FX positions.

Islamic securities: Sukuk (Islamic bonds) require fair valuation under both IFRS and AAOIFI — different measurement approaches.

FSDM: FV_ADJSTMNT (Fair Value Adjustment), HDG (Hedge), HDG_EFCTVNS (Hedge Effectiveness), SCRTY (Security), SCRTY_VLTN (Security Valuation), FX_PSITN (FX Position)""",

    28: """FAIR VALUE & HEDGE ACCOUNTING — Maturity Assessment

Current State (Developing): Government securities valued but process is semi-manual. FX hedging documented but hedge effectiveness testing limited. Islamic securities valuation lacks standardization.

Target State (Innovating): Automated daily mark-to-market, robust hedge effectiveness testing, integrated IFRS + AAOIFI fair value for Islamic securities.

Pakistan Context: PKR 10T+ government securities portfolio makes fair value management a top priority for bank P&L stability.

FSDM: FV_ADJSTMNT (Fair Value), HDG_EFCTVNS (Hedge Effectiveness), SCRTY_VLTN (Security Valuation), FX_PSITN (FX Position)""",

    29: """ASSET VALUATIONS — Capability Detail

PAKISTAN: Bank asset valuations are SBP-regulated. Revaluation surplus on branch properties creates regulatory capital (Tier 2) but must be discounted per SBP capital adequacy rules.

Key valuation areas:
1. Branch properties (16,000+ — significant book value, revaluation every 3-5 years per SBP)
2. Government securities portfolio (daily mark-to-market for AFS/HFT)
3. Lending portfolio (IFRS 9 ECL — collective and individual impairment)
4. Foreclosed collateral (Qarz-e-Hasna for Islamic, foreclosed properties)
5. Equity investments (banking subsidiaries, associates, strategic holdings)

FSDM: ASST (Asset), ASST_VLTN (Asset Valuation), IMPRM (Impairment), RVAL (Revaluation), CLLTRL (Collateral), CLLTRL_VLTN (Collateral Valuation)""",

    30: """ASSET VALUATIONS — Maturity Assessment

Current State (Developing): Branch property revaluation per SBP schedule but manual process. Investment impairment testing semi-automated. Collateral valuations often outdated.

Target State (Innovating): Automated revaluation tracking, real-time impairment testing, dynamic collateral coverage monitoring, complete fixed asset lifecycle management.

Pakistan Context: 16,000+ branch properties represent significant balance sheet value. SBP revaluation surplus rules directly impact Tier 2 capital calculation.

FSDM: ASST_VLTN (Asset Valuation), IMPRM (Impairment), RVAL (Revaluation), CLLTRL_VLTN (Collateral Valuation)""",

    31: """RESERVE ANALYTICS — Capability Detail

PAKISTAN RESERVE/PROVISIONING FRAMEWORK:
SBP requires BOTH general and specific provisioning:
- General: 1% on consumer, 1.5% on SME, 2% on small enterprise, plus IFRS 9 Stage 1/2 ECL
- Specific: 25% (Substandard), 50% (Doubtful), 100% (Loss) — reduced by collateral
- IFRS 9 ECL: PD x LGD x EAD for each loan, with forward-looking macroeconomic scenarios
- SBP allows "regulatory overlay" where SBP provisioning exceeds IFRS 9 ECL

Pakistan-specific challenge: Historical PD data is limited (many banks have <5 years of granular default data). SBP provides sector-level PD benchmarks as fallback.

UBL Context: FSDM star schema supports ECL calculation with FACT_LOAN_PROVISION containing PD, LGD, EAD, Stage, and ECL amount at individual loan level.

FSDM: RSRV (Reserve), ECL (Expected Credit Loss), PD (Probability of Default), LGD (Loss Given Default), EAD (Exposure at Default), PRVSN (Provision), STG (Stage), CLSFCTN (Classification)""",

    32: """RESERVE ANALYTICS — Maturity Assessment

Current State (Developing): IFRS 9 ECL implemented but many banks rely on simplified approaches. PD models lack granular historical data. Forward-looking scenarios limited.

Target State (Innovating): Fully automated ECL at individual loan level, multiple Pakistan macro scenarios, Islamic provisioning under AAOIFI FAS 30, regulatory overlay management.

Pakistan Context: SBP scrutinizes ECL adequacy — banks with insufficient provisions face regulatory action. FSDM-based ECL engine provides defensible, auditable results.

FSDM: ECL (Expected Credit Loss), PD, LGD, EAD, PRVSN (Provision), STG (Stage)""",

    33: """TAX MANAGEMENT & OPTIMISATION — Capability Detail

PAKISTAN BANKING TAX FRAMEWORK:
Banks face the highest effective tax rate in Pakistan:
- Corporate tax: 39% (banks — higher than standard 29%)
- Super tax: 10% on banks (introduced 2022)
- Effective rate: ~49% before other taxes
- WHT on deposit profit: 15% (ATL filers) / 30% (non-filers) — bank acts as collection agent
- Advance tax: Banks deduct advance tax on cash withdrawals >PKR 50K/day (non-filers)
- Minimum tax: Applies if normal tax is below threshold

Automation opportunity:
- Customer filer/non-filer status changes monthly — FBR ATL lookup must be automated
- WHT calculation on millions of deposit profit payments — currently semi-manual at many banks
- Tax provisioning for quarterly/annual filing requires consolidated data from all entities

FSDM: TAX (Tax), WHT (Withholding Tax), TAX_CLCTN (Tax Calculation), TAX_CRDT (Tax Credit), TAX_OBLGTN (Tax Obligation), FLR_STS (Filer Status)""",

    34: """TAX MANAGEMENT & OPTIMISATION — Maturity Assessment

Current State (Developing): WHT calculated but significant manual effort. Filer status verification semi-manual. Group tax planning limited.

Target State (Innovating): Fully automated WHT with ATL integration, consolidated group tax planning, tax optimization within legal framework, real-time tax liability monitoring.

Pakistan Context: 39% corporate + 10% super = ~49% effective rate makes tax optimization a critical capability for bank profitability.

FSDM: TAX (Tax), WHT (Withholding Tax), TAX_CLCTN (Tax Calculation), FLR_STS (Filer Status)""",

    35: """RECONCILIATION, VALIDATION & ADJUSTMENTS — Capability Detail

PAKISTAN RECONCILIATION CHALLENGE:
Pakistan banks spend enormous effort on reconciliation because:
1. Multiple core systems post to GL independently — totals must be reconciled daily
2. Payment system settlements (RAAST, 1Link, SWIFT) must match internal records
3. Card network settlements (Visa, Mastercard) reconcile against card system and GL
4. SBP clearing house (NIFT) cheque settlements must be reconciled
5. Nostro/vostro accounts with correspondent banks must be reconciled daily
6. Inter-branch transactions must be reconciled (Head Office vs. branches)

Estimated: Finance teams at Pakistan banks spend 40-50% of time on reconciliation activities that could be automated.

FSDM: RCNCLTN (Reconciliation), ADJSTMNT (Adjustment), VLDTN (Validation), GL_BLNC (GL Balance), STLMNT (Settlement), NSTRO (Nostro), VSTRO (Vostro)""",

    36: """RECONCILIATION, VALIDATION & ADJUSTMENTS — Maturity Assessment

Current State (Developing): Mostly manual reconciliation. GL-to-sub-ledger breaks common. Payment system reconciliation takes days. Nostro breaks require manual investigation.

Target State (Innovating): Automated daily reconciliation across all systems, exception-based processing, real-time break identification, automated adjustment suggestions.

Pakistan Context: 40-50% of finance team time spent on reconciliation is a massive efficiency opportunity. Automation can free capacity for analysis and advisory.

FSDM: RCNCLTN (Reconciliation), ADJSTMNT (Adjustment), VLDTN (Validation), STLMNT (Settlement)""",

    37: """FINANCIAL CONSOLIDATION — Capability Detail

PAKISTAN CONSOLIDATION CONTEXT:
Typical Pakistan banking group consolidation scope:
- Parent bank (conventional) — e.g., CTL core banking
- Islamic banking subsidiary or window — separate books (IFRS + AAOIFI)
- Insurance/Takaful subsidiary
- Asset management company (mutual funds)
- Brokerage subsidiary
- Modaraba company (if applicable)
- Leasing company (if applicable)
- Microfinance subsidiary (if applicable)

Each may have different ERP, different chart of accounts, different fiscal year end (though SBP requires December 31 for banks).

SBP requires BOTH solo and consolidated regulatory returns.
SECP requires consolidated IFRS financial statements for listed banking groups.

FSDM: CNSLDTN (Consolidation), LGL_ENTTY (Legal Entity), IC_TXN (Inter-Company Transaction), ELMNTN (Elimination), GRP_ENTTY (Group Entity), CNSLDTD_GL (Consolidated GL)""",

    38: """FINANCIAL CONSOLIDATION — Maturity Assessment

Current State (Developing): Manual consolidation using Excel. IC elimination done manually. Different COA across entities requires manual mapping. Solo returns only partially automated.

Target State (Innovating): Automated consolidation across all group entities, automated IC elimination, consistent COA mapping, solo + consolidated returns from unified data.

Pakistan Context: Most banking groups have 6-8 entities with different systems. Manual consolidation adds 5-10 days to close cycle.

FSDM: CNSLDTN (Consolidation), LGL_ENTTY (Legal Entity), IC_TXN (Inter-Company Transaction), ELMNTN (Elimination)""",

    39: """FUNCTIONAL PROFIT & LOSS STATEMENT — Capability Detail

This is the MOST CRITICAL capability for Pakistan banks and directly relates to UBL's Customer Profitability Engine:

P&L DIMENSIONS REQUIRED:
1. BRANCH P&L: Revenue, cost, and profit for each of 16,000+ branches
2. CUSTOMER P&L: Individual customer profitability across all products held
3. PRODUCT P&L: Profitability by product (CASA, Term Deposit, Personal Loan, Credit Card, etc.)
4. SEGMENT P&L: Business segments per IFRS 8 (Corporate, Commercial, Retail, Treasury, Islamic)
5. CHANNEL P&L: Cost-to-serve by channel (branch, ATM, digital, contact center)
6. REGION P&L: Geographic profitability (province, city, area)

P&L COMPONENTS: NII (FTP-adjusted), Non-Interest Income, Direct Costs, Allocated Costs (ABC), Provision Charge (IFRS 9 ECL), Tax.

ISLAMIC P&L: Different line items — "profit earned on financing" not "interest income". Mudaraba profit-sharing ratios replace interest rates.

UBL CONTEXT: FACT_CUSTOMER_PROFITABILITY contains 35+ measures including FTP-adjusted NII, fee income, direct costs, allocated costs, provision charge, and net profit — all at individual customer level.

FSDM: PL_STMNT (P&L Statement), RVNU (Revenue), EXPNS (Expense), NII (Net Interest Income), FEE_INCM (Fee Income), FTP (Funds Transfer Pricing), CST_ALLCTN (Cost Allocation), PRFT (Profit), BRCH (Branch), SGMNT (Segment)""",

    40: """FUNCTIONAL P&L STATEMENT — Maturity Assessment

Current State (Developing): Bank-level P&L only. No branch-level P&L. Limited segment reporting. FTP using flat KIBOR spreads. No customer profitability.

Target State (Innovating): 360-degree profitability: branch x customer x product x segment x channel x region. FTP using KIBOR yield curve. Activity-based costing. Islamic P&L with correct line items.

Pakistan Context: Branch-level P&L available at <15% of banks. Customer profitability calculation virtually non-existent. This is the single highest-value FPM capability for Pakistan banks.

FSDM: PL_STMNT, RVNU, EXPNS, NII, FTP, CST_ALLCTN, PRFT, BRCH, SGMNT""",

    # Part D: EPM Introduction (slides 42-46, 0-indexed 41-45)
    41: """ENTERPRISE PERFORMANCE MANAGEMENT — Why Important

EPM in Pakistan banking must cover both conventional financial metrics AND Islamic banking performance indicators.

Pakistan banking KPIs (industry benchmarks):
- Return on Equity (ROE): 20-30% (top banks), 10-15% (mid-tier)
- Return on Assets (ROA): 1.0-2.0%
- Net Interest Margin (NIM): 3-5%
- Cost-to-Income: 45-65% (range across banks)
- NPL Ratio: 5-10% (varies by bank)
- CASA Ratio: 40-55%
- Capital Adequacy Ratio (CAR): Min 11.5% (SBP), top banks 15-20%
- Advances-to-Deposits Ratio (ADR): 45-55%
- Liquidity Coverage Ratio (LCR): Min 100% (SBP Basel III)

SBP increasingly demands forward-looking analytics (stress testing, scenario analysis) — not just backward-looking reporting.

FSDM: KPI (Key Performance Indicator), PRFMNC (Performance), BDGT (Budget), FRCST (Forecast), TRGTS (Targets)""",

    42: """ENTERPRISE PERFORMANCE MANAGEMENT — Challenges

Pakistan EPM challenges:
- Budget and planning cycles take 3-4 months (manual, Excel-based)
- No rolling forecasts — annual budget is static, outdated within months
- KPI reporting is backward-looking (monthly/quarterly) not forward-looking
- SBP stress testing requirements growing — scenario-based analytics needed
- Islamic banking performance metrics differ from conventional (no NIM, different ratios)
- Profitability measurement at branch/customer/product level not available
- FTP engines use simplified KIBOR spreads — distorts profitability measurement

EPM capabilities detailed in File 08 will address these challenges.

FSDM: BDGT (Budget), FRCST (Forecast), KPI, PRFMNC (Performance), STRSS_TST (Stress Test)""",

    43: """ENTERPRISE PERFORMANCE MANAGEMENT — Are You Able To...

Pakistan EPM assessment questions:
- Can you produce rolling 12-month forecasts updated monthly?
- Can you run stress scenarios (interest rate, FX, GDP, default rate) on P&L and balance sheet?
- Can you measure profitability at branch, customer, product, and segment level?
- Can you track KPIs in real-time (not monthly reporting lag)?
- Can you benchmark performance across branches adjusting for size, location, and market?
- Can you calculate FTP using KIBOR yield curve with maturity matching?
- Can you allocate costs using activity-based costing methodology?

FSDM: BDGT, FRCST, KPI, FTP, ABC, STRSS_TST, PRFTBLTY (Profitability)""",

    44: """ENTERPRISE PERFORMANCE MANAGEMENT — How (Capabilities)

EPM solution approach for Pakistan banks:
1. PROFITABILITY ENGINE: Building on FSDM star schema for branch/customer/product P&L
2. FTP ENGINE: KIBOR yield-curve-based funds transfer pricing
3. ABC FRAMEWORK: Activity-based costing across branch, digital, ATM, agent channels
4. BUDGETING: Automated budget preparation and rolling forecast
5. KPI DASHBOARD: Real-time executive dashboard with drill-down
6. STRESS TESTING: Multi-scenario analysis for SBP regulatory requirements
7. ISLAMIC METRICS: Profit-sharing ratios, Shariah-compliant performance indicators

Detailed capability slides in File 08.

FSDM: PRFTBLTY, FTP, ABC, BDGT, FRCST, KPI, STRSS_TST""",

    45: """ENTERPRISE PERFORMANCE MANAGEMENT — Value

Value of EPM for Pakistan banks:
- BUDGET: From 3-4 month manual cycle to automated 2-week rolling forecast
- PROFITABILITY: From bank-level only to 360-degree branch/customer/product/segment/channel
- FTP: From flat KIBOR spread to yield-curve-based — accurate profitability measurement
- STRESS TESTING: SBP-compliant scenario analysis for capital planning
- DECISION MAKING: Data-driven decisions on pricing, branch network, customer segments

Pakistan banking EPM KPI targets:
- ROE > 20%, ROA > 1.5%, NIM 3-5%, C/I < 50%, NPL < 5%, CASA > 50%, CAR > 15%

Detailed EPM capabilities and value covered in File 08.

FSDM: All EPM entities — delivered through integrated data foundation.""",
}


# Speaker notes for the 3 new slides
NEW_SLIDE_NOTES = {
    'new_a': """FINANCE & PERFORMANCE MANAGEMENT — DOMAIN DASHBOARD

This slide provides a high-level comparison of Finance & Performance Management maturity across global, regional (South Asia & Middle East), and Pakistan banking markets.

Key Takeaways:
- Pakistan monthly close cycle is 15-30 days vs 3-5 day global benchmark
- Financial consolidation automation at <25% vs 80%+ at top global banks
- Activity-based costing adopted by <10% of Pakistan banks vs 60% globally
- IFRS 9 ECL automation at 20-30% (SBP mandate driving adoption)
- Branch-level P&L is rare (<15% of Pakistan banks) vs standard practice globally
- Islamic accounting (AAOIFI) adds unique dual-accounting requirement

This domain covers 35 BVF sub-capabilities with ~200+ FSDM entities across Financial Instrument, Accounting, Party, Product, GL, and Cost Center domains.

Investment Justification: Moving from Developing to Innovating maturity can deliver PKR 5B+ annual savings from improved finance operations, reduced provisioning errors, and optimized capital allocation.

FSDM Domains: Financial Instrument (~100+ entities), Accounting (~80+ entities), GL (~50+ entities), Cost Center (~30+ entities)""",

    'new_b': """PAKISTAN BANKING — FINANCE & PERFORMANCE MANAGEMENT LANDSCAPE

Regulatory Framework:
- SBP: 30+ statutory returns at weekly, monthly, quarterly, annual frequencies
- IFRS 9: ECL provisioning mandate (implemented 2018, ongoing refinement)
- IFRS 17: Insurance contract accounting (effective 2025 for Takaful/insurance)
- AAOIFI: Accounting standards for Islamic banking operations
- Basel III: Capital adequacy (min 11.5%), liquidity ratios (LCR min 100%, NSFR)
- FBR: Tax complexity (WHT 15%/30%, corporate 39%, super tax 10%)

Key Challenges:
- 4-5 core banking systems per bank — fragmented GL/sub-ledger data
- Islamic banking windows require parallel IFRS + AAOIFI accounting
- Monthly close takes 15-30 days vs 3-5 day global benchmark
- Branch-level P&L not available at most banks
- FTP engines use simplified KIBOR spreads, not yield-curve based
- ECIB integration required for IFRS 9 ECL calculations

Pakistan Banking Finance Metrics:
- Banking sector assets: PKR 35T+ (~$125B), Total deposits: PKR 25T+
- Total advances: PKR 13T+, Industry NPLs: ~7.5% (PKR 900B+)
- CASA ratio: ~47%, Number of banks: 33 + 5 Islamic + 11 microfinance
- SBP policy rate: 17.5%, WHT on profit: 15% filers / 30% non-filers""",

    'new_c': """FINANCE & PERFORMANCE MANAGEMENT — IMPLEMENTATION ROADMAP

Phase 1: Accounting Foundation (0-6 months) — Investment: PKR 50-100M
- Unified chart of accounts across all core banking systems
- Sub-ledger integration for top 3 product systems (CASA, lending, cards)
- IFRS 9 ECL data pipeline automation
- Automated GL-to-sub-ledger reconciliation
Quick Win: Close cycle from 30 to 15 days

Phase 2: Performance Analytics (6-18 months) — Investment: PKR 120-250M
- Funds Transfer Pricing engine (KIBOR yield-curve based)
- Activity-based costing framework (branch, digital, ATM, agent channels)
- Branch-level P&L for all 16,000+ branches
- Customer profitability engine (building on UBL FSDM foundation)
Expected: Identify PKR 5B+ in cost optimization opportunities

Phase 3: Enterprise Finance Intelligence (18-36 months) — Investment: PKR 200-400M
- Real-time financial consolidation across all entities
- Islamic accounting hub (dual IFRS + AAOIFI)
- Predictive budgeting and rolling forecasts
- Executive finance dashboard with drill-down to transaction level
Expected: Close in <5 days, 360-degree profitability view

Total 3-Year Investment: PKR 370-750M (~USD 1.3-2.7M)
Estimated Annual Value: PKR 5B+ from optimized finance operations""",
}


# ─────────────────────────────────────────────
# Clean orphaned slides
# ─────────────────────────────────────────────

def clean_pptx(input_path, output_path):
    """Remove orphaned slide XMLs from the PPTX ZIP."""
    import zipfile
    import xml.etree.ElementTree as ET

    with zipfile.ZipFile(input_path, 'r') as zin:
        pres_xml = zin.read('ppt/presentation.xml')
        root = ET.fromstring(pres_xml)
        pns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
               'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

        sld_rids = set()
        for sldId in root.findall('.//p:sldIdLst/p:sldId', pns):
            rid = sldId.attrib.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', '')
            sld_rids.add(rid)

        rels_xml = zin.read('ppt/_rels/presentation.xml.rels')
        rels_root = ET.fromstring(rels_xml)
        referenced_slides = set()
        for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
            rid = rel.attrib.get('Id', '')
            target = rel.attrib.get('Target', '')
            if rid in sld_rids and 'slides/' in target:
                slide_file = target.replace('slides/', '')
                referenced_slides.add(slide_file)

        print(f"  Referenced slides: {len(referenced_slides)}")

        orphaned = set()
        for name in zin.namelist():
            if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                slide_file = name.replace('ppt/slides/', '')
                if slide_file not in referenced_slides:
                    orphaned.add(slide_file.replace('.xml', ''))

        print(f"  Orphaned slides to remove: {len(orphaned)}")

        with zipfile.ZipFile(output_path, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                skip = False
                for orphan in orphaned:
                    if f'slides/{orphan}.xml' in item.filename or f'slides/_rels/{orphan}.xml.rels' in item.filename:
                        skip = True
                        break
                if not skip:
                    data = zin.read(item.filename)
                    zout.writestr(item, data)

    print(f"  Cleaned file saved to {output_path}")


# ─────────────────────────────────────────────
# Main execution
# ─────────────────────────────────────────────

def main():
    print("── Step 0: Cleaning orphaned slides ──")
    clean_pptx(INPUT_FILE, CLEAN_FILE)

    print(f"\nLoading presentation...")
    prs = Presentation(CLEAN_FILE)
    slides = list(prs.slides)
    print(f"  Loaded {len(slides)} slides")

    # ── Step 1: Add speaker notes to ALL 46 slides ──
    print(f"\n── Step 1: Adding speaker notes to all slides ──")
    for idx in range(len(slides)):
        if idx in SPEAKER_NOTES:
            set_notes(slides[idx], SPEAKER_NOTES[idx])
            print(f"  Slide {idx+1}: Speaker notes set")
        else:
            print(f"  Slide {idx+1}: No notes defined (skipped)")

    # ── Step 2: Enrich 10 capability detail slides ──
    print(f"\n── Step 2: Enriching 10 capability detail slides ──")
    for idx, enrichment in CAPABILITY_ENRICHMENTS.items():
        slide = slides[idx]
        table = get_table(slide)
        if not table:
            print(f"  Slide {idx+1}: WARNING - no table found!")
            continue

        # Append to business objectives [1,0]
        biz_lines = '\n'.join(enrichment['biz_obj'])
        append_cell_text(table.cell(1, 0), '\n' + biz_lines)

        # Append to data & solution [3,0]
        data_lines = '\n'.join(enrichment['data_sol'])
        append_cell_text(table.cell(3, 0), '\n' + data_lines)

        # Append to outcome [3,2]
        outcome_lines = '\n'.join(enrichment['outcome'])
        append_cell_text(table.cell(3, 2), '\n' + outcome_lines)

        print(f"  Slide {idx+1}: Detail table enriched")

    # ── Step 3: Enrich 10 maturity slides ──
    print(f"\n── Step 3: Enriching 10 maturity slides ──")
    for idx, enrichment in MATURITY_ENRICHMENTS.items():
        slide = slides[idx]
        table = get_table(slide)
        if not table:
            print(f"  Slide {idx+1}: WARNING - no table found!")
            continue

        # Replace key areas [7,3] and benefits [7,4]
        replace_cell_text(table.cell(7, 3), enrichment['key_areas'])
        replace_cell_text(table.cell(7, 4), enrichment['benefits'])

        # Replace H1 2018
        replaced = replace_h1_2018(table)

        # Remove engagement text from groups
        removed = remove_engagement_text(slide)

        status = []
        if replaced:
            status.append("replaced H1 2018")
        if removed:
            status.append("removed engagement text")
        print(f"  Slide {idx+1}: Maturity enriched + {' + '.join(status)}")

    # ── Step 4: Replace Teradata in notes (slides 6-8) ──
    print(f"\n── Step 4: Replacing Teradata in use case notes ──")
    for idx in [5, 6, 7]:
        slide = slides[idx]
        if slide.has_notes_slide:
            # Notes are already replaced by SPEAKER_NOTES in step 1
            print(f"  Slide {idx+1}: Notes already replaced with Pakistan context")

    # ── Step 5: Replace Teradata text in all slides ──
    print(f"\n── Step 5: Scanning for Teradata text in slide content ──")
    teradata_count = 0
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'Teradata Business Value Framework' in run.text:
                            run.text = run.text.replace('Teradata Business Value Framework', 'Banking Business Value Framework')
                            teradata_count += 1
                        elif 'Teradata' in run.text:
                            run.text = run.text.replace('Teradata', 'Enterprise Analytics Platform')
                            teradata_count += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if 'Teradata Business Value Framework' in run.text:
                                    run.text = run.text.replace('Teradata Business Value Framework', 'Banking Business Value Framework')
                                    teradata_count += 1
                                elif 'Teradata' in run.text:
                                    run.text = run.text.replace('Teradata', 'Enterprise Analytics Platform')
                                    teradata_count += 1
    print(f"  Replaced {teradata_count} Teradata references")

    # ── Step 6: Add 3 new slides ──
    print(f"\n── Step 6: Adding 3 new slides ──")
    # Find a blank layout
    blank_layout = None
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower() or layout.name == 'Blank':
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[-1]  # Use last layout as fallback

    # New Slide A: Dashboard
    slide_a = prs.slides.add_slide(blank_layout)
    _build_dashboard_slide(slide_a, prs)
    set_notes(slide_a, NEW_SLIDE_NOTES['new_a'])
    print(f"  Added: Domain Dashboard (New Slide A)")

    # New Slide B: Pakistan Context
    slide_b = prs.slides.add_slide(blank_layout)
    _build_pakistan_context_slide(slide_b, prs)
    set_notes(slide_b, NEW_SLIDE_NOTES['new_b'])
    print(f"  Added: Pakistan Finance Landscape (New Slide B)")

    # New Slide C: Implementation Roadmap
    slide_c = prs.slides.add_slide(blank_layout)
    _build_roadmap_slide(slide_c, prs)
    set_notes(slide_c, NEW_SLIDE_NOTES['new_c'])
    print(f"  Added: Implementation Roadmap (New Slide C)")

    # ── Step 7: Reorder slides ──
    print(f"\n── Step 7: Reordering slides ──")
    _reorder_slides(prs, len(prs.slides))
    print(f"  Slides reordered: Dashboard, PK Context, original 46, Roadmap (before last 2)")

    # ── Step 8: Save ──
    print(f"\n── Step 8: Saving to {OUTPUT_FILE} ──")
    prs.save(OUTPUT_FILE)

    # Clean up temp file
    if os.path.exists(CLEAN_FILE):
        os.remove(CLEAN_FILE)
        print("  Cleaned up temp file")

    print("  DONE!")


def _build_dashboard_slide(slide, prs):
    """Build the Domain Dashboard slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                shape.text = "Finance & Performance Management — At a Glance"
            elif shape.placeholder_format.idx == 1:
                shape.text = ""

    left = Emu(400000)
    top = Emu(900000)
    width = Emu(8300000)
    height = Emu(3800000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Finance & Performance Management — At a Glance"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    p = tf.add_paragraph()
    p.text = ""

    lines = [
        "                                  Global          South Asia & ME    Pakistan",
        "───────────────────────────────────────────────────────────────────────────────",
        "Monthly close cycle               3-5 days        7-15 days          15-30 days",
        "Financial consolidation auto      80%+ top banks  40-50%             <25%",
        "Activity-based costing            60% large banks 20-30%             <10%",
        "IFRS 9 ECL automation             70%+ EU/UK      30-40%             20-30% (SBP mandate)",
        "FTP engine deployed               80% large banks 40-50%             ~30% (KIBOR-based)",
        "Branch-level P&L                  Standard        Emerging           Rare (<15%)",
        "Islamic accounting (AAOIFI)       N/A             Required Islamic   Dual IFRS + AAOIFI",
        "───────────────────────────────────────────────────────────────────────────────",
    ]
    for line in lines:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(8)
        run.font.name = "Consolas"

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "BVF Sub-capabilities: 35 (11 Acctg Ops + 11 EPM + 7 Treasury + 6 Reporting) | FSDM Entities: ~200+ | BACR: ~150 | Maturity: Developing > Innovating"
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 102, 153)


def _build_pakistan_context_slide(slide, prs):
    """Build the Pakistan Finance Landscape slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                shape.text = "Pakistan Banking — Finance & Performance Management Landscape"
            elif shape.placeholder_format.idx == 1:
                shape.text = ""

    left = Emu(400000)
    top = Emu(850000)
    width = Emu(8300000)
    height = Emu(4000000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Regulatory framework
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Regulatory Framework"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    regs = [
        "SBP: 30+ statutory returns (weekly, monthly, quarterly, annual)",
        "IFRS 9: ECL provisioning mandate (2018+), IFRS 17: Takaful/insurance (2025)",
        "AAOIFI: Islamic banking accounting standards, Basel III: CAR min 11.5%, LCR/NSFR",
        "FBR: Corporate tax 39% + super tax 10%, WHT on deposits 15%/30%",
    ]
    for r in regs:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {r}"
        run.font.size = Pt(9)

    p = tf.add_paragraph()
    p.text = ""

    # Key challenges
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Key Challenges"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    challenges = [
        "4-5 core banking systems per bank — fragmented GL/sub-ledger data",
        "Islamic banking windows require parallel IFRS + AAOIFI accounting",
        "Monthly close 15-30 days (vs 3-5 day global benchmark)",
        "Branch-level P&L not available at most banks (16,000+ branches)",
        "FTP engines use simplified KIBOR spreads, not yield-curve based",
    ]
    for c in challenges:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {c}"
        run.font.size = Pt(9)

    p = tf.add_paragraph()
    p.text = ""

    # Banking metrics
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Pakistan Banking Finance Metrics"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    metrics = [
        "Assets: PKR 35T+ (~$125B) | Deposits: PKR 25T+ | Advances: PKR 13T+",
        "NPLs: ~7.5% (PKR 900B+) | CASA: ~47% | Banks: 33 scheduled + 5 Islamic + 11 MFB",
        "SBP rate: 17.5% | WHT: 15% filers / 30% non-filers | Effective tax: ~49%",
    ]
    for m in metrics:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {m}"
        run.font.size = Pt(9)


def _build_roadmap_slide(slide, prs):
    """Build the Implementation Roadmap slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                shape.text = "Finance & Performance Management — Implementation Roadmap"
            elif shape.placeholder_format.idx == 1:
                shape.text = ""

    left = Emu(400000)
    top = Emu(850000)
    width = Emu(8300000)
    height = Emu(3900000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Title
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Implementation Roadmap"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    p = tf.add_paragraph()
    p.text = ""

    # Phase 1
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Phase 1: Accounting Foundation (0-6 months) — PKR 50-100M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 102, 51)

    phase1 = [
        "Unified chart of accounts across all core banking systems",
        "Sub-ledger integration for top 3 product systems (CASA, lending, cards)",
        "IFRS 9 ECL data pipeline automation",
        "Automated GL-to-sub-ledger reconciliation",
        "Quick Win: Close cycle from 30 to 15 days",
    ]
    for line in phase1:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Phase 2
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Phase 2: Performance Analytics (6-18 months) — PKR 120-250M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 153)

    phase2 = [
        "Funds Transfer Pricing engine (KIBOR yield-curve based)",
        "Activity-based costing framework (branch, digital, ATM, agent)",
        "Branch-level P&L for all 16,000+ branches",
        "Customer profitability engine (building on UBL FSDM foundation)",
        "Expected: Identify PKR 5B+ in cost optimization",
    ]
    for line in phase2:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Phase 3
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Phase 3: Enterprise Finance Intelligence (18-36 months) — PKR 200-400M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(153, 0, 51)

    phase3 = [
        "Real-time financial consolidation across all entities",
        "Islamic accounting hub (dual IFRS + AAOIFI)",
        "Predictive budgeting and rolling forecasts",
        "Executive finance dashboard with drill-down to transaction level",
        "Expected: Close in <5 days, 360-degree profitability view",
    ]
    for line in phase3:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Total
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Total 3-Year Investment: PKR 370-750M (~USD 1.3-2.7M) | Annual Value: PKR 5B+"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)


def _reorder_slides(prs, total_slides):
    """Reorder slides: New A (46), New B (47), Original 0-43, New C (48), Original 44-45.

    Final order (0-indexed source positions):
    [46, 47, 0-43, 48, 44, 45]
    Dashboard, PK Context, all original 46 slides, Roadmap, then EPM Value last 2 slides
    """
    # Actually per the prompt: Roadmap before EPM (last 2 are EPM Value slides 45-46)
    # But the prompt says "insert as second-to-last" which means before the LAST existing slide pair
    # Let's match the final structure from the prompt:
    # 1=Dashboard, 2=PK Context, 3-48=original 46 slides, 49=Roadmap
    # Wait, looking at prompt more carefully:
    # | **49** | **Implementation Roadmap** | **NEW** |
    # So Roadmap is just appended at the end.
    # The reorder is: Dashboard(46), PK Context(47), Original 0-45, Roadmap(48)

    new_order = [46, 47]  # Dashboard, PK Context
    new_order.extend(range(0, 46))  # All 46 original slides
    new_order.append(48)  # Roadmap at end

    pres_ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    pres_elem = prs.part._element
    sldIdLst = pres_elem.find('.//p:sldIdLst', pres_ns)

    sldId_elements = list(sldIdLst)

    for elem in sldId_elements:
        sldIdLst.remove(elem)

    for idx in new_order:
        sldIdLst.append(sldId_elements[idx])

    print(f"  Reordered {len(new_order)} slides")


if __name__ == '__main__':
    main()
