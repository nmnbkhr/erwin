#!/usr/bin/env python3
"""
Update 16_Regulatory_Compliance.pptx — Fill ALL placeholders with Pakistan
banking regulatory context and enhance existing content with SBP/FATF references.
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy

INPUT_FILE = '/mnt/e/erwin/pptout/16_Regulatory_Compliance.pptx'
OUTPUT_FILE = '/mnt/e/erwin/pptout/16_Regulatory_Compliance_UPDATED.pptx'

# ─────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────

def set_notes(slide, notes_text):
    """Set speaker notes for a slide."""
    if not slide.has_notes_slide:
        slide.notes_slide
    notes_tf = slide.notes_slide.notes_text_frame
    while len(notes_tf.paragraphs) > 1:
        p = notes_tf.paragraphs[-1]._p
        p.getparent().remove(p)
    lines = notes_text.strip().split('\n')
    notes_tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        p = notes_tf.add_paragraph()
        p.text = line


def find_shape_by_text(slide, text_fragment):
    """Find a shape containing specific text."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and text_fragment in shape.text_frame.text:
            return shape
    return None


def replace_cell_text(cell, new_text):
    """Replace text in a table cell while preserving basic formatting."""
    # Get format from first run if available
    first_para = cell.text_frame.paragraphs[0]
    font_size = None
    font_bold = None
    font_color = None
    if first_para.runs:
        r = first_para.runs[0]
        font_size = r.font.size
        font_bold = r.font.bold
        if r.font.color and r.font.color.type is not None:
            try:
                font_color = r.font.color.rgb
            except:
                pass

    # Clear all paragraphs except first
    paras = list(cell.text_frame.paragraphs)
    for i in range(len(paras) - 1, 0, -1):
        p = paras[i]._p
        p.getparent().remove(p)

    # Set content
    lines = new_text.strip().split('\n')
    cell.text_frame.paragraphs[0].clear()
    run = cell.text_frame.paragraphs[0].add_run()
    run.text = lines[0]
    if font_size:
        run.font.size = font_size
    if font_bold is not None:
        run.font.bold = font_bold
    if font_color:
        run.font.color.rgb = font_color

    for line in lines[1:]:
        p = cell.text_frame.add_paragraph()
        run = p.add_run()
        run.text = line
        if font_size:
            run.font.size = font_size


def replace_shape_text(shape, new_text, keep_first_para=False):
    """Replace all text in a shape's text_frame."""
    tf = shape.text_frame
    paras = list(tf.paragraphs)

    start_idx = 1 if keep_first_para else 0

    # Remove paragraphs from end
    for i in range(len(paras) - 1, start_idx, -1):
        p = paras[i]._p
        p.getparent().remove(p)

    lines = new_text.strip().split('\n')

    if keep_first_para:
        # Add new content as additional paragraphs
        for line in lines:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(9)
    else:
        # Replace first paragraph
        paras[0].clear()
        run = paras[0].add_run()
        run.text = lines[0]
        run.font.size = Pt(9)
        for line in lines[1:]:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(9)


def fill_use_case_slide(slide, content):
    """Fill a use case slide with content dict mapping header -> list of bullets."""
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        tf = shape.text_frame
        header = tf.paragraphs[0].text.strip() if tf.paragraphs else ''

        if header in content:
            bullets = content[header]
            # Keep header paragraph, replace "Point 1" with real content
            paras = list(tf.paragraphs)
            for i in range(len(paras) - 1, 0, -1):
                p = paras[i]._p
                p.getparent().remove(p)

            for bullet in bullets:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = bullet
                run.font.size = Pt(8)

        # Handle References/Owner field
        if header == 'References' and 'owner' in content:
            paras = list(tf.paragraphs)
            for i in range(len(paras) - 1, 0, -1):
                paras[i]._p.getparent().remove(paras[i]._p)
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"Owner: {content['owner']}"
            run.font.size = Pt(8)


def enhance_use_case_slide(slide, additions):
    """Add Pakistan context to existing use case slides by appending to relevant shapes."""
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        tf = shape.text_frame
        header = tf.paragraphs[0].text.strip() if tf.paragraphs else ''

        if header in additions:
            for text in additions[header]:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = text
                run.font.size = Pt(8)
                run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)


# ─────────────────────────────────────────────
# SLIDE 7 — "Are You Able To…"
# ─────────────────────────────────────────────
def update_slide_7(slide):
    shape = find_shape_by_text(slide, 'Question 1')
    if shape:
        tf = shape.text_frame
        paras = list(tf.paragraphs)
        for i in range(len(paras) - 1, 0, -1):
            paras[i]._p.getparent().remove(paras[i]._p)

        questions = [
            "",
            "Can you produce all 200+ monthly, quarterly, and annual SBP regulatory returns — from prudential classification (CL-1/CL-2), capital adequacy (CAR), foreign exchange position, IFRS 9 ECL staging, to AML/CFT STR reporting — from a single integrated data source with full audit trail and reconciliation to General Ledger within T+5 business days?",
            "",
            "Can you demonstrate to FATF/APG evaluators that your AML/CFT transaction monitoring covers 100% of customer transactions across all channels (branch, ATM, IBFT, RAAST, mobile wallets, trade finance, remittances) with risk-based alert thresholds, sanctions screening, and automated STR generation — all compliant with SBP AML/CFT Guidelines 2020?",
            "",
            "Can you perform Basel III stress testing under SBP-specified macroeconomic scenarios (PKR devaluation, KIBOR spike, GDP contraction, sectoral stress) computing capital adequacy impact at individual facility level with full data lineage from source transaction to capital ratio — and complete this within the 2-week ICAAP submission timeline?",
        ]
        for q in questions:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = q
            run.font.size = Pt(9)
    print("  Slide 7: Filled Are You Able To questions")


# ─────────────────────────────────────────────
# SLIDE 11 — Challenges
# ─────────────────────────────────────────────
def update_slide_11(slide):
    shape = find_shape_by_text(slide, 'Description')
    if shape:
        tf = shape.text_frame
        paras = list(tf.paragraphs)
        for i in range(len(paras) - 1, 0, -1):
            paras[i]._p.getparent().remove(paras[i]._p)

        content = [
            "Building an integrated regulatory compliance data infrastructure that serves multiple regulators (SBP, SECP, FBR, FATF, PTA) from a single source of truth — eliminating the 40-60 separate regulatory data extraction processes that most Pakistani banks maintain, each with different data definitions, reconciliation issues, and manual interventions.",
        ]
        for line in content:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(8)
    print("  Slide 11: Filled challenges")


# ─────────────────────────────────────────────
# SLIDE 16 — Financial Reporting Regulation (table)
# ─────────────────────────────────────────────
def update_slide_16(slide):
    for shape in slide.shapes:
        if type(shape).__name__ == 'GraphicFrame':
            try:
                tbl = shape.table
                # Row [1,0] = Business Objectives content
                replace_cell_text(tbl.cell(1, 0),
                    "Automate financial regulatory reporting covering SBP prudential returns, SECP annual filing requirements, FBR tax compliance, and international standards (IFRS 9, IFRS 16, IFRS 17).\nBuild a single regulatory data warehouse serving all financial reporting obligations with consistent data definitions, automated calculations, and full audit trail.\nReduce reporting cycle from T+15 to T+5 business days while improving accuracy.")

                # Row [3,0] = Data & Solution
                replace_cell_text(tbl.cell(3, 0),
                    "Data:\nGeneral Ledger transaction data\nSub-ledger details (lending, deposits, treasury, trade finance)\nCustomer classification and risk rating data\nIFRS 9 ECL staging and provision data\nCapital adequacy components (RWA by asset class, capital instruments)\nForeign exchange position data\nSector and geographic classification codes (SBP standard)\nOrganizational hierarchy (business units, branches, profit centers)\nFSDM Finance domain entities (300+ entities)\n\nAnalytics:\nAutomated regulatory return generation with validation rules\nReconciliation analytics (GL ↔ sub-ledger ↔ regulatory return)\nData quality monitoring with automated exception detection\nTrend analysis on regulatory metrics (CAR trajectory, NPL trend, provision coverage)\nRegulatory deadline monitoring and workflow management")

                # Row [3,2] = Outcome
                replace_cell_text(tbl.cell(3, 2),
                    "Measurable performance metrics:\nAutomated generation of top 50 SBP regulatory returns from integrated data warehouse\nFull reconciliation trail from source transaction through GL to regulatory line item\nRegulatory calendar with workflow management and deadline tracking\nData quality dashboard showing completeness, accuracy, and timeliness metrics\nRegulatory trend analysis for proactive management of key ratios")
                break
            except:
                pass
    print("  Slide 16: Filled Financial Reporting Regulation")


# ─────────────────────────────────────────────
# SLIDE 17 — Financial Reporting maturity
# ─────────────────────────────────────────────
def update_slide_17(slide):
    for shape in slide.shapes:
        if type(shape).__name__ == 'GraphicFrame':
            try:
                tbl = shape.table
                replace_cell_text(tbl.cell(1, 3),
                    "Fully automated regulatory reporting from integrated data warehouse. Real-time regulatory position monitoring. Automated reconciliation with exception-only manual intervention. Proactive regulatory risk indicators alerting management to approaching threshold breaches. Regulatory sandbox for impact analysis of new regulations before effective date.")
                replace_cell_text(tbl.cell(2, 3),
                    "Most regulatory returns automated with T+5 delivery. Automated reconciliation covering major returns. Regulatory data quality metrics tracked and managed. Regulatory change management process identifies new requirements and maps to data/system changes. Quarterly regulatory position trend analysis for Board.")
                replace_cell_text(tbl.cell(3, 3),
                    "Major regulatory returns semi-automated (automated data extraction, manual formatting and validation). Reconciliation performed manually but systematically. Regulatory deadline tracking in place. Data quality issues identified during preparation and corrected. T+10-15 delivery cycle.")
                replace_cell_text(tbl.cell(4, 3),
                    "Regulatory returns prepared from multiple data extracts with significant manual intervention. Limited reconciliation — breaks identified during audit rather than proactively. Deadline pressure causes quality trade-offs. Multiple spreadsheets and manual processes. T+15-20 delivery.")
                replace_cell_text(tbl.cell(5, 3),
                    "Regulatory reporting entirely manual. Multiple teams prepare overlapping returns from different data sources. No systematic reconciliation. Frequent SBP findings on data quality and timeliness. Regulatory compliance is reactive — preparing returns on deadline rather than maintaining continuous regulatory position visibility.")
                break
            except:
                pass
    print("  Slide 17: Filled Financial Reporting maturity")


# ─────────────────────────────────────────────
# SLIDE 18 — Disclosure (table)
# ─────────────────────────────────────────────
def update_slide_18(slide):
    for shape in slide.shapes:
        if type(shape).__name__ == 'GraphicFrame':
            try:
                tbl = shape.table
                replace_cell_text(tbl.cell(1, 0),
                    "Automate regulatory and public disclosure requirements — including Pillar 3 (Basel III market discipline disclosure), annual report disclosures (IFRS 7 financial instrument disclosures, IFRS 9 credit risk disclosures), SBP quarterly disclosure requirements, and PSX listing compliance.\nBuild consistent disclosure data pipeline ensuring public disclosures reconcile precisely to regulatory returns and financial statements.")

                replace_cell_text(tbl.cell(3, 0),
                    "Data:\nCapital adequacy components for Pillar 3 templates\nRisk exposure data by asset class, geography, sector, maturity\nIFRS 9 stage migration data and ECL movement analysis\nCredit quality data (risk grade distribution, impairment, write-offs)\nMarket risk position data (VaR, stressed VaR)\nRemuneration data for senior management\nCounterparty credit risk data\n\nAnalytics:\nAutomated Pillar 3 template population from regulatory data warehouse\nYear-on-year comparison and movement analysis\nNarrative generation support for qualitative disclosures\nCross-validation between Pillar 3, financial statements, and regulatory returns\nPeer benchmarking on disclosure quality and completeness")

                replace_cell_text(tbl.cell(3, 2),
                    "Measurable performance metrics:\nAutomated Pillar 3 disclosure templates (quantitative sections) produced quarterly\nIFRS 7/9 disclosure tables generated from same data source as financial statements\nDisclosure calendar with workflow management\nCross-validation report showing consistency between disclosure, regulatory returns, and financial statements\nDisclosure quality benchmarking against peer banks")
                break
            except:
                pass
    print("  Slide 18: Filled Disclosure")


# ─────────────────────────────────────────────
# SLIDE 19 — Disclosure maturity
# ─────────────────────────────────────────────
def update_slide_19(slide):
    for shape in slide.shapes:
        if type(shape).__name__ == 'GraphicFrame':
            try:
                tbl = shape.table
                replace_cell_text(tbl.cell(1, 3),
                    "Fully automated quantitative disclosures from integrated data warehouse. AI-assisted qualitative narrative generation. Real-time disclosure position monitoring. Automated cross-validation across all public and regulatory disclosures. Disclosure quality exceeds peer benchmarks and satisfies analyst/investor expectations.")
                replace_cell_text(tbl.cell(2, 3),
                    "Most quantitative disclosures automated. Cross-validation between disclosures and financial statements systematic. Disclosure templates updated within 30 days of new regulatory requirements. Qualitative disclosures drafted by subject matter experts with data support from analytics platform.")
                replace_cell_text(tbl.cell(3, 3),
                    "Key Pillar 3 quantitative tables semi-automated. IFRS disclosure tables prepared with some manual compilation. Basic cross-validation performed. Disclosure updates for regulatory changes handled within reporting cycle. Annual report disclosures meet minimum requirements.")
                replace_cell_text(tbl.cell(4, 3),
                    "Disclosure preparation largely manual using spreadsheet extracts. Limited cross-validation — inconsistencies identified by auditors. Pillar 3 templates populated manually from regulatory returns. Disclosure updates reactive — scrambling when new requirements emerge.")
                replace_cell_text(tbl.cell(5, 3),
                    "Disclosure treated as annual/quarterly compliance exercise with no ongoing data pipeline. Manual compilation from multiple sources. No cross-validation. Minimal Pillar 3 disclosure — meeting letter rather than spirit of requirements. SBP/auditor findings on disclosure gaps.")
                break
            except:
                pass
    print("  Slide 19: Filled Disclosure maturity")


# ─────────────────────────────────────────────
# SLIDE 20 — Consumer Protection intro
# ─────────────────────────────────────────────
def update_slide_20(slide):
    # Fix the "Why is it important?" shape
    shape_why = find_shape_by_text(slide, 'Point 1')
    if shape_why:
        tf = shape_why.text_frame
        paras = list(tf.paragraphs)
        for i in range(len(paras) - 1, 0, -1):
            paras[i]._p.getparent().remove(paras[i]._p)

        points = [
            "",
            "SBP's Fair Treatment of Consumers framework mandates that banks treat customers fairly throughout the product lifecycle — from marketing and sales through servicing and complaint resolution. Compliance requires systematic monitoring of customer outcomes across all products and channels.",
            "",
            "Pakistan's banking sector serves 60M+ account holders including vulnerable populations (low-income, elderly, digitally illiterate) who require additional protections. Financial inclusion goals must be balanced with responsible lending and transparent pricing to prevent predatory outcomes.",
            "",
            "Regulatory penalties for consumer protection violations are increasing — SBP has imposed significant fines for misleading advertising, unfair contract terms, and poor complaint handling. Reputational damage from publicized consumer protection failures impacts customer acquisition and retention in an increasingly competitive market.",
        ]
        for pt in points:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = pt
            run.font.size = Pt(9)

    # Fix the definition shape
    shape_def = find_shape_by_text(slide, 'Revenue Integrity is defined as')
    if shape_def:
        tf = shape_def.text_frame
        paras = list(tf.paragraphs)
        for i in range(len(paras) - 1, -1, -1):
            paras[i]._p.getparent().remove(paras[i]._p) if i > 0 else None
        paras[0].clear()
        run = paras[0].add_run()
        run.text = "Consumer protection in banking encompasses the regulatory framework and analytical capabilities ensuring fair treatment of customers — covering transparent product disclosure, fair pricing, responsible lending, data privacy, complaint resolution, and protection against discriminatory or predatory practices."
        run.font.size = Pt(9)

    print("  Slide 20: Filled Consumer Protection intro")


# ─────────────────────────────────────────────
# SLIDE 21 — Consumer Protection challenges
# ─────────────────────────────────────────────
def update_slide_21(slide):
    shape = find_shape_by_text(slide, 'Description')
    if shape:
        tf = shape.text_frame
        paras = list(tf.paragraphs)
        for i in range(len(paras) - 1, 0, -1):
            paras[i]._p.getparent().remove(paras[i]._p)

        content = [
            "Building systematic, data-driven consumer protection monitoring that demonstrates fair customer outcomes across all products, channels, and customer segments — moving beyond reactive complaint handling to proactive identification and remediation of customer harm.",
        ]
        for line in content:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(8)
    print("  Slide 21: Filled Consumer Protection challenges")


# ─────────────────────────────────────────────
# SLIDES 22-29 — Consumer Protection capability tables and maturity
# ─────────────────────────────────────────────

CAPABILITY_TABLE_CONTENT = {
    22: {  # Truth-in-Advertising
        'objectives': "Monitor marketing and advertising content across all channels (print, digital, social media, branch collateral, SMS/email campaigns) to ensure compliance with SBP Fair Treatment guidelines, SECP advertising regulations, and PBA Code of Conduct.\nDetect misleading claims, hidden charges, unrealistic return promises, and non-compliant product representations before they reach customers or regulators.",
        'data_solution': "Data:\nMarketing campaign content and creative materials\nProduct terms and conditions (actual vs. advertised)\nCustomer complaint data tagged to marketing campaigns\nSocial media monitoring data (customer feedback on advertising claims)\nPricing and fee data (verifying advertised rates match actual charges)\nCompetitor advertising data for benchmarking\n\nAnalytics:\nNLP-based analysis of marketing content for compliance red flags\nAutomated comparison of advertised terms vs. actual product terms\nCustomer outcome monitoring — do customers receive what was advertised\nComplaint root cause analysis linking complaints to specific marketing campaigns\nSocial media sentiment analysis on product advertising",
        'outcome': "Measurable performance metrics:\nPre-publication compliance checking tool for marketing content\nPost-campaign monitoring dashboard tracking customer outcomes vs. advertising claims\nComplaint trend analysis identifying advertising-related customer dissatisfaction\nRegulatory compliance report for SBP advertising guidelines\nMarketing effectiveness analysis integrated with compliance monitoring",
    },
    24: {  # Predatory Practices
        'objectives': "Detect and prevent predatory lending and pricing practices — including excessive interest rates targeting vulnerable borrowers, aggressive upselling of insurance/credit protection products, hidden fees, unauthorized automatic renewals, and lending to customers beyond their repayment capacity.\nEnsure compliance with SBP's responsible lending guidelines and consumer protection framework.",
        'data_solution': "Data:\nLoan pricing data (APR, fees, charges) by customer segment\nCustomer income and debt-service capacity data\nProduct bundling and cross-sell data (insurance attachment to lending)\nCustomer complaint data related to unfair charges\nAccount-level fee revenue analysis\nPayment difficulty indicators (minimum payments, partial payments, restructuring requests)\nDTI (Debt-to-Income) ratios at origination and current\n\nAnalytics:\nCustomer outcome monitoring — identifying segments with disproportionate delinquency, complaint rates, or fee burden\nPricing fairness analysis comparing effective APR across customer demographics\nInsurance attachment rate analysis by channel and sales person (detecting pressure selling)\nDTI monitoring for evidence of over-indebtedness\nFee burden analysis as percentage of customer income/balance",
        'outcome': "Measurable performance metrics:\nPredatory practice monitoring dashboard covering pricing fairness, responsible lending metrics, and product suitability indicators\nAutomated alerts when customer segments show adverse outcome patterns\nSBP compliance report on responsible lending practices\nProduct suitability assessment framework ensuring products match customer needs and capacity",
    },
    26: {  # Privacy Protections
        'objectives': "Build a comprehensive data privacy compliance framework meeting SBP customer data protection guidelines, emerging Pakistan Personal Data Protection legislation, and international standards (GDPR principles where applicable for international customers).\nMonitor data access, consent management, data sharing with third parties, and customer data rights (access, correction, deletion) across all systems.",
        'data_solution': "Data:\nCustomer consent records (what data, what purpose, what duration)\nData access logs across all systems (who accessed what customer data, when, why)\nThird-party data sharing records (what data shared, with whom, under what agreement)\nData breach incident records\nCustomer data rights requests (access, correction, deletion)\nData retention compliance records\n\nAnalytics:\nData access pattern monitoring for unauthorized or unusual access\nConsent coverage analysis (customers without valid consent for data processing)\nData sharing compliance monitoring (ensuring third-party agreements cover shared data)\nPrivacy impact assessment analytics for new products and systems\nCustomer data inventory mapping (what data, where stored, how protected, how long retained)",
        'outcome': "Measurable performance metrics:\nPrivacy compliance dashboard showing consent coverage, data access compliance, and third-party sharing compliance\nAutomated alerts for unusual data access patterns\nCustomer data rights request tracking and response monitoring\nPrivacy impact assessment tool for new product and system launches\nData retention monitoring ensuring compliance with SBP data retention requirements",
    },
    28: {  # Do-Not-Solicit
        'objectives': "Implement comprehensive Do-Not-Solicit (DNS) and communication preference management ensuring customers are not contacted through unwanted channels or for unwanted products.\nComply with SBP Fair Treatment guidelines requiring respect for customer contact preferences, PTA regulations on unsolicited commercial communications, and emerging digital privacy standards.",
        'data_solution': "Data:\nCustomer communication preference records (opt-in/opt-out by channel and product)\nCampaign contact lists and execution records\nCustomer complaint data related to unwanted solicitation\nPTA Do-Not-Call registry data\nRegulatory correspondence on communication violations\nThird-party marketing partner contact records\n\nAnalytics:\nReal-time preference checking against all outbound communication triggers\nCompliance monitoring — detecting contacts made to opted-out customers\nCampaign execution audit comparing contact lists to preference records\nComplaint root cause analysis for communication-related complaints\nOpt-out rate trend analysis by channel and product",
        'outcome': "Measurable performance metrics:\nCentralized communication preference management system integrated with all platforms\nReal-time preference checking preventing non-compliant contacts\nCommunication compliance audit trail for SBP and PTA regulatory inquiries\nCustomer preference analytics informing communication strategy optimization\nZero non-compliant solicitation contacts (target)",
    },
}

MATURITY_CONTENT = {
    23: {  # Truth-in-Advertising maturity
        'leading': "AI-powered pre-publication compliance screening of all marketing content. Real-time social media monitoring for advertising compliance issues. Automated customer outcome verification comparing received vs. advertised terms. Proactive compliance embedded in marketing content creation workflow.",
        'innovating': "Pre-publication compliance review for major campaigns. Post-campaign monitoring of customer outcomes systematic. Social media monitoring for brand and compliance issues operational. Complaint-to-campaign linkage automated.",
        'practicing': "Compliance review of marketing content manual but systematic (review panel for major campaigns). Basic post-campaign monitoring through complaint analysis. Social media monitored periodically. Advertising guidelines documented and communicated to marketing teams.",
        'developing': "Marketing compliance review ad-hoc and inconsistent. Post-campaign monitoring reactive (only when complaints received). No social media monitoring for compliance. Advertising guidelines exist but not systematically enforced.",
        'emerging': "No formal marketing compliance monitoring. Advertising content not reviewed for regulatory compliance. Customer complaints about misleading advertising handled case-by-case without systemic analysis. SBP findings on advertising compliance are frequent.",
    },
    25: {  # Predatory Practices maturity
        'leading': "Comprehensive real-time monitoring of customer outcomes across all products and segments. AI-driven product suitability assessment at origination. Automated detection of emerging predatory patterns before harm materializes. Customer outcome metrics integrated into product development and pricing governance.",
        'innovating': "Systematic monitoring of customer outcomes by segment. Product suitability frameworks operational for major product lines. Pricing fairness analytics conducted quarterly. Insurance attachment monitoring with branch-level reporting. Responsible lending metrics tracked and reported to Board.",
        'practicing': "Annual review of customer outcome metrics. Basic monitoring of complaint trends by product. DTI limits enforced at origination but not monitored ongoing. Insurance attachment guidelines documented. Fee revenue analysis by customer segment conducted periodically.",
        'developing': "Responsible lending guidelines documented but not systematically monitored. Complaint-driven analysis of potential predatory practices. Limited customer outcome analytics. Insurance selling practices monitored through mystery shopping rather than data analytics.",
        'emerging': "No systematic monitoring of predatory practices. Compliance reliance on policy documents rather than outcome monitoring. Customer harm identified only through complaints or regulatory examination. No pricing fairness analytics. No product suitability assessment framework.",
    },
    27: {  # Privacy Protections maturity
        'leading': "Comprehensive data privacy platform with real-time access monitoring, automated consent management, and privacy-by-design embedded in all systems. Customer-facing data portal for self-service privacy management. AI-powered privacy risk detection. Full compliance with Pakistan Personal Data Protection Act.",
        'innovating': "Systematic consent management with automated tracking. Data access monitoring operational across major systems. Third-party data sharing governed by privacy assessments. Privacy impact assessments integrated into product development. Customer data rights fulfilled within regulatory timeframes.",
        'practicing': "Consent collection process documented and implemented. Basic data access logging in place. Third-party agreements include data protection clauses. Privacy impact assessments performed for major new products. Customer data requests handled manually but within timelines.",
        'developing': "Consent collection inconsistent across channels and products. Data access logging limited to critical systems. Third-party data sharing not comprehensively tracked. Privacy impact assessments ad-hoc. Customer data rights requests handled case-by-case without established process.",
        'emerging': "No formal data privacy compliance framework. Consent management absent or minimal. Data access not monitored. Third-party data sharing ungoverned. No privacy impact assessment process. Customer data rights not systematically addressed. Significant regulatory risk exposure.",
    },
    29: {  # Do-Not-Solicit maturity
        'leading': "Real-time communication preference enforcement across all channels with zero non-compliant contacts. Customer self-service preference management portal. AI-optimized contact strategies respecting preferences while maximizing engagement. Full compliance with PTA and SBP communication regulations.",
        'innovating': "Centralized preference management operational across major channels. Real-time preference checking for outbound campaigns. Compliance monitoring automated with exception reporting. Customer preference analytics informing channel strategy.",
        'practicing': "Communication preferences recorded and maintained. Manual preference checking for major campaigns. Periodic compliance audits. Customer complaints about unwanted contacts tracked and resolved. PTA Do-Not-Call registry compliance verified before campaigns.",
        'developing': "Basic opt-out tracking in place but not comprehensively enforced. Preference management varies by channel. Campaign compliance checking inconsistent. Complaint-driven identification of preference violations.",
        'emerging': "No formal communication preference management. Opt-out requests handled ad-hoc. No centralized preference database. Frequent customer complaints about unwanted solicitation. PTA compliance at risk.",
    },
}


# ─────────────────────────────────────────────
# USE CASE SLIDES — Placeholder fills
# ─────────────────────────────────────────────

USE_CASE_CONTENT = {
    33: {  # Basel II
        'owner': 'Regulatory Analytics CoE',
        'Objective / Problem Statement': [
            'Implement Basel II/III Standardized Approach for credit risk capital calculations as mandated by SBP for all scheduled banks in Pakistan',
            'Transition from Basel II foundation to Basel III advanced approaches as SBP roadmap evolves — supporting both Standardized and Foundation IRB where approved',
            'Integrate operational risk capital calculations (Basic Indicator Approach / Standardized Approach) with credit risk RWA for comprehensive capital adequacy reporting',
            'Build data infrastructure supporting future migration to Advanced IRB as Pakistani banking sector matures and SBP approves internal models',
        ],
        'Business Benefit': [
            'Accurate and timely capital adequacy reporting to SBP — meeting minimum CAR of 11.5% (including 1.5% D-SIB buffer for systemically important banks)',
            'Optimize RWA calculations through proper collateral recognition, credit risk mitigation, and exposure classification — potentially freeing PKR 10-20B in capital',
            'Automated Basel III capital calculations reducing manual effort by 60-70% and enabling scenario analysis for capital planning',
            'Foundation for ICAAP (Internal Capital Adequacy Assessment Process) as required by SBP Pillar 2 framework',
        ],
        'Expected Outcome': [
            'Automated daily/monthly RWA calculation by asset class (corporate, retail, SME, sovereign, bank, residential mortgage)',
            'Capital adequacy ratio monitoring dashboard with early warning when CAR approaches minimum thresholds',
            'SBP capital adequacy return (quarterly) auto-generated from integrated data warehouse',
            'Stress testing capability for capital impact of portfolio concentration, sectoral downturn, and macroeconomic scenarios',
        ],
        'POV Success Criteria': [
            'RWA calculations reconcile to SBP quarterly capital adequacy return within 0.1% tolerance',
            'Capital adequacy monitoring available within T+1 of month-end for management reporting',
            'Stress testing scenarios completed within 5 business days of parameter specification',
            'All exposure classifications consistent with SBP Prudential Regulation definitions',
        ],
        'Methodology / Analytic Technique': [
            'Basel III Standardized Approach risk weight assignment by SBP-defined asset classes',
            'Credit risk mitigation (CRM) calculations — collateral haircuts, guarantees, credit derivatives',
            'Operational risk capital (Basic Indicator Approach) — 15% of average gross income',
            'Market risk Standardized Approach — interest rate, equity, FX, commodity position capital charges',
            'Stress testing — scenario-based capital impact modeling using SBP-specified macroeconomic parameters',
        ],
        'Challenges': [
            'Data quality in exposure classification — ensuring consistent SBP asset class mapping across core banking systems',
            'Collateral valuation data often outdated or maintained in separate systems requiring integration for CRM calculations',
            'Multiple core banking systems provide different exposure views requiring reconciliation before capital calculations',
            'SBP regulatory changes require frequent recalibration of capital calculation parameters',
        ],
        'Source Data': [
            'Credit exposure data from core banking systems (4-5 systems) — facility limits, outstanding, accrued interest',
            'Collateral data from collateral management system — type, value, haircut, perfection status',
            'Customer risk ratings (internal and ECIB-sourced)',
            'SBP asset classification data (CL-1 through CL-5) from credit classification system',
            'GL data for operational risk capital calculations (gross income components)',
            'FSDM entities: AGREEMENT, COLLATERAL, PARTY, RISK_RATING, FINANCIAL_STATEMENT_LINE_ITEM',
        ],
    },
    35: {  # IFRS 9
        'owner': 'Finance & Risk Analytics CoE',
        'Objective / Problem Statement': [
            'Implement comprehensive IFRS 9 Expected Credit Loss (ECL) calculation engine compliant with SBP IFRS 9 Implementation Guidelines for all Pakistani scheduled banks',
            'Build integrated ECL model infrastructure covering Stage 1 (12-month ECL), Stage 2 (Lifetime ECL with SICR), and Stage 3 (Credit-Impaired) across all lending portfolios',
            'Automate Significant Increase in Credit Risk (SICR) assessment using both quantitative (PD-based) and qualitative (days past due, restructuring, watchlist) triggers',
            'Incorporate forward-looking macroeconomic scenarios (PKR exchange rate, KIBOR, GDP growth, inflation) with probability-weighted ECL under SBP-specified stress scenarios',
        ],
        'Business Benefit': [
            'Accurate IFRS 9 ECL provisioning compliant with SBP requirements — avoiding regulatory findings and audit qualifications on provision adequacy',
            'Reduced provision volatility through properly calibrated PD/LGD/EAD models with forward-looking overlays',
            'Proactive credit risk management — SICR monitoring identifies deteriorating exposures 30-60 days before migration to Stage 2/3',
            'Board and management reporting on ECL trends, stage migration, and provision coverage — supporting strategic credit decisions worth PKR billions',
        ],
        'Expected Outcome': [
            'Monthly IFRS 9 ECL calculation engine producing facility-level provisions across all portfolios',
            'Stage classification dashboard showing portfolio distribution (Stage 1/2/3), migration trends, and coverage ratios',
            'Forward-looking ECL under 3+ macroeconomic scenarios with probability weights as specified by SBP',
            'IFRS 9 disclosure tables (IFRS 7 credit risk disclosures) auto-generated from ECL engine output',
            'Reconciliation report between IFRS 9 provisions and SBP prudential classification (CL-1 through CL-5)',
        ],
        'POV Success Criteria': [
            'ECL calculations reconcile to audited financial statement provisions within 1% tolerance',
            'SICR identification captures 90%+ of facilities that migrate to Stage 3 within subsequent 6 months',
            'Monthly ECL production completed within T+5 business days',
            'Stage migration analysis available for Board Risk Committee within T+7',
            'Zero SBP inspection findings on IFRS 9 compliance methodology',
        ],
        'Methodology / Analytic Technique': [
            'PD models: logistic regression, survival analysis, and transition matrices calibrated to Pakistan default experience',
            'LGD models: workout LGD using historical recovery data, collateral haircuts per SBP guidelines',
            'EAD models: credit conversion factors for off-balance sheet exposures by product type',
            'SICR assessment: quantitative PD-based triggers + qualitative overlays (DPD >30, restructuring, watchlist)',
            'Macroeconomic scenario modeling: regression of PD/LGD on macro variables (GDP, KIBOR, PKR/USD, inflation, unemployment)',
            'ECL calculation: PD × LGD × EAD × discount factor, summed across time buckets for lifetime ECL',
        ],
        'Challenges': [
            'Limited historical default data for some portfolios (trade finance, agricultural lending) makes model calibration challenging',
            'Macroeconomic scenario specification for Pakistan requires accounting for high-volatility environment (PKR movements, inflation swings)',
            'Reconciliation between IFRS 9 ECL stages and SBP CL-1 to CL-5 classification creates complexity — different triggers and timing',
            'Forward-looking PD models must be regularly updated as Pakistan macroeconomic conditions evolve rapidly',
            'Data quality across 4-5 core banking systems — inconsistent facility dates, missing collateral values, incomplete restructuring history',
        ],
        'Source Data': [
            'Credit exposure data from core banking systems — facility details, disbursement, outstanding, accrued',
            'Historical default and recovery data (minimum 5-7 years, longer for through-the-cycle calibration)',
            'Customer risk ratings and rating migration history from credit risk management system',
            'Collateral valuation data — type, current market value, forced sale value, last valuation date',
            'SBP classification data (CL-1 to CL-5) and days-past-due from collections system',
            'Macroeconomic data from SBP, PBS: GDP, inflation, KIBOR, PKR/USD, unemployment, sector indices',
            'FSDM entities: AGREEMENT, RISK_RATING, COLLATERAL, FINANCIAL_STATEMENT_LINE_ITEM, LOSS_EVENT',
        ],
    },
    41: {  # Regulatory Disclosure Efficiency
        'owner': 'Regulatory Reporting CoE',
        'Objective / Problem Statement': [
            'Automate regulatory disclosure production — Pillar 3 quantitative templates, SBP quarterly disclosures, annual report regulatory notes, and PSX listing compliance disclosures',
            'Build end-to-end disclosure data pipeline from source systems through regulatory data warehouse to formatted disclosure output — eliminating manual spreadsheet-based compilation',
            'Ensure cross-validation between disclosures, regulatory returns, and audited financial statements — preventing inconsistencies that erode stakeholder confidence',
        ],
        'Business Benefit': [
            'Reduce disclosure preparation time by 60-70% through automation — freeing compliance teams for analysis rather than data compilation',
            'Eliminate disclosure inconsistencies — same data source feeding regulatory returns and public disclosures ensures consistency',
            'Improve disclosure quality and completeness — meeting SBP expectations and international peer benchmarks for Pillar 3 disclosure',
            'Faster regulatory response — ad-hoc SBP queries answered from integrated data warehouse within hours rather than days',
        ],
        'Expected Outcome': [
            'Automated Pillar 3 quantitative disclosure templates populated quarterly from regulatory data warehouse',
            'IFRS 7/9 credit risk disclosure tables generated alongside ECL calculations',
            'Cross-validation engine verifying consistency across all public disclosures and regulatory returns',
            'Disclosure calendar with automated workflow management and deadline tracking',
            'Peer benchmarking report comparing disclosure quality with top-5 Pakistani banks',
        ],
        'POV Success Criteria': [
            'Pillar 3 quantitative disclosures produced within T+15 of quarter-end (vs. T+30+ current)',
            'Zero inconsistencies between Pillar 3, annual report, and SBP regulatory returns in external audit',
            'Disclosure completeness score >95% against SBP Pillar 3 disclosure checklist',
            'Annual report regulatory notes cross-referenced to underlying regulatory data warehouse queries',
        ],
        'Methodology / Analytic Technique': [
            'Template mapping engine connecting regulatory data warehouse fields to Pillar 3 template cells',
            'Cross-validation rules engine checking consistency across disclosure documents',
            'Year-on-year movement analysis generating variance explanations automatically',
            'Peer disclosure benchmarking using publicly available annual reports of comparable banks',
            'Regulatory change monitoring scanning SBP circulars for disclosure requirement updates',
        ],
        'Challenges': [
            'Pillar 3 templates require data from multiple regulatory domains (credit risk, market risk, operational risk, capital)',
            'Qualitative disclosure narratives cannot be fully automated — require subject matter expert input',
            'SBP disclosure requirements evolve frequently requiring ongoing template maintenance',
            'Data granularity differences between regulatory returns and disclosure templates create mapping complexity',
        ],
        'Source Data': [
            'Regulatory data warehouse (capital adequacy, risk exposures, IFRS 9 provisions)',
            'General Ledger and financial reporting data for financial statement disclosures',
            'Risk management data (VaR, credit risk models, operational risk incidents)',
            'Corporate governance data (Board composition, committee attendance, remuneration)',
            'FSDM entities: FINANCIAL_STATEMENT_LINE_ITEM, AGREEMENT, RISK_RATING, PARTY',
        ],
    },
    43: {  # Predatory Practices for Credit
        'owner': 'Consumer Protection Analytics',
        'Objective / Problem Statement': [
            'Build analytics capability to detect and prevent predatory lending practices in Pakistan\'s banking sector — targeting excessive pricing, pressure selling of insurance/credit protection, lending beyond repayment capacity, and unfair fee structures',
            'Comply with SBP consumer banking prudential regulations mandating responsible lending, transparent pricing, and fair treatment of borrowers',
            'Implement customer outcome monitoring for consumer lending — ensuring products are suitable for customer circumstances and do not lead to systematic customer harm',
        ],
        'Business Benefit': [
            'Proactive compliance with SBP Fair Treatment of Consumers framework — avoiding regulatory penalties and reputational damage',
            'Reduced non-performing loans through responsible lending (DTI monitoring prevents over-indebtedness)',
            'Improved customer trust and retention — customers receiving suitable products have 40% lower complaint rates',
            'Early detection of predatory patterns prevents systemic customer harm before it triggers regulatory action',
        ],
        'Expected Outcome': [
            'Customer outcome monitoring dashboard tracking delinquency, complaints, and fee burden by customer segment and product',
            'Pricing fairness analysis comparing effective APR (including all fees/charges) across demographics and income segments',
            'Insurance attachment rate monitoring identifying branches or officers with unusually high attachment rates (pressure selling)',
            'Automated DTI monitoring flagging customers whose debt service exceeds SBP-defined thresholds',
        ],
        'POV Success Criteria': [
            'Identify 100% of consumer lending with effective APR exceeding product-specific thresholds',
            'Insurance attachment analysis covers all lending products with branch-level reporting',
            'DTI monitoring operational for all consumer lending originations and portfolio reviews',
            'Customer outcome trends tracked monthly with quarterly Board reporting',
        ],
        'Methodology / Analytic Technique': [
            'Statistical analysis of customer outcomes by product, segment, and origination channel',
            'Regression analysis identifying factors associated with adverse customer outcomes',
            'Outlier detection for pricing, fees, and insurance attachment rates',
            'Time series analysis of complaint trends by product and issue category',
            'Cohort analysis comparing customer outcomes across vintage and product design changes',
        ],
        'Challenges': [
            'Effective APR calculation requires aggregating scheduled and actual fees across multiple systems',
            'Insurance attachment data often maintained by insurance partners rather than bank systems',
            'Defining thresholds for predatory practice indicators requires balancing commercial objectives with consumer protection',
            'Cultural sensitivity in Pakistan — relationship-based banking makes distinguishing advice from pressure selling difficult',
        ],
        'Source Data': [
            'Consumer lending origination data — pricing, fees, terms, customer demographics, income',
            'Insurance attachment records — product, premium, commission, cancellation rates',
            'Collection system data — payment behavior, delinquency, restructuring',
            'Complaint management system — complaint type, product, resolution, customer satisfaction',
            'FSDM entities: AGREEMENT, PARTY, FINANCIAL_TRANSACTION, COMPLAINT_EVENT',
        ],
    },
    44: {  # Privacy Protection (reframe from HIPAA to Pakistan)
        'owner': 'Data Privacy & Compliance',
        'Objective / Problem Statement': [
            'Build data privacy compliance analytics for Pakistan banking — monitoring customer data access, consent management, and third-party data sharing across all systems',
            'Comply with SBP customer data protection guidelines, emerging Pakistan Personal Data Protection Bill, and SBP IT Security Guidelines for data handling',
            'Replace US-centric HIPAA framework with Pakistan-specific privacy requirements — SBP consumer data protection, NADRA data sharing agreements, ECIB reporting compliance, and mobile wallet data privacy',
        ],
        'Business Benefit': [
            'Proactive compliance with emerging Pakistan data protection legislation before enforcement — first-mover advantage in building privacy infrastructure',
            'Reduced data breach risk through systematic monitoring of data access patterns across banking systems',
            'Customer trust through transparent data handling — critical as digital banking adoption grows 40% YoY and customers demand privacy',
            'Regulatory readiness for SBP data protection audits and any future FATF-related data handling requirements',
        ],
        'Expected Outcome': [
            'Data access monitoring dashboard covering all banking systems (core banking, CRM, digital channels, data warehouse)',
            'Customer consent management system tracking consent for data processing, marketing, and third-party sharing',
            'Privacy incident tracking and response management system with automated escalation',
            'Data inventory mapping — what customer data exists, where, with what protection, for how long',
            'Regulatory compliance report for SBP data protection requirements',
        ],
        'POV Success Criteria': [
            'Data access monitoring covers 100% of systems containing customer PII within 12 months',
            'Consent management system operational for all new customer relationships within 6 months',
            'Privacy incident response time under 24 hours for identification, 72 hours for containment',
            'Customer data rights requests (access/correction/deletion) fulfilled within SBP-specified timeframes',
        ],
        'Methodology / Analytic Technique': [
            'Data access pattern analysis using anomaly detection on system access logs',
            'Consent coverage analytics tracking consent rates by channel, product, and data category',
            'Privacy impact assessment scoring model for new products and system changes',
            'Data lineage mapping using metadata analysis across integrated systems',
            'Automated classification of customer data by sensitivity level and retention requirements',
        ],
        'Challenges': [
            'Pakistan\'s data protection legislation still evolving — requirements may change before full implementation',
            'Banking systems (especially legacy core banking) have limited data access logging capabilities',
            'Customer consent collection across multiple channels (branch, digital, call center) requires unified management',
            'Balancing privacy requirements with AML/CFT obligations (which require comprehensive customer data monitoring)',
        ],
        'Source Data': [
            'System access logs from core banking, CRM, data warehouse, and digital channel systems',
            'Customer consent records from account opening, digital registration, and marketing preference systems',
            'Third-party data sharing agreements and data transmission logs',
            'Privacy incident records from information security management system',
            'FSDM entities: PARTY_PRIVACY_PREFERENCE, THIRD_PARTY_CONSENT, LOCATOR, ELECTRONIC_ADDRESS',
        ],
    },
    45: {  # Do-Not-Solicit Compliance
        'owner': 'Consumer Protection Analytics',
        'Objective / Problem Statement': [
            'Build centralized communication preference management ensuring zero non-compliant solicitation contacts across all outbound channels (SMS, email, push notification, WhatsApp, branch, call center)',
            'Comply with SBP Fair Treatment guidelines on customer contact preferences, PTA regulations on unsolicited commercial communications (including DNCR compliance), and emerging digital privacy standards',
            'Detect and prevent unauthorized contact attempts including those from third-party marketing partners and outsourced call centers',
        ],
        'Business Benefit': [
            'Zero regulatory fines from PTA and SBP for unsolicited communication violations',
            'Reduced customer complaints about unwanted contacts — a top-5 complaint category at most Pakistani banks',
            'Improved marketing effectiveness — contacts only to willing recipients improve response rates by 3-5x',
            'Reputational protection — avoiding social media backlash from aggressive marketing practices',
        ],
        'Expected Outcome': [
            'Centralized preference management system integrated with all outbound communication platforms',
            'Real-time preference enforcement preventing non-compliant contacts before execution',
            'Communication compliance audit trail for SBP and PTA regulatory inquiries',
            'Preference analytics dashboard showing opt-in/opt-out trends by channel and product',
            'Third-party marketing partner compliance monitoring',
        ],
        'POV Success Criteria': [
            'Zero non-compliant solicitation contacts within 6 months of implementation',
            'Customer preference sync across all channels within 24 hours of update',
            'PTA DNCR compliance verification automated for all outbound campaigns',
            'Communication-related complaints reduced by 50% within 12 months',
        ],
        'Methodology / Analytic Technique': [
            'Real-time API-based preference checking at point of outbound communication trigger',
            'Campaign execution audit comparing contact lists against preference database',
            'Complaint analysis correlating communication complaints to specific campaigns and channels',
            'Opt-out pattern analysis identifying products/channels driving highest opt-out rates',
            'A/B testing of communication frequency and channel mix within compliant preferences',
        ],
        'Challenges': [
            'Multiple outbound communication platforms (SMS gateway, email platform, WhatsApp Business, call center dialer) require integration with centralized preference system',
            'Customer preferences change frequently — real-time sync across all systems is technically challenging',
            'Third-party marketing partners may not have direct access to preference system — requiring pre-campaign compliance checking',
            'PTA DNCR registry updates require batch processing and may have data quality issues',
        ],
        'Source Data': [
            'Customer preference records from CRM, digital banking, and branch systems',
            'PTA Do-Not-Call Registry (DNCR) data',
            'Campaign management system — contact lists, execution records, response data',
            'Complaint management system — communication-related complaint records',
            'Third-party marketing partner contact records and consent documentation',
            'FSDM entities: PARTY_PRIVACY_PREFERENCE, CAMPAIGN, CAMPAIGN_CELL, DIRECT_CONTACT_EVENT',
        ],
    },
}


# ─────────────────────────────────────────────
# Enhancement content for existing use case slides
# ─────────────────────────────────────────────

USE_CASE_ENHANCEMENTS = {
    31: {  # AML/KYC
        'Objective / Problem Statement': [
            'Pakistan context: Comply with SBP AML/CFT Guidelines 2020, FMU (Financial Monitoring Unit) reporting requirements, and FATF/APG recommendations following Pakistan\'s grey-list experience',
            'Monitor 100% of transactions across RAAST, IBFT, mobile wallets, trade finance, and remittances — Pakistan processes PKR 15T+ annually through digital channels',
        ],
        'Source Data': [
            'CNIC/NADRA verification data for KYC — 13-digit CNIC as primary identifier',
            'RAAST and IBFT transaction data for payment network monitoring',
            'FMU STR/CTR reporting templates and thresholds',
            'FSDM entities: PARTY, INDIVIDUAL, AGREEMENT, FINANCIAL_TRANSACTION, SUSPICIOUS_ACTIVITY_REPORT',
        ],
        'Business Benefit': [
            'Pakistan-specific: Maintain FATF compliance post-grey-list exit — avoiding return to increased monitoring that constrains international banking relationships',
            'Reduce false positive rate in transaction monitoring from typical 90%+ to <70% through behavioral analytics calibrated to Pakistan transaction patterns',
        ],
        'Challenges': [
            'Pakistan-specific: Cash-intensive economy makes transaction monitoring challenging — 85% of transactions are cash-based',
            'Cross-border remittance monitoring (Pakistan receives $30B+ annually) requires integration with correspondent banking data',
        ],
    },
    32: {  # BCBS 239
        'Objective / Problem Statement': [
            'Pakistan context: SBP has issued BCBS 239 compliance requirements for D-SIBs (Domestically Systemically Important Banks) — HBL, UBL, MCB, ABL, NBP must demonstrate compliance with all 14 principles',
        ],
        'Source Data': [
            'Pakistan banking context: 4-5 core banking systems per bank, separate Islamic banking systems, multiple risk and regulatory reporting systems requiring integration',
            'FSDM provides the target data architecture for BCBS 239 compliance — 3,917 entities covering risk, finance, and regulatory domains',
        ],
        'Challenges': [
            'Pakistan-specific: Legacy core banking systems (some 20+ years old) have limited metadata and data lineage capabilities',
            'Separate Islamic banking data infrastructure adds complexity to enterprise risk aggregation',
        ],
    },
    34: {  # CCAR → SBP ICAAP
        'Objective / Problem Statement': [
            'Pakistan context: SBP requires ICAAP (Internal Capital Adequacy Assessment Process) submissions from all banks — equivalent to the CCAR framework but tailored to Pakistan macroeconomic conditions',
            'Stress testing under SBP-specified scenarios: PKR devaluation (30-50%), KIBOR spike (+500bps), GDP contraction (-3%), sectoral stress (agriculture, textile, real estate)',
        ],
        'Business Benefit': [
            'Pakistan context: Demonstrated ICAAP capability builds SBP confidence for approving advanced approaches (Foundation IRB) — reducing capital requirements',
        ],
        'Source Data': [
            'Pakistan macroeconomic data: GDP, inflation, KIBOR, PKR/USD, sectoral indices from SBP/PBS',
            'SBP stress testing scenario specifications',
            'FSDM entities: AGREEMENT, RISK_RATING, COLLATERAL, FINANCIAL_STATEMENT_LINE_ITEM',
        ],
    },
    36: {  # Solvency II → SECP Insurance Rules
        'Objective / Problem Statement': [
            'Pakistan context: While Solvency II is EU-specific, SECP Insurance Rules 2017 govern insurance company solvency in Pakistan. For bancassurance operations, banks must ensure insurance partners meet SECP solvency requirements',
            'Relevant for Pakistani banks\' growing bancassurance businesses — takaful (Islamic insurance) and conventional insurance products',
        ],
    },
    39: {  # Volcker Rule → SBP proprietary trading restrictions
        'Objective / Problem Statement': [
            'Pakistan context: SBP Banking Companies Ordinance 1962 Section 23 restricts banks\' proprietary trading activities. SBP Prudential Regulations limit banks\' equity investments and require marking-to-market of trading portfolios',
            'Pakistan treasury operations primarily involve government securities (T-bills, PIBs), FX trading, and money market — requiring compliance monitoring for position limits and P&L reporting',
        ],
        'Source Data': [
            'Pakistan context: Treasury system data for government securities (MTBs, PIBs), FX position data, money market transactions',
            'SBP weekly/monthly position reporting data',
        ],
    },
    40: {  # Sarbanes-Oxley → Pakistan Companies Act 2017
        'Objective / Problem Statement': [
            'Pakistan context: Pakistan Companies Act 2017 and SECP Code of Corporate Governance require similar internal controls over financial reporting — including Board responsibility for financial statements, internal audit, and external audit requirements',
            'SBP\'s corporate governance guidelines for banks add additional requirements specific to banking — Board composition, risk committee, audit committee, and management reporting',
        ],
        'Source Data': [
            'Pakistan context: GL data from multiple core banking systems requiring integrated financial reporting',
            'Branch-level financial data from 1,000+ branches for internal control monitoring',
        ],
    },
    42: {  # Communications Compliance
        'Objective / Problem Statement': [
            'Pakistan context: SBP Fair Treatment of Consumers guidelines mandate appropriate communication with customers. PTA regulates commercial communications including SMS, calls, and digital messaging',
            'Monitor customer communications across Urdu and English — call center recordings, WhatsApp messages, email, SMS, and branch interaction logs',
        ],
        'Source Data': [
            'Pakistan context: Call center recordings (Urdu/English bilingual analysis), WhatsApp Business messages, SMS campaign records',
            'SBP complaint data from Banking Mohtasib for compliance monitoring',
        ],
        'Challenges': [
            'Pakistan-specific: Bilingual (Urdu/English) text and speech analytics required for comprehensive monitoring',
            'WhatsApp Business emerging as major customer communication channel — requires new monitoring capabilities',
        ],
    },
}


# ─────────────────────────────────────────────
# Speaker Notes for ALL 45 slides
# ─────────────────────────────────────────────

SPEAKER_NOTES = {
    1: """Legal & Regulatory Compliance domain overview — contextualized for Pakistan's banking sector. Two sub-areas: Financial & Regulatory Compliance (AML, Capital Management, Financial Reporting, Disclosure) and Consumer Protection (Truth-in-Advertising, Predatory Practices, Privacy Protection, Do-Not-Solicit).

Pakistan regulatory landscape: SBP (prudential regulations, AML/CFT, Basel III), SECP (corporate governance, securities, insurance), FBR (tax compliance), PTA (communications), FATF/APG (AML/CFT mutual evaluation).

BVF Legal & Regulatory domain: 12 sub-capabilities. FSDM Compliance domain entities support all regulatory data requirements. BACR assessment: Most Pakistani banks at Level 1-2 for regulatory analytics maturity.""",

    2: """Financial & Regulatory Compliance covers: AML/CFT (SBP FMU compliance, FATF requirements), Capital Management (Basel III CAR, ICAAP, stress testing), Financial Reporting Regulation (200+ SBP returns, IFRS 9/16/17), and Disclosure (Pillar 3, annual report, PSX compliance).

Consumer Protection covers: Truth-in-Advertising (SBP Fair Treatment, SECP advertising regs), Predatory Practices (responsible lending, pricing fairness), Privacy Protection (SBP data protection, Pakistan PDP Bill), Do-Not-Solicit (PTA regulations, communication preferences).

FSDM entities supporting regulatory compliance: FINANCIAL_STATEMENT_LINE_ITEM, AGREEMENT, RISK_RATING, PARTY, SUSPICIOUS_ACTIVITY_REPORT, COMPLIANCE_EVENT.""",

    3: """Use Cases page 1 of 2 — Financial & Regulatory Compliance use cases: AML/KYC, BCBS 239, Basel II, CCAR (SBP ICAAP), IFRS 9, Solvency II (SECP Insurance), IFRS 15, Volcker Rule (SBP Section 23), Sarbanes-Oxley (Companies Act 2017), Regulatory Disclosure Efficiency.

Use case color coding: Orange = Traditional Data Warehouse analytics, Red = Big Data/Advanced Analytics, Blue = Non-Analytical Hadoop/Data Lake.

Pakistan prioritization: IFRS 9 ECL (highest priority — provisioning impact), AML/CFT (FATF compliance), Capital Management (Basel III CAR), and SBP Regulatory Reporting (200+ returns) are the top 4 use cases for Pakistani banks.""",

    4: """Use Cases page 2 of 2 — Consumer Protection use cases: Communications Compliance, Predatory Practices for Credit, Privacy Protection, Do-Not-Solicit Compliance.

Pakistan context: Consumer protection is rapidly evolving area for Pakistani banks. SBP Fair Treatment of Consumers framework is being enforced with increasing penalties. Banking Mohtasib complaint volumes increased 25% YoY. Digital banking growth creates new consumer protection challenges around data privacy, responsible digital lending, and transparent pricing.""",

    5: """Legal & Regulatory Compliance definition and importance — the costs of non-compliance in Pakistan are significant: SBP penalties for regulatory breaches, FATF grey-list consequences (constrained international banking), audit qualifications, and reputational damage.

Pakistan's regulatory volume: 200+ SBP regulatory returns, Basel III capital requirements (minimum CAR 11.5% including D-SIB buffer), IFRS 9 ECL provisioning, AML/CFT transaction monitoring, SECP corporate governance, FBR tax compliance.

FSDM provides the integrated data foundation enabling compliance across all regulatory domains from a single source of truth.""",

    6: """Core regulatory areas for Pakistan banking: Finance & Accounting (SBP prudential returns, IFRS 9, GL reconciliation), Capital Management (Basel III CAR, ICAAP, stress testing), AML/CFT (transaction monitoring, STR/CTR reporting, sanctions screening), Consumer Protection (Fair Treatment, responsible lending, data privacy), Trade (trade finance compliance, FX regulations), Labor (SBP employment regulations).

Pakistan-specific regulatory bodies: SBP (banking supervision), SECP (corporate governance, securities), FBR (taxation), FMU (AML/CFT), PTA (telecommunications/communications), Competition Commission of Pakistan.""",

    7: """These three questions cover the most critical regulatory compliance challenges for Pakistani banks:

1. SBP Regulatory Reporting: 200+ returns from integrated data source — most banks use 40-60 separate extraction processes. FSDM Finance domain entities (300+) provide the integrated data model.

2. AML/CFT Compliance: FATF/APG requirements demand 100% transaction coverage — including RAAST/IBFT digital channels growing 40% YoY. FSDM Event and Transaction entities support comprehensive monitoring.

3. Stress Testing/ICAAP: SBP requires bank-specific stress testing under Pakistan macro scenarios — requires end-to-end data lineage that only an integrated data warehouse can provide.

Most Pakistani banks will answer 'No' or 'Partially' to all three — creating the diagnostic opening for BVF-based capability building.""",

    8: """Section divider: Legal and Regulatory Compliance Capabilities. The following slides detail each capability area with maturity frameworks.

Implementation roadmap for Pakistani banks:
Phase 1 (0-6 months): AML/CFT enhancement and IFRS 9 ECL engine
Phase 2 (6-12 months): SBP regulatory reporting automation
Phase 3 (12-18 months): Capital management and stress testing
Phase 4 (18-24 months): Consumer protection analytics and disclosure automation

BACR maturity: Most Pakistani banks at Level 1-2. Target: Level 3 within 18-24 months.""",

    9: """Financial and Regulatory Compliance — Why is it important: Regulators expect consistency, transparency, and timeliness. SBP inspection findings on data quality, late submissions, and reconciliation breaks are among the most common regulatory issues for Pakistani banks.

Pakistan's regulatory intensification: Basel III capital increasing to 12.5% CAR (including D-SIB buffer), IFRS 9 fully effective, FATF compliance continuing, digital banking regulations expanding.

Key insight: Risk and Finance functions share underlying data but often maintain separate data extracts — FSDM provides the common data model enabling both functions from single source of truth.""",

    10: """How - Capabilities: Provide Risk and Finance with common data and analytical infrastructure. Key capabilities include automated regulatory return generation, reconciliation analytics, capital calculations, ECL provisioning, AML monitoring, and disclosure automation.

For Pakistani banks, the priority is building the integrated regulatory data warehouse using FSDM as the logical data model — replacing 40-60 separate regulatory data extraction processes.

FSDM coverage: Finance domain (300+ entities), Risk domain (200+ entities), Compliance domain (100+ entities) — all supporting regulatory reporting requirements.""",

    11: """Pakistan banking regulatory challenges:

1. Regulatory Volume: 200+ SBP returns, multiple SECP/FBR/FATF requirements. Manual preparation consumes 100+ FTE-equivalent at large banks.

2. Data Quality: Returns must reconcile to GL and financial statements — reconciliation breaks are top SBP inspection finding.

3. Evolving Landscape: Basel III increasing, IFRS 9 fully effective, FATF post-grey-list, digital banking regs expanding, ESG/climate emerging.

4. Cross-Regulatory Consistency: Same exposure data must satisfy SBP, SECP, FBR, auditors — each with different classification rules.

FSDM's 3,917 entities provide the integration blueprint eliminating separate data extraction processes.""",

    12: """Anti-Money Laundering capability overview. Pakistan's AML/CFT landscape: SBP AML/CFT Guidelines 2020, FMU (Financial Monitoring Unit) reporting, FATF Mutual Evaluation compliance.

Key challenge: Pakistan's cash-intensive economy (85% cash transactions) combined with rapidly growing digital channels (RAAST, IBFT, mobile wallets) creates complex monitoring requirements.

The integrated data approach uses FSDM Event and Transaction entities to monitor 100% of customer transactions across all channels — branch, ATM, IBFT, RAAST, mobile wallet, trade finance, remittances.""",

    13: """AML maturity assessment. Pakistan context for each level:

Leading: Advanced analytics with RAAST/IBFT network analysis, Urdu text analytics on STR narratives, real-time transaction monitoring across all channels.

Innovating: Path and graph analysis on payment networks, integration of NADRA/ECIB data for enhanced KYC.

Practicing: FATCA/CRS compliance integrated with AML data, comprehensive transaction monitoring with risk-based thresholds.

Developing: Common data platform for AML, fraud, and credit risk — leveraging shared data assets.

Emerging: Basic KYC with CNIC verification, manual STR preparation, threshold-based monitoring.

Most Pakistani banks: Level 2-3 (Developing/Practicing). Target: Level 4 (Innovating) within 24 months.""",

    14: """Capital Management capability overview. Pakistan context: SBP Basel III capital framework requires minimum CAR of 11.5% (including 1.5% D-SIB buffer for systemically important banks).

Current Pakistan banking sector CAR: ~17% average (well above minimum), but this masks significant variation and capital optimization opportunities.

Key insight: Moving from Standardized to Foundation IRB for credit risk (when SBP approves) could free significant capital — estimated PKR 10-20B for top-5 banks. The data infrastructure must be in place first.

FSDM entities: AGREEMENT, RISK_RATING, COLLATERAL, CAPITAL_INSTRUMENT, RWA_CALCULATION.""",

    15: """Capital Management maturity assessment. Pakistan context:

Leading: Dynamic capital allocation with risk/return optimization — only achievable with integrated data warehouse supporting real-time RWA calculations.

Developing (typical Pakistani bank): Capital on Standardized methodology, business lines inherit enterprise-level capital allocation. SBP inspection typically expects Level 3-4 for D-SIBs.

Current state for most Pakistani banks: Level 2 (Developing). Target: Level 3-4 (Practicing/Innovating) within 24-36 months.

Priority: Move from aggregate capital allocation to business line and product-level RAROC — enabling capital-efficient lending decisions worth PKR billions.""",

    16: """Financial Reporting Regulation — Pakistan produces 200+ SBP regulatory returns. The integrated regulatory data warehouse is the highest-priority compliance capability.

Data requirements span FSDM Finance domain (300+ entities): GL transactions, sub-ledger details, customer classification, IFRS 9 staging, capital components, FX positions, sector codes, organizational hierarchy.

Target: Reduce reporting cycle from T+15 to T+5 while improving accuracy and enabling automated reconciliation.""",

    17: """Financial Reporting Regulation maturity. Pakistan bank positioning:

Most banks currently at Developing level (T+15-20, multiple spreadsheets, limited reconciliation).

Target: Innovating level within 24 months (T+5, automated reconciliation, regulatory trend analysis).

Key investment: FSDM-based regulatory data warehouse with automated return generation for top 50 SBP returns.

Quick wins: Automated CL-1/CL-2 classification returns, capital adequacy quarterly return, IFRS 9 provision report.""",

    18: """Disclosure capability — Pakistan context: Pillar 3 quantitative disclosure, SBP quarterly disclosures, annual report IFRS 7/9 disclosures, PSX listing compliance.

Key challenge: Disclosures must reconcile precisely to regulatory returns and financial statements — inconsistencies erode stakeholder confidence and attract SBP/auditor findings.

FSDM provides the single data source ensuring consistency across all disclosure outputs.""",

    19: """Disclosure maturity. Most Pakistani banks at Developing level — manual Pillar 3 preparation, limited cross-validation, auditor-identified inconsistencies.

Target: Practicing/Innovating level within 18-24 months — semi-automated to automated quantitative disclosures with systematic cross-validation.

Key investment: Disclosure data pipeline from regulatory data warehouse to formatted output templates.""",

    20: """Consumer Protection — a rapidly emerging regulatory area for Pakistani banks. SBP's Fair Treatment of Consumers framework mandates fair outcomes across the product lifecycle.

Pakistan context: 60M+ account holders including vulnerable populations, digital lending at high APRs to low-income customers, complex fee structures, and increasing regulatory scrutiny.

Key regulatory frameworks: SBP Fair Treatment guidelines, SBP consumer banking prudential regulations, PTA communication regulations, emerging Pakistan Personal Data Protection Bill.

FSDM entities: PARTY, PARTY_PRIVACY_PREFERENCE, COMPLAINT_EVENT, CAMPAIGN, AGREEMENT.""",

    21: """Consumer Protection challenges for Pakistani banks:

1. Product Complexity: Complex fee structures, embedded charges, KIBOR-linked variable pricing that customers don't fully understand.

2. Digital Channel Risks: Mobile wallets and digital lending serve customers with limited financial literacy. Terms on small screens often not read.

3. Complaint Management: SBP requires complaint resolution within specified timeframes. Most banks treat complaints operationally without analytics.

4. Data Privacy: Emerging Pakistan data protection legislation and SBP guidelines on customer data handling.

Most Pakistani banks have no systematic consumer protection analytics — creating both regulatory risk and opportunity.""",

    22: """Truth-in-Advertising — ensuring marketing content complies with SBP Fair Treatment guidelines, SECP advertising regulations, and PBA Code of Conduct.

Pakistan context: Banks face increasing scrutiny on advertising claims, particularly for digital products and consumer lending. Social media amplifies customer dissatisfaction with misleading advertising.

Key challenge: NLP-based analysis must handle bilingual Urdu/English marketing content.""",

    23: """Truth-in-Advertising maturity. Most Pakistani banks at Emerging/Developing level — no formal marketing compliance monitoring, ad-hoc review.

Target: Practicing level within 12-18 months — systematic compliance review, basic post-campaign monitoring, advertising guidelines enforced.

Quick wins: Complaint-to-campaign linkage analysis, social media monitoring for advertising compliance issues.""",

    24: """Predatory Practices — detecting and preventing predatory lending and pricing in Pakistan banking.

Pakistan context: SBP consumer lending regulations, micro-finance interest rate caps, responsible lending guidelines. Digital lending at 30-40% APR to low-income customers raises particular concerns.

Key analytics: Customer outcome monitoring, pricing fairness analysis, insurance attachment rate analysis, DTI monitoring.""",

    25: """Predatory Practices maturity. Most Pakistani banks at Emerging level — no systematic monitoring, compliance reliance on policy documents.

Target: Practicing level within 18-24 months — annual outcome reviews, DTI limits enforced, complaint trend monitoring.

Key investment: Customer outcome monitoring dashboard covering pricing fairness, responsible lending metrics, and product suitability.""",

    26: """Privacy Protections — data privacy compliance for Pakistani banks.

Pakistan context: SBP customer data protection guidelines, emerging Personal Data Protection Bill, SBP IT Security Guidelines. Unlike HIPAA (US healthcare), Pakistan banking privacy requirements center on SBP regulations and emerging national legislation.

Key challenge: Balancing privacy compliance with AML/CFT obligations that require comprehensive customer data monitoring.""",

    27: """Privacy Protections maturity. Most Pakistani banks at Emerging/Developing level — no formal privacy framework, minimal consent management, data access not monitored.

Target: Practicing level within 18-24 months — consent collection documented, basic data access logging, third-party agreements include data protection.

Key investment: Privacy compliance platform with consent management, data access monitoring, and privacy impact assessment tools.""",

    28: """Do-Not-Solicit — communication preference management for Pakistani banks.

Pakistan context: SBP Fair Treatment guidelines require respecting customer contact preferences. PTA regulates unsolicited commercial communications including Do-Not-Call Registry (DNCR). WhatsApp Business emerging as major channel requiring new preference management.

Target: Zero non-compliant solicitation contacts — critical for both regulatory compliance and customer experience.""",

    29: """Do-Not-Solicit maturity. Most Pakistani banks at Developing level — basic opt-out tracking, inconsistent enforcement, complaint-driven violation identification.

Target: Innovating level within 18-24 months — centralized preference management, real-time checking, automated compliance monitoring.

Quick wins: PTA DNCR compliance automation, centralized preference database, campaign-level compliance checking.""",

    30: """Use Cases section — detailed use case templates for regulatory compliance analytics.

Slides 31-45 provide use case summaries covering: AML/KYC, BCBS 239, Basel II/III, CCAR/ICAAP, IFRS 9, Solvency II/SECP Insurance, IFRS 15, Volcker/SBP Section 23, Sarbanes-Oxley/Companies Act 2017, Regulatory Disclosure, Communications Compliance, Predatory Practices, Privacy Protection, Do-Not-Solicit.

All use cases contextualized for Pakistan banking with SBP regulatory references, FSDM entity mappings, and Pakistan-specific data sources.""",

    31: """AML/KYC use case — enhanced for Pakistan context. Key additions: SBP AML/CFT Guidelines 2020 compliance, FMU STR/CTR reporting, FATF/APG mutual evaluation requirements.

Pakistan-specific challenges: Cash-intensive economy (85% cash), rapidly growing digital channels (RAAST, IBFT), cross-border remittances ($30B+ annually), CNIC-based KYC using NADRA verification.

FSDM entities: PARTY, INDIVIDUAL, AGREEMENT, FINANCIAL_TRANSACTION, SUSPICIOUS_ACTIVITY_REPORT.""",

    32: """BCBS 239 use case — Pakistan context: SBP requires D-SIB compliance with all 14 BCBS 239 principles. Key challenge for Pakistani banks: 4-5 core banking systems, separate Islamic banking infrastructure, limited data lineage capabilities in legacy systems.

FSDM provides the target data architecture for BCBS 239 compliance — 3,917 entities covering risk, finance, and regulatory domains with full data lineage support.""",

    33: """Basel II/III use case — Pakistan Standardized Approach implementation. SBP requires all scheduled banks to calculate capital adequacy using Basel III Standardized Approach.

Key Pakistan context: Minimum CAR 11.5% (including D-SIB buffer), SBP asset class definitions, collateral recognition per SBP Prudential Regulations, future IRB migration roadmap.

FSDM entities: AGREEMENT, COLLATERAL, PARTY, RISK_RATING, FINANCIAL_STATEMENT_LINE_ITEM.""",

    34: """CCAR/ICAAP use case — reframed for SBP ICAAP requirements. Pakistan stress testing scenarios: PKR devaluation, KIBOR spike, GDP contraction, sectoral stress (agriculture, textile, real estate).

Key difference from US CCAR: Pakistan scenarios reflect emerging market volatility — PKR can move 30-50% in adverse scenarios, KIBOR movements of 500+ bps, GDP contraction scenarios more severe.""",

    35: """IFRS 9 use case — critical for Pakistani banks. SBP IFRS 9 Implementation Guidelines require comprehensive ECL calculation covering Stage 1/2/3, SICR assessment, forward-looking scenarios.

Key Pakistan challenges: Limited historical default data, high macroeconomic volatility for forward-looking models, reconciliation between IFRS 9 stages and SBP CL-1 to CL-5 classification.

This is the highest-priority regulatory analytics use case for most Pakistani banks — directly impacts provisioning adequacy and financial statement integrity.""",

    36: """Solvency II use case — reframed for Pakistan context: SECP Insurance Rules 2017 for bancassurance operations, takaful (Islamic insurance) solvency requirements.

Relevant for Pakistani banks' growing bancassurance businesses — banks must ensure insurance partners meet SECP solvency requirements and that bancassurance products are appropriately priced and reserved.""",

    37: """IFRS 15 use case — revenue recognition for banking and telecom. Pakistan context: Banks with telecom partnerships (mobile banking, branchless banking) must apply IFRS 15 to bundled service contracts.

Relevant for revenue recognition on fee-based banking products, subscription-based digital banking services, and partnership revenue sharing arrangements.""",

    38: """IFRS 15 continued — example output using Teradata Value Analyser. Pakistan banking application: Contract-level revenue analysis for subscription banking products, partnership revenue recognition, and fee income allocation.

FSDM Finance entities support contract-level revenue tracking and performance obligation identification.""",

    39: """Volcker Rule use case — reframed for SBP proprietary trading restrictions under Banking Companies Ordinance 1962 Section 23.

Pakistan context: Banks' treasury operations primarily involve government securities (T-bills, PIBs), FX trading, and money market. SBP limits equity investments and requires marking-to-market of trading portfolios. Position limits and P&L reporting requirements apply.""",

    40: """Sarbanes-Oxley use case — reframed for Pakistan Companies Act 2017 and SECP Code of Corporate Governance.

Pakistan context: Board responsibility for financial statements, internal audit requirements, external audit independence, management certification of internal controls. SBP adds banking-specific governance requirements.

Key challenge: GL integration across 4-5 core banking systems and 1,000+ branches for comprehensive internal control monitoring.""",

    41: """Regulatory Disclosure Efficiency use case — automating Pakistan regulatory disclosures: Pillar 3, SBP quarterly disclosures, annual report notes, PSX listing compliance.

Target: Reduce disclosure preparation time by 60-70% through automation, eliminate inconsistencies between disclosures and regulatory returns, improve disclosure quality to peer benchmarks.""",

    42: """Communications Compliance use case — enhanced for Pakistan context: SBP Fair Treatment guidelines, PTA communication regulations, bilingual (Urdu/English) communication monitoring.

Key addition: WhatsApp Business emerging as major customer communication channel requiring new monitoring capabilities. Call center recordings require Urdu/English bilingual text analytics.""",

    43: """Predatory Practices for Credit use case — Pakistan responsible lending context: SBP consumer lending regulations, DTI monitoring, pricing fairness analysis, insurance attachment monitoring.

Key concern: Digital lending at high APRs (30-40%) to low-income customers, aggressive insurance attachment to lending products, complex fee structures in consumer banking.""",

    44: """Privacy Protection use case — ENTIRELY reframed from US HIPAA to Pakistan banking context: SBP customer data protection guidelines, Pakistan Personal Data Protection Bill, SBP IT Security Guidelines.

Key difference: Pakistan's data protection legislation is still evolving — banks building privacy infrastructure now will have first-mover advantage when requirements are enforced.

FSDM entities: PARTY_PRIVACY_PREFERENCE, THIRD_PARTY_CONSENT, LOCATOR, ELECTRONIC_ADDRESS.""",

    45: """Do-Not-Solicit Compliance use case — Pakistan context: SBP Fair Treatment, PTA DNCR (Do-Not-Call Registry), digital communication preference management.

Target: Zero non-compliant solicitation contacts across all channels (SMS, email, push, WhatsApp, branch, call center).

Key challenge: Multiple outbound platforms requiring integration with centralized preference system. PTA DNCR compliance automation for all campaigns.""",
}


# ─────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────

def main():
    print(f"Loading: {INPUT_FILE}")
    prs = Presentation(INPUT_FILE)
    print(f"Total slides: {len(prs.slides)}")

    # ── Framework/Overview slides ──

    # Slide 5 — Add Pakistan context to definition
    print("\nProcessing overview slides...")
    slide5 = prs.slides[4]
    s5_shape = find_shape_by_text(slide5, 'costs of investigations')
    if s5_shape:
        tf = s5_shape.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "In Pakistan, regulatory non-compliance carries severe consequences: SBP penalties for prudential breaches, FATF grey-list constraints on international banking, audit qualifications impacting market confidence, and reputational damage in an increasingly competitive sector with 33 scheduled banks and 5 digital banking licensees."
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
    print("  Slide 5: Added Pakistan regulatory context")

    # Slide 7 — Are You Able To
    update_slide_7(prs.slides[6])

    # Slide 9 — Add Pakistan context to Financial & Regulatory Why
    slide9 = prs.slides[8]
    s9_shape = find_shape_by_text(slide9, 'Regulators, investors')
    if s9_shape:
        tf = s9_shape.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Pakistan context: SBP requires 200+ regulatory returns, Basel III capital adequacy reporting (minimum CAR 11.5%), IFRS 9 ECL provisioning, AML/CFT transaction monitoring across all channels (branch, RAAST, IBFT, mobile wallets), and comprehensive disclosure. Most Pakistani banks use 40-60 separate data extraction processes — FSDM provides the integrated data foundation."
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
    print("  Slide 9: Added Pakistan regulatory context")

    # Slide 10 — Add Pakistan capabilities context
    slide10 = prs.slides[9]
    s10_shape = find_shape_by_text(slide10, 'How - Capabilities')
    if s10_shape:
        tf = s10_shape.text_frame
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "\nPakistan Banking Context: Build integrated regulatory data warehouse using FSDM (3,917 entities) replacing 40-60 separate regulatory extraction processes. Priority capabilities: SBP regulatory return automation, IFRS 9 ECL engine, Basel III capital calculations, AML/CFT transaction monitoring, and consumer protection analytics."
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
    print("  Slide 10: Added Pakistan capabilities context")

    # Slide 11 — Challenges
    update_slide_11(prs.slides[10])

    # ── Capability table slides (16-29) ──
    print("\nProcessing capability slides...")
    update_slide_16(prs.slides[15])
    update_slide_17(prs.slides[16])
    update_slide_18(prs.slides[17])
    update_slide_19(prs.slides[18])
    update_slide_20(prs.slides[19])
    update_slide_21(prs.slides[20])

    # Consumer Protection table slides
    for slide_num, content in CAPABILITY_TABLE_CONTENT.items():
        slide = prs.slides[slide_num - 1]
        for shape in slide.shapes:
            if type(shape).__name__ == 'GraphicFrame':
                try:
                    tbl = shape.table
                    replace_cell_text(tbl.cell(1, 0), content['objectives'])
                    replace_cell_text(tbl.cell(3, 0), content['data_solution'])
                    replace_cell_text(tbl.cell(3, 2), content['outcome'])
                    print(f"  Slide {slide_num}: Filled capability table")
                    break
                except Exception as e:
                    print(f"  Slide {slide_num}: Error filling table: {e}")

    # Consumer Protection maturity slides
    for slide_num, content in MATURITY_CONTENT.items():
        slide = prs.slides[slide_num - 1]
        for shape in slide.shapes:
            if type(shape).__name__ == 'GraphicFrame':
                try:
                    tbl = shape.table
                    replace_cell_text(tbl.cell(1, 3), content['leading'])
                    replace_cell_text(tbl.cell(2, 3), content['innovating'])
                    replace_cell_text(tbl.cell(3, 3), content['practicing'])
                    replace_cell_text(tbl.cell(4, 3), content['developing'])
                    replace_cell_text(tbl.cell(5, 3), content['emerging'])
                    print(f"  Slide {slide_num}: Filled maturity levels")
                    break
                except Exception as e:
                    print(f"  Slide {slide_num}: Error filling maturity: {e}")

    # ── Use case slides — Fill placeholders ──
    print("\nProcessing use case placeholder slides...")
    for slide_num, content in USE_CASE_CONTENT.items():
        slide = prs.slides[slide_num - 1]
        fill_use_case_slide(slide, content)
        print(f"  Slide {slide_num}: Filled use case content")

    # ── Use case slides — Enhance existing ──
    print("\nEnhancing existing use case slides...")
    for slide_num, additions in USE_CASE_ENHANCEMENTS.items():
        slide = prs.slides[slide_num - 1]
        enhance_use_case_slide(slide, additions)
        print(f"  Slide {slide_num}: Enhanced with Pakistan context")

    # ── Update H1 2018 references ──
    print("\nUpdating date references...")
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if type(shape).__name__ == 'GraphicFrame':
                try:
                    tbl = shape.table
                    for row in tbl.rows:
                        for cell in row.cells:
                            if 'H1 2018' in cell.text:
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        if 'H1 2018' in run.text:
                                            run.text = run.text.replace('H1 2018', '2025-2026')
                                print(f"  Slide {idx+1}: Updated H1 2018 → 2025-2026")
                except:
                    pass

    # ── Speaker notes for ALL slides ──
    print("\nSetting speaker notes...")
    for slide_num, notes in SPEAKER_NOTES.items():
        slide = prs.slides[slide_num - 1]
        set_notes(slide, notes)
        print(f"  Slide {slide_num}: Notes set ({len(notes)} chars)")

    # ── Save ──
    print(f"\nSaving to: {OUTPUT_FILE}")
    prs.save(OUTPUT_FILE)
    print("Done!")

    # ── Verify ──
    verify = Presentation(OUTPUT_FILE)
    print(f"\nVerification: {len(verify.slides)} slides")

    notes_count = 0
    placeholder_count = 0
    for idx, slide in enumerate(verify.slides):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notes_count += 1
        # Check for remaining placeholders
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text
                for placeholder in ['Point 1', 'Question 1', 'Question 2', 'Question 3']:
                    if placeholder in text and 'Point 1' not in ['Customer Life Time Value']:
                        placeholder_count += 1
                        print(f"  WARNING: Slide {idx+1} still has placeholder: '{placeholder}' in shape {shape.shape_id}")
                        break
            if type(shape).__name__ == 'GraphicFrame':
                try:
                    tbl = shape.table
                    for row in tbl.rows:
                        for cell in row.cells:
                            ct = cell.text.strip()
                            if ct in ('Space', 'Point 1', 'Description', 'Header', 'Question 1', 'Question 2', 'Question 3'):
                                placeholder_count += 1
                                print(f"  WARNING: Slide {idx+1} table cell still has: '{ct}'")
                except:
                    pass

    print(f"\nSlides with speaker notes: {notes_count}/45")
    print(f"Remaining placeholders found: {placeholder_count}")


if __name__ == '__main__':
    main()
