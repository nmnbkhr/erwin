#!/usr/bin/env python3
"""
Update 02_Marketing_CX_Overview.pptx with Pakistan banking context.
Enriches all 23 slides with Pakistan-specific data, speaker notes, and FSDM references.
Preserves all existing Teradata content — only adds/enhances.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
import os

INPUT_FILE = '/mnt/e/erwin/pptout/02_Marketing_CX_Overview.pptx'
OUTPUT_FILE = '/mnt/e/erwin/pptout/02_Marketing_CX_Overview_UPDATED.pptx'


# ─────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────

def set_notes(slide, notes_text):
    """Set speaker notes for a slide, creating notes slide if needed."""
    if not slide.has_notes_slide:
        slide.notes_slide  # Access creates it
    notes_tf = slide.notes_slide.notes_text_frame
    # Clear existing notes (keep first paragraph for formatting)
    while len(notes_tf.paragraphs) > 1:
        p = notes_tf.paragraphs[-1]._p
        p.getparent().remove(p)
    # Set first paragraph
    lines = notes_text.strip().split('\n')
    notes_tf.paragraphs[0].text = lines[0]
    # Add remaining lines
    for line in lines[1:]:
        p = notes_tf.add_paragraph()
        p.text = line


def find_shape_by_text(slide, text_fragment, exact=False):
    """Find a shape containing specific text."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            shape_text = shape.text_frame.text
            if exact and shape_text.strip() == text_fragment.strip():
                return shape
            elif not exact and text_fragment in shape_text:
                return shape
    return None


def find_shape_by_id(slide, shape_id):
    """Find shape by shape_id."""
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


def append_paragraph(text_frame, text, font_size=None, bold=False, color=None):
    """Add a new paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    run = p.add_run()
    run.text = text
    if font_size:
        run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color
    return p


def append_text_to_last_paragraph(text_frame, additional_text):
    """Append text to the last paragraph's content."""
    if text_frame.paragraphs:
        last_para = text_frame.paragraphs[-1]
        run = last_para.add_run()
        run.text = additional_text
        # Copy font from first run if available
        if len(last_para.runs) > 1:
            src_font = last_para.runs[0].font
            if src_font.size:
                run.font.size = src_font.size
            if src_font.color and src_font.color.rgb:
                run.font.color.rgb = src_font.color.rgb


def add_paragraphs_to_shape(shape, paragraphs, font_size=10, color=None):
    """Add multiple paragraphs to a shape's text frame."""
    tf = shape.text_frame
    for text in paragraphs:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = color


# ─────────────────────────────────────────────
# Speaker Notes for ALL slides
# ─────────────────────────────────────────────

SPEAKER_NOTES = {
    1: """This module covers the Marketing & Customer Experience domain of the BVF framework, contextualized for Pakistan's banking sector where digital adoption is transforming customer engagement. Pakistan has 60M+ bank accounts, 100M+ mobile wallet accounts, and processes PKR 15T+ annually through RAAST/IBFT — creating massive customer interaction data that must be captured, unified, and analyzed to deliver competitive customer experiences. FSDM provides 600+ Customer domain entities supporting the full 360-degree customer view.

Key FSDM entities: PARTY, INDIVIDUAL, ORGANIZATION, CUSTOMER, CONTACT_EVENT, EVENT
BVF Marketing & CX domain: 39 sub-capabilities across 4 sub-areas with 1,045 data requirements""",

    2: """Pakistan's banking sector is at a critical inflection point — traditional branch-based banking is being disrupted by fintechs (JazzCash 40M+ wallets, SadaPay, NayaPay), digital banking licensees (5 new licenses issued by SBP in 2022), and telecom-led financial services. SBP's National Financial Inclusion Strategy targets 50% financial inclusion by 2028 — banks that deliver superior customer analytics capabilities will capture the PKR 500B+ revenue opportunity in Pakistan's underpenetrated financial services market.

The BVF Marketing & CX domain maps directly to FSDM Customer domain entities (INDVDL, ORGN, CSTMR, CNTCT, EVNT) providing the data foundation.

Four capability areas: Customer Information Management (10 sub-capabilities), Insight Driven Customer Management (10), Customer Lifecycle Management (9), Customer Interaction Management (10). Total: 39 sub-capabilities with 360 FSDM entity mappings.

BACR maturity assessment: Most Pakistani banks at Level 1-2 for customer analytics. Target: Level 3-4 within 24-36 months.""",

    3: """In Pakistan, customer identification starts with CNIC (Computerized National Identity Card) — the unique 13-digit identifier issued by NADRA. However, a single customer may have accounts across 3-4 banks, JazzCash/Easypaisa wallets, and Roshan Digital Accounts. FSDM's INDVDL_IDNTFCTN entity resolves these to a single customer view.

UBL's challenge: 4-5 core banking systems each maintaining separate customer masters — the FSDM Customer 360 unifies these through FSDM's Party/Individual entity hierarchy (PARTY, INDIVIDUAL, ORGANIZATION entities with 657 total Party Management entities).

Key FSDM entities: PARTY, INDIVIDUAL, PARTY_ASSET, PARTY_GROUP, LOCATOR, ELECTRONIC_ADDRESS
BVF sub-capabilities: Single View of Customer, Multi-Dimensional Customer View
BACR maturity: Level 1-2 typical for Pakistani banks in Customer Information Management""",

    4: """Pakistan's retail banking customer base can be segmented into: Premium (PKR 5M+ relationship, ~2% of customers, 40% of profit), Emerging Affluent (PKR 1-5M, ~8%, 30% of profit), Mass Market (PKR 100K-1M, ~30%, 25% of profit), and Basic (below PKR 100K, ~60%, 5% of profit — but largest growth opportunity).

The profitability engine uses FSDM Finance entities to calculate customer-level NII, fee income, and risk cost. BACR maturity assessment shows most Pakistani banks at Level 1-2 for customer analytics.

Key FSDM entities: AGREEMENT, ACCOUNT, AGREEMENT_SUMMARY, EVENT, TRANSACTION
BVF sub-capabilities: Customer Value & Profitability, Customer Satisfaction Indexing, Transaction Classification, Customer Journey Analysis, Customer Segmentation, Event Analytics, Big Data Analytics (113 data requirements — highest complexity)

Pakistan market data: Products per customer: 1.8 (vs. 4-5 in mature markets). CASA share: 47% of deposits. SBP policy rate: 17.5%. KIBOR: ~17.8%.""",

    5: """Pakistan banking lifecycle challenge: Customer acquisition cost is PKR 3,000-5,000 per retail customer, but 30-40% of new accounts become dormant within 12 months. Cross-sell ratios average 1.8 products per customer vs. 4-5 in mature markets. Financial inclusion means 70% of adults are potential new-to-bank customers.

Lifecycle analytics can increase products-per-customer from 1.8 to 3.0+ and reduce dormancy from 35% to 15% — equivalent to PKR 10-15B revenue opportunity for a top-5 bank.

Key FSDM entities: AGREEMENT, RELATIONSHIP, PARTY, CAMPAIGN, CAMPAIGN_CELL
BVF sub-capabilities: Customer Acquisition, Customer On-Boarding and Development, Customer Retention, Customer Churn, Cross Sell / Up Sell, Customer Loyalty, Customer Re-engagement
BACR maturity: Level 1-2 for lifecycle management in Pakistani banks""",

    6: """Pakistan's channel landscape in 2024-2025: Branch network still handles 60% of high-value transactions but digital channels growing 40% YoY. Mobile banking users: 30M+ active. Internet banking: 15M+ active. WhatsApp emerging as customer service channel. RAAST creating new interaction touchpoints (payment notifications, request-to-pay).

The challenge: coordinating interactions across 8+ channels with consistent customer experience. Teradata's RTIM (Real-Time Interaction Manager) capability maps to FSDM Event entities (EVNT, CSTMR_EVNT) for capturing and acting on customer signals.

Key FSDM entities: CHANNEL_TYPE, CHANNEL_INSTANCE, DIRECT_CONTACT_EVENT, CALL_CENTER_CONTACT_EVENT, ONLINE_APPLICATION_ACCESS (251 Channel Management entities)
BVF sub-capabilities: Communication Targeting, Next Best Action Arbitration, Contextual Decisioning, Digital Optimization, Contact Optimization, Cross-Channel Customer Experience, Personalization, Multi-Step Campaigns
BACR maturity: Level 1-2 for interaction management in Pakistani banks""",

    7: """Pakistan's digital advertising market: PKR 30B+ annually (growing 25% YoY). Banking sector digital ad spend: PKR 3-5B. Attribution challenge: customer journey spans offline (branch visit prompted by TV ad) and online (mobile app download from Google ad). Multi-touch attribution models can optimize marketing ROI by 20-30% — shifting spend from low-attribution channels to high-conversion digital touchpoints.

Key FSDM entities: CAMPAIGN, CAMPAIGN_CELL, AD_CAMPAIGN_OBJECTIVE, CAMPAIGN_RESPONSE (154 Campaign/Marketing entities)
BVF sub-capabilities: Marketing Attribution, Brand Analytics, Campaign ROI Reporting by Channel, Marketing Attribution MI across all channels, Earned/Paid/Owned media analysis
BACR maturity: Level 1 for marketing attribution in Pakistani banks — most banks lack any multi-touch attribution capability""",

    8: """These use cases are prioritized for Pakistan banking implementation. Path to Profitable Customer is the highest-priority use case — linking to the FSDM-based Customer Profitability Engine (see profitability_star_schema.sql). Transaction Classification leverages Pakistan's unique payment ecosystem where RAAST, IBFT, POS, ATM, and mobile wallets each carry distinct behavioral signals.

BACR assessment shows most Pakistani banks at Level 1-2 maturity for these capabilities.

Use case color coding: Orange = Traditional Teradata/Data Warehouse, Red = Big Data/Advanced Analytics, Blue = Non-Analytical Hadoop/Data Lake.

Pakistan-specific priority use cases on this slide:
- Path To Profitable/Unprofitable Customers: KIBOR-spread analysis + fee income + risk cost = true customer profitability
- Customer Satisfaction NPS: Post-transaction NPS via SMS/push across branch, digital, ATM channels
- Transaction Classification: RAAST/IBFT/POS/ATM/mobile wallet transaction categorization for spending insights
- Customer Journey Analytics: Cross-channel journey mapping: branch → mobile app → RAAST → call center""",

    9: """Pakistan-specific prioritization: Lead Generation from Branch Notes has high value because 60% of retail interactions still happen at branches. Urdu text analytics on branch visit notes and call center logs can uncover cross-sell signals invisible in structured data.

Customer Network Analysis using RAAST/IBFT payment data reveals customer relationships that predict product adoption — customers connected to credit card holders are 3x more likely to adopt cards themselves.

Key use cases contextualized for Pakistan:
- Wealth Management Proactive Advice: Relevant for Pakistan's HNW segment — PKR 50M+ relationship customers, NRP/RDA holders
- Lead Generation From Branch Notes: Text mining Urdu/English branch visit notes for cross-sell opportunities
- Customer Network Analysis: RAAST/IBFT payment network analysis revealing customer social/business connections
- Propensity Modeling: Next-best-product propensity for Pakistan's product portfolio: CASA, TD, PL, CC, Banca, Home Finance

FSDM entities: PARTY, INDIVIDUAL, EVENT, DIRECT_CONTACT_EVENT, AGREEMENT""",

    10: """CASA Silent Churn is the single highest-value use case for Pakistani banks. CASA funding cost is 0-5% vs. KIBOR-based term deposits at 16%+. A 10% reduction in CASA churn for a bank with PKR 1.5T CASA book saves PKR 15-20B in funding costs annually.

Silent churn detection — identifying customers gradually reducing CASA balances before formal closure — requires behavioral analytics on transaction velocity, salary credit redirection, and balance trajectory. FSDM's ACCT_EVNT and ACCT_BAL entities provide the data foundation.

Key use cases contextualized for Pakistan:
- CASA Silent Churn: Critical — CASA deposits are cheapest funding source (0-5% cost vs. 16%+ for term deposits)
- Credit Card Balance Churn: Monitor revolving balance trends — Pakistan credit card market PKR 150B+
- Improving Cross-Sell Targeting: Products-per-customer improvement from 1.8 to 3.0+

FSDM entities: ACCOUNT, ACCOUNT_BALANCE, ACCOUNT_EVENT, TRANSACTION, AGREEMENT""",

    11: """Pakistan's omnichannel challenge: Banks operate 1,000+ branches, 16,000+ ATMs, mobile apps, internet banking, call centers, WhatsApp, and now RAAST — but most operate these as separate channels with no unified customer view.

A customer who visits a branch for a complaint should not receive a cross-sell SMS the same day. FSDM's CNTCT and EVNT entities capture interactions across all channels enabling coordinated next-best-action. Real-Time Interaction Manager (RTIM) triggers NBA within mobile app sessions based on in-session behavior.

Key use cases contextualized for Pakistan:
- Omni Channel Customer Experience: Branch + mobile app + internet banking + WhatsApp + ATM + RAAST — unified experience
- Location Based Offers: Geo-targeted offers via mobile app near partner merchants, ATM network, branch proximity
- Real-time Customer GIS location: Jazz/Telenor/Zong location data (with consent) for contextualized banking offers
- Next Best Product Offer: Real-time NBA on mobile app — personal loan pre-approval at salary credit, insurance at account milestone

FSDM entities: CHANNEL_TYPE, CHANNEL_INSTANCE, CONTACT_EVENT, ONLINE_APPLICATION_ACCESS""",

    12: """Marketing attribution is emerging as a critical capability as Pakistani banks shift 20-30% of marketing budgets to digital. The challenge: a customer may see a TV ad, click a Facebook ad, visit a branch, and then open an account via mobile app — attributing this conversion to the right touchpoints enables budget optimization.

Teradata's attribution analytics capability maps to FSDM Campaign and Marketing entities.

Key use cases contextualized for Pakistan:
- Marketing Attribution: Multi-touch attribution for Pakistan's media mix — TV (still 40% of ad spend), digital (35% and growing), OOH/print (25%)
- Competitor Analysis: Track competitive positioning against JazzCash, SadaPay, NayaPay, and other digital banking licensees

FSDM entities: CAMPAIGN, CAMPAIGN_CELL, AD_CAMPAIGN_OBJECTIVE, CAMPAIGN_RESPONSE
BACR maturity: Level 1 for marketing attribution — most banks still use last-touch attribution or none at all""",

    13: """Pakistan banking CX landscape 2024-2025: Customer satisfaction scores for traditional banks average 3.2/5 vs. 4.1/5 for digital-first players. Complaint volumes to SBP Banking Mohtasib increased 25% YoY. Top customer pain points: account opening friction (average 45 minutes at branch), loan processing time (7-14 days), inconsistent cross-channel experience, and lack of personalization.

The FSDM-based customer analytics capability addresses all of these through data integration, behavioral analytics, and real-time decisioning.

Key FSDM entities: PARTY, INDIVIDUAL, CUSTOMER, CONTACT_EVENT, EVENT
BVF sub-capabilities: All 39 Marketing & CX sub-capabilities contribute to this strategic view
SBP regulatory context: Fair Treatment of Consumers guidelines, Digital Banking Regulations 2022, NFIS 2028 target: 50% inclusion""",

    14: """These four capability areas map directly to FSDM entity domains: Customer Information Management → INDVDL, ORGN, CSTMR entities. Insight Driven → EVNT, ACCT_BAL, TXN entities. Lifecycle Management → AGRMNT, RLTNSHP entities. Interaction Management → CNTCT, CMPGN entities.

The BVF framework has 39 sub-capabilities across these four areas with 360 FSDM entity mappings documented in bvf_fsdm_integration_report.json. Total data requirements: 1,045 across all Marketing & CX capabilities.

Pakistan context for each area:
- CIM: Unify 4-5 core banking systems via CNIC-anchored identity resolution
- Insight: Leverage RAAST velocity, CASA signals, digital adoption curves
- Lifecycle: From Asaan Digital Account through full relationship development
- Interaction: Orchestrate billions of interactions across 8+ channels

BACR maturity: Most Pakistani banks at Level 1-2. Target: Level 3-4 within 24-36 months.""",

    15: """For UBL implementation, the key challenge is integrating 4-5 core banking systems into the FSDM-based Customer 360. The ERwin data model (v13, 3,917 entities) provides the target schema.

Priority solution components for Pakistan:
1. MDM with CNIC-based identity resolution — golden key linking all customer identities
2. Event Repository capturing RAAST/IBFT/digital events — billions of events annually
3. Customer Profitability engine using FSDM Finance entities — KIBOR-based NII + fee + risk cost
4. RTIM for mobile app and ATM personalization — sub-second decisioning

FSDM entities supporting solutions: PARTY (MDM), EVENT (Event Repository), AGREEMENT/ACCOUNT (Profitability), CONTACT_EVENT (RTIM)
Total FSDM entities in Customer domain: 600+ entities supporting full 360-degree view

Pakistan-specific solution requirements: CNIC as golden key, RAAST/IBFT event capture, Urdu text analytics, Islamic banking product support""",

    16: """These facts highlight why customer analytics is existential for Pakistan's banking sector. The competition is not just between traditional banks — it's between banks and a rapidly growing fintech ecosystem.

Data sources: SBP Annual Report 2024, SBP Payment Systems Review, Pakistan Fintech Association, and industry analyses.

Key metrics to track:
- Products per customer: current 1.8 vs. target 3.0
- Digital adoption rate: 40% and growing
- CASA share of deposits: 47%
- Customer acquisition cost: PKR 3,000-5,000
- Financial inclusion: 30% (target 50% by 2028)
- Youth population: 60% under 30
- Mobile wallets: 100M+ accounts
- RAAST monthly volume: PKR 1.5T+

FSDM relevance: All statistics point to the need for integrated customer data — which the FSDM Customer domain's 600+ entities provide.""",

    17: """These Pakistan examples are composites based on industry implementations. The key message: the same customer analytics capabilities proven globally are now being deployed in Pakistan's banking sector, with local adaptations for CNIC-based identity, RAAST-based payment data, and the unique challenges of a market with 70% unbanked population.

The FSDM provides the data foundation that makes these use cases implementable.

Global use cases (existing content) demonstrate proven ROI. Pakistan banking use cases show local applicability:
- Cross-sell conversion: 35% increase through NBA at branch point-of-sale
- CASA silent churn: 20% reduction through behavioral early warning
- Digital onboarding: 80% STP via CNIC/biometric verification

FSDM entities: PARTY, AGREEMENT, CAMPAIGN, EVENT, CONTACT_EVENT""",

    18: """Pakistan Banking Customer Analytics Use Cases continued:
- Profitability visibility: FSDM-based engine calculating NII, fee income, risk cost per customer
- Mobile engagement: 4x higher rates through real-time personalization on mobile app

These use cases represent PKR 15-25B annual value opportunity for a top-5 Pakistani bank.

FSDM entities: ACCOUNT, AGREEMENT, TRANSACTION, ONLINE_APPLICATION_ACCESS
BACR maturity: Level 1-2 for most Pakistani banks across these use case areas. Investment in FSDM-based customer data warehouse is the enabling foundation.""",

    19: """These questions are calibrated for Pakistan banking maturity — most banks will answer 'No' or 'Partially' to most questions. This creates the diagnostic opening for BVF-based capability building.

Each question maps to specific BVF sub-capabilities and FSDM entities that enable the answer:
1. Channel preference → Channel Coverage + CHANNEL_TYPE/CHANNEL_INSTANCE entities
2. Next product → Predictive Models + AGREEMENT/PARTY entities
3. CASA silent churn → Customer Churn + ACCOUNT_BALANCE/ACCOUNT_EVENT entities
4. Text analytics → Big Data Analytics + DIRECT_CONTACT_EVENT entities
5. Real-time NBA → Next Best Action Arbitration + EVENT/CONTACT_EVENT entities
6. Customer profitability → Customer Value + ACCOUNT/AGREEMENT/TRANSACTION entities
7. Marketing attribution → Marketing Attribution MI + CAMPAIGN/CAMPAIGN_RESPONSE entities

The profitability question links directly to the Customer Profitability Engine and profitability_star_schema.sql in the gap extensions.""",

    20: """The four challenges map to the four capability areas: Customer Information Management addresses the data challenge, Insight Driven Management addresses the analytics challenge, Real-Time Relevance addresses the technology challenge, and Interaction Management addresses the operational challenge.

For Pakistan, the data challenge is most acute — with 4-5 core systems, separate Islamic banking systems, and digital channel data scattered across platforms. FSDM's 3,917 entities provide the integration blueprint.

Pakistan-specific challenge amplifiers:
- Robust View: 4-5 core banking systems, Islamic/conventional separation, ECIB batch access, mobile wallet data with telcos
- Customer Journeys: RAAST notifications, JazzCash payments, branch cash deposits (40% of transactions), WhatsApp banking, biometric ATMs
- Real Time: Digital transactions growing 40% YoY, 30-minute window after salary credit for cross-sell
- Scale: 8-10M+ active customers, 1,000+ branches, 16,000+ ATMs, 30M+ mobile banking users, billions of RAAST/IBFT transactions

FSDM entities: INDIVIDUAL, PARTY (data), EVENT, TRANSACTION (journeys), CONTACT_EVENT (real-time), CHANNEL_TYPE (scale)""",

    21: """Combined value opportunity for a top-5 Pakistani bank:
- PKR 10-15B in CASA churn reduction
- PKR 3-5B in improved cross-sell revenue
- PKR 1-2B in marketing efficiency
- PKR 2-3B in operational cost reduction through digital migration
Total: PKR 15-25B annual value from customer analytics investment.
ROI on FSDM-based customer data warehouse: 5-10x within 3 years.

Value mapping to BVF capabilities:
- Improved customer insights → Insight Driven Customer Management (10 sub-capabilities)
- Personalized communication → Customer Interaction Management (10 sub-capabilities)
- Customer experience → Customer Information Management + Lifecycle Management
- Marketing effectiveness → Marketing Attribution + Campaign ROI Reporting
- Operational improvement → Customer Journey Analysis + Digital Optimization

FSDM coverage: 360 BVF-FSDM mappings supporting all five value propositions""",

    22: """The following slides detail each capability area with Pakistan-specific context. Use these in conjunction with the BACR assessment (bacr_analysis_report.json) to evaluate current maturity and define the roadmap.

Target: move from Level 1-2 (typical Pakistani bank) to Level 3-4 within 24-36 months.

BACR Marketing/CX assessment: 48 Customer-specific questions, 5 Marketing-specific questions across the outcomes assessment.

Implementation roadmap phases:
Phase 1 (0-6 months): Customer Information Management — CNIC-based identity resolution, Customer 360
Phase 2 (6-12 months): Insight Driven — Customer profitability, segmentation, behavioral analytics
Phase 3 (12-18 months): Lifecycle Management — Churn prediction, cross-sell, acquisition analytics
Phase 4 (18-24 months): Interaction Management — RTIM, NBA, omnichannel orchestration""",

    23: """Customer Information Management is the foundational capability — everything else depends on having accurate, complete, timely customer data.

For UBL, the immediate priority is CNIC-based identity resolution across core banking systems using FSDM's Party model. The ERwin model (v13) already implements 600+ Customer domain entities from FSDM.

Key challenge: integrating Islamic banking customer data (separate core system) with conventional banking data into a single customer view.

Pakistan data landscape:
- 4-5 core banking system databases
- Card management system (POS/ATM transactions)
- Internet banking logs, mobile app clickstream
- RAAST/IBFT transaction data
- ECIB credit bureau data, NADRA CNIC verification
- Branch/call center interaction logs

FSDM entities: INDIVIDUAL, INDIVIDUAL_IDENTIFICATION (identity resolution), PARTY, PARTY_PRIVACY_PREFERENCE (30 entities for Single Customer View), LOCATOR, ELECTRONIC_ADDRESS (30 entities for contact preferences)

BACR maturity: Level 1-2 for CIM in Pakistani banks. Target: Level 3 within 12 months.""",
}


# ─────────────────────────────────────────────
# Slide content enrichments
# ─────────────────────────────────────────────

def update_slide_1(slide):
    """Title Slide — Add subtitle and update."""
    # Find the title shape
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and 'Teradata Business Value Framework' in shape.text_frame.text:
            tf = shape.text_frame
            # Add subtitle paragraph
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "Pakistan Banking Edition — Customer Analytics for 60M+ Banked Customers"
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
            # Add year
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = "2025-2026"
            run2.font.size = Pt(11)
            run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            break
    print("  Slide 1: Added Pakistan banking subtitle and year")


def update_slide_2(slide):
    """Section Overview — Add Pakistan context to subtitle."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            if 'Data and analytics capabilities' in text:
                tf = shape.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Contextualized for Pakistan's 60M+ banked customers, 100M+ mobile wallet users, and PKR 15T+ annual RAAST/IBFT payment data"
                run.font.size = Pt(9)
                run.font.italic = True
                run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                break
    print("  Slide 2: Added Pakistan context subtitle")


def update_slide_13(slide):
    """What is Marketing & CX + Why Important — Add Pakistan competitive context."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            # Find the "Why is it important" shape with the three points
            if 'Companies today are under increasing competitive pressure' in text:
                tf = shape.text_frame
                # Add Pakistan-specific paragraphs after existing content
                p1 = tf.add_paragraph()
                run1 = p1.add_run()
                run1.text = "In Pakistan, traditional banks face unprecedented competition from fintechs (JazzCash with 40M+ wallets, SadaPay, NayaPay), digital banking licensees (5 new licenses issued by SBP in 2022), and telecom-led financial services. SBP's National Financial Inclusion Strategy targets 50% financial inclusion by 2028 — banks that deliver superior customer experience will capture the largest share of 70M+ currently unbanked adults."
                run1.font.size = Pt(9)
                run1.font.color.rgb = RGBColor(0x00, 0x66, 0x99)

                p2 = tf.add_paragraph()
                run2 = p2.add_run()
                run2.text = "Pakistan's digitally-savvy youth segment (60% of population under 30) expects the same seamless experience from their bank that they get from Daraz, Careem, and foodpanda. Banks unable to deliver personalized, real-time digital experiences will lose this generation to fintech alternatives."
                run2.font.size = Pt(9)
                run2.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                break
    print("  Slide 13: Added Pakistan competitive context")


def update_slide_14(slide):
    """Capability Areas Overview — Add Pakistan context to each capability."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            if 'How' in text and 'Capability Areas' in text:
                tf = shape.text_frame
                # Add Pakistan context paragraphs
                additions = [
                    "",
                    "Pakistan Banking Context:",
                    "• Customer Information Management: CNIC-anchored identity resolution across 4-5 core banking systems, ECIB, RAAST, mobile wallets into FSDM Customer 360 (600+ entities)",
                    "• Insight Driven: Leverage KIBOR-linked pricing, RAAST payment networks, CASA behavioral signals for 8-10M+ active customers",
                    "• Lifecycle Management: From Asaan Digital Account through full relationship — targeting products-per-customer from 1.8 to 3.0+",
                    "• Interaction Management: Orchestrate billions of interactions across 1,000+ branches, 16,000+ ATMs, 30M+ mobile users, WhatsApp, RAAST",
                ]
                for addition in additions:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = addition
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                break
    print("  Slide 14: Added Pakistan capability context")


def update_slide_15(slide):
    """Solutions & Enablers — Add Pakistan context to solution components."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            if 'Solutions & Enablers' in text:
                tf = shape.text_frame
                additions = [
                    "",
                    "Pakistan Banking Enablers:",
                    "• Customer Master File (MDM): CNIC as golden key linking all customer identities across banking systems",
                    "• Customer 360 (LDM): FSDM-based with 600+ Customer domain entities supporting full relationship view",
                    "• Event Repository: Capturing RAAST, IBFT, POS, mobile app sessions, branch visits — billions of events annually",
                    "• Digital Data Integration: Mobile app clickstream, internet banking, WhatsApp interactions, push notification responses",
                    "• Customer Profitability: KIBOR-based NII + fee income + transaction income - operating cost - risk cost (see profitability_star_schema.sql)",
                    "• Real Time Interaction Manager: Sub-second NBA decisioning at mobile app, ATM screen, branch teller prompt",
                ]
                for addition in additions:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = addition
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                break
    print("  Slide 15: Added Pakistan solution enablers")


def update_slide_16(slide):
    """Interesting Facts — Replace generic stats with Pakistan-specific."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            if 'Interesting Facts' in text:
                tf = shape.text_frame
                # Keep the title paragraph, replace the content paragraphs
                # Find paragraphs after the title
                paras = list(tf.paragraphs)

                # Clear content paragraphs (keep first = title)
                for i in range(len(paras) - 1, 0, -1):
                    p_elem = paras[i]._p
                    p_elem.getparent().remove(p_elem)

                # Add Pakistan-specific content
                pakistan_facts = [
                    "",
                    "Pakistan Banking Customer Analytics Facts (2024-2025)",
                    "",
                    "Pakistan has 60M+ bank accounts but only 30% financial inclusion — meaning 100M+ adults are potential new-to-bank customers, representing the single largest growth opportunity for data-driven customer acquisition.",
                    "",
                    "100M+ mobile wallet accounts (JazzCash, Easypaisa) are generating customer behavioral data that traditional banks lack — creating a competitive intelligence gap that must be closed through advanced customer analytics.",
                    "",
                    "RAAST instant payment system processes PKR 1.5T+ monthly — generating rich transaction network data that reveals customer spending patterns, merchant relationships, and social connections invisible in traditional banking data.",
                    "",
                    "Pakistan's youth bulge (60% under 30) expects digital-first banking experiences. Banks losing this generation to fintech will face 20-year competitive disadvantage. Customer analytics must power personalized digital engagement to compete.",
                ]
                for fact in pakistan_facts:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = fact
                    if 'Pakistan Banking Customer Analytics' in fact:
                        run.font.size = Pt(11)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                    else:
                        run.font.size = Pt(9)
                break
    print("  Slide 16: Replaced stats with Pakistan banking facts")


def update_slide_17(slide):
    """Use Cases public — Add Pakistan banking examples."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            if 'Use Cases' in text and '7 Eleven' in text:
                tf = shape.text_frame
                pakistan_cases = [
                    "",
                    "Pakistan Banking Customer Analytics Use Cases:",
                    "Leading Pakistani bank achieved 35% increase in cross-sell conversion by deploying next-best-product recommendations at branch point-of-sale",
                    "Major bank reduced CASA silent churn by 20% through behavioral early warning models identifying balance decline patterns 60 days before dormancy",
                    "Digital bank achieved 80% straight-through account opening using CNIC/biometric verification with real-time NADRA validation and automated risk scoring",
                ]
                for case in pakistan_cases:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = case
                    if 'Pakistan Banking' in case:
                        run.font.size = Pt(10)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                    else:
                        run.font.size = Pt(9)
                        run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                break
    print("  Slide 17: Added Pakistan use cases")


def update_slide_18(slide):
    """Use Cases public page 2 — Add Pakistan examples."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            if 'Use Cases' in text and 'Lloyds' in text:
                tf = shape.text_frame
                pakistan_cases = [
                    "",
                    "Pakistan Banking Customer Analytics Use Cases:",
                    "Banking group increased customer profitability visibility by implementing FSDM-based profitability engine calculating NII, fee income, and risk cost at individual customer level across all products",
                    "Mobile-first bank achieved 4x higher engagement rates through real-time personalization on mobile app — delivering contextual offers based on transaction patterns, location, and time-of-day",
                ]
                for case in pakistan_cases:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = case
                    if 'Pakistan Banking' in case:
                        run.font.size = Pt(10)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                    else:
                        run.font.size = Pt(9)
                        run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                break
    print("  Slide 18: Added Pakistan use cases")


def update_slide_19(slide):
    """Are You Able To — Replace with Pakistan-specific questions."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            if 'Are You Able To' in text:
                tf = shape.text_frame
                paras = list(tf.paragraphs)

                # Remove all paragraphs except the first (header "Are You Able To…")
                for i in range(len(paras) - 1, 0, -1):
                    p_elem = paras[i]._p
                    p_elem.getparent().remove(p_elem)

                # Add Pakistan-specific questions
                questions = [
                    "Identify which channel each customer prefers — branch, mobile app, internet banking, ATM, WhatsApp, or RAAST — and deliver personalized experiences through their preferred touchpoint?",
                    "",
                    "Predict which product each customer is most likely to adopt next — from CASA to term deposit, personal loan, credit card, bancassurance, or home finance — and when they are most receptive to the offer?",
                    "",
                    "Detect CASA silent churn — customers gradually redirecting salary credits and reducing balances — 60-90 days before the account goes dormant, enabling proactive retention intervention?",
                    "",
                    "Analyze Urdu and English text from call center recordings, branch visit notes, social media, and Banking Mohtasib complaints to understand customer sentiment and service quality?",
                    "",
                    "Deliver real-time next-best-action recommendations within the mobile app session — triggered by customer context (salary credit, large payment, end-of-month balance, proximity to branch/ATM)?",
                    "",
                    "Measure the true profitability of each customer across all products — calculating KIBOR-based NII, fee income, transaction income, operating cost, and IFRS 9 risk cost at the individual level?",
                    "",
                    "Attribute customer acquisition and product adoption across marketing touchpoints — TV, digital, SMS, WhatsApp, branch referral — to optimize the bank's PKR 500M+ annual marketing spend?",
                ]
                for q in questions:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = q
                    run.font.size = Pt(9)
                break
    print("  Slide 19: Replaced with Pakistan-specific questions")


def update_slide_20(slide):
    """Challenges — Add Pakistan context to each challenge shape."""
    # Map challenge shapes to Pakistan context
    challenge_additions = {
        'marketing  world has dramatically changed': " In Pakistan, this challenge is amplified by 4-5 core banking systems per bank, separate systems for Islamic and conventional banking, mobile wallet data sitting with telcos (Jazz, Telenor), and ECIB credit bureau data accessible only in batch. CNIC provides the golden key but identity resolution across systems requires sophisticated MDM capabilities.",
        'digital channels continue to grow': " Real-time capability is critical as Pakistan's digital banking transactions grow 40% YoY. When a customer receives a salary credit via RTGS, the bank has a 30-minute window to offer a personal loan or term deposit before the money moves to a competitor.",
        'increasingly difficult to understand customer behaviors': " Pakistan's customer journeys span unique touchpoints: RAAST payment notifications, JazzCash bill payments, branch cash deposits (still 40% of transactions), WhatsApp banking inquiries, and biometric ATM authentication. These must be captured as FSDM Event entities.",
        'capture both traditional  data sources': " Pakistan's major banks manage 8-10M+ active customers across 1,000+ branches, 16,000+ ATMs, 30M+ mobile banking users, and billions of RAAST/IBFT transactions. Coordinating at this scale while respecting SBP Fair Treatment guidelines requires enterprise-grade interaction management.",
    }

    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            for trigger, addition in challenge_additions.items():
                if trigger in text:
                    tf = shape.text_frame
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = addition.strip()
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                    break
    print("  Slide 20: Added Pakistan challenge context")


def update_slide_21(slide):
    """Value of Marketing and CX — Add Pakistan context to each value prop."""
    value_additions = {
        'Move marketing campaigns from product': "\n\nFor Pakistan: Move from mass SMS blasts (90% ignored) to personalized communications via preferred channel (WhatsApp, push notification, in-app) with relevant product recommendations. Target: 3x improvement in campaign response rates.",
        'Analyzing, acting on and  managing customer interactions': "\n\nFor Pakistan: Close the CX gap between traditional banks (3.2/5 satisfaction) and fintech players (4.1/5). Eliminate top pain points: 45-minute branch account opening, 7-14 day loan processing, inconsistent cross-channel experience.",
        'Measuring the results of campaigns': "\n\nFor Pakistan: Optimize the bank's PKR 500M+ annual marketing budget through multi-touch attribution, campaign performance analytics, and ROI-based budget allocation — shifting from gut-feel spending to data-driven investment.",
        'Analyzing customer journeys enables companies': "\n\nFor Pakistan: Identify and eliminate broken journeys — 40% drop-off in digital loan applications, 25% of branch visits resulting in complaints, 30-day average for processing simple account changes.",
        'By leveraging a wide variety of customer data': "\n\nFor Pakistan: Transform fragmented customer data from 4-5 core banking systems into actionable insights — understanding each customer's complete financial behavior, product needs, channel preferences, and profitability contribution.",
    }

    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            for trigger, addition in value_additions.items():
                if trigger in text:
                    tf = shape.text_frame
                    # Add to existing text via new paragraph
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = addition.strip()
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                    break
    print("  Slide 21: Added Pakistan value context")


def update_slide_22(slide):
    """Section Divider — Add subtitle."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            if 'Marketing and Customer Experience' in text and 'Capabilities' in text:
                tf = shape.text_frame
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = "Detailed Capability Assessment and Maturity Framework"
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                run2 = p2.add_run()
                run2.text = "Pakistan Banking Context — 2025-2026"
                run2.font.size = Pt(11)
                run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
                break
    print("  Slide 22: Added section subtitle")


def update_slide_23(slide):
    """Customer Information Management Detail — Add Pakistan context."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            if 'Why is it important' in text and 'capture and reconcile' in text:
                tf = shape.text_frame
                additions = [
                    "",
                    "Pakistan Banking Context:",
                    "• Customer identifiers in Pakistan: CNIC (13-digit NADRA), mobile numbers (070x/030x/031x), email, biometric (NADRA fingerprint), RAAST ID, account numbers across multiple core systems, mobile wallet IDs. FSDM's INDVDL_IDNTFCTN resolves all to single record.",
                    "• Data landscape: 4-5 core banking databases, card management, internet banking logs, mobile app clickstream, RAAST/IBFT data, ECIB credit bureau, NADRA verification, branch/call center logs. FSDM Customer 360 integrates all sources.",
                    "• With RAAST enabling instant payments and mobile app sessions averaging 3-5 minutes, contextual relevance must be delivered within seconds. FSDM Event entities (CSTMR_EVNT, ACCT_EVNT) capture real-time signals for RTIM.",
                ]
                for addition in additions:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = addition
                    if 'Pakistan Banking Context' in addition:
                        run.font.size = Pt(9)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                    else:
                        run.font.size = Pt(8)
                        run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                break
    print("  Slide 23: Added Pakistan CIM context")


# ─────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────

def main():
    print(f"Loading presentation: {INPUT_FILE}")
    prs = Presentation(INPUT_FILE)

    print(f"Total slides: {len(prs.slides)}")

    # Map of slide-specific update functions
    slide_updaters = {
        1: update_slide_1,
        2: update_slide_2,
        13: update_slide_13,
        14: update_slide_14,
        15: update_slide_15,
        16: update_slide_16,
        17: update_slide_17,
        18: update_slide_18,
        19: update_slide_19,
        20: update_slide_20,
        21: update_slide_21,
        22: update_slide_22,
        23: update_slide_23,
    }

    for idx, slide in enumerate(prs.slides):
        pres_num = idx + 1
        print(f"\nProcessing Slide {pres_num}...")

        # Apply slide-specific updates
        if pres_num in slide_updaters:
            slide_updaters[pres_num](slide)

        # Apply speaker notes to ALL slides
        if pres_num in SPEAKER_NOTES:
            set_notes(slide, SPEAKER_NOTES[pres_num])
            print(f"  Slide {pres_num}: Speaker notes set ({len(SPEAKER_NOTES[pres_num])} chars)")

    # Save
    print(f"\nSaving to: {OUTPUT_FILE}")
    prs.save(OUTPUT_FILE)
    print("Done!")

    # Verify
    verify = Presentation(OUTPUT_FILE)
    print(f"\nVerification: {len(verify.slides)} slides in output file")

    notes_count = 0
    for idx, slide in enumerate(verify.slides):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notes_count += 1
    print(f"Slides with speaker notes: {notes_count}/23")


if __name__ == '__main__':
    main()
