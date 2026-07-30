#!/usr/bin/env python3
"""
Update File 11 — Product Management Use Cases
Enriches 17 slides with Pakistan banking product context, rewrites 2 insurance slides,
replaces Discover case study, removes Teradata/IC branding, adds 3 new slides.
"""

import os
import copy
import re
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

INPUT_FILE = './11_Product_Management_Use_Cases.pptx'
CLEAN_FILE = './pptout/11_Product_Management_Use_Cases_CLEAN.pptx'
OUTPUT_FILE = './pptout/11_Product_Management_Use_Cases_UPDATED.pptx'


# ─────────────────────────────────────────────
# Field label mapping
# ─────────────────────────────────────────────

FIELD_LABELS = {
    'objective':        'Objective / Problem Statement',
    'business_benefit': 'Business Benefit',
    'source_data':      'Source Data',
    'methodology':      'Methodology / Analytic Technique',
    'outcome':          'Expected Outcome',
    'challenges':       'Challenges',
    'success_criteria': 'Success Criteria',
}

# IC / owner names to remove from References shapes
IC_NAMES = [
    "Andrew Johnston", "Dominic Ligot", "Vince Leat", "Tim Dickson",
    "Yasmeen Ahmed", "Nick Campbell", "Simon Axon", "Terence Kong",
]

# Text replacements for Teradata tool references
TOOL_REPLACEMENTS = {
    'on Aster': 'on advanced analytics platform',
    ' Aster': ' advanced analytics platform',
    'nPath': 'path analysis',
    'Npath': 'Path analysis',
    'NPath': 'Path Analysis',
}


# ─────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────

def set_notes(slide, text):
    """Set or replace the speaker notes on a slide."""
    if not slide.has_notes_slide:
        slide.notes_slide  # creates notes slide
    notes_tf = slide.notes_slide.notes_text_frame
    for i in range(len(notes_tf.paragraphs) - 1, 0, -1):
        p = notes_tf.paragraphs[i]._element
        p.getparent().remove(p)
    notes_tf.paragraphs[0].text = text


def find_shape_by_label(slide, label):
    """Find a text box shape by its first paragraph (label) text."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            paras = shape.text_frame.paragraphs
            if paras and paras[0].text.strip() == label:
                return shape
    return None


def get_content_font_size(shape):
    """Get the font size used in content paragraphs (P1+)."""
    if len(shape.text_frame.paragraphs) > 1:
        for run in shape.text_frame.paragraphs[1].runs:
            if run.font.size:
                return run.font.size
    return Pt(8)


def append_to_shape(shape, lines):
    """Append content lines to a shape's text frame after existing content."""
    font_size = get_content_font_size(shape)
    for line in lines:
        p = shape.text_frame.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = font_size


def replace_shape_content(shape, lines):
    """Replace all content in a shape, keeping the P0 label paragraph."""
    tf = shape.text_frame
    font_size = get_content_font_size(shape)

    # Remove P1+ paragraphs
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._element
        p.getparent().remove(p)

    # Add new content
    for line in lines:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = font_size


def replace_tool_references(slide):
    """Replace Aster/nPath/Teradata tool references in all text on a slide."""
    replaced = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    original = run.text
                    for find_text, replace_text in TOOL_REPLACEMENTS.items():
                        if find_text in run.text:
                            run.text = run.text.replace(find_text, replace_text)
                    if run.text != original:
                        replaced += 1
    return replaced


def remove_ic_names(slide):
    """Remove IC owner names and related references from References shapes."""
    removed = 0
    ref_shape = find_shape_by_label(slide, 'References')
    if not ref_shape:
        return removed

    tf = ref_shape.text_frame
    paras_to_remove = []

    for i, para in enumerate(tf.paragraphs):
        if i == 0:  # Skip label
            continue
        text = para.text.strip()
        # Check for IC names
        should_remove = False
        if text.startswith('IC ') or text.startswith('DS:') or text.startswith('Owner:'):
            should_remove = True
        for name in IC_NAMES:
            if name in text:
                should_remove = True
                break
        if should_remove:
            paras_to_remove.append(para)
            removed += 1

    for para in paras_to_remove:
        para._element.getparent().remove(para._element)

    # Check if only label remains — add replacement
    remaining_content = [p.text.strip() for p in tf.paragraphs[1:] if p.text.strip()]
    if not remaining_content:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "Banking Industry Best Practice"
        run.font.size = Pt(9)

    return removed


def remove_race_text(slide):
    """Remove 'RACE Relevant' and 'Not eligible for RACE' from Success Criteria."""
    removed = 0
    sc_shape = find_shape_by_label(slide, 'Success Criteria')
    if not sc_shape:
        return removed

    tf = sc_shape.text_frame
    paras_to_remove = []

    for i, para in enumerate(tf.paragraphs):
        if i == 0:
            continue
        text = para.text.strip()
        if 'RACE' in text:
            paras_to_remove.append(para)
            removed += 1

    for para in paras_to_remove:
        para._element.getparent().remove(para._element)

    return removed


def clean_pptx(input_path, output_path):
    """Remove orphaned slide XMLs from the PPTX zip to prevent save corruption."""
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
          'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
          'rel': 'http://schemas.openxmlformats.org/package/2006/relationships'}

    with zipfile.ZipFile(input_path, 'r') as zin:
        pres_xml = zin.read('ppt/presentation.xml')
        root = ET.fromstring(pres_xml)
        sldIdLst = root.findall('.//p:sldIdLst/p:sldId', ns)
        active_rids = set()
        for s in sldIdLst:
            rid = s.attrib.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            if rid:
                active_rids.add(rid)

        rels_xml = zin.read('ppt/_rels/presentation.xml.rels')
        rels_root = ET.fromstring(rels_xml)

        active_slides = set()
        for rel in rels_root.findall('.//rel:Relationship', ns):
            if rel is None:
                continue
            rid = rel.get('Id', '')
            target = rel.get('Target', '')
            if rid in active_rids and 'slides/slide' in target:
                active_slides.add('ppt/' + target.lstrip('/'))

        all_slides = {n for n in zin.namelist()
                      if n.startswith('ppt/slides/slide') and n.endswith('.xml')}
        orphans = all_slides - active_slides
        orphan_rels = {s.replace('ppt/slides/', 'ppt/slides/_rels/') + '.rels' for s in orphans}
        exclude = orphans | orphan_rels

        print(f"  Active slides: {len(active_slides)}, Orphaned: {len(orphans)}")

        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                if item not in exclude:
                    zout.writestr(item, zin.read(item))

    print(f"  Cleaned PPTX saved to {output_path}")


# ─────────────────────────────────────────────
# ENRICHMENTS — Pakistan content to APPEND to existing use case slides
# Keyed by slide index (0-based)
# ─────────────────────────────────────────────

ENRICHMENTS = {
    1: {  # Product Development and Packaging
        'objective': [
            "Design new banking products aligned with Pakistan market needs: digital-first accounts, Islamic product variants, financial inclusion products, Roshan Digital Account extensions, green banking products."
        ],
        'source_data': [
            "FSDM customer/transaction data, SBP financial inclusion survey, mobile app usage patterns, RAAST adoption data, competitor product analysis, Islamic banking demand indicators"
        ],
        'methodology': [
            "Product usage clustering (identify unmet needs), propensity modeling for new product adoption, cannibalization analysis (will new product erode existing?)"
        ],
        'outcome': [
            "Data-driven product pipeline: 3-5 new products/year based on analytics, not intuition. Reduced product failure rate from ~40% to <20%."
        ],
        'challenges': [
            "SBP product approval timeline (4-8 weeks), Islamic Shariah board approval, core banking configuration across 4-5 systems"
        ],
    },
    2: {  # Relationship Mgt. and Optimization
        'objective': [
            "Analyze relationships between banking objects — customer, product, channel, branch, RM, transaction — using graph analytics to identify optimization opportunities.",
            "Pakistan banking: understand which RM behaviors drive higher product adoption. Identify why certain branches have 3x products/customer vs. others."
        ],
        'source_data': [
            "FSDM: customer-product holdings, transaction patterns, channel usage, RM assignment, branch allocation, product lifecycle events"
        ],
        'methodology': [
            "Graph analysis on FSDM relationship data, RM productivity benchmarking, product affinity network analysis"
        ],
        'outcome': [
            "Optimize RM-customer allocation for maximum cross-sell, identify successful RM behaviors for training/replication"
        ],
    },
    3: {  # Merchant Acquiring Insights
        'objective': [
            "Leverage card transaction data to provide merchant acquirers (banks operating POS networks) with customer insights. Pakistan's POS network is growing rapidly — 100,000+ terminals.",
            "New revenue opportunity: monetize aggregated card transaction analytics for merchant business insights (footfall patterns, average ticket size, competitive share-of-wallet)"
        ],
        'source_data': [
            "Card transaction data from FSDM, merchant category codes (MCC), merchant location data, 1Link POS switch data, customer demographics"
        ],
        'outcome': [
            "New fee income from merchant analytics services, increased POS terminal deployment, merchant loyalty through value-added insights"
        ],
    },
    4: {  # Credit/Debit Card Spend Stimulation — CRITICAL
        'objective': [
            "Pakistan has 30M+ debit cards but <30% are actively used for POS purchases (most used only for ATM cash withdrawal). Stimulate card spend to increase interchange revenue and reduce cash handling cost.",
            "Identify dormant/low-usage cardholders with high spending potential. Match with relevant merchant offers."
        ],
        'source_data': [
            "Card transaction data (ATM vs. POS vs. e-commerce split), customer FSDM profile, salary credit data (spending capacity indicator), merchant category preferences, previous offer response data"
        ],
        'methodology': [
            "Collaborative filtering (customers with similar profiles who spend more — target non-spenders), propensity to activate POS usage, merchant offer matching"
        ],
        'outcome': [
            "Target: activate 30% of dormant POS cards, increase average POS spend by 25% in targeted segment. Revenue impact: PKR 2-5B additional interchange industry-wide."
        ],
    },
    5: {  # Individualized Merchants Recommendation
        'objective': [
            "Recommend relevant merchants to individual cardholders based on spending patterns, location, preferences. Increase card usage frequency and average spend."
        ],
        'source_data': [
            "FSDM card transaction data by MCC and merchant, customer location (branch, home address), spending pattern history, merchant location data"
        ],
        'outcome': [
            "Personalized merchant recommendations via mobile app push notification, SMS, or in-app widget. Target: 10-15% increase in card transaction frequency."
        ],
    },
    6: {  # Improvement in Risk Scoring Models — CRITICAL
        'objective': [
            "Enhance ECIB-based credit scoring with transactional behavior data from FSDM. Pakistan's ECIB has limited data for many customers (thin-file or new-to-bureau). Adding behavioral indicators from transactions dramatically improves predictive power."
        ],
        'source_data': [
            "ECIB bureau data, FSDM: salary credit regularity, balance volatility, RAAST/IBFT outflow patterns, merchant spend categories, overdraft frequency, cheque bounce history, mobile app login frequency"
        ],
        'methodology': [
            "Build enhanced scorecards: traditional credit variables + behavioral variables from FSDM. Test with out-of-sample validation. Champion/challenger deployment."
        ],
        'outcome': [
            "Improve Gini coefficient by 10-20%, reduce false rejection rate by 15-25%, enable lending to 'thin-file' customers previously declined"
        ],
    },
    7: {  # Multivariate Testing (page 1)
        'objective': [
            "A/B and multivariate testing for Pakistan digital banking: test mobile app product screens, lending application flow, deposit rate display, card offer placement, onboarding funnel optimization.",
            "Pakistan context: digital banking is early stage — A/B testing can yield 20-50% improvement in conversion rates because baseline is unoptimized."
        ],
        'source_data': [
            "Mobile app event data, internet banking clickstream, onboarding funnel data, digital sales conversion data"
        ],
        'outcome': [
            "Evidence-based digital UX optimization, reduced onboarding dropoff (target: from 40-60% to <20%), improved digital product cross-sell conversion"
        ],
    },
    9: {  # Competitor Analysis
        'objective': [
            "Monitor 33+ commercial banks, 5 full Islamic banks, 11 microfinance banks, and fintechs (JazzCash, Easypaisa, SadaPay, NayaPay, TAG) — their products, pricing, features, marketing, and customer sentiment."
        ],
        'source_data': [
            "Competitor websites, app store reviews (Google Play), SBP published data (rates, financials), social media (Twitter, Facebook), job postings (indicate strategic direction), annual reports"
        ],
        'methodology': [
            "NLP/text analytics on competitor reviews and social media, rate comparison dashboards, feature gap analysis, market share trend tracking (SBP data)"
        ],
        'outcome': [
            "Quarterly competitor intelligence report, real-time rate comparison alerts, product gap identification, fintech threat assessment"
        ],
    },
    12: {  # Credit Risk Models for New-to-Lending — CRITICAL
        'objective': [
            "Pakistan has 150M+ adults with no formal credit history. Build credit models using alternative data to enable lending to the unbanked and underbanked — SBP financial inclusion mandate."
        ],
        'source_data': [
            "Mobile wallet transaction data (JazzCash/Easypaisa), utility payment history (K-Electric, SSGC, SNGPL), telecom usage/payment data (Jazz, Telenor, Zong), RAAST P2P patterns, NADRA demographic data, employer verification, social media (with consent)"
        ],
        'methodology': [
            "Alternative data credit scoring: ML models trained on default outcomes of thin-file customers who were subsequently lent to. Feature engineering from mobile wallet and utility payment regularity."
        ],
        'outcome': [
            "Expand addressable lending market from 30M (current bank customers) to 80M+ (include mobile wallet users). Enable nano-lending (PKR 5K-50K) via digital channel."
        ],
    },
    13: {  # Understanding Sales Across All Channels (page 1)
        'objective': [
            "Unified view of product sales across all Pakistan banking channels: branch (16,000+), mobile app, internet banking, call center, agent network (branchless), ATM/kiosk, employer partnership (salary accounts), RAAST onboarding."
        ],
        'source_data': [
            "FSDM: product origination channel tag, branch sales register, mobile app funnel data, internet banking applications, call center lead data, agent network transaction data"
        ],
        'outcome': [
            "Channel-level product sales dashboard: which products sell on which channels? Where are the dropoffs? Which channels are cannibalizing (intended or not)?"
        ],
    },
    15: {  # Compare Portfolio to Benchmark
        'objective': [
            "Pakistan's wealth management is nascent — bank financial advisors manage customer investment portfolios (mutual funds, government securities, equities, insurance). Compare customer portfolio allocation to age/risk-appropriate benchmarks.",
            "Pakistan context: National Savings Certificates dominate retail investment — banks need to offer competitive alternatives. SECP-regulated mutual fund industry growing (PKR 1.5T+ AUM)."
        ],
        'source_data': [
            "Customer investment holdings (FSDM), age/risk profile, benchmark portfolio models, market data (PSX, mutual fund NAVs), National Savings rates"
        ],
        'outcome': [
            "Advisor-facing portfolio comparison tool, automated alerts for out-of-benchmark portfolios, product recommendation engine"
        ],
    },
}


# ─────────────────────────────────────────────
# REWRITES — slides 10, 11 get FULL content replacement
# ─────────────────────────────────────────────

REWRITES = {
    10: {  # Advanced Risk and Pricing Insights → Advanced Risk-Based Pricing Analytics
        'title': "Advanced Risk-Based Pricing Analytics",
        'objective': [
            "Use granular risk data to price lending products more accurately than standard KIBOR + flat spread",
            "Differentiate pricing by customer risk profile, collateral quality, sector risk, and geographic risk"
        ],
        'business_benefit': [
            "Improved risk-adjusted returns on lending portfolio",
            "Reduced adverse selection — retain low-risk customers, price appropriately for high-risk",
            "20-50bps NIM improvement potential"
        ],
        'source_data': [
            "ECIB credit scores, FSDM behavioral risk indicators",
            "Collateral valuation data, sector performance data (SBP sector-wise NPL data)",
            "Geographic risk data (flood zone, security zone), macroeconomic forward indicators"
        ],
        'methodology': [
            "Risk-based pricing models combining PD (from enhanced scorecard) x LGD (from collateral and recovery data) x EAD",
            "Set minimum risk-adjusted price per customer",
            "Overlay with competitive pricing intelligence"
        ],
        'outcome': [
            "Risk-adjusted lending rates that properly compensate for credit risk",
            "Reduced adverse selection",
            "Improved portfolio risk-return trade-off"
        ],
        'challenges': [
            "Building granular PD/LGD models requires sufficient default history",
            "Core banking system limitations for individual customer pricing",
            "SBP rate regulations and minimum spread requirements"
        ],
        'success_criteria': [
            "20-50bps improvement in portfolio NIM",
            "Reduced adverse selection measurable through improved loss ratios",
            "Competitive retention of high-quality borrowers"
        ],
    },
    11: {  # Behavioral-Based Pricing with Telematics → Behavioral-Based Product Pricing
        'title': "Behavioral-Based Product Pricing",
        'objective': [
            "Price banking products based on actual customer behavior rather than static attributes",
            "Reward positive financial behavior with better pricing",
            "Pakistan context: reward customers who maintain stable balances (lower deposit cost), who transact digitally (lower cost-to-serve), who use products across categories (higher lifetime value)"
        ],
        'business_benefit': [
            "Better customer retention through fairness-based pricing",
            "Reduced cost-to-serve by incentivizing digital behavior",
            "Improved NIM through behavioral segmentation"
        ],
        'source_data': [
            "FSDM behavioral data: balance stability (coefficient of variation), digital transaction ratio",
            "Product breadth, payment punctuality, salary credit consistency"
        ],
        'methodology': [
            "Behavioral scoring models that create dynamic pricing tiers",
            "Updated monthly based on trailing 6-month behavior"
        ],
        'outcome': [
            "Better customer retention through fairness-based pricing",
            "Reduced cost-to-serve by incentivizing digital behavior",
            "Improved NIM through behavioral segmentation"
        ],
        'challenges': [
            "SBP minimum rate constraints",
            "Customer communication of dynamic pricing",
            "System capability for frequent rate changes"
        ],
        'success_criteria': [
            "Measurable improvement in customer retention for high-value segments",
            "15-20% increase in digital transaction ratio",
            "Reduced voluntary attrition in top-tier customers"
        ],
    },
}


# ─────────────────────────────────────────────
# DISCOVER REPLACEMENT — slide 16 content
# ─────────────────────────────────────────────

DISCOVER_REPLACEMENT = {
    'title': "Pakistan Product Analytics — Success Metrics",
    'content': [
        "PRODUCT MANAGEMENT ANALYTICS IMPACT (Pakistan Banking Targets):",
        "",
        "Products per Customer:        1.8 → 2.5        (+39% revenue per customer)",
        "Card POS Activation:          <30% → 50%       (PKR 3-5B additional interchange)",
        "Digital Product Sales:        15% → 40%         (of total product originations)",
        "Product Launch Cycle:         16 weeks → 8      (with data-driven development)",
        "Cross-sell Response Rate:     <5% → 10%         (through analytics targeting)",
        "New-to-Lending Coverage:      20M → 40M         (through alternative data models)",
        "Product Failure Rate:         ~40% → <20%       (through business impact prediction)",
        "Risk Scoring Accuracy:        +10-20% Gini      (through FSDM behavioral data)",
        "",
        "INVESTMENT:  PKR 200-400M over 3 years",
        "RETURN:      PKR 10-20B annually through revenue uplift + cost reduction",
    ],
}


# ─────────────────────────────────────────────
# SPEAKER NOTES — for ALL 17 slides
# ─────────────────────────────────────────────

SPEAKER_NOTES = {
    0: (  # Section Divider
        "Product Management Use Cases — 15 analytical use cases for banking product lifecycle management.\n"
        "Covers: Product Development (7 use cases), Product Introduction (1), Pricing & Promotion (3), Product Performance (3).\n"
        "Pakistan Focus: Card spend stimulation, risk scoring improvement, new-to-lending credit models, digital product testing.\n"
        "BVF Reference: File 10 (Product Management capability detail + maturity), File 06 (CIM use cases — related), File 08-09 (Profitability — pricing)"
    ),
    1: (  # Product Development and Packaging
        "PAKISTAN PRODUCT DEVELOPMENT PRIORITIES:\n"
        "1. DIGITAL-FIRST: Account opening via selfie + CNIC, instant nano-lending (PKR 5K-50K), QR merchant payments\n"
        "2. ISLAMIC VARIANTS: Every conventional product needs Islamic equivalent — many gaps exist\n"
        "3. FINANCIAL INCLUSION: SBP mandate for Basic Banking Account, Asaan Mobile Account, agent banking products\n"
        "4. ROSHAN DIGITAL: Extensions for overseas Pakistanis (Naya Pakistan Certificate, stock market access, home finance)\n"
        "5. SME PRODUCTS: Digital onboarding, supply chain finance, invoice discounting\n"
        "6. GREEN BANKING: SBP green taxonomy products — green loans, green deposits\n\n"
        "FSDM: PRDCT (Product), PRDCT_FTRE (Product Feature), PRDCT_DSGN (Product Design), MKT_ND (Market Need), CSTMR_SGMNT (Customer Segment)\n"
        "BVF Capability: Product Development > Research New Product Opportunities (File 10, Slides 21-22)"
    ),
    2: (  # Relationship Mgt. and Optimization
        "PAKISTAN CONTEXT:\n"
        "RM effectiveness varies enormously in Pakistan banks:\n"
        "- Top 10% of RMs generate 3-5x products/customer vs. average\n"
        "- Graph analysis can reveal: which product sequencing works, which RM behaviors drive adoption, which branch configurations produce best results\n"
        "- Challenge: RM data (assignments, interactions) often not captured systematically in core banking systems — FSDM can integrate CRM + core banking + call center data\n\n"
        "FSDM: IP_RLTNP (Inter-Party Relationship), CSTMR_PRDCT (Customer-Product), RM_ALLCTN (RM Allocation), GRP_ANLYS (Graph Analysis)"
    ),
    3: (  # Merchant Acquiring Insights
        "PAKISTAN MERCHANT ACQUIRING:\n"
        "Pakistan POS landscape:\n"
        "- 100,000+ POS terminals (growing rapidly with SBP push for digital payments)\n"
        "- Major acquirers: HBL, MCB, UBL, Allied, Bank Alfalah\n"
        "- Visa/Mastercard interchange: 1.5-2.5% — bank earns acquiring share\n"
        "- Opportunity: Sell anonymized, aggregated spending insights to merchants (like Visa Analytics, MC SpendingPulse)\n"
        "- Example: 'Your restaurant average ticket is PKR 2,500, 30% below area average — opportunity to upsell'\n\n"
        "FSDM: MRCHNT (Merchant), MRCHNT_TXN (Merchant Transaction), MCC (Merchant Category Code), INTRCHNG (Interchange), POS_TRML (POS Terminal)"
    ),
    4: (  # Credit/Debit Card Spend Stimulation — CRITICAL
        "PAKISTAN CARD SPEND:\n"
        "The card spend opportunity in Pakistan is massive:\n"
        "- 30M+ debit cards issued, but >70% used ONLY for ATM cash withdrawal\n"
        "- Credit card base ~3M with ~60% active — low by global standards\n"
        "- Cash-on-delivery still dominates e-commerce (70%+)\n"
        "- SBP pushing digital payments: RAAST QR, POS incentives\n\n"
        "Stimulation strategies:\n"
        "1. Cashback on first POS transaction (activate dormant cards)\n"
        "2. Merchant-specific offers based on spending patterns (collaborative filtering)\n"
        "3. Salary-day spend stimulation (within 48 hours of salary credit)\n"
        "4. Category spend challenges (spend PKR 10K at restaurants → get PKR 500 cashback)\n"
        "5. EMI conversion offers (convert large purchases to installments — increases revolving revenue)\n\n"
        "Data requirements: Full card transaction history, merchant data, customer demographics, salary patterns.\n\n"
        "FSDM: CRD (Card), CRD_TXN (Card Transaction), POS_TXN (POS Transaction), ATM_TXN (ATM Transaction), MRCHNT_OFFR (Merchant Offer), INTRCHNG (Interchange)"
    ),
    5: (  # Individualized Merchants Recommendation
        "PAKISTAN MERCHANT RECOMMENDATION:\n"
        "Opportunity: Pakistan mobile app downloads growing 30%+ YoY — push notifications for merchant offers are a viable channel.\n"
        "Current state: Most banks send generic merchant offers to all card customers. No personalization.\n"
        "Target: Use collaborative filtering on card transaction data to recommend merchants each customer is likely to visit.\n"
        "Key challenge: Merchant data quality — POS terminals often mapped to generic MCCs, not specific merchant names.\n\n"
        "FSDM: CRD_TXN (Card Transaction), MRCHNT (Merchant), MCC (Merchant Category Code), CSTMR_PREF (Customer Preference)"
    ),
    6: (  # Improvement in Risk Scoring Models — CRITICAL
        "PAKISTAN RISK SCORING:\n"
        "Current state: Most Pakistan banks use ECIB score + basic internal rating.\n"
        "Challenge: ECIB has limited history for 70%+ of population (new-to-bureau or thin-file).\n\n"
        "FSDM-enhanced variables that improve scoring:\n"
        "1. SALARY REGULARITY: Is salary credited consistently? Amount trend up/down?\n"
        "2. BALANCE BEHAVIOR: Average balance trend, minimum balance, balance volatility\n"
        "3. TRANSACTION PATTERNS: Regular bill payments (disciplined), frequent ATM withdrawals (cash-dependent)\n"
        "4. RAAST/IBFT: Outflow pattern — regular payments suggest obligations, erratic suggest instability\n"
        "5. MERCHANT SPEND: Spend categories reveal lifestyle/income level\n"
        "6. DIGITAL ENGAGEMENT: App login frequency, bill payment regularity\n"
        "7. RELATIONSHIP DEPTH: Products held, tenure, cross-sell response history\n"
        "8. OVERDRAFT BEHAVIOR: Frequency, amount, recovery speed\n\n"
        "These variables are available in FSDM but NOT in ECIB bureau data — unique competitive advantage for banks with integrated data warehouses.\n\n"
        "FSDM: RSK_SCR (Risk Score), BHVRL_SCR (Behavioral Score), ECIB_SCR (ECIB Score), PD (Probability of Default), SCRCRD (Scorecard), MDLNG (Modeling)\n"
        "BVF Reference: File 12-13 Risk Management, File 08-09 Profitability (risk-adjusted pricing)"
    ),
    7: (  # Multivariate Testing (page 1)
        "PAKISTAN DIGITAL TESTING:\n"
        "Pakistan digital banking is in early stages — most bank apps have never been A/B tested.\n"
        "Opportunity: Low baseline conversion rates mean simple tests can yield 20-50% improvement.\n"
        "Key test areas: Account opening funnel, loan application flow, deposit rate display, card activation, bill payment UX.\n"
        "Example: One bank reduced account opening dropoff from 55% to 22% by simplifying the CNIC upload screen.\n\n"
        "FSDM: DGTL_EVNT (Digital Event), APP_SSSN (App Session), CNVRSN_EVNT (Conversion Event), FNNL_STEP (Funnel Step)"
    ),
    8: (  # Multivariate Testing (page 2 — analytics)
        "PAKISTAN DIGITAL TESTING (continued):\n"
        "Analytics approach: Session-level path analysis on mobile app and internet banking events.\n"
        "Key metrics: Conversion rate per funnel step, session duration, error rates, A/B test lift.\n"
        "Pakistan-specific: Test Urdu vs. English UI, test CNIC photo vs. manual entry, test biometric vs. PIN verification.\n\n"
        "FSDM: DGTL_EVNT (Digital Event), A_B_TEST (A/B Test), TST_VRNT (Test Variant), CNVRSN_RT (Conversion Rate)"
    ),
    9: (  # Competitor Analysis
        "PAKISTAN COMPETITIVE LANDSCAPE:\n"
        "Direct competitors: HBL, MCB, UBL, Allied, Bank Alfalah, Meezan (Islamic), Bank Islami\n"
        "Fintech disruptors: JazzCash (50M+ users), Easypaisa (40M+), SadaPay, NayaPay, TAG\n"
        "International: Standard Chartered, Citibank (exited retail)\n\n"
        "Key competitive dimensions:\n"
        "1. DEPOSIT RATES: Savings rate, FD rates by tenor — published by SBP\n"
        "2. LENDING RATES: Personal loan markup, home finance rate, credit card rate\n"
        "3. DIGITAL EXPERIENCE: App rating (Google Play), feature set, onboarding speed\n"
        "4. BRANCH NETWORK: Coverage, quality, operating hours\n"
        "5. ISLAMIC PRODUCTS: Range, Shariah compliance reputation\n"
        "6. CUSTOMER SERVICE: NPS, complaint resolution (SBP publishes complaint data)\n\n"
        "FSDM: CMPTTR (Competitor), CMPTTR_PRDCT (Competitor Product), MKT_SHR (Market Share), BNCHMRK (Benchmark)"
    ),
    10: (  # Advanced Risk-Based Pricing Analytics (REWRITTEN)
        "PAKISTAN RISK-BASED PRICING:\n"
        "Currently: Most Pakistan banks price lending as KIBOR + flat spread by segment. A corporate customer with AA rating pays similar rate to BBB customer in same segment.\n\n"
        "Target: Granular risk-based pricing where:\n"
        "- Rate = Cost of Funds (FTP) + Operating Cost (ABC) + Expected Loss (PD x LGD) + Capital Charge + Target Margin\n"
        "- Each component calculated per individual customer\n"
        "- Result: Low-risk customers get competitive rates (retain against competitors), high-risk customers are priced for risk (or declined)\n\n"
        "Impact: 20-50bps NIM improvement, reduced adverse selection, better risk-adjusted returns.\n\n"
        "FSDM: RSK_PRCNG (Risk-Based Pricing), PD (PD), LGD (LGD), EAD (EAD), RSK_PREMM (Risk Premium), FTP_RT (FTP Rate)"
    ),
    11: (  # Behavioral-Based Product Pricing (REWRITTEN)
        "PAKISTAN BEHAVIORAL PRICING:\n"
        "Example applications:\n"
        "1. SAVINGS RATE: Customers who maintain stable monthly balance get 50bps premium on savings rate\n"
        "2. LENDING RATE: Customers with 12+ months perfect repayment get 100bps reduction on next loan\n"
        "3. CARD FEES: Heavy digital spenders get annual fee waiver (lower cost-to-serve justifies it)\n"
        "4. ACCOUNT FEES: Customers who go fully digital (no branch visits) get zero maintenance fee\n"
        "5. FD RATES: Customers with 3+ products get loyalty premium on FD rate (retention value justifies cost)\n\n"
        "Behavioral variables from FSDM:\n"
        "- Balance stability index (low volatility = lower funding risk = reward)\n"
        "- Digital transaction ratio (higher digital = lower cost = reward)\n"
        "- Product breadth score (more products = lower churn risk = reward)\n"
        "- Payment discipline score (on-time payments = lower credit risk = reward)\n"
        "- Relationship tenure (longer = more predictable = reward)\n\n"
        "FSDM: BHVRL_SCR (Behavioral Score), BHVRL_PRCNG (Behavioral Pricing), CSTMR_BHVR (Customer Behavior), PRCNG_TR (Pricing Tier)"
    ),
    12: (  # Credit Risk Models for New-to-Lending — CRITICAL
        "PAKISTAN NEW-TO-LENDING:\n"
        "This is THE most impactful use case for Pakistan financial inclusion:\n\n"
        "THE PROBLEM:\n"
        "- 150M+ adults, only ~60M have bank accounts, only ~20M have credit bureau history\n"
        "- ECIB data is insufficient to score 80%+ of the population\n"
        "- Result: Banks lend only to salaried/documented sector (20-30% of workforce)\n"
        "- 70% informal sector is excluded from formal credit\n\n"
        "THE SOLUTION:\n"
        "Alternative data sources available in Pakistan:\n"
        "1. MOBILE WALLETS: JazzCash (50M+) and Easypaisa (40M+) have transaction histories — payment regularity, balance patterns, peer-to-peer behavior\n"
        "2. TELECOM: 190M+ mobile connections — top-up patterns, call behavior, data usage correlate with income stability\n"
        "3. UTILITY PAYMENTS: K-Electric, SSGC, SNGPL, WASA — regular payment = financial discipline\n"
        "4. RAAST: Growing P2P transaction data — can infer income, spending patterns\n"
        "5. NADRA: Demographic data (age, location, family size) — can inform risk segmentation\n\n"
        "MODEL APPROACH:\n"
        "- Train on existing customers where outcome (default/non-default) is known\n"
        "- Engineer features from alternative data\n"
        "- Validate on out-of-sample population\n"
        "- Deploy with small initial credit limits, expand based on repayment behavior\n"
        "- Champion/challenger: traditional model vs. alternative data model\n\n"
        "IMPACT: If even 10% of unbanked adults get access to PKR 50K credit, that's PKR 750B+ in new lending opportunity.\n\n"
        "FSDM: ALTNTV_DATA (Alternative Data), CRDTH_SCR (Credit Score), MOBL_WLT (Mobile Wallet), UTLTY_PMT (Utility Payment), THNFL (Thin File)"
    ),
    13: (  # Understanding Sales Across All Channels (page 1)
        "PAKISTAN OMNI-CHANNEL PRODUCT SALES:\n"
        "Pakistan banking channels: 16,000+ branches, growing mobile apps, internet banking, call centers, agent banking, ATM/kiosk.\n"
        "Challenge: Product sales tracked in separate systems per channel — no unified view of which products sell where.\n"
        "Opportunity: Identify channel-specific conversion rates and optimize product placement by channel.\n"
        "Example: Savings accounts open 70% in branches, but digital-only accounts could shift 40% to mobile if funnel is optimized.\n\n"
        "FSDM: CHNL (Channel), PRDCT_ORGN (Product Origination), CHNL_SLS (Channel Sales), BRCH_SLS (Branch Sales)"
    ),
    14: (  # Understanding Sales Across All Channels (page 2)
        "PAKISTAN OMNI-CHANNEL (continued):\n"
        "Analytics approach: Path analysis on customer product journey — which touchpoints lead to conversion?\n"
        "Key insight: Many Pakistan bank customers research online but transact in branch (ROPO effect).\n"
        "Agent banking channel growing rapidly — 500,000+ agents — but product range limited to basic transactions.\n\n"
        "FSDM: CHNL_JRNY (Channel Journey), TUCHPNT (Touchpoint), CNVRSN_EVNT (Conversion Event), AGNT_NTW (Agent Network)"
    ),
    15: (  # Compare Portfolio to Benchmark
        "PAKISTAN WEALTH MANAGEMENT:\n"
        "National Savings Certificates: PKR 4T+ outstanding — dominant retail investment vehicle.\n"
        "Mutual fund industry: PKR 1.5T+ AUM under SECP regulation, growing.\n"
        "Bank wealth management: Early stage — most advisory is basic asset allocation.\n"
        "Opportunity: Build portfolio comparison tools that benchmark against age/risk-appropriate models.\n"
        "Challenge: Customer investment data fragmented across bank, NCCPL (stocks), AMCs (mutual funds), National Savings.\n\n"
        "FSDM: INVSTMNT_HLDNG (Investment Holding), PRTFL (Portfolio), BNCHMRK_PRTFL (Benchmark Portfolio), ADVSRY (Advisory)"
    ),
    16: (  # Pakistan Product Analytics — Success Metrics (REPLACED)
        "PAKISTAN PRODUCT ANALYTICS IMPACT:\n"
        "These targets represent the combined impact of implementing all 15 product management use cases.\n"
        "Investment of PKR 200-400M over 3 years yields estimated PKR 10-20B annual return.\n"
        "Key drivers: Card activation, digital sales, new-to-lending expansion, risk scoring improvement.\n"
        "Cross-reference: File 10 capability maturity shows current Product Management analytics maturity at Developing level."
    ),
}


# ─────────────────────────────────────────────
# NEW SLIDE NOTES — for 3 new slides
# ─────────────────────────────────────────────

NEW_SLIDE_NOTES = {
    'new_a': (  # Dashboard
        "PRODUCT MANAGEMENT USE CASES — PORTFOLIO DASHBOARD:\n"
        "15 use cases organized by product lifecycle stage: Development, Introduction, Pricing & Promotion, Performance.\n"
        "Priority matrix based on Pakistan market impact and implementation readiness.\n"
        "Critical use cases: Card Spend Stimulation, Risk Scoring Improvement, New-to-Lending Credit Models.\n"
        "Cross-reference: File 10 (Product Management capabilities), File 06 (CIM use cases), File 08-09 (Profitability)."
    ),
    'new_b': (  # Pakistan Context
        "PAKISTAN PRODUCT USE CASE PRIORITIES:\n"
        "Top 5 use cases with highest Pakistan market impact, with specific data sources and expected outcomes.\n"
        "Key insight: Pakistan's large unbanked population (100M+) and low card POS activation (<30%) represent massive untapped opportunities.\n"
        "FSDM provides the integrated data foundation required for all 15 use cases.\n"
        "SBP financial inclusion mandate creates regulatory tailwind for New-to-Lending and digital product use cases."
    ),
    'new_c': (  # Roadmap
        "IMPLEMENTATION ROADMAP:\n"
        "3-phase approach: Quick Wins (0-6 months), Product Intelligence (6-18 months), Advanced Analytics (18-36 months).\n"
        "Phase 1 focuses on 5 highest-impact use cases with existing data: card spend, risk scoring, new-to-lending, competitor analysis, omni-channel.\n"
        "Total investment: PKR 200-400M over 3 years. Expected return: PKR 10-20B annually.\n"
        "Dependencies: FSDM data integration (Phase 1), advanced analytics platform (Phase 2), graph analytics (Phase 3)."
    ),
}


# ─────────────────────────────────────────────
# Main processing
# ─────────────────────────────────────────────

def main():
    print("=" * 70)
    print("File 11 — Product Management Use Cases — Update Script")
    print("=" * 70)

    # ── Step 1: Clean orphaned slides ──
    print(f"\n── Step 1: Cleaning orphaned slides ──")
    clean_pptx(INPUT_FILE, CLEAN_FILE)
    prs = Presentation(CLEAN_FILE)
    print(f"  Loaded: {len(prs.slides)} slides")

    # ── Step 2: Replace Teradata tool references ──
    print(f"\n── Step 2: Replacing tool references (Aster/nPath) ──")
    total_replaced = 0
    for i, slide in enumerate(prs.slides):
        if i in REWRITES:
            continue  # Rewrite slides get completely replaced
        if i == 16:
            continue  # Discover slide gets completely replaced
        count = replace_tool_references(slide)
        if count:
            print(f"  Slide {i}: replaced {count} tool reference(s)")
            total_replaced += count
    print(f"  Total tool references replaced: {total_replaced}")

    # ── Step 3: Remove IC names and RACE text ──
    print(f"\n── Step 3: Removing IC names, Owner lines, and RACE text ──")
    total_ic = 0
    total_race = 0
    for i, slide in enumerate(prs.slides):
        if i == 16:
            continue
        ic_count = remove_ic_names(slide)
        race_count = remove_race_text(slide)
        if ic_count:
            print(f"  Slide {i}: removed {ic_count} IC/Owner/DS reference(s)")
            total_ic += ic_count
        if race_count:
            print(f"  Slide {i}: removed {race_count} RACE reference(s)")
            total_race += race_count
    print(f"  Total IC/Owner removed: {total_ic}, RACE removed: {total_race}")

    # ── Step 4: Enrich use case slides (append Pakistan content) ──
    print(f"\n── Step 4: Enriching use case slides ──")
    for slide_idx, enrichment in ENRICHMENTS.items():
        slide = prs.slides[slide_idx]
        fields_enriched = 0
        for field_key, lines in enrichment.items():
            label = FIELD_LABELS.get(field_key)
            if not label:
                continue
            shape = find_shape_by_label(slide, label)
            if shape:
                append_to_shape(shape, lines)
                fields_enriched += 1
            else:
                print(f"  WARNING: Slide {slide_idx} — field '{label}' not found")
        print(f"  Slide {slide_idx}: enriched {fields_enriched} field(s)")

    # ── Step 5: Rewrite insurance slides (10, 11) ──
    print(f"\n── Step 5: Rewriting insurance/telematics slides ──")
    for slide_idx, rewrite in REWRITES.items():
        slide = prs.slides[slide_idx]
        fields_rewritten = 0

        # Replace title
        new_title = rewrite.get('title')
        if new_title:
            # Find the title shape — it's the large text that isn't a field label
            for shape in slide.shapes:
                if shape.has_text_frame:
                    p0 = shape.text_frame.paragraphs[0].text.strip()
                    # Title shapes: "Advanced Risk and  Pricing Insights" or "Behavioral-Based Pricing..."
                    if ('Risk' in p0 and 'Pricing' in p0 and p0 not in FIELD_LABELS.values()) or \
                       ('Behavioral' in p0 and 'Pricing' in p0 and p0 not in FIELD_LABELS.values()):
                        for run in shape.text_frame.paragraphs[0].runs:
                            run.text = ''
                        shape.text_frame.paragraphs[0].runs[0].text = new_title
                        print(f"  Slide {slide_idx}: title → '{new_title}'")
                        break

        # Replace each field
        for field_key, lines in rewrite.items():
            if field_key == 'title':
                continue
            label = FIELD_LABELS.get(field_key)
            if not label:
                continue
            shape = find_shape_by_label(slide, label)
            if shape:
                replace_shape_content(shape, lines)
                fields_rewritten += 1
            else:
                print(f"  WARNING: Slide {slide_idx} — field '{label}' not found for rewrite")
        print(f"  Slide {slide_idx}: rewrote {fields_rewritten} field(s)")

        # Also clean IC names on rewritten slides
        remove_ic_names(slide)

    # ── Step 6: Replace Discover case study (slide 16) ──
    print(f"\n── Step 6: Replacing Discover case study (slide 16) ──")
    discover_slide = prs.slides[16]
    for shape in discover_slide.shapes:
        if shape.has_text_frame:
            p0 = shape.text_frame.paragraphs[0].text.strip()
            # Title shape: "Driving Further Innovation at Discover"
            if 'Discover' in p0 and 'Innovation' in p0:
                for run in shape.text_frame.paragraphs[0].runs:
                    run.text = ''
                shape.text_frame.paragraphs[0].runs[0].text = DISCOVER_REPLACEMENT['title']
                shape.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
                print(f"  Replaced title → '{DISCOVER_REPLACEMENT['title']}'")
            # Company name: "Discover"
            elif p0 == 'Discover':
                for run in shape.text_frame.paragraphs[0].runs:
                    run.text = ''
                print(f"  Cleared company name")
            # URL shape
            elif p0.startswith('http'):
                for run in shape.text_frame.paragraphs[0].runs:
                    run.text = ''
                print(f"  Cleared URL")
            # Main content: "Challenge" → replace with metrics
            elif p0 == 'Challenge' or 'centralized repository' in shape.text_frame.text:
                # Clear all content
                tf = shape.text_frame
                for i in range(len(tf.paragraphs) - 1, -1, -1):
                    if i > 0:
                        tf.paragraphs[i]._element.getparent().remove(tf.paragraphs[i]._element)
                    else:
                        for run in tf.paragraphs[0].runs:
                            run.text = ''

                # Add new content
                first = True
                for line in DISCOVER_REPLACEMENT['content']:
                    if first:
                        run = tf.paragraphs[0].add_run()
                        run.text = line
                        run.font.size = Pt(9)
                        run.font.bold = True if 'IMPACT' in line or 'INVESTMENT' in line or 'RETURN' in line else False
                        first = False
                    else:
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = line
                        run.font.size = Pt(9)
                        if 'INVESTMENT' in line or 'RETURN' in line:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(0, 102, 51)
                        elif line and not line.startswith(' ') and ':' in line and '→' in line:
                            run.font.name = "Consolas"
                            run.font.size = Pt(8)
                print(f"  Replaced content with Pakistan metrics")

    # ── Step 7: Set speaker notes ──
    print(f"\n── Step 7: Setting speaker notes ──")
    for slide_idx, notes_text in SPEAKER_NOTES.items():
        set_notes(prs.slides[slide_idx], notes_text)
    print(f"  Set notes on {len(SPEAKER_NOTES)} slides")

    # ── Step 8: Add 3 new slides ──
    print(f"\n── Step 8: Adding new slides ──")
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name == 'Blank':
            blank_layout = layout
            break
    if not blank_layout:
        blank_layout = prs.slide_layouts[-1]
        print(f"  WARNING: 'Blank' layout not found, using '{blank_layout.name}'")

    # New Slide A: Dashboard
    slide_a = prs.slides.add_slide(blank_layout)
    _build_dashboard_slide(slide_a)
    set_notes(slide_a, NEW_SLIDE_NOTES['new_a'])
    print(f"  Added: Use Case Portfolio Dashboard (New Slide A)")

    # New Slide B: Pakistan Context
    slide_b = prs.slides.add_slide(blank_layout)
    _build_pakistan_context_slide(slide_b)
    set_notes(slide_b, NEW_SLIDE_NOTES['new_b'])
    print(f"  Added: Pakistan Product Use Case Priorities (New Slide B)")

    # New Slide C: Implementation Roadmap
    slide_c = prs.slides.add_slide(blank_layout)
    _build_roadmap_slide(slide_c)
    set_notes(slide_c, NEW_SLIDE_NOTES['new_c'])
    print(f"  Added: Implementation Roadmap (New Slide C)")

    # ── Step 9: Reorder slides ──
    print(f"\n── Step 9: Reordering slides ──")
    _reorder_slides(prs)
    print(f"  Slides reordered: Dashboard, PK Context, original 17, Roadmap")

    # ── Step 10: Save ──
    print(f"\n── Step 10: Saving to {OUTPUT_FILE} ──")
    prs.save(OUTPUT_FILE)

    if os.path.exists(CLEAN_FILE):
        os.remove(CLEAN_FILE)
        print("  Cleaned up temp file")

    final_prs = Presentation(OUTPUT_FILE)
    file_size = os.path.getsize(OUTPUT_FILE)
    print(f"  DONE! Output: {len(final_prs.slides)} slides ({file_size // 1024:,} KB)")


# ─────────────────────────────────────────────
# Builder functions for new slides
# ─────────────────────────────────────────────

def _build_dashboard_slide(slide):
    """Build the Use Case Portfolio Dashboard slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx in (0, 1):
                shape.text = ""

    left = Emu(400000)
    top = Emu(850000)
    width = Emu(8300000)
    height = Emu(4000000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Title
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Product Management Use Cases — Portfolio Dashboard"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    p = tf.add_paragraph()
    p.text = ""

    # Section summary
    sections = [
        ("Product Development (7 UCs)", "Development & Packaging, Relationship Optimization, Merchant Acquiring, Card Spend Stimulation, Merchant Recommendation, Risk Scoring Improvement, Multivariate Testing", "Critical: Card Spend, Risk Scoring"),
        ("Product Introduction (1 UC)", "Competitor Analysis", "High"),
        ("Pricing & Promotion (3 UCs)", "Risk-Based Pricing, Behavioral Pricing, New-to-Lending Credit Models", "Critical: New-to-Lending"),
        ("Product Performance (3 UCs)", "Omni-Channel Sales, Portfolio Benchmark", "High: Omni-Channel"),
    ]

    for section, use_cases, priority in sections:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = section
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 51, 102)

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {use_cases}"
        run.font.size = Pt(8)

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  Priority: {priority}"
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(153, 0, 51)

    p = tf.add_paragraph()
    p.text = ""

    # Priority matrix
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Pakistan Priority Matrix:"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    matrix = [
        "IMMEDIATE (0-6 mo):   Card Spend Stimulation, Risk Scoring, New-to-Lending, Competitor Analysis, Omni-Channel Sales",
        "HIGH (6-18 mo):       Merchant Acquiring, Merchant Recommendation, Behavioral Pricing, Product Development, Multivariate Testing",
        "STRATEGIC (18-36 mo): Relationship Optimization (Graph), Portfolio Benchmark (Wealth), Advanced Risk-Based Pricing",
    ]
    for line in matrix:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(8)
        run.font.name = "Consolas"

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "TOTAL: 15 Use Cases (1 Discover case study removed) | BVF Reference: File 10 (Product Management Capabilities)"
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 102, 153)


def _build_pakistan_context_slide(slide):
    """Build the Pakistan Product Use Case Priorities slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx in (0, 1):
                shape.text = ""

    left = Emu(400000)
    top = Emu(850000)
    width = Emu(8300000)
    height = Emu(4000000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Title
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Pakistan Banking — Product Use Case Priorities"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "USE CASES WITH HIGHEST PAKISTAN IMPACT:"
    run.font.size = Pt(10)
    run.font.bold = True

    use_cases = [
        ("1. CREDIT RISK MODELS FOR NEW-TO-LENDING",
         "Why: 70% of Pakistan population has no credit bureau history — cannot access formal lending",
         "Data: Mobile wallet patterns (JazzCash/Easypaisa), utility payment history, RAAST P2P, NADRA demographic data",
         "Impact: Expand lending to 50M+ previously unbanked — SBP financial inclusion priority"),
        ("2. CARD SPEND STIMULATION",
         "Why: Pakistan has 30M+ debit cards but <30% are actively used for POS purchases",
         "Data: Card transaction data (POS, ATM, e-commerce), merchant category codes, 1Link switch data, customer FSDM profile",
         "Impact: Shift cash transactions to card — increase interchange revenue PKR 5-10B industry-wide"),
        ("3. RISK SCORING IMPROVEMENT",
         "Why: Current scoring models use limited variables. Adding transactional behavior improves Gini by 10-20%",
         "Data: ECIB bureau data + FSDM transactional patterns (salary regularity, spending behavior, balance trends, RAAST usage)",
         "Impact: Reduce false rejections, expand lending to thin-file customers"),
        ("4. COMPETITOR ANALYSIS",
         "Why: Intense competition from 33+ banks + 5 Islamic + JazzCash/Easypaisa/SadaPay/NayaPay",
         "Data: Competitor product features, rate sheets, app store reviews, social media sentiment",
         "Impact: Real-time competitive intelligence for product and pricing decisions"),
        ("5. OMNI-CHANNEL SALES",
         "Why: Product sales fragmented across branch, mobile, internet, agent, call center — no unified view",
         "Data: FSDM channel data + mobile app analytics + branch sales register + agent network",
         "Impact: Optimize product-channel mix, reduce conversion dropoffs"),
    ]

    for title, why, data, impact in use_cases:
        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = title
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 51, 102)

        for detail in [why, data, impact]:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"  {detail}"
            run.font.size = Pt(7)


def _build_roadmap_slide(slide):
    """Build the Implementation Roadmap slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx in (0, 1):
                shape.text = ""

    left = Emu(400000)
    top = Emu(850000)
    width = Emu(8300000)
    height = Emu(3900000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Title
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Product Management Use Cases — Implementation Roadmap"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    p = tf.add_paragraph()
    p.text = ""

    # Phase 1
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Phase 1: Quick Wins (0-6 months) — 5 use cases — PKR 40-80M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 102, 51)

    phase1 = [
        "Card Spend Stimulation, Risk Scoring Improvement, Credit Models for New-to-Lending,",
        "Competitor Analysis, Omni-Channel Sales Understanding",
        "Data: FSDM + ECIB + card transaction data",
        "Quick Win: 15% increase in card activation rate, 10% expansion of lending eligibility",
    ]
    for line in phase1:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Phase 2
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Phase 2: Product Intelligence (6-18 months) — 5 use cases — PKR 80-150M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 153)

    phase2 = [
        "Merchant Acquiring Insights, Individualized Merchant Recommendation,",
        "Behavioral-Based Pricing, Product Development & Packaging, Multivariate Testing",
        "Data: + merchant data + digital funnel + A/B testing framework",
        "Expected: 20% improvement in card interchange revenue, data-driven product pipeline",
    ]
    for line in phase2:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Phase 3
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Phase 3: Advanced Product Analytics (18-36 months) — 5 use cases — PKR 100-200M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(153, 0, 51)

    phase3 = [
        "Relationship Optimization (Graph), Portfolio Benchmark (Wealth),",
        "Advanced Risk-Based Pricing, + new use cases from product roadmap",
        "Data: + graph analytics + market data feeds",
        "Expected: Dynamic product pricing, predictive product lifecycle management",
    ]
    for line in phase3:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Total
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Total 3-Year Investment: PKR 220-430M | Annual Value: PKR 10-20B through revenue uplift + cost reduction"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)


def _reorder_slides(prs):
    """Reorder slides: Dashboard(17), PK Context(18), Original 0-16, Roadmap(19).

    After adding 3 new slides, indices are:
    0-16: original 17 slides
    17: Dashboard (New A)
    18: Pakistan Context (New B)
    19: Roadmap (New C)

    Target order: [17, 18, 0, 1, 2, ..., 16, 19]
    """
    new_order = [17, 18]  # Dashboard, PK Context
    new_order.extend(range(0, 17))  # All 17 original slides
    new_order.append(19)  # Roadmap at end

    pres_ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    pres_elem = prs.part._element
    sldIdLst = pres_elem.find('.//p:sldIdLst', pres_ns)

    sldId_elements = list(sldIdLst)

    for elem in sldId_elements:
        sldIdLst.remove(elem)

    for idx in new_order:
        sldIdLst.append(sldId_elements[idx])

    print(f"  Reordered {len(new_order)} slides")


if __name__ == '__main__':
    main()
