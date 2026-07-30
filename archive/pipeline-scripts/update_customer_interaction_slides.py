#!/usr/bin/env python3
"""
Update 05_Customer_Interaction_Management.pptx with Pakistan banking context.
Enriches all 30 slides (15 capabilities × 2 each), adds 3 new slides,
removes Teradata branding, and adds comprehensive speaker notes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy
import os

INPUT_FILE = '/mnt/e/erwin/pptout/05_Customer_Interaction_Management.pptx'
CLEAN_FILE = '/mnt/e/erwin/pptout/05_Customer_Interaction_Management_CLEAN.pptx'
OUTPUT_FILE = '/mnt/e/erwin/pptout/05_Customer_Interaction_Management_UPDATED.pptx'

NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# ─────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────

def set_notes(slide, notes_text):
    """Set speaker notes for a slide, creating notes slide if needed."""
    if not slide.has_notes_slide:
        slide.notes_slide
    notes_tf = slide.notes_slide.notes_text_frame
    while len(notes_tf.paragraphs) > 1:
        p = notes_tf.paragraphs[-1]._p
        p.getparent().remove(p)
    lines = notes_text.strip().split('\n')
    notes_tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        p = notes_tf.add_paragraph()
        p.text = line


def replace_cell_text(cell, new_text):
    """Replace text in a table cell preserving basic formatting."""
    first_para = cell.text_frame.paragraphs[0]
    font_size = None
    font_bold = None
    font_color = None
    if first_para.runs:
        r = first_para.runs[0]
        font_size = r.font.size
        font_bold = r.font.bold
        if r.font.color and r.font.color.type is not None:
            try:
                font_color = r.font.color.rgb
            except:
                pass

    paras = list(cell.text_frame.paragraphs)
    for i in range(len(paras) - 1, 0, -1):
        p = paras[i]._p
        p.getparent().remove(p)

    lines = new_text.strip().split('\n')
    cell.text_frame.paragraphs[0].clear()
    run = cell.text_frame.paragraphs[0].add_run()
    run.text = lines[0]
    if font_size:
        run.font.size = font_size
    if font_bold is not None:
        run.font.bold = font_bold
    if font_color:
        run.font.color.rgb = font_color

    for line in lines[1:]:
        p = cell.text_frame.add_paragraph()
        run = p.add_run()
        run.text = line
        if font_size:
            run.font.size = font_size


def append_cell_text(cell, additional_lines):
    """Append bullet lines to existing cell content."""
    first_para = cell.text_frame.paragraphs[0]
    font_size = None
    if first_para.runs:
        font_size = first_para.runs[0].font.size
    for line in additional_lines:
        p = cell.text_frame.add_paragraph()
        run = p.add_run()
        run.text = line
        if font_size:
            run.font.size = font_size


def remove_smart_text(slide):
    """Remove 'smart' orphan text from GroupShapes on a slide."""
    for shape in slide.shapes:
        if type(shape).__name__ == 'GroupShape':
            for t_elem in shape._element.findall('.//a:t', NSMAP):
                if t_elem.text and t_elem.text.strip().lower() == 'smart':
                    t_elem.text = ''


def remove_maturity_engagement_text(slide):
    """Remove 'For use in Maturity Assessment & Roadmap Engagements' from GroupShapes."""
    for shape in slide.shapes:
        if type(shape).__name__ == 'GroupShape':
            for t_elem in shape._element.findall('.//a:t', NSMAP):
                if t_elem.text and 'Maturity Assessment' in t_elem.text:
                    t_elem.text = ''
                if t_elem.text and 'Engagements' in t_elem.text:
                    t_elem.text = ''


def get_table(slide):
    """Get the first table shape from a slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'has_table') and shape.has_table:
            return shape.table
    return None


def replace_h1_2018(slide):
    """Replace 'H1 2018' with 'H1 2025 — H2 2026' in maturity table."""
    tbl = get_table(slide)
    if tbl:
        cell = tbl.cell(7, 5)
        if 'H1 2018' in cell.text:
            replace_cell_text(cell, 'H1 2025 —\nH2 2026')


# ─────────────────────────────────────────────
# Pakistan context to add to capability detail table cells (odd slides)
# Each entry: (slide_index, [business_obj_additions], [data_solution_additions], [outcome_additions])
# ─────────────────────────────────────────────

CAPABILITY_ENRICHMENTS = {
    # Slide 1: Communication Targeting
    0: {
        'biz_obj': [
            "Build targeting across Pakistan's channel mix: SMS (95% reach), push, WhatsApp, in-app, branch, agent",
            "Define opt-in/opt-out rules per SBP Consumer Protection guidelines",
        ],
        'data_sol': [
            "Pre-approved segments: salary holders, RAAST-active, dormant, Islamic-pref, diaspora, youth, HNW",
            "CNIC-based household targeting for family product bundles",
        ],
        'outcome': [
            "Urdu and English message variants for all targeted communications",
            "Compliant targeting via SBP consent management",
        ],
    },
    # Slide 3: Contact Optimization
    2: {
        'biz_obj': [
            "Solve Pakistan's SMS fatigue — customers receive 5-10 promotional SMS daily across banks",
            "Channel-specific frequency limits: SMS max 2/week, push max 3/week, WhatsApp max 1/week",
        ],
        'data_sol': [
            "Longer suppression after customer complaint (SBP Consumer Protection requirement)",
            "Differentiate contact rules for Ramadan (higher spiritual content, lower hard-sell)",
        ],
        'outcome': [
            "Coordinate outbound marketing with SBP-mandated service notifications (OTP, alerts)",
            "30% reduction in SMS opt-outs via contact policy framework",
        ],
    },
    # Slide 5: Response Optimization
    4: {
        'biz_obj': [
            "Optimize response across mobile-first audience — 85% of digital banking via smartphone",
            "Test Urdu vs English vs bilingual variants — response rates differ 30%+ by segment",
        ],
        'data_sol': [
            "Optimize timing: salary-day (1st-5th), Eid-season, Ramadan, back-to-school (August)",
            "Measure response beyond click — track through to product purchase/activation",
        ],
        'outcome': [
            "Build A/B testing culture — most Pakistan banks do zero controlled testing",
            "2x response rates via channel, timing, and language optimization",
        ],
    },
    # Slide 7: Cross-Channel Customer Experience
    6: {
        'biz_obj': [
            "Seamless experience across 7+ channels: branch (16K+), ATM, mobile, internet, CC, USSD, agent (500K+)",
            "Enable branch-to-digital handoff: start application in branch, complete on mobile app",
        ],
        'data_sol': [
            "Integrate RAAST/IBFT as interaction channels — they generate customer events",
            "Design for Pakistan's digital divide: full digital urban, assisted Tier 2/3, agent-led rural",
        ],
        'outcome': [
            "Unified interaction view across all Pakistan banking channels",
            "Cross-channel continuity: agent never asks customer to repeat information",
        ],
    },
    # Slide 9: Contextual Decisioning
    8: {
        'biz_obj': [
            "Use RAAST/IBFT transaction context to trigger real-time offers (salary > loan within 1 hour)",
            "Mobile app context: browsing home finance page triggers RM callback or WhatsApp message",
        ],
        'data_sol': [
            "Branch context: customer queuing triggers relevant product on branch display",
            "Islamic calendar context: Ramadan triggers Zakat calculator + Islamic investment offers",
        ],
        'outcome': [
            "Location-based: near branch > walk-in notification, at mall > POS cashback offer",
            "Real-time contextual offers vs current batch-driven, context-free SMS blasts",
        ],
    },
    # Slide 11: Call/Contact Center Optimization
    10: {
        'biz_obj': [
            "Optimize Urdu/English bilingual contact center (some banks add Sindhi, Pashto, Punjabi IVR)",
            "Provide agents with customer 360 view: RAAST/IBFT transactions, app activity, branch visits",
        ],
        'data_sol': [
            "Implement sentiment analysis on call recordings — detect frustrated customers in real-time",
            "Pre-load customer context when CNIC/mobile number identified via IVR",
        ],
        'outcome': [
            "Route Islamic banking queries to Shariah-trained agents",
            "Reduce average handle time by 2-3 minutes via pre-loaded customer context",
        ],
    },
    # Slide 13: Digital Optimization
    12: {
        'biz_obj': [
            "Optimize mobile banking app for Pakistan — 85% Android, low-bandwidth conditions",
            "Track and fix digital onboarding dropoff (40-60% at NADRA biometric step)",
        ],
        'data_sol': [
            "A/B test Urdu vs English app interface — significant impact on feature adoption",
            "Optimize RAAST/IBFT transfer flow: reduce taps from 5 to 3",
        ],
        'outcome': [
            "Monitor internet banking low-adoption features and redesign or remove",
            "App load time <3s on 3G; key action in <5 taps",
        ],
    },
    # Slide 15: Search Engine Optimization
    14: {
        'biz_obj': [
            "Optimize for Urdu + English bilingual search — Pakistan is 60% English, 40% Urdu/Roman Urdu",
            "Target product search terms: 'bank account kholna', 'car loan Pakistan', 'best rate Pakistan'",
        ],
        'data_sol': [
            "Optimize mobile web experience — 85% of Pakistan internet access is mobile",
            "Target financial literacy searches: 'RAAST kya hai', 'Zakat calculator', 'SBP rate today'",
        ],
        'outcome': [
            "Capture organic acquisition at zero marginal cost — first bank to invest wins",
            "Islamic banking keyword capture: halal investment, Shariah-compliant savings",
        ],
    },
    # Slide 17: Personalization
    16: {
        'biz_obj': [
            "Move from 'Dear Customer' SMS blasts to personalized offers based on individual behavior",
            "Personalize by Islamic preference — detected from transaction patterns (Zakat, Islamic FD)",
        ],
        'data_sol': [
            "Language personalization: Urdu for mass market, English for corporate/affluent, bilingual hybrid",
            "Life-stage personalization: youth/student, working professional, family, retired, diaspora",
        ],
        'outcome': [
            "Personalize digital experience — relevant products on mobile app home screen",
            "4-level personalization ladder: name > segment > behavioral > individual real-time",
        ],
    },
    # Slide 19: Next Best Action Arbitration
    18: {
        'biz_obj': [
            "NBA arbitrates: cross-sell, retention, service recovery, risk alerts, compliance notifications",
            "NBA must consider Islamic preference — never recommend conventional to Islamic-pref customer",
        ],
        'data_sol': [
            "Prioritize by: customer value, propensity, recency of last contact, channel context, urgency",
            "Deploy NBA at: branch counter, mobile app home, contact center agent screen, ATM message",
        ],
        'outcome': [
            "Branch cross-sell: from <5% to 15-20% response with data-driven recommendations",
            "Contact center: from zero cross-sell to 5-8% conversion on inbound calls",
        ],
    },
    # Slide 21: Product Recommendation
    20: {
        'biz_obj': [
            "Pakistan product affinity: CASA > Debit > Credit Card > Personal Loan > Auto > Home > Insurance > Investment",
            "Factor in Islamic product alternatives for every conventional recommendation",
        ],
        'data_sol': [
            "Leverage transaction classification for relevant recs (frequent remitter > forex card, travel insurance)",
            "Recommend government savings products (NSS, PIBs) for conservative/senior segments",
        ],
        'outcome': [
            "Islamic alternative in every recommendation (conventional + Shariah option)",
            "Products-per-customer from 1.8 to 3.0+ via data-driven recommendations",
        ],
    },
    # Slide 23: Multi-Step Campaigns
    22: {
        'biz_obj': [
            "Islamic-calendar-aware sequences: Ramadan awareness, Zakat, Islamic investment, Eid greeting",
            "Salary-day sequence: Day 1 savings nudge, Day 7 investment, Day 14 lending, Day 28 reminder",
        ],
        'data_sol': [
            "Onboarding sequence: welcome > app download > first txn > card activation > cross-sell > 90-day check",
            "Design sequences that work across SMS to WhatsApp to in-app to branch handoff",
        ],
        'outcome': [
            "Seasonal campaign lift: 3-5x response during Ramadan/Eid windows",
            "Onboarding completion from 30% to 60%+ via automated drip sequences",
        ],
    },
    # Slide 25: Marketing Effectiveness
    24: {
        'biz_obj': [
            "Implement control groups for ALL campaigns — currently rare in Pakistan banking",
            "Measure incremental value, not just response count — additional revenue vs control",
        ],
        'data_sol': [
            "Track beyond 7-day window — Pakistan customers take 14-30 days to act on financial offers",
            "Compare channel effectiveness: cost-per-conversion across SMS, WhatsApp, in-app, branch, agent",
        ],
        'outcome': [
            "Build CMO dashboard: campaign ROI, acquisition cost by channel, response rates by segment",
            "Data-driven budget allocation vs current gut-feel marketing spend decisions",
        ],
    },
    # Slide 27: Marketing Attribution
    26: {
        'biz_obj': [
            "Attribution across Pakistan's mix: ATL (TV, billboards) + digital + direct (SMS, branch) + earned (referral)",
            "Pakistan-specific: attribute remittance-driven account openings to diaspora word-of-mouth",
        ],
        'data_sol': [
            "Track offline-to-online: billboard > Google search > website > app download > account opening",
            "Measure agent network effectiveness — which agents drive activations, not just openings",
        ],
        'outcome': [
            "UTM tracking + unique promo codes + RAAST referral tracking per channel/campaign",
            "Multi-touch attribution model starting with digital channels",
        ],
    },
    # Slide 29: Brand Analytics
    28: {
        'biz_obj': [
            "Measure brand health vs fintech challengers — SadaPay, NayaPay have stronger millennial perception",
            "Track brand perception: trust, digital innovation, service quality, Islamic credentials, accessibility",
        ],
        'data_sol': [
            "Monitor social media brand mentions: Twitter/X, Facebook, Instagram, TikTok (growing bank presence)",
            "Measure brand impact of CSR activities — education, disaster relief, financial literacy",
        ],
        'outcome': [
            "Track employer brand strength — important for salary account acquisition",
            "Close NPS gap: Pakistan bank avg 20-30 vs fintech 55-70",
        ],
    },
}


# ─────────────────────────────────────────────
# Maturity table enrichments (even slides)
# Key areas for improvement [7,3] and benefits [7,4]
# ─────────────────────────────────────────────

MATURITY_ENRICHMENTS = {
    # Slide 2: Communication Targeting maturity
    1: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Move from untargeted bulk SMS to segment-based targeting\n"
            "Build SBP opt-in/opt-out compliant targeting framework\n"
            "CNIC-based household segmentation for family product bundles\n"
            "Seasonal segment library: Ramadan, Eid, Hajj, tax season\n"
            "WhatsApp Business API integration for high-engagement targeting"
        ),
        'benefits': (
            "2-3x campaign response rate from targeted vs broadcast\n"
            "30% reduction in SMS costs via precision targeting\n"
            "SBP compliance for consumer protection opt-in rules\n"
            "Higher ROI from seasonal campaign segments"
        ),
    },
    # Slide 4: Contact Optimization maturity
    3: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Unified contact policy across SMS, call, push, WhatsApp, branch, agent\n"
            "Frequency caps: SMS max 2/week, push max 3/week, WhatsApp max 1/week\n"
            "Time-gating: no promos before 9AM/after 9PM; Ramadan sehri/iftar blocks\n"
            "Comply with PTA DND registry — auto-suppress registered numbers\n"
            "30-day marketing blackout after any customer complaint (SBP requirement)"
        ),
        'benefits': (
            "30% reduction in contact fatigue complaints\n"
            "SBP consumer protection compliance + PTA DND registry compliance\n"
            "Higher engagement: customers receive fewer, more relevant messages\n"
            "Reduced SMS cost from elimination of wasted contacts"
        ),
    },
    # Slide 6: Response Optimization maturity
    5: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Build A/B testing culture — most Pakistan banks do zero controlled testing\n"
            "Test Urdu vs English vs bilingual variants — 30%+ response difference\n"
            "Optimize timing: salary-day (1st-5th), Ramadan, Eid, back-to-school\n"
            "Measure beyond click — track through to product purchase/activation\n"
            "2-week measurement window minimum (Pakistan customers take longer to act)"
        ),
        'benefits': (
            "2x response rates from optimized content and timing\n"
            "WhatsApp 3-5x response vs SMS for same offer\n"
            "40% higher response with Urdu SMS in mass market segment\n"
            "Data-driven campaign timing for Pakistan-specific patterns"
        ),
    },
    # Slide 8: Cross-Channel Customer Experience maturity
    7: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Unified interaction view across 7 channels (branch, ATM, mobile, USSD, CC, agent, WhatsApp)\n"
            "Branch-to-digital handoff: start in branch, complete on mobile\n"
            "USSD optimization for 30-40% feature phone users\n"
            "Agent-to-bank transition path for financial inclusion customers\n"
            "WhatsApp conversational banking integration"
        ),
        'benefits': (
            "Eliminate channel silos: RM sees full digital + branch history\n"
            "Increased digital adoption through smooth branch handoff\n"
            "Financial inclusion retention: agent > branch > digital journey\n"
            "Reduced service cost via WhatsApp deflection of routine queries"
        ),
    },
    # Slide 10: Contextual Decisioning maturity
    9: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Salary credit event trigger for instant pre-approved offers\n"
            "Ramadan contextual engine: Islamic product promotion, post-iftar timing\n"
            "Branch walk-in context: RM tablet with 360° view + recommended action\n"
            "RAAST/IBFT large transfer triggers: wealth management or retention follow-up\n"
            "Post-complaint suppression rules (14-day marketing blackout)"
        ),
        'benefits': (
            "10-15% offer response when contextually triggered vs 1-3% batch\n"
            "Salary-day offers: highest conversion window in Pakistan banking\n"
            "Branch RM productivity: data-driven recommendations vs intuition\n"
            "Customer satisfaction from relevant, timely contextual offers"
        ),
    },
    # Slide 12: Call/Contact Center Optimization maturity
    11: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Urdu IVR redesign: 3-level menu, natural language, voice biometric auth\n"
            "WhatsApp chatbot for Tier 1 queries: balance, mini-statement, branch locator\n"
            "Urdu NLP for sentiment analysis on call recordings\n"
            "Agent assist: real-time screen pop with customer 360° and recommendations\n"
            "SBP complaint SLA tracking: 2-day ack, 15-day resolution"
        ),
        'benefits': (
            "20-30% call volume reduction via WhatsApp deflection\n"
            "First call resolution: from 55-65% to 75-85%\n"
            "5-8% cross-sell conversion on inbound calls with NBA\n"
            "SBP complaint compliance and reduced regulatory risk"
        ),
    },
    # Slide 14: Digital Optimization maturity
    13: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Urdu-first mobile UI (not translated English) for 65%+ of market\n"
            "Low-bandwidth optimization: <30MB app, works on 3G/Edge\n"
            "3 persona UX tracks: Digital Native (18-30), Adopter (30-50), Hesitant (50+)\n"
            "USSD menu streamlining for feature phone banking\n"
            "Session-to-conversion optimization: reduce 80% browse, <5% complete gap"
        ),
        'benefits': (
            "3x digital adoption with Urdu-first interface\n"
            "Competitive parity with SadaPay/NayaPay polished UX\n"
            "Rural/Tier 3 coverage via USSD channel optimization\n"
            "Revenue uplift from improved conversion funnel"
        ),
    },
    # Slide 16: Search Engine Optimization maturity
    15: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Bilingual SEO: Google.pk in Urdu + English content\n"
            "Islamic banking keyword strategy: halal, Shariah-compliant, Islamic finance\n"
            "Google My Business for 16,000+ branch locations (bank near me)\n"
            "Diaspora SEO: Roshan Digital Account, NRP banking, send money to Pakistan\n"
            "Seasonal SEO: Hajj savings, Ramadan banking, Eid loans"
        ),
        'benefits': (
            "Capture 30%+ YoY growing banking search demand\n"
            "First-mover advantage in Urdu banking content SEO\n"
            "Local branch discovery: high-intent 'bank near me' searches\n"
            "Diaspora account acquisition via targeted search keywords"
        ),
    },
    # Slide 18: Personalization maturity
    17: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Move from broadcast to 5-layer personalization: name > segment > behavioral > context > predictive\n"
            "Islamic preference personalization: Shariah-compliant product recommendations\n"
            "Language personalization: Urdu mass, English HNW, bilingual youth\n"
            "Life-stage personalization: student > professional > home buyer > retiree\n"
            "Islamic calendar-aware content: Ramadan, Eid, Shab-e-Qadr"
        ),
        'benefits': (
            "3-5x lift from segment-based vs broadcast marketing\n"
            "15-20% higher response with Urdu personalization in Tier 2/3\n"
            "Islamic preference routing captures Shariah-conscious market share\n"
            "Life-stage relevance increases product adoption rates"
        ),
    },
    # Slide 20: Next Best Action Arbitration maturity
    19: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "NBA deployment at branch (RM tablet), mobile app, contact center\n"
            "Arbitration engine: value x propensity x offer value x business priority x contact policy\n"
            "Islamic product alternative in every arbitration set\n"
            "Real-time NBA for RAAST/IBFT transaction triggers\n"
            "SBP cooling-off period and suitability constraints in arbitration"
        ),
        'benefits': (
            "Branch cross-sell: <5% to 15-20% response with data-driven NBA\n"
            "Contact center: zero to 5-8% cross-sell on inbound calls\n"
            "Mobile app: <1% banner click to 5-10% personalized engagement\n"
            "Holistic customer treatment: right offer, right channel, right time"
        ),
    },
    # Slide 22: Product Recommendation maturity
    21: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Product affinity pipeline: CASA > debit > credit card > loan > auto > home > insurance\n"
            "Salary-triggered recommendations within 24 hours of salary credit\n"
            "Life-event detection: marriage, child, Hajj intent, retirement triggers\n"
            "Diaspora pipeline: remittance > Roshan Digital > NPC > property\n"
            "Islamic alternative in every product recommendation set"
        ),
        'benefits': (
            "Higher product holdings per customer (from 1.5 to 3+ products)\n"
            "Salary-day conversion: highest response window in Pakistan banking\n"
            "Islamic product capture: 40%+ customers prefer Shariah-compliant\n"
            "Diaspora monetization: high-value Roshan Digital products"
        ),
    },
    # Slide 24: Multi-Step Campaigns maturity
    23: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Ramadan 30-day campaign sequence: Zakat > invest > Sadaqah > Eid\n"
            "Salary-day monthly sequence: savings nudge > bill pay > invest > lending\n"
            "90-day onboarding drip: welcome > app > first RAAST > card > cross-sell\n"
            "Win-back campaign: dormancy > re-engage > RM call > offer > notice\n"
            "Cross-channel sequences: SMS > WhatsApp > In-app > Branch handoff"
        ),
        'benefits': (
            "3-5x campaign lift during Ramadan/Eid seasonal windows\n"
            "60% onboarding completion vs 30% without drip sequence\n"
            "Reduced dormancy through automated win-back sequences\n"
            "Calendar-driven revenue from Islamic seasonal opportunities"
        ),
    },
    # Slide 26: Marketing Effectiveness maturity
    25: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Implement control groups for ALL campaigns — currently rare in Pakistan\n"
            "Measure incremental value vs control, not just response count\n"
            "Track beyond 7-day window — Pakistan customers take 14-30 days to act\n"
            "Channel cost-per-conversion: SMS vs WhatsApp vs in-app vs branch vs agent\n"
            "CMO dashboard: campaign ROI, acquisition cost, response rates by segment"
        ),
        'benefits': (
            "5-10x target ROI vs current unmeasured 2-3x industry average\n"
            "Budget reallocation from low-ROI to high-ROI channels\n"
            "Answer basic questions most Pakistan banks cannot answer today\n"
            "Data-driven marketing budget justification to board"
        ),
    },
    # Slide 28: Marketing Attribution maturity
    27: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Build attribution across: ATL (TV 35%, billboards 20%, print 15%), digital (15%), direct (10%), agent (5%)\n"
            "UTM tracking on all digital campaigns + unique promo codes per channel\n"
            "Track offline-to-online: billboard > Google search > website > app > account\n"
            "Pakistan-specific: attribute diaspora openings to word-of-mouth\n"
            "Measure agent network effectiveness — activations, not just account openings"
        ),
        'benefits': (
            "Accurate channel ROI: move from last-click to multi-touch model\n"
            "Fair attribution across branch, digital, and agent channels\n"
            "Pakistan marketing mix optimization with spend-to-conversion tracking\n"
            "12+ months cross-channel data enables ML-based attribution"
        ),
    },
    # Slide 30: Brand Analytics maturity
    29: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Brand health tracking: Trust, Innovation, Service, Islamic, Digital dimensions\n"
            "Benchmark vs fintechs (SadaPay, NayaPay) and Islamic banks (Meezan, BankIslami)\n"
            "Social media monitoring: Twitter/X, Facebook, Instagram, TikTok (growing)\n"
            "App store rating management (Google Play reviews)\n"
            "Track employer brand strength for salary account acquisition"
        ),
        'benefits': (
            "Close NPS gap: bank avg 20-30 vs fintech 55-70\n"
            "Real-time brand crisis detection via social monitoring\n"
            "Competitive brand positioning vs Big 5, Islamic, fintech, digital banks\n"
            "Talent attraction and salary account wins from employer brand"
        ),
    },
}


# ─────────────────────────────────────────────
# Speaker notes for all 33 slides (30 existing + 3 new)
# ─────────────────────────────────────────────

SPEAKER_NOTES = {
    # Slide 1: Communication Targeting (detail)
    0: """COMMUNICATION TARGETING — Capability Detail

PAKISTAN CHANNEL TARGETING MATRIX:
| Segment | Primary Channel | Secondary | Tertiary |
| Youth (<30) | Push notification / In-app | WhatsApp | SMS |
| Mass Market | SMS | Agent outreach | Branch |
| Affluent | WhatsApp / RM call | In-app | Email |
| HNW/Private | RM personal call | WhatsApp | Branch invitation |
| Diaspora | Email / WhatsApp | In-app (RDA) | SMS |
| Rural/Agent | Agent network | SMS (USSD) | Branch |
| Islamic-pref | SMS (Urdu) | Branch (Islamic) | WhatsApp |
| Corporate | Email | RM call | Branch visit |

Key Pakistan regulations:
- SBP requires explicit opt-in for marketing communications
- PTA restricts bulk SMS sending — DND (Do Not Disturb) registry
- PEMRA governs above-the-line advertising content
- Data privacy considerations (no formal law yet, but SECP guidelines)

FSDM: CMPGN (Campaign), CMPGN_TRGTS (Campaign Targets), SGMNT (Segment), ELIG (Eligibility), CHNL (Channel), OPT_IN (Opt-In), CNTCT_PREF (Contact Preference)""",

    # Slide 2: Communication Targeting (maturity)
    1: """COMMUNICATION TARGETING — Maturity Assessment

Current State (Developing): Pakistan banks use basic segments — product ownership, channel registration, standard exclusions. Most campaigns target broad customer groups with minimal analytics.

Target State (Innovating): Rich target segments built from offline + advanced analytics. Event-based and predictive targeting. Reusable segment library for rapid campaign deployment.

Pakistan Maturity Context:
- Leading (0% of PK banks): ML-based real-time micro-targeting, self-optimizing segments
- Innovating (0%): Advanced analytics-based segments, event triggers
- Practicing (<5%): Comprehensive segments with change control processes
- Developing (~15%): Basic product/channel segments for most campaigns
- Emerging (~80%): Each campaign built from scratch, no segment reuse

Key Gaps: No reusable segment library, no CNIC-based household targeting, no seasonal segment automation, limited analytics resources for target optimization.

FSDM: CMPGN (Campaign), CMPGN_TRGTS (Campaign Targets), SGMNT (Segment), ELIG (Eligibility), CHNL (Channel), OPT_IN (Opt-In), CNTCT_PREF (Contact Preference)""",

    # Slide 3: Contact Optimization (detail)
    2: """CONTACT OPTIMIZATION — Capability Detail

PAKISTAN CONTACT FATIGUE CRISIS:
Average Pakistan bank customer receives:
- 3-5 promotional SMS per day from their own bank
- 2-3 promotional SMS from other banks (data sharing via third parties)
- 5+ SMS from telcos, retailers, and service providers
Total: 10-15+ promotional SMS daily = severe fatigue

Result: <2% response rates on SMS (was 5-8% five years ago), rising opt-out rates, customer complaints to SBP.

Contact policy framework priorities for Pakistan:
1. FREQUENCY CAP: Max 2 promotional SMS per week per customer (not per campaign)
2. RECENCY RULE: Min 48-hour gap between promotional contacts to same customer
3. COMPLAINT SUPPRESSION: 30-day marketing blackout after any complaint
4. CHANNEL PREFERENCE: Respect stated preference (many customers prefer WhatsApp over SMS)
5. RAMADAN RULES: Adjust frequency and timing — no hard-sell during Sehri/Iftar, focus on Zakat/Islamic products
6. DO-NOT-DISTURB: Comply with PTA DND registry — auto-suppress registered numbers
7. TRANSACTIONAL vs. PROMOTIONAL: Never mix promotional content in transactional notifications (SBP regulation)

FSDM: CNTCT_PLCY (Contact Policy), CNTCT_FREQ (Contact Frequency), CNTCT_RCNCY (Contact Recency), SPRSN (Suppression), OPT_OUT (Opt-Out), CHNL_PREF (Channel Preference), CNTCT_HSTRY (Contact History)""",

    # Slide 4: Contact Optimization (maturity)
    3: """CONTACT OPTIMIZATION — Maturity Assessment

Current State (Developing): Basic frequency caps exist at some banks but enforcement is inconsistent. No Ramadan/Friday timing rules. No unified contact policy across channels.

Target State (Innovating): Enterprise contact policy engine governing all channels with cultural/regulatory timing rules, customer preference management, and SBP compliance.

Pakistan Maturity Context:
- Developing (~15%): Basic SMS frequency caps, product-level limits
- Emerging (~85%): No contact policy, unlimited promotional SMS, no time-gating

Critical Gap: Pakistan banks' lack of contact governance causes massive customer fatigue, leading to SMS blocking that impacts even critical transactional messages (OTP, fraud alerts).

FSDM: CNTCT_PLCY (Contact Policy), CNTCT_FREQ (Contact Frequency), CNTCT_RCNCY (Contact Recency), SPRSN (Suppression), OPT_OUT (Opt-Out), CHNL_PREF (Channel Preference), CNTCT_HSTRY (Contact History)""",

    # Slide 5: Response Optimization (detail)
    4: """RESPONSE OPTIMIZATION — Capability Detail

PAKISTAN RESPONSE OPTIMIZATION:
Current state: Most Pakistan banks measure only "SMS delivered" and "call completed" — not actual response or conversion.

Optimization opportunities:
1. TIMING: Salary credit day (1st-5th) generates 3x response for lending products. Ramadan Zakat period generates 5x response for Islamic products. Friday afternoon generates lowest response (prayer time).
2. LANGUAGE: Urdu SMS generates 40% higher response in mass market segment; English better for urban affluent.
3. CHANNEL: WhatsApp generates 3-5x response vs. SMS for same offer (due to richer content + less fatigue).
4. OFFER TYPE: Cashback offers outperform interest rate offers 2:1 in Pakistan's price-sensitive market.
5. PERSONALIZATION: Using customer name + product they hold + specific benefit generates 2x vs. generic.

A/B testing framework needed:
- Test one variable at a time (channel, timing, language, offer, creative)
- Minimum 10K control group per test
- Measure through to activation, not just response
- 2-week measurement window minimum (Pakistani customers take longer to act)
- Document learnings in campaign knowledge base

FSDM: CMPGN_RSPNS (Campaign Response), CNVRSN (Conversion), A_B_TST (A/B Test), CMPGN_KPI (Campaign KPI), RSP_RT (Response Rate), CMPGN_RSLT (Campaign Result)""",

    # Slide 6: Response Optimization (maturity)
    5: """RESPONSE OPTIMIZATION — Maturity Assessment

Current State (Developing): Most Pakistan banks measure only "SMS delivered" — not actual response or conversion. No A/B testing culture. No systematic optimization.

Target State (Innovating): Continuous A/B testing with automated optimization. Response analysis by segment, channel, timing, and language. 2-week measurement windows.

Pakistan Benchmarks: SMS response 1-3% (industry), WhatsApp 3-5x higher for same offer. Massive optimization opportunity from untested campaigns.

FSDM: CMPGN_RSPNS (Campaign Response), CNVRSN (Conversion), A_B_TST (A/B Test), CMPGN_KPI (Campaign KPI), RSP_RT (Response Rate), CMPGN_RSLT (Campaign Result)""",

    # Slide 7: Cross-Channel Customer Experience (detail)
    6: """CROSS-CHANNEL CUSTOMER EXPERIENCE — Capability Detail

PAKISTAN CHANNEL ECOSYSTEM:
| Channel | Reach | Cost/Interaction | Trend |
| Branch (16,000+) | Universal | PKR 150-300 | Stable (declining share) |
| ATM (16,000+) | 60M+ cardholders | PKR 30-50 | Stable |
| Mobile App | 30M+ downloads | PKR 5-15 | Rapidly growing (40% YoY) |
| Internet Banking | 15M+ users | PKR 10-20 | Growing |
| Contact Center | Universal | PKR 80-150 per call | Stable |
| USSD | 100M+ mobile users | PKR 2-5 | Stable (feature phones) |
| Agent Network | 500,000+ agents | PKR 20-40 | Growing |
| WhatsApp Business | Emerging | PKR 3-8 | New — high potential |
| RAAST/IBFT | 50M+ monthly | PKR 2-5 | Explosive growth |

Cross-channel challenge: Most Pakistan banks have NO unified interaction history. A customer who calls the contact center about a failed mobile transaction is asked to repeat everything because the agent cannot see the mobile app session.

FSDM: CHNL (Channel), CHNL_INTN (Channel Interaction), CHNL_PRFNC (Channel Performance), CSTMR_CHNL_PREF (Customer Channel Preference), OMNI_CHNL (Omni-Channel), CHNL_INTGRTN (Channel Integration)""",

    # Slide 8: Cross-Channel Customer Experience (maturity)
    7: """CROSS-CHANNEL CUSTOMER EXPERIENCE — Maturity Assessment

Current State (Developing): Channels operate in silos. No unified view. No handoff capability. Agent and branch systems completely disconnected from digital.

Target State (Innovating): Unified interaction view, seamless channel handoff, consistent experience regardless of touchpoint, journey-aware interactions.

Pakistan Reality: 9 distinct channel systems with zero integration. Customer must repeat information at every touchpoint. Branch has no knowledge of digital engagement.

FSDM: CHNL (Channel), CHNL_INTN (Channel Interaction), CHNL_PRFNC (Channel Performance), CSTMR_CHNL_PREF (Customer Channel Preference), OMNI_CHNL (Omni-Channel), CHNL_INTGRTN (Channel Integration)""",

    # Slide 9: Contextual Decisioning (detail)
    8: """CONTEXTUAL DECISIONING — Capability Detail

PAKISTAN CONTEXTUAL OPPORTUNITIES:
Real-time context signals available in Pakistan banking:
1. TRANSACTION: Salary credit, large deposit, large withdrawal, bill payment, international remittance
2. APP BEHAVIOR: Product page browse, calculator use, loan eligibility check, FD rate comparison
3. LOCATION: Branch proximity (GPS), merchant proximity, airport/travel detection
4. TIME: Salary day (1st-5th), bill payment day, Ramadan/Eid, weekend vs. weekday
5. LIFECYCLE: Just opened account, just activated card, just made first RAAST payment
6. RISK: Failed transaction, login from new device, suspected fraud alert

Current Pakistan reality: Almost zero contextual decisioning deployed. Offers are pre-generated in batch and pushed regardless of customer context. No bank is doing real-time NBA in Pakistan as of 2025.

FSDM: CNTXT (Context), DCSNNG (Decisioning), RL_TM_EVNT (Real-Time Event), OFFR_ARBTTN (Offer Arbitration), TRGR (Trigger), CNTXT_DCSNNG (Contextual Decisioning)""",

    # Slide 10: Contextual Decisioning (maturity)
    9: """CONTEXTUAL DECISIONING — Maturity Assessment

Current State (Developing): Batch-driven campaigns with no contextual awareness. No event triggers. No real-time decisioning.

Target State (Innovating): Event-triggered contextual offers across multiple touchpoints with real-time decisioning and personalization.

Pakistan Gap: Almost zero contextual decisioning deployed. Salary credit, RAAST, and branch walk-in contexts are completely untapped.

FSDM: CNTXT (Context), DCSNNG (Decisioning), RL_TM_EVNT (Real-Time Event), OFFR_ARBTTN (Offer Arbitration), TRGR (Trigger), CNTXT_DCSNNG (Contextual Decisioning)""",

    # Slide 11: Call/Contact Center Optimization (detail)
    10: """CALL/CONTACT CENTER OPTIMIZATION — Capability Detail

PAKISTAN CONTACT CENTER LANDSCAPE:
Top-5 banks handle 1M+ calls/month each. Typical contact center:
- 200-500 agents per large bank
- IVR handles 40-60% of calls (balance inquiry, mini-statement, card block)
- Average handle time: 5-8 minutes (high — agents lack customer context)
- First call resolution: 60-70% (low — requires callback or branch visit)
- Languages: Urdu (70%), English (25%), Regional (5%)

Key optimization opportunities:
1. CUSTOMER 360 ON SCREEN: Show agent last 10 transactions, products held, pending complaints, segment, value tier, last campaign offered — saves 2-3 minutes per call
2. IVR PERSONALIZATION: Greet by name, offer top 3 likely reasons for call based on recent activity
3. SENTIMENT DETECTION: Real-time voice analysis to detect frustration, auto-escalate to supervisor
4. CROSS-SELL AT RESOLUTION: After resolving issue, agent can offer relevant product
5. COMPLAINT ROUTING: Auto-categorize per SBP Consumer Protection categories

FSDM: CNTCT_CNTR (Contact Center), IVR (Interactive Voice Response), AGNT (Agent), CLL (Call), CLL_RSLN (Call Resolution), SNTMNT (Sentiment), CMPLNT (Complaint), ESCLTN (Escalation)""",

    # Slide 12: Call/Contact Center Optimization (maturity)
    11: """CALL/CONTACT CENTER OPTIMIZATION — Maturity Assessment

Current State (Developing): Basic IVR (5+ levels), Urdu quality poor, no WhatsApp integration, no sentiment analysis, manual complaint tracking.

Target State (Innovating): Simplified Urdu IVR, WhatsApp chatbot deflection, NLP sentiment analysis, agent NBA screen pop, automated SBP SLA compliance.

Pakistan Benchmarks: 60-70% FCR vs 75-85% global. 5-8 min AHT vs 3-4 min. Massive efficiency opportunity.

FSDM: CNTCT_CNTR (Contact Center), IVR (Interactive Voice Response), AGNT (Agent), CLL (Call), CLL_RSLN (Call Resolution), SNTMNT (Sentiment), CMPLNT (Complaint), ESCLTN (Escalation)""",

    # Slide 13: Digital Optimization (detail)
    12: """DIGITAL OPTIMIZATION — Capability Detail

PAKISTAN DIGITAL BANKING METRICS:
Mobile App:
- Downloads: 30M+ across top-10 banks
- Monthly Active Users: ~40% of downloaders
- Top features: Balance check (80%), fund transfer (60%), bill payment (40%)
- Lowest adoption: Investment (5%), insurance (3%), loan application (2%)
- OS split: Android 85%, iOS 15%
- Key friction: NADRA biometric verification for high-value transactions

Optimization priorities:
1. SPEED: App load time <3 seconds on 3G (many users in low-bandwidth areas)
2. SIMPLICITY: Reduce steps for top 3 features (transfer, balance, bill pay)
3. LANGUAGE: Toggle between Urdu/English without app restart
4. BIOMETRIC: Fingerprint/face login for faster access
5. ONBOARDING: Fix NADRA biometric dropoff — offer alternatives (video KYC, branch fallback)

FSDM: DGTL_CHNL (Digital Channel), APP_EVNT (App Event), SSSN (Session), PG_VW (Page View), CLCK_STRM (Clickstream), FNNEL (Funnel), DRPFF (Dropoff)""",

    # Slide 14: Digital Optimization (maturity)
    13: """DIGITAL OPTIMIZATION — Maturity Assessment

Current State (Developing): English-only apps, heavy (60-100MB), slow on low-end devices. No USSD optimization. High browse-to-conversion dropoff.

Target State (Innovating): Urdu-first responsive apps, lightweight, 3-persona UX tracks, USSD banking, data-driven conversion optimization.

Pakistan Reality: SadaPay/NayaPay have set new UX benchmark. Traditional banks losing digital-first millennials.

FSDM: DGTL_CHNL (Digital Channel), APP_EVNT (App Event), SSSN (Session), PG_VW (Page View), CLCK_STRM (Clickstream), FNNEL (Funnel), DRPFF (Dropoff)""",

    # Slide 15: Search Engine Optimization (detail)
    14: """SEARCH ENGINE OPTIMIZATION — Capability Detail

PAKISTAN SEO CONTEXT:
Pakistan has 120M+ internet users (mobile-first). Banking product search volume growing 30% YoY.

Top banking search queries in Pakistan:
1. "Best savings account Pakistan" / "Best bank Pakistan"
2. "Car loan calculator" / "Home loan rate Pakistan"
3. "Online account opening Pakistan"
4. "RAAST payment" / "How to use RAAST"
5. "Zakat calculator" / "Islamic banking Pakistan"

SEO is immature in Pakistan banking — most bank websites have poor Urdu content, slow mobile load times, and thin product pages. First bank to invest captures organic acquisition at zero marginal cost.

FSDM: WB_SRCH (Web Search), SRCH_QRY (Search Query), LNDNG_PG (Landing Page), CNVRSN_RT (Conversion Rate), KYWD (Keyword)""",

    # Slide 16: Search Engine Optimization (maturity)
    15: """SEARCH ENGINE OPTIMIZATION — Maturity Assessment

Current State (Developing): English-only SEO, no Urdu content optimization, basic Google My Business for some branches, no Islamic banking keyword strategy.

Target State (Innovating): Bilingual SEO, Islamic banking keyword capture, comprehensive Google My Business for 16K+ branches, diaspora content strategy.

Pakistan Opportunity: 120M+ internet users, banking searches +30% YoY, Urdu search growing rapidly, Islamic banking keywords have high intent and low competition.

FSDM: WB_SRCH (Web Search), SRCH_QRY (Search Query), LNDNG_PG (Landing Page), CNVRSN_RT (Conversion Rate), KYWD (Keyword)""",

    # Slide 17: Personalization (detail)
    16: """PERSONALIZATION — Capability Detail

PAKISTAN PERSONALIZATION MATURITY:
Current state: <5% of Pakistan bank communications are personalized beyond customer name.

Personalization ladder for Pakistan:
Level 1 (Current): Name + product held in SMS
Level 2 (Next): Segment-based offer (youth gets different offer than retiree)
Level 3 (Target): Behavioral — based on actual transaction patterns, channel usage, life events
Level 4 (Aspirational): Individual real-time — dynamic content based on context + history

Quick win variables:
- Salary amount tier determines lending offer amount
- Transaction patterns drive spending category offers
- Islamic product holding triggers Shariah-compliant alternatives
- RAAST usage frequency drives digital-first feature promotions
- Branch visit frequency drives digital migration nudges

FSDM: PRSNLZTN (Personalization), CSTMR_PRFL (Customer Profile), BHVR_SGM (Behavioral Segment), DYNMC_CNTNT (Dynamic Content), OFFR_PRSNLZTN (Offer Personalization)""",

    # Slide 18: Personalization (maturity)
    17: """PERSONALIZATION — Maturity Assessment

Current State (Developing): Name insertion at best. <5% personalized beyond name. No language, Islamic, or value-tier personalization.

Target State (Innovating): 4-level personalization ladder: name > segment > behavioral > real-time individual. Language, Islamic preference, and life-stage aware.

Pakistan Reality: No bank at Level 3+. Massive gap between broadcast SMS reality and 1:1 personalization aspiration.

FSDM: PRSNLZTN (Personalization), CSTMR_PRFL (Customer Profile), BHVR_SGM (Behavioral Segment), DYNMC_CNTNT (Dynamic Content), OFFR_PRSNLZTN (Offer Personalization)""",

    # Slide 19: Next Best Action Arbitration (detail)
    18: """NEXT BEST ACTION ARBITRATION — Capability Detail

PAKISTAN NBA FRAMEWORK:
1. RETAIN (highest priority): If churn risk > threshold, retention offer first
2. GROW (cross-sell): If eligible for next product in affinity chain, cross-sell offer
3. SERVE: If recent complaint or low satisfaction, service recovery action
4. MIGRATE: If branch-heavy, digital adoption nudge
5. COMPLY: If KYC expiring or CDD refresh needed, compliance notification

Pakistan-specific NBA rules:
- NEVER offer conventional lending to Islamic-preference customers
- ALWAYS check ECIB before lending offers
- RESPECT Ramadan timing — spiritual/charity offers preferred
- DIASPORA customers get Roshan Digital Account / NPC offers
- SALARY-DAY timing: Cross-sell within 24 hours of salary credit

Current state: NO Pakistan bank has deployed NBA. All interactions are product-push based on marketing calendar.

FSDM: NBA (Next Best Action), ARBTTN (Arbitration), OFFR_PRTY (Offer Priority), PRPNSTY (Propensity), CSTMR_VALU (Customer Value), DCSNNG_ENGN (Decisioning Engine)""",

    # Slide 20: Next Best Action Arbitration (maturity)
    19: """NEXT BEST ACTION ARBITRATION — Maturity Assessment

Current State (Developing): No real-time NBA. RM uses intuition. Generic banners in app. No cross-sell in contact center.

Target State (Innovating): Real-time NBA across branch, app, contact center with arbitration engine, propensity models, and feedback loop.

Pakistan Reality: Zero banks have NBA capability. All interactions are product-push based on marketing calendar, not data-driven.

FSDM: NBA (Next Best Action), ARBTTN (Arbitration), OFFR_PRTY (Offer Priority), PRPNSTY (Propensity), CSTMR_VALU (Customer Value), DCSNNG_ENGN (Decisioning Engine)""",

    # Slide 21: Product Recommendation (detail)
    20: """PRODUCT RECOMMENDATION — Capability Detail

PAKISTAN PRODUCT AFFINITY MAP:
1. CASA to Debit Card to ATM to Mobile App (universal starter)
2. Salary Account to Credit Card to Personal Loan (employed segment)
3. CASA to Fixed Deposit to Mutual Fund to PIB/T-Bill (savings-focused)
4. CASA to Islamic Savings to Diminishing Musharaka to Takaful (Islamic track)
5. Current Account to Trade Finance to LC/LG to Working Capital (business track)
6. Remittance to RDA to NPC to Property Investment (diaspora track)

FSDM: PRDCT_RCMNDTN (Product Recommendation), PRDCT_AFFNTY (Product Affinity), RCMNDTN_ENGN (Recommendation Engine), PRDCT (Product), PRDCT_BNDL (Product Bundle)""",

    # Slide 22: Product Recommendation (maturity)
    21: """PRODUCT RECOMMENDATION — Maturity Assessment

Current State (Developing): Static product push based on RM intuition. No data-driven recommendations. No trigger-based system.

Target State (Innovating): Automated product recommendation engine with trigger-based rules, affinity pipeline, Islamic alternatives, and diaspora path.

Pakistan Opportunity: Product-per-customer average 1.5 vs 3+ at leading banks. 6 defined affinity tracks (universal, employed, savings, Islamic, business, diaspora).

FSDM: PRDCT_RCMNDTN (Product Recommendation), PRDCT_AFFNTY (Product Affinity), RCMNDTN_ENGN (Recommendation Engine), PRDCT (Product), PRDCT_BNDL (Product Bundle)""",

    # Slide 23: Multi-Step Campaigns (detail)
    22: """MULTI-STEP CAMPAIGNS — Capability Detail

PAKISTAN MULTI-STEP CAMPAIGN EXAMPLES:

SEQUENCE 1: Ramadan Campaign (30 days)
Day 1: Ramadan Mubarak + Zakat calculator link
Day 5: Islamic savings account offer
Day 15: Ramadan spending insights + cashback on groceries
Day 20: Eid shopping pre-approved credit card limit increase
Day 25: Zakat auto-deduction setup reminder
Eid: Eid Mubarak greeting + Eid-ul-Adha savings plan
Post-Eid: Financial health check + investment offer

SEQUENCE 2: Salary Day Campaign (monthly)
Day 0: Salary credited notification
Day 1: Auto-savings nudge (transfer 10% to savings)
Day 3: Bill payment reminders with auto-pay setup
Day 7: Investment opportunity (mutual fund SIP)
Day 14: Personal loan pre-approval (if eligible)
Day 28: Month-end balance summary + spending insights

SEQUENCE 3: New Customer Onboarding (90 days)
Day 0: Welcome + app download link
Day 1: Tutorial video (Urdu/English)
Day 3: First RAAST transfer incentive
Day 7: Debit card activation reminder
Day 14: Internet banking registration nudge
Day 30: First cross-sell (based on behavior so far)
Day 60: Satisfaction survey
Day 90: Loyalty program enrollment

FSDM: MLT_STP_CMPGN (Multi-Step Campaign), CMPGN_STP (Campaign Step), CMPGN_SQNC (Campaign Sequence), TRGR (Trigger), RSPNS (Response), STPNG_LGIC (Stepping Logic)""",

    # Slide 24: Multi-Step Campaigns (maturity)
    23: """MULTI-STEP CAMPAIGNS — Maturity Assessment

Current State (Developing): One-off campaigns. No automated sequences. No calendar-driven campaign library.

Target State (Innovating): Automated multi-step sequences for lifecycle, seasonal, and win-back scenarios with conditional branching.

Pakistan Calendar Advantage: Ramadan (30-day), Salary Day (monthly), and Onboarding (90-day) create natural multi-step campaign opportunities that few banks exploit systematically.

FSDM: MLT_STP_CMPGN (Multi-Step Campaign), CMPGN_STP (Campaign Step), CMPGN_SQNC (Campaign Sequence), TRGR (Trigger), RSPNS (Response), STPNG_LGIC (Stepping Logic)""",

    # Slide 25: Marketing Effectiveness (detail)
    24: """MARKETING EFFECTIVENESS — Capability Detail

PAKISTAN MARKETING EFFECTIVENESS GAP:
Most Pakistan banks cannot answer basic questions:
- "Which of our 50 campaigns last month generated the most revenue?" — Unknown
- "What is our cost per acquisition by channel?" — Estimated at best
- "Should we spend more on SMS or WhatsApp?" — No data to decide
- "Did the Ramadan campaign increase deposits?" — Anecdotal only

Minimum viable marketing effectiveness:
1. CONTROL GROUPS: 10% holdout on every campaign — no exceptions
2. RESPONSE TRACKING: Link campaign ID to product application to activation to 30-day usage
3. ROI CALCULATION: (Incremental revenue from campaign group minus control) / Campaign cost
4. CHANNEL COMPARISON: Same offer, same segment, different channel — A/B test
5. DASHBOARD: Weekly automated report to marketing leadership

FSDM: CMPGN_EFCTVNS (Campaign Effectiveness), ROI (Return on Investment), CNTRL_GRP (Control Group), INCRMNTL_VALU (Incremental Value), CMPGN_KPI (Campaign KPI)""",

    # Slide 26: Marketing Effectiveness (maturity)
    25: """MARKETING EFFECTIVENESS — Maturity Assessment

Current State (Developing): Basic sent/response tracking. No ROI measurement. No control groups. Most banks cannot answer basic effectiveness questions.

Target State (Innovating): Full marketing effectiveness framework: control groups on every campaign, CPA, ROI, incremental lift, LTV impact, channel efficiency, with CMO dashboard.

Pakistan Gap: Marketing budgets (1-2% of revenue) are smallest globally. Control groups and effectiveness measurement would justify budget increases and optimize allocation.

FSDM: CMPGN_EFCTVNS (Campaign Effectiveness), ROI (Return on Investment), CNTRL_GRP (Control Group), INCRMNTL_VALU (Incremental Value), CMPGN_KPI (Campaign KPI)""",

    # Slide 27: Marketing Attribution (detail)
    26: """MARKETING ATTRIBUTION — Capability Detail

PAKISTAN MARKETING MIX:
- TV advertising: 35% of bank marketing spend
- Billboards/outdoor: 20% of spend
- Newspaper/print: 15% of spend
- Digital: 15% of spend (growing fast)
- Direct (SMS, branch): 10% of spend
- Agent incentives: 5% of spend

First steps for Pakistan attribution:
- UTM tracking on all digital campaigns
- Unique promo codes per channel/campaign
- Ask "How did you hear about us?" at account opening
- RAAST referral tracking (refer a friend via RAAST)
- Build multi-touch model starting with digital channels only

FSDM: MKTG_ATRBTN (Marketing Attribution), TCHPNT (Touchpoint), CNVRSN_PTH (Conversion Path), CHNL_CNTRBN (Channel Contribution), ATRBTN_MDL (Attribution Model)""",

    # Slide 28: Marketing Attribution (maturity)
    27: """MARKETING ATTRIBUTION — Maturity Assessment

Current State (Developing): Last-touch attribution only. Branch gets all credit. No multi-touch model. No digital-to-branch tracking.

Target State (Innovating): Multi-touch attribution across ATL (TV 35%, billboards 20%, print 15%), digital (15%), direct (10%), and agent (5%) with UTM tracking and time-decay weighting.

Pakistan Challenge: 6 channel categories, long sales cycles, data silos. Need UTM tracking + promo codes + 12 months cross-channel data before ML attribution is feasible.

FSDM: MKTG_ATRBTN (Marketing Attribution), TCHPNT (Touchpoint), CNVRSN_PTH (Conversion Path), CHNL_CNTRBN (Channel Contribution), ATRBTN_MDL (Attribution Model)""",

    # Slide 29: Brand Analytics (detail)
    28: """BRAND ANALYTICS — Capability Detail

PAKISTAN BRAND LANDSCAPE:
| Bank Type | Trust | Innovation | Service | Islamic | Digital |
| Big 5 (HBL, UBL, MCB, ABL, NBP) | High | Medium | Medium | Medium | Medium |
| Islamic (Meezan, BankIslami) | High | Medium | Medium | Very High | Medium |
| Fintech (SadaPay, NayaPay) | Medium | Very High | High | Low | Very High |
| Digital Banks (new licensees) | Low | High | Unknown | Varies | High |

Brand analytics priorities:
1. NPS TRACKING: Quarterly NPS by segment vs. fintech benchmarks
2. SOCIAL LISTENING: Monitor Urdu + English social media
3. APP STORE RATINGS: Track Google Play rating and reviews
4. EMPLOYER BRAND: Survey HR managers on preferred salary banking partner
5. ISLAMIC BRAND: Measure Shariah compliance perception

FSDM: BRND (Brand), BRND_HLTH (Brand Health), NPS (Net Promoter Score), BRND_PRCP (Brand Perception), SCIAL_MDIA (Social Media), SNTMNT (Sentiment)""",

    # Slide 30: Brand Analytics (maturity)
    29: """BRAND ANALYTICS — Maturity Assessment

Current State (Developing): Annual brand survey only. No real-time social monitoring. No app store analysis. No competitive brand benchmarking.

Target State (Innovating): Real-time brand health dashboard, social sentiment monitoring, competitive benchmarking across Big 5, Islamic, fintech, and digital bank categories, employer brand tracking.

Pakistan Reality: Fintechs (SadaPay, NayaPay) rate Very High on Innovation and Digital. Traditional banks rate High on Trust only. Gap in brand perception is a key competitive vulnerability.

FSDM: BRND (Brand), BRND_HLTH (Brand Health), NPS (Net Promoter Score), BRND_PRCP (Brand Perception), SCIAL_MDIA (Social Media), SNTMNT (Sentiment)""",
}


# Speaker notes for the 3 new slides (indexed as 'new_a', 'new_b', 'new_c')
NEW_SLIDE_NOTES = {
    'new_a': """CUSTOMER INTERACTION MANAGEMENT — DOMAIN DASHBOARD

This slide provides a high-level comparison of Customer Interaction Management maturity across global, regional (South Asia & Middle East), and Pakistan banking markets.

Key Takeaways:
- Pakistan lags significantly in marketing automation (<15% adoption vs 75% global)
- Real-time personalization exists at <5% of Pakistan banks vs 50%+ of digital banks globally
- NBA (Next Best Action) deployed at <3% of Pakistan banks vs 40% of top global banks
- Marketing attribution in place at <5% of Pakistan banks vs 55% globally
- Campaign ROI measurement is rare — most lack control groups

This domain covers 15 BVF sub-capabilities with ~100+ FSDM entities across Campaign, Channel, Interaction, and Contact domains. BACR assessment covers ~80 maturity questions.

Investment Justification: Moving from Developing to Innovating maturity across these 15 capabilities can deliver 3-5x marketing ROI improvement, 2x campaign response rates, and products/customer from 1.8 to 3.0+.

FSDM Domains: Campaign (~100+ entities), Channel (~250+ entities), Interaction/Event (~700 entities)""",

    'new_b': """PAKISTAN — CUSTOMER INTERACTION LANDSCAPE

Channel Reach & Effectiveness:
- SMS: 95%+ reach (mobile penetration), 2-5% response rate, OVERUSED
- Push Notification: 30% reach (app installed), 8-15% open rate, GROWING
- WhatsApp Business: 20% reach, 40-60% open rate, EMERGING — highest engagement
- Email: 15% reach (urban/corporate only), 5-10% open rate, LOW priority
- In-App Messages: 25% reach, 15-25% engagement, HIGHEST conversion
- Branch Face-to-Face: 100% for walk-ins, high conversion, HIGH cost
- Contact Center: Urdu/English bilingual, 1M+ calls/month (top-5 banks), CRITICAL for complaints
- Agent Network: 500,000+ agents, rural reach, ESSENTIAL for financial inclusion

Key Challenges:
1. SMS FATIGUE: Customers receive 5-10 bank SMS daily — regulatory crackdown looming
2. NO CDP: <5% of banks have Customer Data Platform — interactions siloed by channel
3. BATCH-ONLY: Most campaigns run on monthly batch — no real-time personalization
4. NO ATTRIBUTION: Marketing spend not tracked to conversion — budget allocation is gut-feel
5. LANGUAGE: Must support Urdu, English, and increasingly Sindhi, Pashto, Punjabi
6. COMPLIANCE: SBP Consumer Protection requires opt-in for marketing + complaint response SLAs

Opportunity: Pakistan's unique combination of high mobile penetration, growing smartphone adoption, and emerging digital payment infrastructure (RAAST) creates a window for banks to leapfrog to omni-channel marketing.""",

    'new_c': """CUSTOMER INTERACTION MANAGEMENT — IMPLEMENTATION ROADMAP

Phase 1: Foundation (0-6 months) — Investment: PKR 40-70M
- Implement contact policy framework (prevent SMS fatigue)
- Build unified customer contact history across channels
- Launch WhatsApp Business channel for top 20% customers
- A/B testing framework for SMS and push notifications
Quick Win: 30% reduction in SMS opt-outs

Phase 2: Intelligence (6-18 months) — Investment: PKR 100-180M
- Customer Data Platform (CDP) implementation
- Multi-step campaign automation (Ramadan, Eid, salary-day sequences)
- Personalization engine — segment-level then individual-level
- Call center optimization with customer 360 view for agents
Expected: 2x campaign response rates

Phase 3: Real-Time (18-36 months) — Investment: PKR 150-280M
- Next Best Action engine across branch, mobile, contact center
- Real-time contextual decisioning (location, transaction, app behavior)
- Marketing attribution model (digital + branch + ATL integration)
- AI-powered product recommendation at every touchpoint
Expected: 3x marketing ROI, products/customer from 1.8 to 3.0+

Total 3-Year Investment: PKR 290-530M (~USD 1-2M)
Estimated Annual Revenue Impact: PKR 500M-1B from improved marketing effectiveness""",
}


# ─────────────────────────────────────────────
# Main execution
# ─────────────────────────────────────────────

def clean_pptx(input_path, output_path):
    """Remove orphaned slide XMLs from the PPTX ZIP to prevent duplicate name conflicts."""
    import zipfile
    import xml.etree.ElementTree as ET

    with zipfile.ZipFile(input_path, 'r') as zin:
        # Parse presentation.xml to find referenced slide filenames
        pres_xml = zin.read('ppt/presentation.xml')
        root = ET.fromstring(pres_xml)
        pns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
               'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

        # Get all rIds from sldIdLst
        sld_rids = set()
        for sldId in root.findall('.//p:sldIdLst/p:sldId', pns):
            rid = sldId.attrib.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', '')
            sld_rids.add(rid)

        # Parse presentation.xml.rels to find which slide files are referenced
        rels_xml = zin.read('ppt/_rels/presentation.xml.rels')
        rels_root = ET.fromstring(rels_xml)
        referenced_slides = set()
        for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
            rid = rel.attrib.get('Id', '')
            target = rel.attrib.get('Target', '')
            if rid in sld_rids and 'slides/' in target:
                slide_file = target.replace('slides/', '')
                referenced_slides.add(slide_file)

        print(f"  Referenced slides: {len(referenced_slides)}")

        # Build list of orphaned slide files
        orphaned = set()
        for name in zin.namelist():
            if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                slide_file = name.replace('ppt/slides/', '')
                if slide_file not in referenced_slides:
                    orphaned.add(slide_file.replace('.xml', ''))

        print(f"  Orphaned slides to remove: {len(orphaned)}")
        if orphaned:
            print(f"    {sorted(orphaned)}")

        # Write cleaned ZIP
        with zipfile.ZipFile(output_path, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                # Skip orphaned slide files and their rels
                skip = False
                for orphan in orphaned:
                    if f'ppt/slides/{orphan}.xml' == item.filename:
                        skip = True
                    if f'ppt/slides/_rels/{orphan}.xml.rels' == item.filename:
                        skip = True
                    # Also skip orphaned notesSlides if they exist
                    if f'ppt/notesSlides/{orphan}' in item.filename:
                        skip = True
                if not skip:
                    data = zin.read(item.filename)
                    zout.writestr(item, data)

    print(f"  Cleaned file saved to {output_path}")


def main():
    # Step 0: Clean orphaned slides from PPTX
    print("── Step 0: Cleaning orphaned slides ──")
    clean_pptx(INPUT_FILE, CLEAN_FILE)

    print("\nLoading presentation...")
    prs = Presentation(CLEAN_FILE)
    print(f"  Loaded {len(prs.slides)} slides")

    # ── Step 1: Enrich existing 30 slides ──
    print("\n── Step 1: Enriching existing slides ──")

    for slide_idx in range(30):
        slide = prs.slides[slide_idx]
        is_detail = (slide_idx % 2 == 0)  # Odd slides (0-indexed even) = capability detail
        is_maturity = (slide_idx % 2 == 1)  # Even slides (0-indexed odd) = maturity table

        slide_num = slide_idx + 1
        print(f"  Slide {slide_num}: ", end="")

        if is_detail:
            # ── Capability Detail Slide ──
            # 1. Remove "smart" orphan text
            remove_smart_text(slide)

            # 2. Enrich table cells with Pakistan context
            if slide_idx in CAPABILITY_ENRICHMENTS:
                enrich = CAPABILITY_ENRICHMENTS[slide_idx]
                tbl = get_table(slide)
                if tbl:
                    # Add Pakistan bullets to Business Objectives [1,0]
                    if enrich.get('biz_obj'):
                        append_cell_text(tbl.cell(1, 0), enrich['biz_obj'])

                    # Add Pakistan bullets to Data & Solution [3,0]
                    if enrich.get('data_sol'):
                        append_cell_text(tbl.cell(3, 0), enrich['data_sol'])

                    # Add Pakistan bullets to Outcome [3,2]
                    if enrich.get('outcome'):
                        append_cell_text(tbl.cell(3, 2), enrich['outcome'])

            print(f"Detail - enriched table + removed 'smart'")

        else:
            # ── Maturity Table Slide ──
            # 1. Replace "H1 2018" with "H1 2025 — H2 2026"
            replace_h1_2018(slide)

            # 2. Remove "For use in Maturity Assessment & Roadmap Engagements"
            remove_maturity_engagement_text(slide)

            # 3. Enrich key areas and benefits cells
            if slide_idx in MATURITY_ENRICHMENTS:
                enrich = MATURITY_ENRICHMENTS[slide_idx]
                tbl = get_table(slide)
                if tbl:
                    if enrich.get('key_areas'):
                        replace_cell_text(tbl.cell(7, 3), enrich['key_areas'])
                    if enrich.get('benefits'):
                        replace_cell_text(tbl.cell(7, 4), enrich['benefits'])

            print(f"Maturity - enriched table + replaced H1 2018 + removed engagement text")

        # 3. Set speaker notes for ALL slides
        if slide_idx in SPEAKER_NOTES:
            set_notes(slide, SPEAKER_NOTES[slide_idx])
            print(f"    + Speaker notes set")

    # ── Step 2: Replace Teradata branding in header shapes ──
    print("\n── Step 2: Replacing Teradata branding ──")
    teradata_count = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'Teradata' in (run.text or ''):
                            run.text = run.text.replace('Teradata Business Value Framework', 'Banking Business Value Framework')
                            run.text = run.text.replace('Teradata', 'Enterprise Analytics Platform')
                            teradata_count += 1
    print(f"  Replaced {teradata_count} Teradata references")

    # ── Step 3: Add 3 new slides ──
    print("\n── Step 3: Adding 3 new slides ──")

    # We need to add slides by duplicating an existing slide's layout and inserting
    # Since python-pptx doesn't support insert_slide_at(), we'll add slides
    # at the end and then reorder via XML manipulation

    # Get slide layouts from existing slides
    detail_layout = prs.slides[0].slide_layout  # 2_Title and Subtitle Only
    maturity_layout = prs.slides[1].slide_layout  # 4_Title and Subtitle Only

    # Add New Slide A: Domain Dashboard
    slide_a = prs.slides.add_slide(detail_layout)
    _build_dashboard_slide(slide_a, prs)
    set_notes(slide_a, NEW_SLIDE_NOTES['new_a'])
    print("  Added: Domain Dashboard (New Slide A)")

    # Add New Slide B: Pakistan Market Context
    slide_b = prs.slides.add_slide(detail_layout)
    _build_pakistan_context_slide(slide_b, prs)
    set_notes(slide_b, NEW_SLIDE_NOTES['new_b'])
    print("  Added: Pakistan Interaction Landscape (New Slide B)")

    # Add New Slide C: Implementation Roadmap
    slide_c = prs.slides.add_slide(detail_layout)
    _build_roadmap_slide(slide_c, prs)
    set_notes(slide_c, NEW_SLIDE_NOTES['new_c'])
    print("  Added: Implementation Roadmap (New Slide C)")

    # ── Step 4: Reorder slides ──
    # New slides are at indices 30, 31, 32
    # We want: [30, 31, 0..28, 32, 29] → i.e., Dashboard, PK Context, then original 1-29, then Roadmap, then original 30
    # Actually per prompt: New A = slide 1, New B = slide 2, originals 1-30 = slides 3-32, New C = second to last (slide 32)
    # Final order: [new_a(30), new_b(31), orig_0..orig_27(0-27), new_c(32), orig_28, orig_29]
    print("\n── Step 4: Reordering slides ──")
    _reorder_slides(prs, len(prs.slides))
    print("  Slides reordered: Dashboard, PK Context, 15 capabilities, Roadmap, Brand Analytics")

    # ── Step 5: Save ──
    print(f"\n── Step 5: Saving to {OUTPUT_FILE} ──")
    prs.save(OUTPUT_FILE)

    # Clean up temp file
    if os.path.exists(CLEAN_FILE):
        os.remove(CLEAN_FILE)
        print("  Cleaned up temp file")

    print("  DONE!")


def _build_dashboard_slide(slide, prs):
    """Build the Domain Dashboard slide."""
    # Clear any inherited shapes from layout that have text
    # Set the title placeholder
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:  # Title
                shape.text = "Customer Interaction Management — At a Glance"
            elif shape.placeholder_format.idx == 1:  # Subtitle
                shape.text = ""

    # Add a text box with the dashboard content
    from pptx.util import Inches, Pt, Emu
    left = Emu(400000)
    top = Emu(900000)
    width = Emu(8300000)
    height = Emu(3600000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Title line
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Banking Business Value Framework — Customer Interaction Management"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    # Blank line
    p = tf.add_paragraph()
    p.text = ""

    # Dashboard table as text
    lines = [
        "                              Global          South Asia & ME    Pakistan",
        "──────────────────────────────────────────────────────────────────────────",
        "Marketing automation          75%+ top banks  30-40%             <15%",
        "Real-time personalization     50%+ digital    15-20%             <5%",
        "Omnichannel campaign exec     65% mature      25%                <10%",
        "NBA (Next Best Action)        40% top banks   10%                <3%",
        "Marketing attribution         55% top banks   15%                <5%",
        "Campaign ROI measurement      Standard        Emerging           Rare — most lack control groups",
        "──────────────────────────────────────────────────────────────────────────",
    ]
    for line in lines:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9)
        run.font.name = "Consolas"

    # Info bar
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "BVF Sub-capabilities: 15 | FSDM Entities: ~100+ (Campaign, Channel, Interaction, Contact) | BACR Questions: ~80 | Maturity: Developing > Innovating"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 102, 153)


def _build_pakistan_context_slide(slide, prs):
    """Build the Pakistan Market Context slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                shape.text = "Pakistan — Customer Interaction Landscape"
            elif shape.placeholder_format.idx == 1:
                shape.text = ""

    left = Emu(400000)
    top = Emu(900000)
    width = Emu(8300000)
    height = Emu(3800000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Channel Reach header
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Channel Reach and Effectiveness"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    channel_lines = [
        "SMS: 95%+ reach (mobile penetration), 2-5% response rate, OVERUSED",
        "Push Notification: 30% reach (app installed), 8-15% open rate, GROWING",
        "WhatsApp Business: 20% reach, 40-60% open rate, EMERGING — highest engagement",
        "Email: 15% reach (urban/corporate only), 5-10% open rate, LOW priority",
        "In-App Messages: 25% reach, 15-25% engagement, HIGHEST conversion",
        "Branch Face-to-Face: 100% for walk-ins, high conversion, HIGH cost",
        "Contact Center: Urdu/English bilingual, 1M+ calls/month (top-5), CRITICAL",
        "Agent Network: 500,000+ agents, rural reach, ESSENTIAL for inclusion",
    ]
    for line in channel_lines:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    p = tf.add_paragraph()
    p.text = ""

    # Key Challenges
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Key Challenges"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    challenges = [
        "SMS FATIGUE: Customers receive 5-10 bank SMS daily — regulatory crackdown looming",
        "NO CDP: <5% of banks have Customer Data Platform — interactions siloed by channel",
        "BATCH-ONLY: Most campaigns run on monthly batch — no real-time personalization",
        "NO ATTRIBUTION: Marketing spend not tracked to conversion — budget is gut-feel",
        "LANGUAGE: Must support Urdu, English, and increasingly Sindhi, Pashto, Punjabi",
        "COMPLIANCE: SBP requires opt-in for marketing + complaint response SLAs",
    ]
    for c in challenges:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {c}"
        run.font.size = Pt(9)


def _build_roadmap_slide(slide, prs):
    """Build the Implementation Roadmap slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                shape.text = "Customer Interaction Management — Implementation Roadmap"
            elif shape.placeholder_format.idx == 1:
                shape.text = ""

    left = Emu(400000)
    top = Emu(850000)
    width = Emu(8300000)
    height = Emu(3900000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Phase 1
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Phase 1: Foundation (0-6 months) — PKR 40-70M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 102, 51)

    phase1 = [
        "Implement contact policy framework (prevent SMS fatigue)",
        "Build unified customer contact history across channels",
        "Launch WhatsApp Business channel for top 20% customers",
        "A/B testing framework for SMS and push notifications",
        "Quick Win: 30% reduction in SMS opt-outs",
    ]
    for line in phase1:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Phase 2
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Phase 2: Intelligence (6-18 months) — PKR 100-180M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 153)

    phase2 = [
        "Customer Data Platform (CDP) implementation",
        "Multi-step campaign automation (Ramadan, Eid, salary-day sequences)",
        "Personalization engine — segment-level then individual-level",
        "Call center optimization with customer 360 view for agents",
        "Expected: 2x campaign response rates",
    ]
    for line in phase2:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Phase 3
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Phase 3: Real-Time (18-36 months) — PKR 150-280M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(153, 0, 51)

    phase3 = [
        "Next Best Action engine across branch, mobile, contact center",
        "Real-time contextual decisioning (location, transaction, app behavior)",
        "Marketing attribution model (digital + branch + ATL integration)",
        "AI-powered product recommendation at every touchpoint",
        "Expected: 3x marketing ROI, products/customer from 1.8 to 3.0+",
    ]
    for line in phase3:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Total
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Total 3-Year Investment: PKR 290-530M (~USD 1-2M) | Annual Revenue Impact: PKR 500M-1B"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)


def _reorder_slides(prs, total_slides):
    """Reorder slides: New A (30), New B (31), Original 0-27, New C (32), Original 28-29.

    Final order (0-indexed source positions):
    [30, 31, 0, 1, 2, ..., 27, 32, 28, 29]
    """
    new_order = [30, 31]  # Dashboard, Pakistan Context
    new_order.extend(range(0, 28))  # Original slides 1-28 (capabilities 1-14)
    new_order.append(32)  # Implementation Roadmap (second to last)
    new_order.extend([28, 29])  # Brand Analytics (original last 2)

    # Access the slide ID list via the presentation XML element
    pres_ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    pres_elem = prs.part._element
    sldIdLst = pres_elem.find('.//p:sldIdLst', pres_ns)

    sldId_elements = list(sldIdLst)

    for elem in sldId_elements:
        sldIdLst.remove(elem)

    for idx in new_order:
        sldIdLst.append(sldId_elements[idx])

    print(f"  Reordered {len(new_order)} slides")


if __name__ == '__main__':
    main()
