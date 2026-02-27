"""
Split BVF Finance Industry 1.pptx into logical sub-presentations.
"""

import os
from lxml import etree
from pptx import Presentation

SOURCE = "/mnt/e/erwin/BVF Finance Industry 1.pptx"
OUTDIR = "/mnt/e/erwin/pptout"

SECTIONS = [
    (1,   12,  "01_Introduction_Industry_Challenges"),
    (13,  35,  "02_Marketing_CX_Overview"),
    (36,  65,  "03_Customer_Information_Insight_Analytics"),
    (66,  83,  "04_Customer_Lifecycle_Management"),
    (84,  113, "05_Customer_Interaction_Management"),
    (114, 182, "06_Marketing_CX_Use_Cases"),
    (183, 228, "07_Finance_PM_Accounting_Operations"),
    (229, 286, "08_Finance_EPM_Treasury_Reporting"),
    (287, 318, "09_Finance_PM_Use_Cases"),
    (319, 372, "10_Product_Management"),
    (373, 389, "11_Product_Management_Use_Cases"),
    (390, 449, "12_Risk_Management"),
    (450, 473, "13_Risk_Management_Use_Cases"),
    (474, 533, "14_Security_Fraud_Overview"),
    (534, 574, "15_Security_Fraud_Use_Cases"),
    (575, 619, "16_Regulatory_Compliance"),
]


def extract_slides(source_path, output_path, start_idx, end_idx):
    """
    Extract slides [start_idx..end_idx] (0-based inclusive) by removing all others.
    """
    prs = Presentation(source_path)
    slide_list = list(prs.slides)

    # Get the presentation XML element
    prs_xml = prs.part._element
    nsmap = {
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

    # Get sldIdLst element
    sldIdLst = prs_xml.find('.//p:sldIdLst', nsmap)

    # Collect indices to remove (reverse order)
    indices_to_remove = [i for i in range(len(slide_list)) if i < start_idx or i > end_idx]

    for i in reversed(indices_to_remove):
        slide = slide_list[i]

        # Find the relationship ID for this slide
        rId = None
        for rel_key, rel in prs.part.rels.items():
            if rel.target_part == slide.part:
                rId = rel_key
                break

        if rId:
            # Remove from sldIdLst
            for sldId in list(sldIdLst):
                r_id_attr = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if r_id_attr == rId:
                    sldIdLst.remove(sldId)
                    break

            # Drop the relationship
            prs.part.drop_rel(rId)

    prs.save(output_path)


def main():
    os.makedirs(OUTDIR, exist_ok=True)

    prs = Presentation(SOURCE)
    total_slides = len(prs.slides)
    print(f"Source file has {total_slides} slides")
    print(f"Splitting into {len(SECTIONS)} files...\n")

    for start, end, name in SECTIONS:
        start_idx = start - 1
        end_idx = end - 1
        slide_count = end - start + 1
        output_path = os.path.join(OUTDIR, f"{name}.pptx")

        print(f"  {name}.pptx (slides {start}-{end}, {slide_count} slides)...", end=" ", flush=True)
        extract_slides(SOURCE, output_path, start_idx, end_idx)

        # Verify
        verify = Presentation(output_path)
        actual = len(verify.slides)
        size_mb = os.path.getsize(output_path) / (1024 * 1024)
        status = "OK" if actual == slide_count else f"MISMATCH (expected {slide_count}, got {actual})"
        print(f"{actual} slides, {size_mb:.1f} MB - {status}")

    print(f"\nDone! All files saved to {OUTDIR}/")


if __name__ == "__main__":
    main()
