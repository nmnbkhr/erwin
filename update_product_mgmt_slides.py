#!/usr/bin/env python3
"""
Update File 10 — Product Management
Enriches 54 slides with Pakistan banking product context, rewrites retail/telco
language as banking, adds 3 new slides, removes Teradata branding, adds speaker notes.
Output: ./pptout/10_Product_Management_UPDATED.pptx (57+ slides)
"""

import os
import copy
import re
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.shapes.group import GroupShape

INPUT_FILE = './10_Product_Management.pptx'
CLEAN_FILE = './pptout/10_Product_Management_CLEAN.pptx'
OUTPUT_FILE = './pptout/10_Product_Management_UPDATED.pptx'


# ─────────────────────────────────────────────
# Industry Language Translation Map
# ─────────────────────────────────────────────

LANGUAGE_REPLACEMENTS = [
    # Order matters — more specific patterns first
    ('SKU', 'banking product'),
    ('SKUs', 'banking products'),
    ('merchandise', 'banking product'),
    ('merchandising', 'product management'),
    ('assortment', 'product portfolio'),
    ('assortments', 'product portfolios'),
    ('product mix', 'product portfolio'),
    ('store ', 'branch '),
    ('stores', 'branches'),
    ('Store ', 'Branch '),
    ('Stores', 'Branches'),
    ('point of sale', 'point of service'),
    ('POS', 'branch/digital channel'),
    ('ARPU', 'revenue per customer'),
    ('shelf space', 'cross-sell placement'),
    ('adjacencies', 'product affinity'),
    ('basket size', 'products per customer'),
    ('UPT', 'products per customer'),
    ('markdown', 'rate adjustment'),
    ('markdowns', 'rate adjustments'),
    ('clearance', 'product sunset'),
    ('inventory', 'product capacity'),
    ('stock ', 'product pipeline '),
    ('supplier', 'product manufacturer'),
    ('suppliers', 'product manufacturers'),
    ('vendor', 'technology provider'),
    ('consumer', 'customer'),
    ('consumers', 'customers'),
    ('Consumer', 'Customer'),
    ('Consumers', 'Customers'),
    ('shopper', 'customer'),
    ('shoppers', 'customers'),
    ('Shopper', 'Customer'),
    ('wireframe', 'digital channel'),
    ('wire line', 'digital channel'),
    ('wireline', 'digital channel'),
    ('Wireframe', 'Digital Channel'),
    ('content viewership', 'mobile app usage'),
    ('Content viewership', 'Mobile app usage'),
    ('churn', 'product closure rate'),
    ('retail ', 'banking '),
    ('Retail ', 'Banking '),
    ('H1 2018', 'H1 2025 — H2 2026'),
]


# ─────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────

def set_notes(slide, text):
    """Set or replace the speaker notes on a slide."""
    if not slide.has_notes_slide:
        slide.notes_slide  # creates notes slide
    notes_tf = slide.notes_slide.notes_text_frame
    # Clear existing
    for i in range(len(notes_tf.paragraphs) - 1, 0, -1):
        p = notes_tf.paragraphs[i]._element
        p.getparent().remove(p)
    notes_tf.paragraphs[0].text = text


def replace_cell_text(cell, new_text):
    """Replace cell text while preserving first run formatting."""
    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
        fmt = cell.text_frame.paragraphs[0].runs[0].font
        font_size = fmt.size
        font_bold = fmt.bold
        font_name = fmt.name
    else:
        font_size = Pt(10)
        font_bold = False
        font_name = None

    for i in range(len(cell.text_frame.paragraphs) - 1, 0, -1):
        p = cell.text_frame.paragraphs[i]._element
        p.getparent().remove(p)

    lines = new_text.split('\n')
    for idx, line in enumerate(lines):
        if idx == 0:
            p = cell.text_frame.paragraphs[0]
            p.clear()
        else:
            p = cell.text_frame.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.bold = font_bold
        if font_name:
            run.font.name = font_name


def append_cell_text(cell, additional_text):
    """Append text to existing cell content."""
    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
        fmt = cell.text_frame.paragraphs[0].runs[0].font
        font_size = fmt.size
        font_name = fmt.name
    else:
        font_size = Pt(10)
        font_name = None

    for line in additional_text.split('\n'):
        p = cell.text_frame.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        if font_name:
            run.font.name = font_name


def remove_engagement_text(slide):
    """Remove 'For use in Maturity Assessment & Roadmap Engagements' text from groups."""
    removed = False
    for shape in list(slide.shapes):
        if isinstance(shape, GroupShape):
            for s in shape.shapes:
                if s.has_text_frame:
                    txt = s.text_frame.text.strip()
                    if 'Engagements' in txt and 'Maturity Assessment' in txt:
                        for para in s.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = ''
                        removed = True
    return removed


def get_table(slide):
    """Return the first table shape on a slide."""
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def replace_h1_2018(table):
    """Replace 'H1 2018' with 'H1 2025 — H2 2026' in table cells."""
    replaced = False
    for row in table.rows:
        for cell in row.cells:
            if 'H1 2018' in cell.text:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if 'H1 2018' in run.text:
                            run.text = run.text.replace('H1 2018', 'H1 2025 —\nH2 2026')
                            replaced = True
    return replaced


def apply_language_replacements(text):
    """Apply retail/telco → banking language translations."""
    result = text
    for old, new in LANGUAGE_REPLACEMENTS:
        result = result.replace(old, new)
    return result


def find_shape_by_text(slide, text_fragment, exact=False):
    """Find a shape containing specific text."""
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            shape_text = shape.text_frame.text
            if exact and shape_text.strip() == text_fragment.strip():
                return shape
            elif not exact and text_fragment in shape_text:
                return shape
    return None


def append_paragraph(text_frame, text, font_size=None, bold=False, color=None):
    """Add a new paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    run = p.add_run()
    run.text = text
    if font_size:
        run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color
    return p


# ─────────────────────────────────────────────
# Capability enrichments for 13 detail slides
# Table structure: 4x3
#   [0,0]=Business Objectives header, [1,0]=content
#   [2,0]=Data & Solution header, [2,2]=Outcome header
#   [3,0]=data content, [3,2]=outcome content
# ─────────────────────────────────────────────

CAPABILITY_ENRICHMENTS = {
    # Slide 21 (0-indexed 20): Research New Product Development Opportunities
    20: {
        'biz_obj': [
            "Research unmet needs in Pakistan's 220M population — only 30% banked, 60% have mobile phones",
            "Analyze RAAST transaction patterns to identify product opportunities (P2P senders → savings product, frequent billers → auto-debit product)",
            "Identify Islamic product gaps — Shariah-compliant alternatives for every conventional product",
            "Products for financially excluded: basic accounts, micro-savings, nano-lending via mobile",
        ],
        'data_sol': [
            "FSDM customer/transaction data, mobile app usage analytics",
            "RAAST transaction patterns, SBP financial inclusion survey data",
            "Competitor product monitoring, social media sentiment analysis",
            "NADRA demographic data for market sizing by segment",
        ],
        'outcome': [
            "Data-driven product pipeline aligned with market need",
            "SBP financial inclusion target alignment (50% inclusion by 2028)",
            "Islamic product gap analysis: identify missing Shariah-compliant equivalents",
            "Digital-first product pipeline for Pakistan's mobile-savvy youth (60% under 30)",
        ],
    },
    # Slide 23 (0-indexed 22): Predict Business Impact (Discovery)
    22: {
        'biz_obj': [
            "Predict impact of new product launch on existing product balances (cannibalization)",
            "Forecast channel capacity impact: branch/call center load from new product inquiries",
            "Estimate revenue impact: NII + fees from new product vs. cannibalization of existing",
            "Pakistan example: digital-only savings — will it cannibalize branch savings? Expected adoption curve?",
        ],
        'data_sol': [
            "Historical product launch data (past 5 years), similar product performance at competitors",
            "Customer segment adoption patterns from FSDM, channel capacity models",
            "FSDM: PRDCT_DVLPMNT, MKT_RSRCH, CSTMR_ND, PRDCT_PPLN",
        ],
        'outcome': [
            "Evidence-based product launch business case with cannibalization estimate",
            "Channel capacity plan for new product support",
            "Risk assessment: operational, regulatory (SBP approval), market risk for new products",
        ],
    },
    # Slide 25 (0-indexed 24): Develop and Test Product Prototypes
    24: {
        'biz_obj': [
            "Test product features, pricing, and packaging through A/B testing and pilot programs",
            "Pakistan banking: pilot new products in select branches/regions before national rollout",
            "Test pricing variants (rate sensitivity), feature variants (min balance, fee structure), channel variants",
            "A/B testing for digital products, regional pilot for branch products, champion/challenger for pricing",
        ],
        'data_sol': [
            "FSDM transaction data for pilot vs. control groups",
            "Mobile app engagement data, branch sales data, customer feedback",
            "Call center query data during pilot period",
            "NADRA biometric success rates for digital product onboarding",
        ],
        'outcome': [
            "Evidence-based product launch decisions, optimized features and pricing before full investment",
            "Reduced product failure rate through data-driven pilot validation",
            "SBP-ready product performance data supporting NOC application",
        ],
    },
    # Slide 29 (0-indexed 28): Plan Product Investments / Go-To-Market
    28: {
        'biz_obj': [
            "Plan GTM for Pakistan banking: branch staff training (16,000+ branches)",
            "Core banking system configuration (4-5 systems), SBP reporting setup",
            "Marketing collateral (Urdu + English), Islamic Shariah compliance documentation",
            "Forecast resource requirements: call center staffing, branch teller training, digital UX readiness",
        ],
        'data_sol': [
            "Historical product launch data, branch capacity data",
            "Call center staffing model, digital channel capacity",
            "Card production lead times, NADRA integration capacity",
            "SBP NOC submission requirements and timeline data",
        ],
        'outcome': [
            "Fully resourced product launch plan with channel-specific readiness checklist",
            "Reduced time-to-market through parallel workstream management",
            "SBP approval timeline integrated into launch planning",
        ],
    },
    # Slide 31 (0-indexed 30): Improve Product Onboarding
    30: {
        'biz_obj': [
            "Pakistan digital onboarding has 40-60% dropoff rate at NADRA biometric verification step",
            "Optimize product onboarding: account opening (7-14 days → <1 day), loan origination (21-45 days → <7 days)",
            "Card issuance optimization: 14-30 days → <5 days",
            "Track onboarding funnel: application → document submission → NADRA → KYC → activation → first transaction",
        ],
        'data_sol': [
            "Digital onboarding funnel data, NADRA verification success/failure rates",
            "Document upload completion rates, first transaction timing",
            "Call center queries during onboarding, branch walk-in data",
            "Competitor onboarding benchmarks (SadaPay, NayaPay — minutes not days)",
        ],
        'outcome': [
            "Reduced digital dropoff from 40-60% to <20% through funnel optimization",
            "Faster time-to-first-transaction: from days to hours for digital products",
            "Onboarding cost reduction through digital-first channel steering",
        ],
    },
    # Slide 35 (0-indexed 34): Analyze Competitive Pricing
    34: {
        'biz_obj': [
            "Monitor competitive pricing: deposit rates (SBP publishes), lending rates (KIBOR + spreads)",
            "Track card fees (annual, late payment, FX markup), trade finance pricing (LC commission rates)",
            "Digital transaction fees (RAAST/IBFT charges), account maintenance fees",
            "Pakistan-specific: SBP minimum deposit rate creates floor; KIBOR creates lending benchmark",
        ],
        'data_sol': [
            "SBP published rate data, competitor rate sheets (publicly available)",
            "KIBOR curves, card fee comparison across peers",
            "Customer rate sensitivity data from FSDM",
            "Digital banking fee comparison (JazzCash, Easypaisa, SadaPay, NayaPay)",
        ],
        'outcome': [
            "Competitive pricing dashboard updated daily for deposit and lending rates",
            "Fee optimization opportunity identification across product portfolio",
            "Market share impact analysis from pricing changes",
        ],
    },
    # Slide 37 (0-indexed 36): Manage Dynamic Pricing
    36: {
        'biz_obj': [
            "Dynamic pricing in Pakistan banking: adjust deposit rates for large depositors (relationship pricing)",
            "Risk-based lending: KIBOR + variable risk premium based on ECIB/internal rating",
            "Promotional card rates: balance transfer offers, spend-based fee waivers",
            "Constraints: SBP minimum deposit rate, KIBOR-linked lending floors, SBP credit card rate caps (under discussion)",
        ],
        'data_sol': [
            "Customer deposit balance and tenor data from FSDM",
            "ECIB credit bureau scores, internal risk ratings",
            "Transaction and spending pattern data for card pricing",
            "Competitor pricing intelligence and market elasticity data",
        ],
        'outcome': [
            "Optimize NII by dynamically adjusting deposit gathering cost vs. lending margin",
            "Retain large depositors through competitive relationship pricing while protecting margin",
            "Risk-adjusted lending pricing maximizing risk-return trade-off",
        ],
    },
    # Slide 41 (0-indexed 40): Determine Product Migration
    40: {
        'biz_obj': [
            "Sunset legacy products: old savings schemes, discontinued FD tenors",
            "Migrate customers from conventional to Islamic products (growing demand)",
            "Consolidate product variants across merged entities (200+ variants → 50-80)",
            "Regulatory sunset: SBP discontinued products require customer migration",
        ],
        'data_sol': [
            "Product-level balance and customer data from FSDM",
            "Customer communication preference, channel usage data",
            "Product switching history, customer satisfaction scores",
            "Regulatory timeline data for mandated product changes",
        ],
        'outcome': [
            "Smooth migration preserving customer relationships, minimal attrition during sunset",
            "Optimized product portfolio: fewer variants, lower operational complexity",
            "Proactive conventional → Islamic migration for willing customers",
        ],
    },
    # Slide 45 (0-indexed 44): Analyze Cross Channel Performance
    44: {
        'biz_obj': [
            "Track product sales and usage across all Pakistan banking channels",
            "Branch (16,000+), ATM (16,000+), mobile app, internet banking, USSD, call center, agent network, RAAST",
            "Key insight: which products sell best on which channels?",
            "Personal loans: 80% via branch but digital origination growing; CASA opens increasingly digital",
        ],
        'data_sol': [
            "FSDM transaction data by channel, product sales by channel",
            "Mobile app product funnel data, branch sales register",
            "Agent banking transaction data, RAAST/IBFT channel data",
            "Channel cost-to-serve data for profitability by channel",
        ],
        'outcome': [
            "Channel-optimized product distribution strategy",
            "Digital migration roadmap: identify products ready for digital-first distribution",
            "Channel ROI by product — cost-to-serve optimization",
        ],
    },
    # Slide 47 (0-indexed 46): Compare to Product Benchmarks
    46: {
        'biz_obj': [
            "Benchmark Pakistan banking products against SBP published industry data",
            "Peer group performance (top 5/10 banks), international benchmarks (Gulf, SE Asian banks)",
            "Key product benchmarks: CASA ratio (~47%), card activation rate, loan-to-deposit ratio",
            "Fee income per product, digital adoption rate, products per customer (1.8 vs 4-6 global)",
        ],
        'data_sol': [
            "SBP statistical bulletins, annual reports of peer banks",
            "Industry surveys (PBA), FSDM internal product performance data",
            "International benchmarking data (Gulf banks, SE Asian banks)",
            "Digital banking benchmarks from fintech competitors",
        ],
        'outcome': [
            "Product performance scorecard vs. industry benchmarks",
            "Gap analysis: identify underperforming products vs. peers",
            "Actionable improvement targets for each product category",
        ],
    },
    # Slide 49 (0-indexed 48): Manage Product Profitability — CRITICAL
    48: {
        'biz_obj': [
            "Calculate true profitability: FTP-adjusted NII + fee income - direct costs - ABC-allocated costs - IFRS 9 provision",
            "Cross-reference with File 08-09 Profitability Analytics — product profitability is one dimension",
            "Pakistan insight: some 'popular' products may be unprofitable (basic savings PKR 100 min, free debit cards)",
            "Analytics identifies cross-subsidy flows between product lines",
        ],
        'data_sol': [
            "FSDM Customer Profitability Engine output aggregated to product level",
            "FTP rates by product, ABC unit costs by product",
            "IFRS 9 provision allocation by product category",
            "Revenue and cost data: NII, fee income, interchange, operating costs",
        ],
        'outcome': [
            "True product-level profitability revealing cross-subsidy flows",
            "Identify bottom 10% unprofitable products for repricing or sunset",
            "Product pricing optimization based on true cost-to-serve",
        ],
    },
    # Slide 51 (0-indexed 50): Produce Behavioral Segmentations
    50: {
        'biz_obj': [
            "Segment by product usage: heavy transactors (high RAAST/ATM), savers (FD-focused), borrowers (multiple lending)",
            "Digital-first (mobile app primary), branch-loyal (branch-only), dormant (no activity 90+ days)",
            "Islamic preference segmentation: Islamic-only vs. mixed vs. conventional-only customers",
            "Product affinity segmentation: which product combinations predict high lifetime value?",
        ],
        'data_sol': [
            "FSDM transaction data, product holding data, channel usage data",
            "Islamic/conventional preference flag, mobile app engagement metrics",
            "Product switching history, cross-sell acceptance/rejection data",
            "Demographic data linked to product usage patterns",
        ],
        'outcome': [
            "Behavioral segments driving targeted product recommendations",
            "Islamic product opportunity sizing by customer segment",
            "Product bundle optimization: right product combinations for each segment",
        ],
    },
    # Slide 53 (0-indexed 52): Predict Product/Service Retention
    52: {
        'biz_obj': [
            "Predict which products each customer is likely to close: FD non-renewal, CASA closure, card cancellation",
            "Pakistan-specific: salary account switching is highest product-level churn risk",
            "When employer switches bank, entire product bundle at risk (salary + CASA + card + loan)",
            "Loan prepayment prediction: refinancing at competitors offering lower KIBOR spreads",
        ],
        'data_sol': [
            "FSDM: product tenure, balance trend, transaction frequency change",
            "Salary credit continuity monitoring, competitor rate data",
            "Customer life event indicators, ECIB inquiry frequency (rate shopping signal)",
            "Product usage decline patterns: ATM/POS decline, mobile app inactivity",
        ],
        'outcome': [
            "Product retention alerts 30 days before closure",
            "Targeted retention offers by product type (rate match, fee waiver, upgrade)",
            "Salary account switching early warning system",
        ],
    },
}


# ─────────────────────────────────────────────
# Maturity enrichments for 13 maturity slides
# ─────────────────────────────────────────────

MATURITY_ENRICHMENTS = {
    # Slide 22 (0-indexed 21): Research New Product Development maturity
    21: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "RAAST/IBFT transaction pattern analysis for product opportunity identification\n"
            "Digital product demand analysis from mobile app usage analytics\n"
            "Islamic product gap analysis (Shariah-compliant alternatives for every conventional product)\n"
            "Financial inclusion product research (basic accounts, nano-lending, micro-savings)\n"
            "Competitor product monitoring dashboard (bank + fintech landscape)"
        ),
        'benefits': (
            "Data-driven product pipeline aligned with 220M population needs\n"
            "SBP financial inclusion target compliance through analytically designed products\n"
            "Islamic product innovation addressing 25%+ deposit market\n"
            "Reduced product failure rate through evidence-based development"
        ),
    },
    # Slide 24 (0-indexed 23): Predict Business Impact maturity
    23: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Product cannibalization modeling (digital vs. branch product impact)\n"
            "Channel capacity forecasting for new product launches\n"
            "Revenue prediction: NII + fee income impact with KIBOR sensitivity\n"
            "SBP approval timeline integration in launch impact modeling\n"
            "Islamic product impact assessment (Shariah board + SBP dual approval)"
        ),
        'benefits': (
            "Accurate new product revenue forecasting reducing launch risk\n"
            "Channel capacity planning avoiding service degradation at launch\n"
            "Cannibalization quantification enabling informed portfolio decisions\n"
            "Regulatory timeline accuracy improving launch date predictability"
        ),
    },
    # Slide 26 (0-indexed 25): Develop and Test Product Prototypes maturity
    25: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "A/B testing framework for digital banking products (rate, feature, UX variants)\n"
            "Regional pilot methodology for branch-based products\n"
            "Champion/challenger pricing testing (deposit rates, lending spreads)\n"
            "Focus group infrastructure for Islamic product design validation\n"
            "NADRA biometric integration testing for digital onboarding flows"
        ),
        'benefits': (
            "Reduced product launch risk through evidence-based pilot validation\n"
            "Optimized product features and pricing before full-scale investment\n"
            "Faster iteration on digital products through A/B testing capability\n"
            "Improved SBP NOC applications supported by pilot performance data"
        ),
    },
    # Slide 30 (0-indexed 29): Plan Product Investments / GTM maturity
    29: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Multi-channel GTM planning: 16,000+ branches, digital, ATM, agent network\n"
            "SBP NOC timeline integration in product launch planning\n"
            "Bilingual marketing collateral management (Urdu + English)\n"
            "Core banking configuration management across 4-5 systems\n"
            "Islamic product launch checklist: Shariah board + SBP + marketing + CBS config"
        ),
        'benefits': (
            "Reduced time-to-market through parallel workstream management\n"
            "Complete channel readiness at launch (branch training, digital UX, CBS config)\n"
            "SBP approval integrated into standard product launch methodology\n"
            "Consistent product experience across all channels from day one"
        ),
    },
    # Slide 32 (0-indexed 31): Improve Product Onboarding maturity
    31: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "NADRA biometric verification optimization (30-40% first-attempt failure rate)\n"
            "Digital onboarding funnel analytics (identify and fix dropout points)\n"
            "Account opening time reduction: branch (7-14 → <1 day), digital (<30 min)\n"
            "Loan origination streamlining: 21-45 days → <7 days (personal), <30 days (home)\n"
            "Card issuance acceleration: 14-30 days → <5 days"
        ),
        'benefits': (
            "Digital dropoff reduction from 40-60% to <20%\n"
            "Faster time-to-first-transaction increasing product activation rates\n"
            "Reduced onboarding cost through digital-first channel steering\n"
            "Competitive parity with fintechs on onboarding speed (SadaPay: minutes)"
        ),
    },
    # Slide 36 (0-indexed 35): Analyze Competitive Pricing maturity
    35: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Real-time competitive pricing dashboard: deposits, lending, cards, trade finance\n"
            "SBP rate data integration for deposit floor monitoring\n"
            "KIBOR curve analytics for lending spread optimization\n"
            "Digital banking fee benchmarking (RAAST, IBFT, mobile wallet fees)\n"
            "Cross-product pricing correlation analysis"
        ),
        'benefits': (
            "Daily competitive intelligence on pricing across all product categories\n"
            "Optimized deposit gathering cost through rate-vs-volume analytics\n"
            "Lending spread optimization maximizing NII within competitive constraints\n"
            "Fee income optimization from card and digital product pricing"
        ),
    },
    # Slide 38 (0-indexed 37): Manage Dynamic Pricing maturity
    37: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Relationship-based deposit pricing for large depositors (tiered rate engine)\n"
            "Risk-adjusted lending pricing: KIBOR + variable premium per ECIB/internal rating\n"
            "Card promotional pricing automation (balance transfer, spend rewards)\n"
            "SBP constraint management: minimum deposit rates, lending caps\n"
            "Real-time pricing authority delegation to branch/RM level"
        ),
        'benefits': (
            "NII optimization through dynamic deposit cost management\n"
            "Large depositor retention through competitive relationship pricing\n"
            "Risk-return optimization on lending portfolio\n"
            "Automated compliance with SBP pricing constraints"
        ),
    },
    # Slide 42 (0-indexed 41): Determine Product Migration maturity
    41: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Legacy product sunset methodology with customer communication plan\n"
            "Post-merger product rationalization framework (200+ → 50-80 variants)\n"
            "Conventional → Islamic product migration for willing customers\n"
            "Branch → digital product migration for eligible customers\n"
            "SBP-mandated product change management (regulatory sunset)"
        ),
        'benefits': (
            "Minimal customer attrition during product sunset (target <5% loss)\n"
            "Optimized product portfolio reducing operational complexity\n"
            "Proactive Islamic migration capturing growing demand segment\n"
            "Regulatory compliance for SBP-mandated product changes"
        ),
    },
    # Slide 46 (0-indexed 45): Analyze Cross Channel Performance maturity
    45: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Cross-channel product performance dashboard: branch, ATM, mobile, internet, agent, RAAST\n"
            "Channel-product affinity analytics (which products sell best where)\n"
            "Digital migration tracking for traditionally branch-based products\n"
            "Agent banking product performance monitoring (branchless banking)\n"
            "Channel cost-to-serve integration for product distribution ROI"
        ),
        'benefits': (
            "Optimized product distribution across 8+ banking channels\n"
            "Digital migration acceleration reducing branch dependency\n"
            "Channel-specific product promotion effectiveness measurement\n"
            "Agent banking expansion guided by product performance data"
        ),
    },
    # Slide 48 (0-indexed 47): Compare to Product Benchmarks maturity
    47: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "SBP industry data integration for automated product benchmarking\n"
            "Peer group comparison (top 5/10 banks) across all product categories\n"
            "International benchmarking: Gulf, SE Asian, global digital banks\n"
            "Fintech competitor benchmarking (JazzCash, Easypaisa, SadaPay, NayaPay)\n"
            "CASA ratio, card activation, digital adoption, products per customer benchmarks"
        ),
        'benefits': (
            "Automated product performance scorecard vs. industry benchmarks\n"
            "Identified underperformance gaps with actionable improvement targets\n"
            "Competitive positioning clarity across all product categories\n"
            "Board-ready benchmark reports supporting strategic product decisions"
        ),
    },
    # Slide 50 (0-indexed 49): Manage Product Profitability maturity
    49: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "True product profitability: FTP-adjusted NII + fees - direct costs - ABC costs - IFRS 9 provision\n"
            "Cross-subsidy flow identification between product lines\n"
            "Islamic product profitability parity analysis (Islamic vs. conventional)\n"
            "Digital product profitability measurement (very low cost-to-serve)\n"
            "Integration with File 08-09 multi-dimensional profitability model"
        ),
        'benefits': (
            "True product P&L revealing hidden unprofitable products\n"
            "Informed product pricing decisions based on actual cost-to-serve\n"
            "Product portfolio optimization reducing cross-subsidy leakage\n"
            "Board-level product profitability visibility for strategic decisions"
        ),
    },
    # Slide 52 (0-indexed 51): Produce Behavioral Segmentations maturity
    51: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Product usage behavioral segmentation: transactors, savers, borrowers, digital-first, dormant\n"
            "Islamic preference segmentation for targeted product offers\n"
            "Product affinity modeling: which product combinations predict high LTV\n"
            "RAAST/digital behavioral signals for segment identification\n"
            "Youth segment product preferences (60% under 30)"
        ),
        'benefits': (
            "Targeted product recommendations by behavioral segment\n"
            "Islamic product cross-sell opportunity sizing\n"
            "Optimized product bundles increasing products per customer from 1.8\n"
            "Digital-first segment identification for fintech-competitive offers"
        ),
    },
    # Slide 54 (0-indexed 53): Predict Product/Service Retention maturity
    53: {
        'key_areas': (
            "Pakistan priority improvements:\n"
            "Product closure prediction: FD non-renewal, CASA closure, card cancellation, loan prepayment\n"
            "Salary account switching early warning system\n"
            "Rate shopping detection via ECIB inquiry frequency monitoring\n"
            "Product usage decline pattern recognition\n"
            "Competitor product launch impact assessment on retention"
        ),
        'benefits': (
            "Product retention alerts 30 days before closure\n"
            "Targeted retention offers by product type (rate match, fee waiver, upgrade)\n"
            "Salary account switching prevention saving entire product bundles\n"
            "Reduced product attrition by 15-25% through proactive intervention"
        ),
    },
}


# ─────────────────────────────────────────────
# Speaker notes for all 54 existing slides
# ─────────────────────────────────────────────

SPEAKER_NOTES = {
    # Part A: Overview (slides 1-18, 0-indexed 0-17)
    0: """PRODUCT MANAGEMENT — Navigation
This domain covers 6 sub-domains with 12 BVF sub-capabilities:
1. Product Development (3 capabilities)
2. Product Introduction (2 capabilities)
3. Price and Promotion (2 capabilities)
4. Product End-of-Life (1 capability)
5. Product Performance (5 capabilities — largest sub-domain)

PAKISTAN CONTEXT: Product management in banking is the revenue engine — every PKR of income comes from a product. Pakistan banks average 1.8 products per customer vs. 4-6 globally — massive cross-sell opportunity.

FSDM: PRDCT (Product), PRDCT_TYP (Product Type), AGRMNT (Agreement), PRDCT_FTRE (Product Feature), PRDCT_PRCNG (Product Pricing), PRDCT_PRFMNC (Product Performance)""",

    1: """PRODUCT MANAGEMENT — Functional Insights (6 Areas)
The 6 functional insight areas cover the complete product lifecycle from development through end-of-life. In Pakistan banking:
- Product Development: New product design (CASA, lending, cards, Islamic, digital)
- Product Introduction: Go-to-market with SBP approval integration
- Price and Promotion: KIBOR-based pricing, SBP rate constraints
- Product End-of-Life: Legacy product sunset, merger rationalization
- Product Performance: Cross-channel analytics, benchmarking, profitability

FSDM: PRDCT (Product), AGRMNT (Agreement), PRCNG (Pricing), PRFMNC (Performance), LFCYCL (Lifecycle)""",

    2: """PRODUCT MANAGEMENT — Capabilities Page 1
These capabilities map directly to Pakistan banking product lifecycle needs. Each capability has a detail slide and maturity assessment slide.

Pakistan banks should prioritize: Product Profitability (most banks cannot calculate true product P&L), Cross-Channel Performance (8+ channels with no unified view), and Product Onboarding (40-60% digital dropoff rate).

BACR Assessment: ~60 product management questions across all capabilities.
FSDM: 12 sub-capabilities with FSDM entity mappings across Product, Agreement, Party, Transaction, Channel, Pricing domains.""",

    3: """PRODUCT MANAGEMENT — Capabilities Page 2
Continued capability listing. Note that Product Performance has the most sub-capabilities (5) reflecting the critical importance of measuring how products perform across channels, segments, and profitability dimensions.

Pakistan-specific priority: Behavioral Segmentation (Islamic vs. conventional preference) and Product Retention (salary switching is #1 churn risk).

FSDM: PRDCT_SGMNT (Product Segment), PRDCT_RTNTN (Product Retention), BHVRL_SGMNT (Behavioral Segmentation)""",

    4: """PRODUCT MANAGEMENT BVF (1 of 2) — Use Case Map
Use cases mapped to Product Management capabilities. Pakistan priority use cases:
- Product Profitability Analytics (from FSDM Customer Profitability Engine)
- Cross-Channel Product Performance Dashboard
- Product Onboarding Journey Optimization (NADRA biometric fix)
- Islamic Product Innovation Pipeline
- Dynamic Pricing Engine (risk-adjusted lending, competitive deposit pricing)

FSDM: Use cases span Product, Agreement, Transaction, Channel, and Pricing entity domains.""",

    5: """PRODUCT MANAGEMENT BVF (2 of 2) — Use Case Map
Continued use case mapping. Additional Pakistan priority use cases:
- Product Retention Prediction (salary switching early warning)
- Behavioral Segmentation (Islamic preference, digital-first, branch-loyal)
- Product Benchmark Analytics (vs. SBP industry data, peer banks, Gulf/SE Asian banks)
- New Product Development Research (RAAST patterns, financial inclusion)

FSDM: PRDCT_RTNTN, BHVRL_SGMNT, PRDCT_BNCHMARK, MKT_RSRCH""",

    6: """(Spacer slide — retained for navigation structure)""",

    7: """PRODUCTS/SERVICES VARY BY INDUSTRY — Banking Context
In banking, a 'product' encompasses: deposit accounts (CASA, term deposits), lending facilities (personal loans, home finance, corporate working capital), cards (debit, credit, prepaid), trade finance instruments (LC, LG, SBLC), treasury products (government securities, FX), digital services (mobile banking, QR payments, RAAST), and Islamic equivalents (Murabaha, Ijarah, Diminishing Musharaka).

Pakistan banking product portfolio for a typical large bank: 100-200+ product variants across deposit, lending, card, trade, treasury, digital, and Islamic categories.

FSDM: PRDCT (Product), PRDCT_TYP (Product Type), PRDCT_CTGRY (Product Category), PRDCT_FTRE (Product Feature)""",

    8: """EXAMPLES: WHAT IS A PRODUCT OR SERVICE — Banking Products
Pakistan banking product examples:
DEPOSITS: Current Account (zero interest CASA), Savings Account (SBP min rate ~11-13%), Fixed Deposit (3M-5Y), Foreign Currency Account, Roshan Digital Account, Basic Banking Account, Islamic: Mudaraba Savings, Mudaraba Term Deposit, Wakalah Pool
LENDING: Corporate working capital, term finance, project finance, trade finance (LC/LG/SBLC), personal loan, home finance, auto finance, credit card revolving, SME running/term finance, agriculture finance, Islamic: Murabaha, Ijarah, Diminishing Musharaka
CARDS: Debit (30M+), Credit (3M+), Prepaid (travel, remittance, payroll)
DIGITAL: Mobile banking, internet banking, RAAST P2P/P2M, QR payments, digital wallet integration

FSDM: PRDCT, PRDCT_TYP, DPST_PRDCT, LNDNG_PRDCT, CARD_PRDCT, DGTL_PRDCT""",

    9: """PRODUCT MANAGEMENT — Why Important
Product management is the engine of banking revenue — every PKR of income comes from a product.
Pakistan banks average 1.8 products per customer vs. 4-6 globally — massive cross-sell opportunity.
Islamic product innovation is critical — 25%+ of deposits are now in Islamic products, growing 20%+ annually.
Digital product adoption determines competitive survival against fintechs (JazzCash 50M+, Easypaisa 40M+ users).

PAKISTAN BANKING PRODUCT LANDSCAPE — Revenue by product category (typical large bank):
CASA: 30-40% of NII (low/zero cost deposits deployed at KIBOR)
Term Deposits: 15-20% of NII (spread between FD rate and deployment)
Corporate Lending: 20-25% of total income (NII + trade finance fees)
Consumer Lending: 10-15% (personal loans, cards, home/auto finance)
Cards: 8-12% (interchange, revolving interest, annual fees)
Trade Finance: 8-12% (LC/LG commissions — high-margin)
Treasury: 10-15% (government securities yield)

FSDM: PRDCT, PRDCT_TYP, AGRMNT, PRDCT_FTRE, PRDCT_PRCNG, PRDCT_PRFMNC""",

    10: """PRODUCT MANAGEMENT — How (Capabilities)
Product management capabilities enable data-driven decisions across the product lifecycle. The approach:
1. PRODUCT INTELLIGENCE: Unified product performance data from all channels and systems
2. PRICING ANALYTICS: Dynamic pricing within SBP constraints (deposit floors, lending caps)
3. LIFECYCLE MANAGEMENT: From concept through sunset with analytics at every stage
4. PROFITABILITY: True product P&L with FTP, ABC, and IFRS 9 provision allocation
5. BEHAVIORAL INSIGHTS: Product usage segmentation driving personalized recommendations
6. ISLAMIC PRODUCTS: Shariah-compliant product analytics with dual measurement

FSDM: All Product domain entities plus Agreement, Transaction, Channel, Pricing, and Performance entities.""",

    11: """PRODUCT MANAGEMENT — Solutions
Solution components for Pakistan banking product management:
- Product Performance Dashboard: Real-time sales, balances, revenue, cost by product and channel
- Pricing Analytics Engine: Competitive monitoring + dynamic pricing within SBP constraints
- Onboarding Optimization: Funnel analytics for NADRA-integrated digital onboarding
- Product Profitability: FTP-adjusted NII + fees - costs - provisions at product level
- Cross-Sell Engine: Product affinity models driving NBA recommendations
- Islamic Product Analytics: Shariah-compliant performance measurement

FSDM: PRDCT_PRFMNC (Performance), PRDCT_PRCNG (Pricing), ONBRDNG (Onboarding), PRDCT_PRFTBLTY (Profitability)""",

    12: """PRODUCT MANAGEMENT — Interesting Facts
Pakistan Banking Product Facts:
- Products per customer: 1.8 (vs. 4-6 global benchmark) — PKR 15-25B cross-sell opportunity for top-5 bank
- SBP product approval (NOC) adds 4-8 weeks to launch. Islamic: Shariah board + SBP = dual gate
- 30M+ debit cards, 3M+ credit cards — but card activation rates only 50-60% (vs. 80%+ global)
- Digital product adoption: 30M+ mobile banking users, growing 40% YoY
- Islamic products: 25%+ of deposits, growing 20% annually — every conventional product needs Islamic equivalent
- CASA ratio: ~47% — zero/low cost deposits are the single most profitable product category
- Digital onboarding dropoff: 40-60% at NADRA biometric step — fixing this is highest-ROI product initiative

FSDM: All Product domain entities. BVF: 12 sub-capabilities across 6 sub-domains.""",

    13: """PRODUCT MANAGEMENT FOR BANKING
Pakistan banking product management priorities:
1. DIGITAL-FIRST: Account opening via mobile (NADRA biometric), instant nano-lending, QR payments
2. ISLAMIC PRODUCTS: Islamic credit card (commodity Murabaha-based), Islamic home finance (Diminishing Musharaka), Islamic digital savings
3. FINANCIAL INCLUSION: Asaan Mobile Account (simplified KYC), agent banking products, micro-insurance (Takaful)
4. ROSHAN DIGITAL: Products for overseas Pakistanis — Naya Pakistan Certificate, Roshan Apna Ghar
5. GREEN BANKING: SBP green taxonomy — green lending products, green deposits (linked to ESG)
6. SME PRODUCTS: SBP refinance schemes (concessional rates), supply chain finance, digital SME onboarding

Product growth areas: Digital (fastest), Islamic (20%+ annual), Roshan Digital ($7B+ since 2020), SME (SBP priority)

FSDM: PRDCT, PRDCT_TYP, DGTL_PRDCT, ISLMC_PRDCT, INCLN_PRDCT""",

    14: """PRODUCT MANAGEMENT — Are You Able To...
Pakistan banking product management assessment questions:
- Can you calculate true profitability for each product in your portfolio?
- Can you track product performance across all channels (branch, ATM, digital, agent)?
- Can you identify which products drive cross-sell and which cannibalize?
- Can you optimize deposit pricing (rate) vs. volume trade-off analytically?
- Can you measure product onboarding success rate by channel?
- Can you predict which products a customer is likely to close next?
- Can you launch a new product variant in under 4 weeks?

Each question maps to BVF sub-capabilities and FSDM entities. Most Pakistan banks will answer 'No' or 'Partially' — this defines the capability gap.

FSDM: PRDCT_PRFTBLTY, PRDCT_PRFMNC, ONBRDNG, PRDCT_RTNTN, PRDCT_DVLPMNT""",

    15: """PRODUCT MANAGEMENT — Business Challenges
Pakistan-specific product management challenges:
- SBP product approval process adds 4-8 weeks to launch timeline
- Islamic products require dual approval: Shariah board + SBP
- 4-5 core banking systems make product configuration inconsistent
- Pricing constrained: deposit floor rate (SBP), lending ceiling (KIBOR + risk), card rates (SBP caps under discussion)
- Product data fragmented across core banking, cards, trade finance, Islamic sub-systems
- Products per customer: 1.8 (vs 4-6 global) — cross-sell capability severely underdeveloped
- Digital onboarding: 40-60% dropoff at NADRA verification
- Product profitability: <15% of banks can calculate true product-level P&L

FSDM: Integration challenges span Product, Agreement, Transaction, Channel, Pricing domains.""",

    16: """VALUE OF PRODUCT MANAGEMENT
Value proposition for Pakistan banks:
- CROSS-SELL: Increase products per customer from 1.8 to 2.5+ (PKR 5-10B revenue opportunity)
- PROFITABILITY: Identify and fix unprofitable products (repricing saves PKR 2-5B)
- PRICING: Dynamic pricing optimizing NII (deposit cost vs. lending margin — PKR 3-5B)
- ONBOARDING: Fix digital dropoff (40% → 20%) — reduces acquisition cost by PKR 1-2B
- ISLAMIC: Islamic product parity analytics — capture 25%+ and growing deposit market
- RETENTION: Product closure prediction — 15-25% reduction in product attrition

Estimated annual value: PKR 15-25B for a top-5 Pakistan bank.
ROI on product analytics investment: 5-10x within 3 years.

FSDM: Value delivered through integrated Product domain data across all 12 sub-capabilities.""",

    17: """FUNCTIONAL INSIGHTS AND CAPABILITIES — Section Divider
Transitioning from Product Management domain overview to detailed capability slides.

The following slides cover 6 sub-domains:
1. Product Development (Slides 19-26): 3 capability pairs
2. Product Introduction (Slides 27-32): 2 capability pairs
3. Price and Promotion (Slides 33-38): 2 capability pairs
4. Product End-of-Life (Slides 39-42): 1 capability pair
5. Product Performance (Slides 43-54): 5 capability pairs

Each pair = Capability Detail + Maturity Assessment.
Pakistan banks should use BACR assessment to evaluate current maturity and define the roadmap.

FSDM: All 12 sub-capabilities with FSDM entity mappings.""",

    # Part B: Product Development (slides 19-26, 0-indexed 18-25)
    18: """PRODUCT DEVELOPMENT — Why Important
Product development in banking means: designing new deposit variants, structuring new lending products, creating card propositions, launching digital services, and innovating Islamic products.

Pakistan context: SBP encourages financial inclusion products (Basic Banking Account, Asaan Mobile Account, digital lending), Islamic product innovation, and green banking products.

Key product development opportunities in Pakistan:
1. DIGITAL-FIRST: Account opening via mobile, instant nano-lending (PKR 5K-50K)
2. ISLAMIC: Islamic credit card (commodity Murabaha), Islamic home finance (DM), Islamic digital savings
3. FINANCIAL INCLUSION: Asaan Mobile Account, agent banking, micro-insurance (Takaful)
4. ROSHAN DIGITAL: Products for overseas Pakistanis ($7B+ since 2020)
5. GREEN BANKING: SBP green taxonomy — green lending, green deposits

SBP product approval process:
1. Internal product committee approval → 2. Shariah board (if Islamic) → 3. SBP NOC submission → 4. SBP review (4-8 weeks) → 5. CBS configuration → 6. Launch

FSDM: PRDCT_DVLPMNT, MKT_RSRCH, CSTMR_ND, PRDCT_PPLN""",

    19: """PRODUCT DEVELOPMENT — How
Solution approach for product development analytics:
1. MARKET RESEARCH: RAAST/transaction pattern analysis, competitor monitoring, customer need identification
2. DEMAND MODELING: Segment-level demand forecasting, financial inclusion gap analysis
3. BUSINESS CASE: Revenue projection, cannibalization assessment, channel capacity planning
4. PROTOTYPING: A/B testing for digital, regional pilot for branch, champion/challenger for pricing
5. REGULATORY: SBP NOC preparation with data-driven product documentation

FSDM: PRDCT_DVLPMNT, MKT_RSRCH, CSTMR_ND, DMND_MDLNG""",

    20: """RESEARCH NEW PRODUCT DEVELOPMENT OPPORTUNITIES — Capability Detail
Research unmet needs in Pakistan's 220M population — only 30% banked, 60% have mobile phones.
Analyze RAAST transaction patterns to identify product opportunities.
Products for financially excluded: basic accounts, micro-savings, nano-lending via mobile.
Islamic product gap analysis — Shariah-compliant alternatives for every conventional product.

FSDM: PRDCT_DVLPMNT, MKT_RSRCH, CSTMR_ND, PRDCT_PPLN, RAAST_TXN, DGTL_BHVR""",

    21: """RESEARCH NEW PRODUCT DEVELOPMENT — Maturity Assessment
Current State (Developing): Ad hoc product research. New products based on competitor copying rather than data-driven customer need analysis.
Target State (Innovating): RAAST/transaction pattern mining, social media sentiment, competitor monitoring, financial inclusion gap analytics driving product pipeline.
Pakistan Context: With 70% unbanked, the product development opportunity is enormous but requires analytical understanding of unmet needs.

FSDM: PRDCT_DVLPMNT, MKT_RSRCH, CSTMR_ND, PRDCT_PPLN""",

    22: """PREDICT BUSINESS IMPACT (DISCOVERY) — Capability Detail
Predict impact of new product launch on: existing product balances (cannibalization), channel capacity (branch/call center load), revenue (NII + fees), operational risk, and customer experience.
Pakistan example: launching a digital-only savings account — will it cannibalize branch savings? What's the expected adoption curve? Impact on call center?
Source Data: Historical product launch data (past 5 years), similar product performance at competitors, customer segment adoption patterns from FSDM.

FSDM: PRDCT_IMPCT, CNNBLZTN, CHNL_CPCTY, RVNU_FRCST""",

    23: """PREDICT BUSINESS IMPACT — Maturity Assessment
Current State (Developing): Product launches based on qualitative judgment. No systematic cannibalization analysis. Channel capacity not modeled.
Target State (Innovating): Quantitative impact modeling with cannibalization, channel capacity, and revenue forecasting for every new product.
Pakistan Context: With SBP NOC adding 4-8 weeks, accurate business case preparation is essential — cannot afford to launch and fail.

FSDM: PRDCT_IMPCT, CNNBLZTN, CHNL_CPCTY, RVNU_FRCST""",

    24: """DEVELOP AND TEST PRODUCT PROTOTYPES — Capability Detail
Test product features, pricing, and packaging through A/B testing and pilot programs before full-scale launch.
Pakistan banking: pilot new products in select branches/regions before national rollout.
Test pricing variants (rate sensitivity), feature variants (minimum balance, fee structure), and channel variants (branch-only vs. digital).
Data: FSDM transaction data for pilot vs. control groups, mobile app engagement, branch sales, customer feedback.
Methodology: A/B testing for digital products, regional pilot for branch products, champion/challenger for pricing.

FSDM: PRDCT_TST, AB_TST, PLT_PRGRM, CMPGN_CNTRL""",

    25: """DEVELOP AND TEST PRODUCT PROTOTYPES — Maturity Assessment
Current State (Developing): No systematic product testing. Products launched nationally without pilot. Pricing set by committee without testing.
Target State (Innovating): A/B testing framework for digital, regional pilots for branch, champion/challenger for pricing. Evidence-based launch decisions.
Pakistan Context: With 40-60% digital onboarding dropoff, A/B testing capability is essential to optimize product features before full investment.

FSDM: PRDCT_TST, AB_TST, PLT_PRGRM""",

    # Part C: Product Introduction (slides 27-32, 0-indexed 26-31)
    26: """PRODUCT INTRODUCTION — Why Important
Product introduction in Pakistan banking faces unique challenges: SBP approval (4-8 weeks), dual approval for Islamic products (Shariah board + SBP), multi-system configuration (4-5 core banking systems), multi-channel training (16,000+ branches + digital + ATM).
Speed-to-market determines competitive success against fintechs that can launch products in days vs. banks' 8-16 weeks.
Pakistan digital product launch speed: 8-16 weeks (SBP approval + IT) vs. 2-4 weeks at global digital banks.

FSDM: PRDCT_INTRO, GTM, CHNL_RDNSS, ONBRDNG""",

    27: """PRODUCT INTRODUCTION — How
Solution approach:
1. GTM PLANNING: Multi-channel launch plan with branch training, digital UX, and CBS configuration
2. SBP INTEGRATION: NOC timeline embedded in product launch methodology
3. CHANNEL READINESS: Training, collateral (Urdu + English), system configuration across all channels
4. ONBOARDING ANALYTICS: Funnel tracking from application to first transaction
5. LAUNCH MONITORING: Real-time product adoption dashboard tracking sales by channel

FSDM: PRDCT_INTRO, GTM, CHNL_RDNSS, ONBRDNG, LNCH_MNTRNG""",

    28: """PLAN PRODUCT INVESTMENTS / GO-TO-MARKET — Capability Detail
Plan GTM for Pakistan banking: branch staff training (16,000+ branches), core banking system configuration (4-5 systems), SBP reporting setup, marketing collateral (Urdu + English), Islamic Shariah compliance documentation.
Forecast resource requirements: call center staffing for inquiries, branch teller training, digital UX readiness, card production pipeline.
Source Data: Historical product launch data, branch capacity data, call center staffing model, digital channel capacity, card production lead times.

FSDM: GTM, CHNL_RDNSS, RSRC_PLN, LNCH_PLN""",

    29: """PLAN PRODUCT INVESTMENTS / GTM — Maturity Assessment
Current State (Developing): Product launches ad hoc. Branch training inconsistent. CBS configuration delayed. No integrated launch methodology.
Target State (Innovating): Integrated multi-channel launch methodology with SBP timeline, parallel workstreams, and real-time readiness dashboard.
Pakistan Context: With 16,000+ branches and 4-5 CBS, product introduction requires systematic project management — most banks do this manually.

FSDM: GTM, CHNL_RDNSS, RSRC_PLN, LNCH_PLN""",

    30: """IMPROVE PRODUCT ONBOARDING — Capability Detail
PAKISTAN PRODUCT ONBOARDING:
Current onboarding timelines (average):
- Current account: 3-7 days (branch), 1-3 days (digital — if NADRA works)
- Savings account: Same as current
- Personal loan: 7-21 days (documentation + credit check + approval)
- Home finance: 30-90 days (property valuation + legal + SBP guidelines)
- Credit card: 7-14 days (ECIB check + card production + courier)
- Digital onboarding: 10-30 minutes (if NADRA biometric succeeds on first attempt)

Key dropout points:
1. NADRA biometric verification failure (30-40% first-attempt failure rate)
2. Document upload issues (photo quality, format)
3. Minimum balance requirement confusion
4. Waiting for physical card/cheque book delivery
5. First transaction complexity

FSDM: ONBRDNG, ONBRDNG_JRNY, DRPOFF, ACTVTN, FRST_TXN""",

    31: """IMPROVE PRODUCT ONBOARDING — Maturity Assessment
Current State (Developing): Digital onboarding exists but with 40-60% dropoff. No funnel analytics. Loan origination still paper-based at many banks.
Target State (Innovating): Optimized digital onboarding with <20% dropoff, real-time funnel analytics, automated loan origination <7 days.
Pakistan Context: SadaPay and NayaPay can onboard in minutes. Banks must match this speed or lose the digital-first segment.

FSDM: ONBRDNG, ONBRDNG_JRNY, DRPOFF, ACTVTN""",

    # Part D: Price and Promotion (slides 33-38, 0-indexed 32-37)
    32: """PRICE AND PROMOTION — Why Important
Pricing in Pakistan banking operates within a constrained environment:
- Deposit floor: SBP minimum deposit rate (~11-13% on savings)
- Lending benchmark: KIBOR + risk spread (currently KIBOR ~17-18%)
- Card rates: SBP caps under discussion
- Fee income: Subject to market competition and SBP fair pricing guidelines

Competition is on spread, fees, and service quality — not on base rates.
Dynamic pricing opportunity: large depositor relationship pricing, risk-based lending, promotional card pricing.

FSDM: PRCNG, PRDCT_PRCNG, CMPTTV_PRCNG, DYNMC_PRCNG""",

    33: """PRICE AND PROMOTION — How
Solution approach:
1. COMPETITIVE MONITORING: Real-time dashboard of deposit rates, lending rates, card fees across peers
2. PRICING OPTIMIZATION: Deposit rate-vs-volume trade-off analytics within SBP constraints
3. RISK-BASED LENDING: KIBOR + variable risk premium using ECIB/internal ratings
4. DYNAMIC PRICING: Relationship pricing for large depositors, promotional card rates
5. FEE ANALYTICS: Fee income optimization across products and channels

FSDM: PRCNG, CMPTTV_PRCNG, RSK_PRCNG, DYNMC_PRCNG, FEE_ANLYTCS""",

    34: """ANALYZE COMPETITIVE PRICING — Capability Detail
Monitor competitive pricing across Pakistan banking: deposit rates (SBP publishes), lending rates (KIBOR + spreads), card fees (annual, late payment, FX markup), trade finance pricing (LC commission rates), digital transaction fees (RAAST/IBFT charges).
Pakistan-specific: SBP minimum deposit rate creates floor. KIBOR creates lending benchmark. Competition is on spread, fees, and service quality.
Source Data: SBP published rate data, competitor rate sheets, KIBOR curves, card fee comparison, customer rate sensitivity from FSDM.

FSDM: CMPTTV_PRCNG, MKT_RT, KIBOR, DPST_RT, LNDNG_SPRD""",

    35: """ANALYZE COMPETITIVE PRICING — Maturity Assessment
Current State (Developing): Manual competitive monitoring. Rate sheets collected periodically. No systematic fee benchmarking.
Target State (Innovating): Real-time competitive pricing dashboard with automated alerting, rate sensitivity analytics, and fee optimization recommendations.
Pakistan Context: With SBP publishing deposit rates and KIBOR setting lending benchmarks, the competition is on spread optimization and fee income — requiring granular pricing analytics.

FSDM: CMPTTV_PRCNG, MKT_RT, KIBOR""",

    36: """MANAGE DYNAMIC PRICING — Capability Detail
Dynamic pricing in Pakistan banking: adjust deposit rates for large depositors (relationship pricing), risk-based lending (KIBOR + variable risk premium based on ECIB/internal rating), promotional card rates (balance transfer offers, spend-based fee waivers).
Constraints: SBP minimum deposit rate, SBP maximum markup on credit cards (under discussion), KIBOR-linked lending floors.
Outcome: Optimize NII by dynamically adjusting deposit gathering cost vs. lending margin, retain large depositors through competitive pricing while protecting margin.

FSDM: DYNMC_PRCNG, RLTNSHP_PRCNG, RSK_PRCNG, PROMO_PRCNG""",

    37: """MANAGE DYNAMIC PRICING — Maturity Assessment
Current State (Developing): Pricing set by committee. Large depositor rates negotiated manually. No systematic risk-based lending pricing.
Target State (Innovating): Automated relationship pricing, real-time risk-adjusted lending rates, promotional pricing engine with eligibility automation.
Pakistan Context: Banks losing large depositors to competitors offering +25-50bps. Dynamic pricing can retain while optimizing NII.

FSDM: DYNMC_PRCNG, RLTNSHP_PRCNG, RSK_PRCNG""",

    # Part E: Product End-of-Life (slides 39-42, 0-indexed 38-41)
    38: """PRODUCT END-OF-LIFE — Why Important
PAKISTAN PRODUCT LIFECYCLE:
Product sunset scenarios in Pakistan banking:
1. LEGACY PRODUCTS: Old savings schemes with rates no longer competitive — migrate to current
2. POST-MERGER RATIONALIZATION: Two banks merge, combined 300+ product variants → rationalize to 80-100
3. REGULATORY CHANGE: SBP changes product rules (e.g., PLS savings rate) — migrate affected customers
4. CONVENTIONAL → ISLAMIC: Growing demand for Islamic — proactive migration of willing customers
5. BRANCH → DIGITAL: Sunset branch-only products, launch digital equivalents
6. CARD PRODUCT UPGRADE: Migrate classic cardholders to gold/platinum based on usage patterns

Key risk: Product sunset can trigger customer attrition if poorly managed. Must offer equivalent/better alternative and communicate proactively.

FSDM: PRDCT_LFCYCL, PRDCT_SNST, PRDCT_MGRTN, CSTMR_MGRTN""",

    39: """PRODUCT END-OF-LIFE — How
Solution approach:
1. PRODUCT LIFECYCLE ANALYTICS: Track product from launch through growth, maturity, decline to sunset
2. MIGRATION PLANNING: Data-driven migration plan with customer communication strategy
3. IMPACT ASSESSMENT: Estimate attrition risk, channel impact, and revenue impact from sunset
4. REGULATORY COMPLIANCE: SBP-mandated product change management
5. POST-MIGRATION: Monitor migrated customers for satisfaction and attrition signals

FSDM: PRDCT_LFCYCL, PRDCT_SNST, PRDCT_MGRTN, CSTMR_MGRTN, ATTRTN""",

    40: """DETERMINE PRODUCT MIGRATION — Capability Detail
Product migration in Pakistan banking: sunset legacy products (old savings schemes, discontinued FD tenors), migrate customers from conventional to Islamic products (growing demand), consolidate product variants across merged entities.
Pakistan context: Bank mergers create product rationalization need — merged bank may have 200+ product variants that need consolidation to 50-80.
Regulatory sunset: SBP discontinued products (e.g., old foreign currency schemes) require customer migration.
Source Data: Product-level balance and customer data from FSDM, customer communication preference, channel usage, product switching history.
Outcome: Smooth migration preserving customer relationships, minimal attrition during sunset, optimized product portfolio.

FSDM: PRDCT_MGRTN, CSTMR_MGRTN, PRDCT_SNST, PRDCT_RTNLZTN""",

    41: """DETERMINE PRODUCT MIGRATION — Maturity Assessment
Current State (Developing): Product rationalization done manually post-merger. No systematic migration methodology. Customer communication reactive.
Target State (Innovating): Analytics-driven migration with attrition prediction, proactive communication, and automated account migration.
Pakistan Context: With ongoing bank consolidation (SBP encouraging mergers), product rationalization capability is increasingly critical.

FSDM: PRDCT_MGRTN, CSTMR_MGRTN, PRDCT_SNST""",

    # Part F: Product Performance (slides 43-54, 0-indexed 42-53)
    42: """PRODUCT PERFORMANCE — Why Important
Product performance analytics is the most data-intensive sub-domain with 5 capability pairs.
Pakistan priority: Most banks cannot answer basic questions about product performance:
- Which products are profitable and which are cross-subsidized?
- Which channels drive the best product sales at lowest cost?
- How do we compare to peer banks on key product metrics?
- Which customer segments are most receptive to which products?
- Which products are at risk of customer closure?

Without product performance analytics, PKR 15-25B annual value opportunity remains untapped.

FSDM: PRDCT_PRFMNC, PRDCT_PRFTBLTY, CHNL_PRFMNC, BNCHMARK, SGMNTTN, RTNTN""",

    43: """PRODUCT PERFORMANCE — How
Solution approach:
1. PERFORMANCE DASHBOARD: Product sales, balances, revenue, cost by product and channel — real-time
2. PROFITABILITY ENGINE: FTP-adjusted NII + fees - costs - provisions at product level
3. CROSS-CHANNEL: Which products perform best on which channels? Cost-to-serve by channel
4. BENCHMARKING: vs. SBP industry data, peer banks, international benchmarks
5. SEGMENTATION: Product usage behavioral segments driving targeted recommendations
6. RETENTION: Product closure prediction with targeted retention offers

FSDM: PRDCT_PRFMNC, PRDCT_PRFTBLTY, CHNL_PRFMNC, BNCHMARK, SGMNTTN, RTNTN""",

    44: """ANALYZE CROSS CHANNEL PERFORMANCE — Capability Detail
Track product sales and usage across all Pakistan banking channels: branch (16,000+), ATM (16,000+), mobile app, internet banking, USSD, call center, agent network (branchless banking), RAAST.
Key insight: Which products sell best on which channels? Personal loans sell 80% via branch but digital origination growing. CASA opens increasingly digital. Cards are branch-initiated but digitally activated.
Source Data: FSDM transaction data by channel, product sales by channel, mobile app product funnel, branch sales register, agent banking data.

FSDM: CHNL_PRFMNC, PRDCT_SLS, CHNL_CST, DGTL_MGRTN""",

    45: """ANALYZE CROSS CHANNEL PERFORMANCE — Maturity Assessment
Current State (Developing): Product sales tracked by channel but not integrated. No unified cross-channel view. Digital vs. branch attribution unclear.
Target State (Innovating): Unified cross-channel product performance dashboard with cost-to-serve, digital migration tracking, and channel optimization recommendations.
Pakistan Context: 8+ channels with no unified view — branches, ATMs, mobile, internet, USSD, call center, agent, RAAST. Unified analytics is foundational.

FSDM: CHNL_PRFMNC, PRDCT_SLS, CHNL_CST""",

    46: """COMPARE TO PRODUCT BENCHMARKS — Capability Detail
Benchmark Pakistan banking products against: SBP published industry data (deposit rates, advance rates, NIM), peer group performance (top 5/10 banks), international benchmarks (Gulf banks, SE Asian banks).
Key product benchmarks: CASA ratio (industry ~47%), card activation rate, loan-to-deposit ratio, fee income per product, digital adoption rate.
Source Data: SBP statistical bulletins, annual reports of peer banks, industry surveys (PBA), FSDM internal product performance data.

FSDM: BNCHMARK, PRDCT_KPI, INDSTY_DATA, PEER_CMP""",

    47: """COMPARE TO PRODUCT BENCHMARKS — Maturity Assessment
Current State (Developing): Annual report comparison only. No systematic benchmarking. International comparison absent.
Target State (Innovating): Automated benchmark dashboard with SBP data feed, peer analysis, and international comparison. Underperformance alerts.
Pakistan Context: SBP publishes comprehensive industry data making automated benchmarking feasible — but few banks have implemented it.

FSDM: BNCHMARK, PRDCT_KPI, INDSTY_DATA""",

    48: """MANAGE PRODUCT PROFITABILITY — Capability Detail — CRITICAL
PAKISTAN PRODUCT PROFITABILITY:
Typical profitability profile:
| Product | Revenue | Cost | Profit | Notes |
| Corporate working capital | High NII + fees | Low per-unit | Very profitable | Relationship-based |
| Trade finance (LC/LG) | High fee income | Moderate | Highly profitable | Pakistan's #1 fee source |
| Term deposits (large) | Negative NII (high rate) | Low | Negative | Cost of funds |
| CASA | High NII (low cost) | Moderate | Very profitable | Foundation of banking profit |
| Personal loans | High NII | Moderate | Profitable (if NPL controlled) | ECIB scoring critical |
| Credit cards | Interchange + interest | High (fraud, rewards) | Mixed | Top 20% customers profitable |
| Basic savings | Minimal NII | High (branch service) | Unprofitable | Financial inclusion mandate |
| Home finance | Moderate NII | High (origination) | Break-even → profitable | Long tenure needed |
| Digital products | Fee income | Very low | Highly profitable per txn | Scale is key |

Key insight: Product profitability varies enormously. Without analytics, banks cross-subsidize unknowingly.

FSDM: PRDCT_PRFTBLTY, PRDCT_RVNU, PRDCT_CST, FTP
BVF Reference: File 08 EPM > Profitability Modelling, File 09 > Profitability Analytics""",

    49: """MANAGE PRODUCT PROFITABILITY — Maturity Assessment
Current State (Developing): No product-level P&L. Revenue tracked at product level but costs not allocated. FTP not product-specific.
Target State (Innovating): True product profitability with FTP-adjusted NII, ABC cost allocation, IFRS 9 provision allocation, revealing cross-subsidy flows.
Pakistan Context: <15% of banks can calculate true product profitability. This is the single highest-value product analytics capability.

FSDM: PRDCT_PRFTBLTY, FTP, ABC, PRVSN""",

    50: """PRODUCE BEHAVIORAL SEGMENTATIONS — Capability Detail
Segment banking customers by product usage behavior:
- Heavy transactors: High RAAST/ATM usage, multiple daily transactions
- Savers: FD-focused, low transaction frequency, high balance
- Borrowers: Multiple lending products, active credit utilization
- Digital-first: Mobile app primary channel, rarely visits branch
- Branch-loyal: Branch-only, no digital adoption
- Dormant: No activity 90+ days — reactivation opportunity
- Islamic-only: Only uses Shariah-compliant products

Islamic preference segmentation: customers who only use Islamic products vs. mixed vs. conventional-only.
Source Data: FSDM transaction data, product holding data, channel usage, Islamic/conventional preference, mobile app engagement.

FSDM: BHVRL_SGMNT, PRDCT_USGN, CHNL_PRFRNC, ISLMC_PRFRNC""",

    51: """PRODUCE BEHAVIORAL SEGMENTATIONS — Maturity Assessment
Current State (Developing): Basic demographic segmentation only. No behavioral product usage segmentation. Islamic/conventional preference not tracked.
Target State (Innovating): Multi-dimensional behavioral segmentation driving targeted product recommendations, Islamic product opportunity sizing, and product bundle optimization.
Pakistan Context: Islamic preference segmentation is unique to Pakistan/ME markets and critical — 25%+ of deposits are Islamic.

FSDM: BHVRL_SGMNT, PRDCT_USGN, ISLMC_PRFRNC""",

    52: """PREDICT PRODUCT/SERVICE RETENTION — Capability Detail
Predict which products each customer is likely to close: FD non-renewal (rate shopping), CASA closure (salary switch), card cancellation (dormancy), loan prepayment (refinancing).
Pakistan-specific: salary account switching is highest product-level churn risk. When employer switches bank, entire product bundle is at risk.
Source Data: FSDM — product tenure, balance trend, transaction frequency change, salary credit continuity, competitor rate data, customer life event indicators.
Outcome: Product retention alerts 30 days before closure, targeted retention offers by product type.

FSDM: PRDCT_RTNTN, CHURN_PRDCTN, SLRY_MNTRNG, RT_SHPPNG""",

    53: """PREDICT PRODUCT/SERVICE RETENTION — Maturity Assessment
Current State (Developing): Reactive — customer closes product, then bank reacts. No early warning. No product-level churn prediction.
Target State (Innovating): Predictive retention with 30-day early warning, automated retention offers, salary switching detection, rate shopping alerts.
Pakistan Context: Salary account switching can trigger loss of 3-5 products per customer. Early detection saves entire product bundles.

FSDM: PRDCT_RTNTN, CHURN_PRDCTN, SLRY_MNTRNG""",
}


# Speaker notes for the 3 new slides
NEW_SLIDE_NOTES = {
    'new_a': """PRODUCT MANAGEMENT — DOMAIN DASHBOARD

This slide provides a high-level comparison of Product Management maturity across global, regional (South Asia & Middle East), and Pakistan banking markets.

Key Takeaways:
- Pakistan banks average 1.8 products per customer vs. 4-6 globally — massive cross-sell opportunity
- Product profitability analytics available at <15% of Pakistan banks vs. 70%+ of large global banks
- Dynamic pricing capability at <10% of Pakistan banks vs. 50%+ at digital banks globally
- Product A/B testing standard at digital banks globally but <5% of Pakistan banks
- Islamic product suite is a unique Pakistan/ME requirement — 25%+ of deposits
- Digital product launch speed: 8-16 weeks in Pakistan vs. 2-4 weeks at global digital banks

BVF Sub-capabilities: 12 across 6 sub-domains
FSDM Entities: Product, Agreement, Party, Transaction, Channel, Pricing
BACR Questions: ~60 | Maturity Focus: Developing → Innovating

Investment Justification: Moving from Developing to Innovating can deliver PKR 15-25B annual value through cross-sell improvement, pricing optimization, and product portfolio rationalization.""",

    'new_b': """PAKISTAN BANKING — PRODUCT MANAGEMENT LANDSCAPE

Product Portfolio (typical large Pakistan bank):
DEPOSIT PRODUCTS: Current Account (zero interest, CASA), Savings Account (SBP min rate ~11-13%), Fixed Deposit (3M-5Y), Foreign Currency Account, Roshan Digital Account, Basic Banking Account, Islamic: Mudaraba Savings, Mudaraba Term Deposit, Wakalah Pool

LENDING PRODUCTS: Corporate (working capital, term finance, project finance, trade finance LC/LG/SBLC), Consumer (personal loan, home finance, auto finance, credit card revolving), SME (running finance, term finance), Agriculture (crop, livestock, dairy), Islamic (Murabaha, Ijarah, Diminishing Musharaka, Running Musharaka)

CARD PRODUCTS: Debit (30M+ cards), Credit (3M+ cards), Prepaid (travel, remittance, payroll)

DIGITAL PRODUCTS: Mobile banking, Internet banking, RAAST P2P/P2M, QR payments, Digital wallet (JazzCash/Easypaisa), Agent banking

Key Challenges:
- Products per customer: 1.8 (vs. 4-6 global)
- SBP product approval: 4-8 weeks for new products
- Islamic dual approval: Shariah board + SBP
- Pricing rigidity: deposit floor (SBP), lending ceiling (KIBOR)
- Product data fragmented across 4-5 core banking systems""",

    'new_c': """PRODUCT MANAGEMENT — IMPLEMENTATION ROADMAP

Phase 1: Product Intelligence Foundation (0-6 months) — PKR 40-80M
- Product profitability analytics (from FSDM Customer Profitability Engine)
- Product performance dashboard (sales, balances, revenue, cost by product)
- Cross-channel product sales tracking
- Product benchmark analysis (vs. SBP data)
Quick Win: Identify bottom 10% unprofitable products

Phase 2: Product Optimization (6-18 months) — PKR 80-150M
- Competitive pricing analytics (deposit rates, lending spreads, card fees)
- Product onboarding journey optimization (reduce dropoff 40-60% to <20%)
- Cross-sell propensity models by product affinity
- Islamic product performance parity analysis
Expected: Increase products/customer from 1.8 to 2.5

Phase 3: Advanced Product Management (18-36 months) — PKR 120-250M
- Dynamic pricing engine (risk-adjusted lending, competitive deposit pricing)
- A/B testing framework for product features and pricing
- Product lifecycle management (launch → growth → maturity → sunset)
- Predictive product retention (which customers will close which products?)
Expected: 15-25% improvement in product revenue per customer

Total 3-Year Investment: PKR 240-480M (~USD 0.9-1.7M)
Estimated Annual Value: PKR 15-25B from cross-sell, pricing, and portfolio optimization""",
}


# ─────────────────────────────────────────────
# Clean orphaned slides
# ─────────────────────────────────────────────

def clean_pptx(input_path, output_path):
    """Remove orphaned slide XMLs from the PPTX ZIP."""
    with zipfile.ZipFile(input_path, 'r') as zin:
        pres_xml = zin.read('ppt/presentation.xml')
        root = ET.fromstring(pres_xml)
        pns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
               'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

        sld_rids = set()
        for sldId in root.findall('.//p:sldIdLst/p:sldId', pns):
            rid = sldId.attrib.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', '')
            sld_rids.add(rid)

        rels_xml = zin.read('ppt/_rels/presentation.xml.rels')
        rels_root = ET.fromstring(rels_xml)
        referenced_slides = set()
        for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
            rid = rel.attrib.get('Id', '')
            target = rel.attrib.get('Target', '')
            if rid in sld_rids and 'slides/' in target:
                slide_file = target.replace('slides/', '')
                referenced_slides.add(slide_file)

        print(f"  Referenced slides: {len(referenced_slides)}")

        orphaned = set()
        for name in zin.namelist():
            if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                slide_file = name.replace('ppt/slides/', '')
                if slide_file not in referenced_slides:
                    orphaned.add(slide_file.replace('.xml', ''))

        print(f"  Orphaned slides to remove: {len(orphaned)}")

        with zipfile.ZipFile(output_path, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                skip = False
                for orphan in orphaned:
                    if f'slides/{orphan}.xml' in item.filename or f'slides/_rels/{orphan}.xml.rels' in item.filename:
                        skip = True
                        break
                if not skip:
                    data = zin.read(item.filename)
                    zout.writestr(item, data)

    print(f"  Cleaned file saved to {output_path}")


# ─────────────────────────────────────────────
# Main execution
# ─────────────────────────────────────────────

def main():
    print("=" * 60)
    print("FILE 10: PRODUCT MANAGEMENT UPDATE")
    print("=" * 60)

    # ── Step 0: Clean orphaned slides ──
    print("\n── Step 0: Cleaning orphaned slides ──")
    clean_pptx(INPUT_FILE, CLEAN_FILE)

    print(f"\nLoading presentation...")
    prs = Presentation(CLEAN_FILE)
    slides = list(prs.slides)
    print(f"  Loaded {len(slides)} slides")

    # ── Step 1: Apply industry language replacements ──
    print(f"\n── Step 1: Applying retail/telco → banking language replacements ──")
    lang_count = 0
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        new_text = apply_language_replacements(run.text)
                        if new_text != run.text:
                            lang_count += 1
                            run.text = new_text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                new_text = apply_language_replacements(run.text)
                                if new_text != run.text:
                                    lang_count += 1
                                    run.text = new_text
    print(f"  Applied {lang_count} language replacements")

    # ── Step 2: Replace Teradata branding ──
    print(f"\n── Step 2: Replacing Teradata branding ──")
    teradata_count = 0
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'Teradata Business Value Framework' in run.text:
                            run.text = run.text.replace('Teradata Business Value Framework', 'Banking Business Value Framework')
                            teradata_count += 1
                        elif 'Teradata' in run.text:
                            run.text = run.text.replace('Teradata', 'Enterprise Analytics Platform')
                            teradata_count += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if 'Teradata Business Value Framework' in run.text:
                                    run.text = run.text.replace('Teradata Business Value Framework', 'Banking Business Value Framework')
                                    teradata_count += 1
                                elif 'Teradata' in run.text:
                                    run.text = run.text.replace('Teradata', 'Enterprise Analytics Platform')
                                    teradata_count += 1
        # Also check and replace in notes
        if slide.has_notes_slide:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                for run in para.runs:
                    if 'Teradata Business Value Framework' in run.text:
                        run.text = run.text.replace('Teradata Business Value Framework', 'Banking Business Value Framework')
                        teradata_count += 1
                    elif 'Teradata' in run.text:
                        run.text = run.text.replace('Teradata', 'Enterprise Analytics Platform')
                        teradata_count += 1
    print(f"  Replaced {teradata_count} Teradata references")

    # ── Step 3: Remove owner names and outdated references ──
    print(f"\n── Step 3: Removing owner names and outdated references ──")
    owner_names = ['Tidell', 'Kong', 'Axon', 'Erdelyi', 'Marcom Central']
    owner_count = 0
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for name in owner_names:
                            if name in run.text:
                                run.text = run.text.replace(name, '')
                                owner_count += 1
                        if 'EKN Research' in run.text:
                            run.text = run.text.replace('EKN Research', 'Pakistan Banking Industry Data')
                            owner_count += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                for name in owner_names:
                                    if name in run.text:
                                        run.text = run.text.replace(name, '')
                                        owner_count += 1
                                if 'EKN Research' in run.text:
                                    run.text = run.text.replace('EKN Research', 'Pakistan Banking Industry Data')
                                    owner_count += 1
    print(f"  Removed {owner_count} owner/reference items")

    # ── Step 4: Add speaker notes to ALL 54 slides ──
    print(f"\n── Step 4: Adding speaker notes to all slides ──")
    for idx in range(len(slides)):
        if idx in SPEAKER_NOTES:
            set_notes(slides[idx], SPEAKER_NOTES[idx])
            print(f"  Slide {idx+1}: Speaker notes set")
        else:
            print(f"  Slide {idx+1}: No notes defined (skipped)")

    # ── Step 5: Enrich 13 capability detail slides ──
    print(f"\n── Step 5: Enriching 13 capability detail slides ──")
    for idx, enrichment in CAPABILITY_ENRICHMENTS.items():
        slide = slides[idx]
        table = get_table(slide)
        if not table:
            print(f"  Slide {idx+1}: WARNING - no table found!")
            continue

        nrows = len(table.rows)
        ncols = len(table.columns)

        # Standard layout: 4x3 — biz_obj [1,0], data [3,0], outcome [3,2]
        # Variant layout: 3x3 — biz_obj [1,0], data [2,0], outcome [2,2]
        if nrows >= 4:
            biz_row, data_row = 1, 3
        else:
            biz_row, data_row = 1, 2

        # Append to business objectives
        biz_lines = '\n'.join(enrichment['biz_obj'])
        append_cell_text(table.cell(biz_row, 0), '\n' + biz_lines)

        # Append to data & solution
        data_lines = '\n'.join(enrichment['data_sol'])
        append_cell_text(table.cell(data_row, 0), '\n' + data_lines)

        # Append to outcome
        outcome_lines = '\n'.join(enrichment['outcome'])
        append_cell_text(table.cell(data_row, 2), '\n' + outcome_lines)

        print(f"  Slide {idx+1}: Detail table enriched ({nrows}x{ncols})")

    # ── Step 6: Enrich 13 maturity slides ──
    print(f"\n── Step 6: Enriching 13 maturity slides ──")
    for idx, enrichment in MATURITY_ENRICHMENTS.items():
        slide = slides[idx]
        table = get_table(slide)
        if not table:
            print(f"  Slide {idx+1}: WARNING - no table found!")
            continue

        # Replace key areas [7,3] and benefits [7,4]
        try:
            replace_cell_text(table.cell(7, 3), enrichment['key_areas'])
            replace_cell_text(table.cell(7, 4), enrichment['benefits'])
        except IndexError:
            print(f"  Slide {idx+1}: WARNING - table cell index out of range, trying alternative cells")
            try:
                # Try finding the last row with key areas
                nrows = len(table.rows)
                ncols = len(table.columns)
                # Try last row, columns 3 and 4
                if nrows > 5 and ncols > 4:
                    replace_cell_text(table.cell(nrows-1, 3), enrichment['key_areas'])
                    replace_cell_text(table.cell(nrows-1, 4), enrichment['benefits'])
            except Exception as e:
                print(f"  Slide {idx+1}: ERROR - {e}")

        # Replace H1 2018
        replaced = replace_h1_2018(table)

        # Remove engagement text from groups
        removed = remove_engagement_text(slide)

        status = []
        if replaced:
            status.append("replaced H1 2018")
        if removed:
            status.append("removed engagement text")
        extra = f" + {' + '.join(status)}" if status else ""
        print(f"  Slide {idx+1}: Maturity enriched{extra}")

    # ── Step 7: Add Pakistan context to key overview slides ──
    print(f"\n── Step 7: Adding Pakistan context to overview slides ──")

    # Slide 8-9 (0-indexed 7-8): Products vary by industry / Examples
    for idx in [7, 8]:
        slide = slides[idx]
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text
                if any(kw in text for kw in ['product', 'Product', 'service', 'Service', 'industry', 'Industry']):
                    tf = shape.text_frame
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = ("In banking, a 'product' encompasses: deposit accounts (CASA, term deposits), "
                               "lending facilities (personal loans, home finance, corporate working capital), "
                               "cards (debit, credit, prepaid), trade finance (LC, LG, SBLC), treasury products, "
                               "digital services (mobile banking, RAAST), and Islamic equivalents "
                               "(Murabaha, Ijarah, Diminishing Musharaka)")
                    run.font.size = Pt(9)
                    run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                    print(f"  Slide {idx+1}: Added banking product context")
                    break

    # Slide 10 (0-indexed 9): Why Important — add banking context
    slide = slides[9]
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text.lower()
            if 'important' in text or 'why' in text:
                tf = shape.text_frame
                additions = [
                    "Product management is the engine of banking revenue — every PKR of income comes from a product",
                    "Pakistan banks average 1.8 products per customer vs. 4-6 globally — massive cross-sell opportunity",
                    "Islamic product innovation critical — 25%+ of deposits in Islamic products, growing 20%+ annually",
                    "Digital product adoption determines competitive survival against fintechs (JazzCash 50M+, Easypaisa 40M+)",
                ]
                for addition in additions:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = addition
                    run.font.size = Pt(9)
                    run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                print(f"  Slide 10: Added Pakistan banking context to Why Important")
                break

    # Slide 14 (0-indexed 13): Product Management for Retail → for Banking
    slide = slides[13]
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if 'Retail' in run.text:
                        run.text = run.text.replace('Retail', 'Banking')
                    if 'retail' in run.text:
                        run.text = run.text.replace('retail', 'banking')

    # Slide 15 (0-indexed 14): Are You Able To — add Pakistan questions
    slide = slides[14]
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            if 'Able To' in text or 'able to' in text or 'Are You' in text:
                tf = shape.text_frame
                pk_questions = [
                    "Can you calculate true profitability for each product in your portfolio?",
                    "Can you track product performance across all channels (branch, ATM, digital, agent)?",
                    "Can you identify which products drive cross-sell and which cannibalize?",
                    "Can you optimize deposit pricing (rate) vs. volume trade-off analytically?",
                    "Can you measure product onboarding success rate by channel?",
                    "Can you predict which products a customer is likely to close next?",
                    "Can you launch a new product variant in under 4 weeks?",
                ]
                for q in pk_questions:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = q
                    run.font.size = Pt(9)
                    run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                print(f"  Slide 15: Added Pakistan assessment questions")
                break

    # Slide 16 (0-indexed 15): Business Challenges — add Pakistan challenges
    slide = slides[15]
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text
            if 'Challenge' in text or 'challenge' in text:
                tf = shape.text_frame
                pk_challenges = [
                    "SBP product approval process adds 4-8 weeks to launch timeline",
                    "Islamic products require dual approval: Shariah board + SBP",
                    "4-5 core banking systems make product configuration inconsistent",
                    "Pricing constrained: deposit floor (SBP), lending ceiling (KIBOR + risk), card rates (SBP caps under discussion)",
                    "Product data fragmented across core banking, cards, trade finance, Islamic sub-systems",
                ]
                for c in pk_challenges:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = c
                    run.font.size = Pt(9)
                    run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
                print(f"  Slide 16: Added Pakistan business challenges")
                break

    # ── Step 8: Add 3 new slides ──
    print(f"\n── Step 8: Adding 3 new slides ──")
    # Find a blank layout
    blank_layout = None
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower() or layout.name == 'Blank':
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[-1]

    # New Slide A: Dashboard
    slide_a = prs.slides.add_slide(blank_layout)
    _build_dashboard_slide(slide_a, prs)
    set_notes(slide_a, NEW_SLIDE_NOTES['new_a'])
    print(f"  Added: Domain Dashboard (New Slide A)")

    # New Slide B: Pakistan Context
    slide_b = prs.slides.add_slide(blank_layout)
    _build_pakistan_context_slide(slide_b, prs)
    set_notes(slide_b, NEW_SLIDE_NOTES['new_b'])
    print(f"  Added: Pakistan Product Management Context (New Slide B)")

    # New Slide C: Implementation Roadmap
    slide_c = prs.slides.add_slide(blank_layout)
    _build_roadmap_slide(slide_c, prs)
    set_notes(slide_c, NEW_SLIDE_NOTES['new_c'])
    print(f"  Added: Implementation Roadmap (New Slide C)")

    # ── Step 9: Reorder slides ──
    print(f"\n── Step 9: Reordering slides ──")
    total = len(prs.slides)
    _reorder_slides(prs, total)
    print(f"  Slides reordered: Dashboard, PK Context, original 54, Roadmap")

    # ── Step 10: Save ──
    print(f"\n── Step 10: Saving to {OUTPUT_FILE} ──")
    prs.save(OUTPUT_FILE)

    # Clean up temp file
    if os.path.exists(CLEAN_FILE):
        os.remove(CLEAN_FILE)
        print("  Cleaned up temp file")

    final_count = len(prs.slides)
    print(f"\n{'=' * 60}")
    print(f"DONE! {final_count} slides saved to {OUTPUT_FILE}")
    print(f"{'=' * 60}")


# ─────────────────────────────────────────────
# New slide builders
# ─────────────────────────────────────────────

def _build_dashboard_slide(slide, prs):
    """Build the Domain Dashboard slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                shape.text = "Product Management — At a Glance"
            elif shape.placeholder_format.idx == 1:
                shape.text = ""

    left = Emu(400000)
    top = Emu(900000)
    width = Emu(8300000)
    height = Emu(3800000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Product Management — At a Glance"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    p = tf.add_paragraph()
    p.text = ""

    lines = [
        "                                  Global          South Asia & ME    Pakistan",
        "───────────────────────────────────────────────────────────────────────────────",
        "Products per customer             4-6 top banks   2-3                1.8 (massive opportunity)",
        "Product profitability analytics   70%+ large banks 30-40%            <15%",
        "Dynamic pricing capability        50%+ digital    15-25%             <10%",
        "Product A/B testing               Standard digital 15-20%            <5%",
        "Product lifecycle analytics       60% large banks 25-30%             <15%",
        "Islamic product suite             N/A (where appl) Growing           Required (25%+ deposits)",
        "Digital product launch speed      2-4 weeks       4-8 weeks          8-16 weeks (SBP + IT)",
        "───────────────────────────────────────────────────────────────────────────────",
    ]
    for line in lines:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(8)
        run.font.name = "Consolas"

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = ("BVF Sub-capabilities: 12 across 6 sub-domains | FSDM: Product, Agreement, Party, Transaction, Channel, Pricing\n"
                "BACR: ~60 questions | Maturity: Developing → Innovating\n"
                "Products: CASA, Term Deposits, Personal Loans, Home Finance, Auto Finance, Credit Cards, Debit Cards,\n"
                "Trade Finance (LC/LG/SBLC), Islamic (Murabaha, Ijarah, Diminishing Musharaka, Musharaka, Wakalah, Takaful)")
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 102, 153)


def _build_pakistan_context_slide(slide, prs):
    """Build the Pakistan Product Management Landscape slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                shape.text = "Pakistan Banking — Product Management Landscape"
            elif shape.placeholder_format.idx == 1:
                shape.text = ""

    left = Emu(400000)
    top = Emu(850000)
    width = Emu(8300000)
    height = Emu(4200000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Product Portfolio
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Product Portfolio (Typical Large Pakistan Bank)"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    sections = {
        "DEPOSIT PRODUCTS": [
            "Current Account (zero interest, CASA), Savings Account (SBP min ~11-13%)",
            "Fixed Deposit (3M-5Y), Foreign Currency Account, Roshan Digital Account",
            "Basic Banking Account, Islamic: Mudaraba Savings, Mudaraba Term Deposit, Wakalah Pool",
        ],
        "LENDING PRODUCTS": [
            "Corporate: Working capital, term finance, project finance, trade finance (LC/LG/SBLC)",
            "Consumer: Personal loan, home finance, auto finance, credit card revolving",
            "SME/Agri: Running finance, crop/livestock finance",
            "Islamic: Murabaha, Ijarah (auto/home), Diminishing Musharaka, Running Musharaka",
        ],
        "CARD & DIGITAL": [
            "Debit (30M+ cards), Credit (3M+ cards), Prepaid (travel, remittance, payroll)",
            "Mobile banking, Internet banking, RAAST P2P/P2M, QR payments, Agent banking",
        ],
    }

    for header, items in sections.items():
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = header
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 51, 102)
        for item in items:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"  {item}"
            run.font.size = Pt(8)

    p = tf.add_paragraph()
    p.text = ""

    # Key Challenges
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Key Challenges"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    challenges = [
        "Products per customer: 1.8 (vs. 4-6 global benchmark)",
        "SBP product approval: New products require NOC — 4-8 weeks to launch",
        "Islamic product complexity: Shariah board + SBP approval = dual gate",
        "Pricing rigidity: Deposit rates (SBP minimum), lending (KIBOR), card rates (SBP caps)",
        "Cross-sell: <5% response rate on product campaigns",
        "Product cannibalization: Digital products cannibalizing branch products (unmanaged)",
    ]
    for c in challenges:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {c}"
        run.font.size = Pt(8)


def _build_roadmap_slide(slide, prs):
    """Build the Implementation Roadmap slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                shape.text = "Product Management — Implementation Roadmap"
            elif shape.placeholder_format.idx == 1:
                shape.text = ""

    left = Emu(400000)
    top = Emu(850000)
    width = Emu(8300000)
    height = Emu(3900000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Title
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Implementation Roadmap"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    p = tf.add_paragraph()
    p.text = ""

    # Phase 1
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Phase 1: Product Intelligence Foundation (0-6 months) — PKR 40-80M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 102, 51)

    phase1 = [
        "Product profitability analytics (from FSDM Customer Profitability Engine)",
        "Product performance dashboard (sales, balances, revenue, cost by product)",
        "Cross-channel product sales tracking",
        "Product benchmark analysis (vs. industry averages from SBP data)",
        "Quick Win: Identify bottom 10% unprofitable products",
    ]
    for line in phase1:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Phase 2
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Phase 2: Product Optimization (6-18 months) — PKR 80-150M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 153)

    phase2 = [
        "Competitive pricing analytics (deposit rates, lending spreads, card fees)",
        "Product onboarding journey optimization (reduce dropoff from 40-60% to <20%)",
        "Cross-sell propensity models by product affinity",
        "Islamic product performance parity analysis (Islamic vs. conventional)",
        "Expected: Increase products/customer from 1.8 to 2.5",
    ]
    for line in phase2:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Phase 3
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Phase 3: Advanced Product Management (18-36 months) — PKR 120-250M"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(153, 0, 51)

    phase3 = [
        "Dynamic pricing engine (risk-adjusted lending, competitive deposit pricing)",
        "A/B testing framework for product features and pricing",
        "Product lifecycle management (launch → growth → maturity → sunset)",
        "Predictive product retention (which customers will close which products?)",
        "Real-time product recommendations integrated with all channels",
        "Expected: 15-25% improvement in product revenue per customer",
    ]
    for line in phase3:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {line}"
        run.font.size = Pt(9)

    # Total
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Total 3-Year Investment: PKR 240-480M (~USD 0.9-1.7M) | Annual Value: PKR 15-25B"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)


def _reorder_slides(prs, total_slides):
    """Reorder slides: Dashboard (54), PK Context (55), Original 0-53, Roadmap (56).

    Final order: [54, 55, 0-53, 56]
    Dashboard, PK Context, all original 54 slides, Roadmap at end
    """
    new_order = [54, 55]  # Dashboard, PK Context
    new_order.extend(range(0, 54))  # All 54 original slides
    new_order.append(56)  # Roadmap at end

    pres_ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    pres_elem = prs.part._element
    sldIdLst = pres_elem.find('.//p:sldIdLst', pres_ns)

    sldId_elements = list(sldIdLst)

    if len(sldId_elements) != len(new_order):
        print(f"  WARNING: Expected {len(new_order)} slides but found {len(sldId_elements)}")
        # Adjust: use actual count
        n_orig = len(sldId_elements) - 3  # 3 new slides
        new_order = [n_orig, n_orig + 1]
        new_order.extend(range(0, n_orig))
        new_order.append(n_orig + 2)

    for elem in list(sldId_elements):
        sldIdLst.remove(elem)

    for idx in new_order:
        sldIdLst.append(sldId_elements[idx])

    print(f"  Reordered {len(new_order)} slides")


if __name__ == '__main__':
    main()
